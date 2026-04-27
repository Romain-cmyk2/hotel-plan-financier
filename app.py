"""
Plan Financier Hotel 5* - Application Streamlit
Reproduction interactive du classeur Excel Plan_Financier_Hotel.xlsx
"""

import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import date
import json
import os
import copy
import hashlib
from pathlib import Path
from calculs import (
    params_defaut, projection_complete, indicateurs_annuels,
    calc_tableau_pret, calc_prix_moyen_pondere, jours_dans_mois,
)
import github_sync

# ─── Sauvegarde / Chargement des plans ────────────────────────────────────────

PLANS_DIR = Path(__file__).parent / "plans"
PLANS_DIR.mkdir(exist_ok=True)


def _serialiser_params(p):
    """Convertit les params en dict JSON-compatible."""
    d = {}
    for k, v in p.items():
        if isinstance(v, date):
            d[k] = v.isoformat()
        else:
            d[k] = v
    return d


def _deserialiser_params(d):
    """Restaure les types depuis un dict JSON."""
    if "date_ouverture" in d and isinstance(d["date_ouverture"], str):
        d["date_ouverture"] = date.fromisoformat(d["date_ouverture"])
    return d


def _sync_plans_from_github_once():
    """
    Au premier run de la session, telecharge les plans depuis GitHub pour
    contourner le filesystem ephemere de Streamlit Cloud, puis indexe les
    hashs locaux pour eviter les pushs vides ulterieurs. Idempotent.
    """
    if st.session_state.get("_plans_synced", False):
        return
    if github_sync.is_enabled():
        try:
            github_sync.sync_directory_from_github(PLANS_DIR, "plans")
        except Exception:
            pass
        # Pre-remplir les hashs avec l'etat actuel du filesystem,
        # de sorte que le 1er auto-save sans vraie modif ne pousse pas.
        for f in PLANS_DIR.glob("*.json"):
            try:
                contenu = f.read_text(encoding="utf-8")
                st.session_state[f"_last_pushed_hash::{f.stem}"] = (
                    hashlib.md5(contenu.encode("utf-8")).hexdigest()
                )
            except OSError:
                pass
    st.session_state["_plans_synced"] = True


def sauvegarder_plan(nom, p, local_only=False):
    """
    Sauvegarde locale + push GitHub (sauf si local_only=True).
    Le push GitHub n'est fait que sur sauvegarde explicite (bouton). Les
    auto-save ecrivent uniquement en local pour eviter de polluer l'historique
    avec les rerun Streamlit (chaque interaction declenche _auto_save).
    """
    fichier = PLANS_DIR / f"{nom}.json"
    contenu = json.dumps(_serialiser_params(p), ensure_ascii=False, indent=2)
    with open(fichier, "w", encoding="utf-8") as f:
        f.write(contenu)
    if local_only or not github_sync.is_enabled():
        return
    # Hash check rapide en session pour ne pas re-pousser un contenu
    # deja synchronise lors de la session courante.
    new_hash = hashlib.md5(contenu.encode("utf-8")).hexdigest()
    hash_key = f"_last_pushed_hash::{nom}"
    if st.session_state.get(hash_key) == new_hash:
        return
    ok, msg = github_sync.push_file(
        f"plans/{nom}.json", contenu, f"Sauvegarde plan: {nom}"
    )
    if ok:
        st.session_state[hash_key] = new_hash
    else:
        try:
            st.toast(f"Sauvegarde locale OK, GitHub KO : {msg}", icon="⚠️")
        except Exception:
            pass


def charger_plan(nom):
    """Charge les parametres depuis un fichier sauvegarde."""
    fichier = PLANS_DIR / f"{nom}.json"
    with open(fichier, "r", encoding="utf-8") as f:
        d = json.load(f)
    return _deserialiser_params(d)


def lister_plans():
    """Liste les noms des plans sauvegardes (sync GitHub au premier appel)."""
    _sync_plans_from_github_once()
    return sorted(f.stem for f in PLANS_DIR.glob("*.json"))


def renommer_plan(ancien_nom, nouveau_nom):
    """Renomme un plan localement et sur GitHub."""
    ancien = PLANS_DIR / f"{ancien_nom}.json"
    nouveau = PLANS_DIR / f"{nouveau_nom}.json"
    if not (ancien.exists() and not nouveau.exists()):
        return False
    contenu = ancien.read_text(encoding="utf-8")
    ancien.rename(nouveau)
    if github_sync.is_enabled():
        ok, msg = github_sync.rename_file(
            f"plans/{ancien_nom}.json",
            f"plans/{nouveau_nom}.json",
            contenu,
            f"Renomme plan: {ancien_nom} -> {nouveau_nom}",
        )
        if not ok:
            try:
                st.toast(f"Renomme local OK, GitHub KO : {msg}", icon="⚠️")
            except Exception:
                pass
    return True


def supprimer_plan(nom):
    """Supprime un plan localement et sur GitHub."""
    fichier = PLANS_DIR / f"{nom}.json"
    if fichier.exists():
        fichier.unlink()
    if github_sync.is_enabled():
        ok, msg = github_sync.delete_file(
            f"plans/{nom}.json", f"Suppression plan: {nom}"
        )
        if not ok:
            try:
                st.toast(f"Supprime local OK, GitHub KO : {msg}", icon="⚠️")
            except Exception:
                pass

st.set_page_config(
    page_title="Plan Financier Hotel 5*",
    page_icon="\U0001F3E8",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ─── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1rem; border-radius: 12px; color: white; text-align: center;
    }
    .metric-card h3 { margin: 0; font-size: 0.8rem; opacity: 0.85; }
    .metric-card h1 { margin: 0.2rem 0 0 0; font-size: 1.5rem; }
    .metric-red { background: linear-gradient(135deg, #f5576c 0%, #ff6b6b 100%); }
    .metric-green { background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); }
    .metric-blue { background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); }
    .metric-orange { background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); }
    .stTabs [data-baseweb="tab-list"] { gap: 4px; }
    .stTabs [data-baseweb="tab"] {
        padding: 8px 16px; border-radius: 8px 8px 0 0;
        font-weight: 600; background-color: #e8e8e8;
    }
    .stTabs [aria-selected="true"] { background-color: #667eea; color: white; }
    .section-header {
        background: #f0f2f6; padding: 0.6rem 1rem; border-radius: 8px;
        margin: 0.5rem 0; font-weight: 700; font-size: 1.05rem;
    }
    @media print {
        /* Empecher les graphiques d'etre coupes */
        .js-plotly-plot, .plotly, .plot-container,
        [data-testid="stPlotlyChart"],
        [data-testid="stExpander"],
        .stPlotlyChart {
            overflow: visible !important;
            break-inside: avoid !important;
            page-break-inside: avoid !important;
        }
        /* Masquer les boutons et controles */
        button, [data-testid="stButton"],
        [data-testid="stSidebar"],
        header, footer,
        .stDeployButton {
            display: none !important;
        }
        /* Largeur pleine */
        .main .block-container {
            max-width: 100% !important;
            padding: 0 1cm !important;
        }
        /* Eviter les coupures sur les sections */
        .section-header, h2, h3 {
            break-after: avoid !important;
            page-break-after: avoid !important;
        }
    }
</style>
""", unsafe_allow_html=True)


def fmt_eur(val):
    if abs(val) >= 1_000_000:
        return f"{val/1_000_000:,.2f} M\u20ac"
    elif abs(val) >= 1_000:
        return f"{val/1_000:,.0f} K\u20ac"
    return f"{val:,.0f} \u20ac"


def fmt_pct(val):
    return f"{val:.1f}%"


def metric_card(title, value, css_class=""):
    st.markdown(f"""
    <div class="metric-card {css_class}">
        <h3>{title}</h3>
        <h1>{value}</h1>
    </div>""", unsafe_allow_html=True)


def _x_labels_annees(annual):
    """Génère des labels d'axe X avec juste l'année calendaire (ex: '2029')."""
    if "annee" in annual.columns:
        return [str(int(row["annee"])) for _, row in annual.iterrows()]
    return [f"A{int(a)}" for a in annual["annee_exploitation"]]


# Largeur minimale par barre/point pour garantir la lisibilité
_PX_PAR_BARRE = 55
_PX_PAR_MOIS = 30


def _scrollable_chart(fig, n_points, px_par_point, height, key_suffix=""):
    """Affiche un graphique avec scroll horizontal si trop de points."""
    import streamlit.components.v1 as components
    import json as _json

    min_width = n_points * px_par_point
    if min_width > 900:
        fig.update_layout(width=min_width, height=height)
        fig_json = fig.to_json()
        wrapper = f"""
        <script src="https://cdn.plot.ly/plotly-2.35.2.min.js"></script>
        <div style="overflow-x:auto; max-width:100%; -webkit-overflow-scrolling:touch;
                    border:1px solid #e0e0e0; border-radius:6px;">
            <div id="plotly-chart"></div>
        </div>
        <script>
            var figData = {fig_json};
            Plotly.newPlot('plotly-chart', figData.data, figData.layout, {{displayModeBar: false}});
        </script>
        """
        components.html(wrapper, height=height + 30, scrolling=False)
    else:
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


def _make_bar_fig(annual, col, title, color, height, en_milliers, show_labels, x_labels=None):
    """Cree un graphique a barres avec options d'affichage."""
    y_vals = annual[col]
    if en_milliers:
        y_vals = y_vals / 1000
        suffix = " K\u20ac"
        tick_fmt = ",.0f"
    else:
        suffix = " \u20ac"
        tick_fmt = ","

    if x_labels is None:
        x_labels = _x_labels_annees(annual)

    fig = go.Figure()
    text_vals = [f"{v:,.0f}{suffix}" for v in y_vals] if show_labels else None
    fig.add_trace(go.Bar(
        x=x_labels, y=y_vals,
        marker_color=color, opacity=0.85,
        text=text_vals,
        textposition="outside" if show_labels else "none",
        textfont=dict(size=13),
        hovertemplate="%{x}<br>%{y:,.2f}<extra></extra>",
    ))
    fig.update_layout(
        title=dict(text=title, font=dict(size=13)),
        height=height, margin=dict(l=60, r=10, t=35, b=30),
        xaxis=dict(title="", showticklabels=True, type="category",
                   tickfont=dict(size=12), tickangle=0),
        yaxis=dict(title="", tickformat=tick_fmt, showgrid=True),
        showlegend=False,
    )
    return fig


def _make_monthly_bar_fig(df_mensuel, col, title, color, height, en_milliers):
    """Cree un graphique a barres mensuel."""
    y_vals = df_mensuel[col]
    if en_milliers:
        y_vals = y_vals / 1000
        tick_fmt = ",.0f"
    else:
        tick_fmt = ","

    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df_mensuel["date"], y=y_vals,
        marker_color=color, opacity=0.85,
        hovertemplate="%{x}<br>%{y:,.2f}<extra></extra>",
    ))
    fig.update_layout(
        title=dict(text=title, font=dict(size=13)),
        height=height, margin=dict(l=60, r=10, t=35, b=30),
        xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                   tickfont=dict(size=11), tickangle=-45),
        yaxis=dict(title="", tickformat=tick_fmt, showgrid=True),
        showlegend=False,
    )
    return fig


def _chart_with_zoom(df, col, title, color, key_id):
    """Affiche un graphique dans un expander avec choix annuel/mensuel et unité."""
    annual = df.groupby("annee")[col].sum().reset_index()
    n_annees = len(annual)
    n_mois = len(df)

    with st.expander(f"{title}", expanded=False):
        c1, c2 = st.columns(2)
        with c1:
            unite = st.radio("Unite", ["\u20ac", "K\u20ac (milliers)"], key=f"unite_{key_id}", horizontal=True)
        with c2:
            periode = st.radio("Periode", ["Annuel", "Mensuel"], key=f"periode_{key_id}", horizontal=True)
        en_k = unite.startswith("K")

        if periode == "Annuel":
            fig_big = _make_bar_fig(annual, col, title, color, 450, en_k, True)
            _scrollable_chart(fig_big, n_annees, _PX_PAR_BARRE, 450)
        else:
            fig_big = _make_monthly_bar_fig(df, col, title, color, 450, en_k)
            _scrollable_chart(fig_big, n_mois, _PX_PAR_MOIS, 450)


def _build_waterfall_moyens_besoins(st_mod, go_mod, prets, fonds_propres, total_moyens,
                                     total_invest, besoin_tresorerie, solde, key_suffix="",
                                     pret_intra=0, pret_intra_label="Pret a Argenteau"):
    """Graphique waterfall Moyens / Besoins reutilisable (Chateau et Rocher)."""
    wf_labels = []
    wf_values = []
    wf_measures = []

    # Moyens — dedupliquer les noms de prets
    _nom_counts = {}
    for pr in prets:
        _n = pr["nom"]
        _nom_counts[_n] = _nom_counts.get(_n, 0) + 1
    _nom_idx = {}
    for pr in prets:
        _n = pr["nom"]
        if _nom_counts[_n] > 1:
            _nom_idx[_n] = _nom_idx.get(_n, 0) + 1
            label = f"{_n} ({_nom_idx[_n]})"
        else:
            label = _n
        wf_labels.append(label)
        wf_values.append(pr["montant"])
        wf_measures.append("relative")
    wf_labels.append("Fonds propres")
    wf_values.append(fonds_propres)
    wf_measures.append("relative")

    # Sous-total moyens
    wf_labels.append("TOTAL MOYENS")
    wf_values.append(total_moyens)
    wf_measures.append("total")

    # Besoins (negatif) — distinguer investissements et pret intra-groupe si applicable
    _invest_hors_pret = total_invest - pret_intra
    wf_labels.append("Investissements")
    wf_values.append(-_invest_hors_pret)
    wf_measures.append("relative")

    if pret_intra > 0:
        wf_labels.append(pret_intra_label)
        wf_values.append(-pret_intra)
        wf_measures.append("relative")

    if besoin_tresorerie > 0:
        wf_labels.append("Besoin tresorerie")
        wf_values.append(-besoin_tresorerie)
        wf_measures.append("relative")

    # Solde final
    wf_labels.append("SOLDE")
    wf_values.append(solde)
    wf_measures.append("total")

    _wf_vals_k = [v / 1000 for v in wf_values]
    # Positionner les labels: inside pour les grandes barres, outside pour les petites
    _max_abs = max(abs(v) for v in _wf_vals_k) if _wf_vals_k else 1
    _threshold = _max_abs * 0.08
    _positions = ["inside" if abs(v) > _threshold else "outside" for v in _wf_vals_k]
    _text_colors = []
    for i, v in enumerate(_wf_vals_k):
        if abs(v) > _threshold:
            _text_colors.append("white")
        else:
            _text_colors.append("#1a1a2e")
    fig_wf = go_mod.Figure(go_mod.Waterfall(
        x=wf_labels, y=_wf_vals_k,
        measure=wf_measures,
        connector=dict(line=dict(color="rgba(0,0,0,0.3)", width=1)),
        increasing=dict(marker=dict(color="#38ef7d")),
        decreasing=dict(marker=dict(color="#f5576c")),
        totals=dict(marker=dict(color="#667eea")),
        textposition=_positions,
        text=[f"<b>{v:,.0f} K\u20ac</b>" for v in _wf_vals_k],
        textfont=dict(size=12),
        customdata=[[v] for v in _wf_vals_k],
        hovertemplate="%{x}<br><b>%{customdata[0]:,.0f} K\u20ac</b><extra></extra>",
    ))
    # Appliquer les couleurs de texte par barre
    fig_wf.update_traces(textfont_color=_text_colors)
    fig_wf.update_layout(
        title=dict(text="Cascade Moyens / Besoins (K\u20ac)", font=dict(size=16, color="#1a1a2e")),
        height=600,
        xaxis=dict(tickfont=dict(size=10, color="#333"), tickangle=-30),
        yaxis=dict(tickformat=",.0f", title="K\u20ac", tickfont=dict(size=12)),
        showlegend=False,
        margin=dict(t=60, b=110),
        uniformtext=dict(minsize=9, mode="show"),
    )
    st_mod.plotly_chart(fig_wf, use_container_width=True, config={"displayModeBar": False},
                        key=f"wf_moyens_besoins_{key_suffix}")


def _render_hypotheses_charts(df, indic, params_hyp=None):
    """Affiche les mini-graphiques dans l'onglet hypotheses."""
    if "chart_containers" not in st.session_state:
        return

    # Colonne brasserie hors PDJ
    if "ca_brasserie" in df.columns and "ca_pdj" in df.columns:
        df["ca_brasserie_hors_pdj"] = df["ca_brasserie"] - df["ca_pdj"]

    # Colonne mariages total (location + commission catering)
    if "ca_mariages" in df.columns and "ca_commission_catering" in df.columns:
        df["ca_mariages_total"] = df["ca_mariages"] + df["ca_commission_catering"]

    containers = st.session_state["chart_containers"]

    chart_configs = {
        "hebergement": [
            ("ca_hebergement", "CA Hebergement / an", "#4facfe"),
            ("ca_divers", "CA Divers (mini-bar, ...) / an", "#f093fb"),
            ("nuitees", "Nuitees / an", "#764ba2"),
        ],
        "brasserie": [
            ("ca_pdj", "CA Petit-dejeuner / an", "#38ef7d"),
            ("ca_brasserie_hors_pdj", "CA Brasserie (midi & soir) / an", "#f5576c"),
        ],
        "bar": [
            ("ca_bar", "CA Bar / an", "#ffcc00"),
        ],
        "spa": [
            ("ca_spa", "CA Spa / an", "#11998e"),
        ],
        "seminaires": [
            ("ca_seminaires", "CA Seminaires / an", "#667eea"),
        ],
        "mariages": [
            ("ca_mariages_total", "CA Mariages (location + catering) / an", "#f093fb"),
        ],
        "salles_chateau": [
            ("ca_salles_chateau", "CA Salles chateau / an", "#a0522d"),
        ],
        "loyer_restaurant": [
            ("ca_loyer_restaurant", "CA Loyer restaurant / an", "#ff8c00"),
        ],
        "total": [
            ("ca_total", "CA Total / an", "#4facfe"),
            ("ca_hebergement", "Hebergement", "#667eea"),
            ("ca_brasserie", "Brasserie", "#f5576c"),
            ("ca_bar", "Bar", "#ffcc00"),
            ("ca_spa", "Spa", "#11998e"),
            ("ca_salles", "Salles & Evenements", "#a0522d"),
            ("ca_loyer_restaurant", "Loyer restaurant", "#ff8c00"),
        ],
    }

    for key, container in containers.items():
        if key not in chart_configs:
            continue
        configs = chart_configs[key]
        with container:
            if key == "total":
                agg_cols = {c[0]: "sum" for c in configs}
                annual = df.groupby("annee").agg(agg_cols).reset_index()
                x_labels = [str(int(a)) for a in annual["annee"]]
                n_annees = len(annual)
                n_mois = len(df)

                with st.expander("CA Total par service", expanded=False):
                    tab_bar, tab_pie = st.tabs(["\U0001F4CA Evolution", "\U0001F967 Ventilation"])

                    with tab_bar:
                        c1, c2 = st.columns(2)
                        with c1:
                            unite_t = st.radio("Unite", ["\u20ac", "K\u20ac (milliers)"], key="unite_total", horizontal=True)
                        with c2:
                            periode_t = st.radio("Periode", ["Annuel", "Mensuel"], key="periode_total", horizontal=True)
                        en_k = unite_t.startswith("K")
                        diviseur = 1000 if en_k else 1
                        sfx = " K\u20ac" if en_k else " \u20ac"

                        fig_big = go.Figure()
                        _legend_cfg = dict(orientation="h", y=1.08, font=dict(size=14),
                                           xanchor="center", x=0.5,
                                           bgcolor="rgba(255,255,255,0.8)",
                                           bordercolor="rgba(0,0,0,0.1)", borderwidth=1)
                        if periode_t == "Annuel":
                            for col_name, label, color in configs[1:]:
                                vals = annual[col_name] / diviseur
                                fig_big.add_trace(go.Bar(
                                    x=x_labels, y=vals,
                                    name=label, marker_color=color,
                                    hovertemplate="%{x} - " + label + "<br>%{y:,.2f}<extra></extra>",
                                ))
                            # Total en K€ au-dessus de chaque barre
                            totals_raw = annual[configs[0][0]] / diviseur
                            totals_k = annual[configs[0][0]] / 1000
                            max_total = totals_raw.max() if len(totals_raw) > 0 else 1
                            fig_big.add_trace(go.Scatter(
                                x=x_labels,
                                y=[v * 1.02 for v in totals_raw],
                                mode="text",
                                text=[f"<b>{v:,.0f} K\u20ac</b>" for v in totals_k],
                                textposition="top center",
                                textfont=dict(size=10, color="black"),
                                showlegend=False, hoverinfo="skip",
                                cliponaxis=False,
                            ))
                            fig_big.update_layout(
                                title=dict(text="CA Total par service / an", font=dict(size=15)),
                                height=550, barmode="stack",
                                xaxis=dict(type="category", title="", tickfont=dict(size=12)),
                                yaxis=dict(tickformat=",.0f", title="",
                                           range=[0, max_total * 1.15]),
                                showlegend=True, legend=_legend_cfg,
                            )
                            st.plotly_chart(fig_big, use_container_width=True, config={"displayModeBar": False})
                        else:
                            for col_name, label, color in configs[1:]:
                                vals = df[col_name] / diviseur
                                fig_big.add_trace(go.Bar(
                                    x=df["date"], y=vals,
                                    name=label, marker_color=color,
                                    hovertemplate="%{x} - " + label + "<br>%{y:,.2f}<extra></extra>",
                                ))
                            # Total en K€ au-dessus de chaque barre mensuelle
                            totals_m_raw = df[configs[0][0]] / diviseur
                            totals_m_k = df[configs[0][0]] / 1000
                            max_total_m = totals_m_raw.max() if len(totals_m_raw) > 0 else 1
                            fig_big.add_trace(go.Scatter(
                                x=df["date"],
                                y=[v * 1.02 for v in totals_m_raw],
                                mode="text",
                                text=[f"<b>{v:,.0f} K\u20ac</b>" for v in totals_m_k],
                                textposition="top center",
                                textfont=dict(size=8, color="black"),
                                showlegend=False, hoverinfo="skip",
                                cliponaxis=False,
                            ))
                            fig_big.update_layout(
                                title=dict(text="CA Total par service / mois", font=dict(size=15)),
                                height=550, barmode="stack",
                                xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                                           tickfont=dict(size=11), tickangle=-45),
                                yaxis=dict(tickformat=",.0f", title="",
                                           range=[0, max_total_m * 1.15]),
                                showlegend=True, legend=_legend_cfg,
                            )
                            st.plotly_chart(fig_big, use_container_width=True, config={"displayModeBar": False})

                    with tab_pie:
                        annees_dispo = [str(int(a)) for a in annual["annee"]]
                        toutes = st.checkbox("Toutes les annees", value=True, key="pie_all_years")
                        if toutes:
                            annees_sel = annees_dispo
                        else:
                            annees_sel = st.multiselect(
                                "Annees", annees_dispo, default=[annees_dispo[0]],
                                key="pie_annees_total",
                            )
                        if annees_sel:
                            mask = annual["annee"].astype(int).astype(str).isin(annees_sel)
                            filtered = annual[mask]
                            labels_pie = [c[1] for c in configs[1:]]
                            colors_pie = [c[2] for c in configs[1:]]
                            values_pie = [filtered[c[0]].sum() for c in configs[1:]]
                            fig_pie = go.Figure(data=[go.Pie(
                                labels=labels_pie, values=values_pie,
                                marker=dict(colors=colors_pie),
                                textinfo="label+percent",
                                textfont=dict(size=14),
                                hovertemplate="%{label}<br>%{value:,.0f} \u20ac<br>%{percent}<extra></extra>",
                            )])
                            periode_label = annees_sel[0] if len(annees_sel) == 1 else f"{annees_sel[0]} - {annees_sel[-1]}"
                            fig_pie.update_layout(
                                title=dict(text=f"Repartition CA par service ({periode_label})", font=dict(size=15)),
                                height=450,
                                showlegend=True,
                                legend=dict(orientation="h", y=-0.05, font=dict(size=14),
                                            xanchor="center", x=0.5),
                            )
                            st.plotly_chart(fig_pie, use_container_width=True, config={"displayModeBar": False})
                        else:
                            st.info("Selectionnez au moins une annee.")

            elif len(configs) > 1:
                for i, (col_name, label, color) in enumerate(configs):
                    _chart_with_zoom(df, col_name, label, color, f"{key}_{i}")
            else:
                col_name, label, color = configs[0]
                _chart_with_zoom(df, col_name, label, color, key)

    # ── Graphiques pour les autres onglets ──────────────────────────────

    # 2. Frais variables
    if "cv_detail_heberg" in containers:
        with containers["cv_detail_heberg"]:
            _chart_with_zoom(df, "cv_hebergement", "CV Hebergement / an", "#4facfe", "cv_cv_hebergement")

    for cv_key, cv_col, cv_label, cv_color in [
        ("cv_brasserie", "cv_brasserie", "CV Brasserie / an", "#f5576c"),
        ("cv_bar", "cv_bar", "CV Bar / an", "#ffcc00"),
        ("cv_spa", "cv_spa", "CV Spa / an", "#11998e"),
        ("cv_seminaires", "cv_seminaires", "CV Seminaires / an", "#667eea"),
        ("cv_mariages", "cv_mariages", "CV Mariages / an", "#f093fb"),
        ("cv_salles_chateau", "cv_salles_chateau", "CV Salles chateau / an", "#a0522d"),
    ]:
        if cv_key in containers:
            with containers[cv_key]:
                _chart_with_zoom(df, cv_col, cv_label, cv_color, f"cv_{cv_col}")

    if "cv_total" in containers:
        with containers["cv_total"]:
            cv_configs = [
                ("cv_total", "CV Total", "#f5576c"),
                ("cv_hebergement", "Hebergement", "#4facfe"),
                ("cv_brasserie", "Brasserie", "#f5576c"),
                ("cv_bar", "Bar", "#ffcc00"),
                ("cv_spa", "Spa", "#11998e"),
                ("cv_seminaires", "Seminaires", "#667eea"),
                ("cv_mariages", "Mariages", "#f093fb"),
                ("cv_salles_chateau", "Salles chateau", "#a0522d"),
            ]
            agg_cv = {c[0]: "sum" for c in cv_configs}
            annual_cv = df.groupby("annee").agg(agg_cv).reset_index()
            x_labels_cv = [str(int(a)) for a in annual_cv["annee"]]
            n_annees_cv = len(annual_cv)

            with st.expander("Total charges variables par service", expanded=False):
                tab_cv_bar, tab_cv_pie = st.tabs(["\U0001F4CA Evolution", "\U0001F967 Ventilation"])

                with tab_cv_bar:
                    c1, c2 = st.columns(2)
                    with c1:
                        unite_cv = st.radio("Unite", ["\u20ac", "K\u20ac (milliers)"], key="unite_cv_total", horizontal=True)
                    with c2:
                        periode_cv = st.radio("Periode", ["Annuel", "Mensuel"], key="periode_cv_total", horizontal=True)
                    en_k_cv = unite_cv.startswith("K")
                    diviseur_cv = 1000 if en_k_cv else 1

                    _legend_cv = dict(orientation="h", y=1.08, font=dict(size=14),
                                      xanchor="center", x=0.5,
                                      bgcolor="rgba(255,255,255,0.8)",
                                      bordercolor="rgba(0,0,0,0.1)", borderwidth=1)

                    fig_cv = go.Figure()
                    if periode_cv == "Annuel":
                        for col_name, label, color in cv_configs[1:]:
                            vals = annual_cv[col_name] / diviseur_cv
                            fig_cv.add_trace(go.Bar(
                                x=x_labels_cv, y=vals,
                                name=label, marker_color=color,
                                hovertemplate="%{x} - " + label + "<br>%{y:,.2f}<extra></extra>",
                            ))
                        totals_cv_raw = annual_cv[cv_configs[0][0]] / diviseur_cv
                        totals_cv_k = annual_cv[cv_configs[0][0]] / 1000
                        max_cv = totals_cv_raw.max() if len(totals_cv_raw) > 0 else 1
                        fig_cv.add_trace(go.Scatter(
                            x=x_labels_cv,
                            y=[v * 1.02 for v in totals_cv_raw],
                            mode="text",
                            text=[f"<b>{v:,.0f} K\u20ac</b>" for v in totals_cv_k],
                            textposition="top center",
                            textfont=dict(size=10, color="black"),
                            showlegend=False, hoverinfo="skip",
                            cliponaxis=False,
                        ))
                        fig_cv.update_layout(
                            title=dict(text="Charges variables par service / an", font=dict(size=15)),
                            height=550, barmode="stack",
                            xaxis=dict(type="category", title="", tickfont=dict(size=12)),
                            yaxis=dict(tickformat=",.0f", title="",
                                       range=[0, max_cv * 1.15]),
                            showlegend=True, legend=_legend_cv,
                        )
                        st.plotly_chart(fig_cv, use_container_width=True, config={"displayModeBar": False})
                    else:
                        for col_name, label, color in cv_configs[1:]:
                            vals = df[col_name] / diviseur_cv
                            fig_cv.add_trace(go.Bar(
                                x=df["date"], y=vals,
                                name=label, marker_color=color,
                                hovertemplate="%{x} - " + label + "<br>%{y:,.2f}<extra></extra>",
                            ))
                        totals_m_cv_raw = df[cv_configs[0][0]] / diviseur_cv
                        totals_m_cv_k = df[cv_configs[0][0]] / 1000
                        max_cv_m = totals_m_cv_raw.max() if len(totals_m_cv_raw) > 0 else 1
                        fig_cv.add_trace(go.Scatter(
                            x=df["date"],
                            y=[v * 1.02 for v in totals_m_cv_raw],
                            mode="text",
                            text=[f"<b>{v:,.0f} K\u20ac</b>" for v in totals_m_cv_k],
                            textposition="top center",
                            textfont=dict(size=8, color="black"),
                            showlegend=False, hoverinfo="skip",
                            cliponaxis=False,
                        ))
                        fig_cv.update_layout(
                            title=dict(text="Charges variables par service / mois", font=dict(size=15)),
                            height=550, barmode="stack",
                            xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                                       tickfont=dict(size=11), tickangle=-45),
                            yaxis=dict(tickformat=",.0f", title="",
                                       range=[0, max_cv_m * 1.15]),
                            showlegend=True, legend=_legend_cv,
                        )
                        st.plotly_chart(fig_cv, use_container_width=True, config={"displayModeBar": False})

                with tab_cv_pie:
                    annees_cv_dispo = [str(int(a)) for a in annual_cv["annee"]]
                    toutes_cv = st.checkbox("Toutes les annees", value=True, key="pie_all_years_cv")
                    if toutes_cv:
                        annees_cv_sel = annees_cv_dispo
                    else:
                        annees_cv_sel = st.multiselect(
                            "Annees", annees_cv_dispo, default=[annees_cv_dispo[0]],
                            key="pie_annees_cv",
                        )
                    if annees_cv_sel:
                        mask_cv = annual_cv["annee"].astype(int).astype(str).isin(annees_cv_sel)
                        filtered_cv = annual_cv[mask_cv]
                        labels_cv_pie = [c[1] for c in cv_configs[1:]]
                        colors_cv_pie = [c[2] for c in cv_configs[1:]]
                        values_cv_pie = [filtered_cv[c[0]].sum() for c in cv_configs[1:]]
                        fig_cv_pie = go.Figure(data=[go.Pie(
                            labels=labels_cv_pie, values=values_cv_pie,
                            marker=dict(colors=colors_cv_pie),
                            textinfo="label+percent",
                            textfont=dict(size=14),
                            hovertemplate="%{label}<br>%{value:,.0f} \u20ac<br>%{percent}<extra></extra>",
                        )])
                        lbl_cv = annees_cv_sel[0] if len(annees_cv_sel) == 1 else f"{annees_cv_sel[0]} - {annees_cv_sel[-1]}"
                        fig_cv_pie.update_layout(
                            title=dict(text=f"Repartition CV par service ({lbl_cv})", font=dict(size=15)),
                            height=450,
                            showlegend=True,
                            legend=dict(orientation="h", y=-0.05, font=dict(size=14),
                                        xanchor="center", x=0.5),
                        )
                        st.plotly_chart(fig_cv_pie, use_container_width=True, config={"displayModeBar": False})
                    else:
                        st.info("Selectionnez au moins une annee.")

    # 3. Frais fixes directs
    if "cf_directs" in containers:
        with containers["cf_directs"]:
            _chart_with_zoom(df, "cf_personnel_direct", "Personnel direct / an", "#667eea", "cf_personnel_dir")
            _chart_with_zoom(df, "cf_directs_total", "Total charges fixes directes / an", "#764ba2", "cf_directs_chart")

    # 4. Frais fixes indirects
    if "cf_indirects" in containers:
        with containers["cf_indirects"]:
            _chart_with_zoom(df, "cf_autres", "Charges fixes indirectes / an", "#f093fb", "cf_autres_chart")

    # 5. Investissements & Financement
    if "invest_finance" in containers:
        with containers["invest_finance"]:
            _chart_with_zoom(df, "amortissement", "Amortissements / an", "#667eea", "amort_chart")

    # Equilibre Moyens / Besoins
    if "moyens_besoins" in containers and params_hyp is not None:
        with containers["moyens_besoins"]:
            total_inv_mb = sum(i["montant"] for i in params_hyp.get("investissements", []))
            prets_mb = params_hyp.get("prets", [])
            total_prets_mb = sum(pr["montant"] for pr in prets_mb)
            fonds_propres_mb = params_hyp.get("fonds_propres_initial", 0)

            # Besoin de tresorerie : cash flow cumule minimum (point le plus bas)
            cf_cumul = df["cash_flow"].cumsum()
            tresorerie_min = cf_cumul.min()
            besoin_tresorerie = abs(min(0, tresorerie_min))

            # Mois pour atteindre l'equilibre
            mois_equilibre = None
            for idx_m, val in enumerate(cf_cumul):
                if val > 0 and idx_m > 0:
                    mois_equilibre = idx_m + 1
                    break

            total_moyens = total_prets_mb + fonds_propres_mb
            total_besoins = total_inv_mb + besoin_tresorerie
            solde = total_moyens - total_besoins

            # Tableau comparatif
            col_m, col_b = st.columns(2)

            with col_m:
                st.markdown("#### Moyens (ressources)")
                for pr in prets_mb:
                    st.markdown(f"- {pr['nom']} : **{pr['montant']:,.0f} \u20ac**")
                st.markdown(f"- Fonds propres : **{fonds_propres_mb:,.0f} \u20ac**")
                st.markdown(f"**TOTAL MOYENS : {total_moyens:,.0f} \u20ac**")

            with col_b:
                st.markdown("#### Besoins (emplois)")
                st.markdown(f"- Investissements initiaux : **{total_inv_mb:,.0f} \u20ac**")
                st.markdown(f"- Besoin tresorerie activite : **{besoin_tresorerie:,.0f} \u20ac**")
                st.markdown(f"**TOTAL BESOINS : {total_besoins:,.0f} \u20ac**")

            # Solde
            if solde >= 0:
                st.success(f"**Solde : +{solde:,.0f} \u20ac** — Les moyens couvrent les besoins"
                           + (f" avec une marge de {solde/total_besoins*100:.1f}%." if total_besoins > 0 else "."))
            else:
                st.error(f"**Deficit : {solde:,.0f} \u20ac** — Il manque **{abs(solde):,.0f} \u20ac** "
                         f"pour couvrir les besoins.")

            if mois_equilibre:
                st.caption(f"Point d'equilibre de tresorerie atteint au mois {mois_equilibre} "
                           f"(~{mois_equilibre // 12} an(s) et {mois_equilibre % 12} mois apres ouverture)")

            # Graphique en cascade (waterfall)
            _build_waterfall_moyens_besoins(
                st, go, prets_mb, fonds_propres_mb, total_moyens,
                total_inv_mb, besoin_tresorerie, solde, "Chateau"
            )

    # 6. Fiscalite
    if "fiscalite" in containers:
        with containers["fiscalite"]:
            _chart_with_zoom(df, "impot", "Impot des societes (ISOC) / an", "#f5576c", "impot_chart")
            _chart_with_zoom(df, "resultat_avant_impot", "Resultat avant impot / an", "#38ef7d", "rai_chart")


# ─── Onglet Hypotheses (corps principal, pleine largeur) ─────────────────────

def _personnel_table(p, personnel_list, key_prefix, cp, nb_chambres=None, ratio_base=None, ratio_label="ETP/10", allow_delete=True):
    """Affiche un tableau personnel editable avec ajout/suppression.
    Si ratio_base est fourni, affiche le ratio ETP / (ratio_base/10) par ligne.
    nb_chambres est un alias pour ratio_base (retro-compatibilite)."""
    if ratio_base is None and nb_chambres is not None:
        ratio_base = nb_chambres
    # Init session state for this table
    state_key = f"_pers_{key_prefix}"
    version_key = f"_pers_ver_{key_prefix}"
    if state_key not in st.session_state:
        st.session_state[state_key] = list(personnel_list)
    if version_key not in st.session_state:
        st.session_state[version_key] = 0
    _ver = st.session_state[version_key]

    current = st.session_state[state_key]
    new_list = []
    to_delete = []

    for i, pers in enumerate(current):
        if ratio_base and allow_delete:
            c1, c2, c2b, c3, c4a, c4b, c5 = st.columns([3, 1.5, 1.0, 2, 1.5, 1.5, 0.5])
        elif ratio_base:
            c1, c2, c2b, c3, c4a, c4b = st.columns([3, 1.5, 1.0, 2, 1.5, 1.5])
        elif allow_delete:
            c1, c2, c3, c4a, c4b, c5 = st.columns([3, 1.5, 2, 1.5, 1.5, 0.5])
        else:
            c1, c2, c3, c4a, c4b = st.columns([3, 1.5, 2, 1.5, 1.5])
        with c1:
            poste = st.text_input("p", pers["poste"], key=f"{key_prefix}_p_{_ver}_{i}", label_visibility="collapsed")
        with c2:
            etp = st.number_input("e", 0.0, 50.0, float(pers["etp"]), step=0.5, key=f"{key_prefix}_e_{_ver}_{i}", label_visibility="collapsed")
        if ratio_base:
            with c2b:
                ratio = etp / (ratio_base / 10) if ratio_base > 0 else 0
                st.markdown(f"*{ratio:.2f}*")
        with c3:
            cout = st.number_input("c", 0, 250_000, int(pers["cout_brut"]), step=1_000, key=f"{key_prefix}_c_{_ver}_{i}", label_visibility="collapsed")
        with c4a:
            _cout_an_etp = cout * (1 + cp)
            st.markdown(f'<div style="text-align:center"><b>{_cout_an_etp:,.0f} \u20ac</b></div>', unsafe_allow_html=True)
        with c4b:
            _cout_an_total = cout * (1 + cp) * etp
            st.markdown(f'<div style="text-align:center"><b>{_cout_an_total:,.0f} \u20ac</b></div>', unsafe_allow_html=True)
        if allow_delete:
            with c5:
                _visu = st.session_state.get("_visu_mode", False)
                if st.button("\U0001F5D1", key=f"{key_prefix}_del_{_ver}_{i}", help="Supprimer", disabled=_visu):
                    to_delete.append(i)
        new_list.append({"poste": poste, "etp": etp, "cout_brut": cout})

    # Supprimer les entrées marquées
    if to_delete:
        # Reconstruire depuis current (source de verite), pas depuis new_list (valeurs widget)
        cleaned = [e for idx, e in enumerate(current) if idx not in to_delete]
        st.session_state[state_key] = cleaned
        st.session_state[version_key] = _ver + 1
        st.rerun()

    # Bouton ajouter
    _visu = st.session_state.get("_visu_mode", False)
    if st.button(f"\u2795 Ajouter un poste", key=f"{key_prefix}_add", disabled=_visu):
        new_list.append({"poste": "Nouveau poste", "etp": 1.0, "cout_brut": 30_000})
        st.session_state[state_key] = new_list
        st.rerun()

    st.session_state[state_key] = new_list
    return new_list


def _cf_grid(p, cf_dict, key_prefix, nb_cols=4, charts=None, protected_keys=None):
    """Affiche une grille de charges fixes editables avec ajout/suppression.
    'Energie fixe' a un graphique retractable, 'Precompte immobilier' un popover annees.
    protected_keys : set de noms de charges qui ne peuvent pas etre supprimees."""
    if protected_keys is None:
        protected_keys = set()
    state_key = f"_cf_{key_prefix}"
    version_key = f"_cf_ver_{key_prefix}"
    if state_key not in st.session_state:
        st.session_state[state_key] = dict(cf_dict)
    if version_key not in st.session_state:
        st.session_state[version_key] = 0
    _ver = st.session_state[version_key]

    current = st.session_state[state_key]
    new_cf = {}
    to_delete = []
    items = list(current.items())

    for i, (k, v) in enumerate(items):
        is_energie = "nergie" in k and "fixe" in k.lower()
        is_precompte = "recompte" in k or "fonciere" in k.lower() or "fonci" in k.lower()

        c1, c2, c3, c4 = st.columns([3, 2, 0.5, 0.5])
        with c1:
            new_label = st.text_input("l", k, key=f"{key_prefix}_l_{_ver}_{i}", label_visibility="collapsed")
        with c2:
            new_val = st.number_input("v", 0, 1_000_000, int(v), step=1_000, key=f"{key_prefix}_v_{_ver}_{i}", label_visibility="collapsed")
        with c3:
            _is_protected = k in protected_keys
            _visu = st.session_state.get("_visu_mode", False)
            if _is_protected:
                st.markdown("")  # placeholder vide
            elif st.button("\U0001F5D1", key=f"{key_prefix}_del_{_ver}_{i}", help="Supprimer", disabled=_visu):
                to_delete.append(i)
        with c4:
            if is_energie:
                _ek = f"_energie_expand_{key_prefix}"
                if _ek not in st.session_state:
                    st.session_state[_ek] = False
                st.session_state[_ek] = st.checkbox("\u26A1 Detail",
                    value=st.session_state[_ek],
                    key=f"{key_prefix}_enrg_{_ver}_{i}",
                    help="Afficher/masquer le graphique energie fixe + variable")
            elif is_precompte:
                with st.popover("\U0001F4C5", help="Annees d'application"):
                    st.markdown("**Annees d'application de la taxe fonciere**")
                    nb_annees_proj = p.get("nb_mois_projection", 204) // 12 + 1
                    annees_list = list(range(1, nb_annees_proj + 1))
                    default_actives = p.get("precompte_annees_actives", annees_list)
                    tout_select = st.checkbox("Toutes les annees",
                        value=len(default_actives) >= len(annees_list), key="precompte_all_pop")
                    if tout_select:
                        p["precompte_annees_actives"] = annees_list
                    else:
                        date_ouv = p.get("date_ouverture")
                        if hasattr(date_ouv, "year"):
                            an_labels = [f"An {a} ({date_ouv.year + a - 1})" for a in annees_list]
                        else:
                            an_labels = [f"An {a}" for a in annees_list]
                        selected = st.multiselect("Annees", options=annees_list,
                            default=[a for a in default_actives if a in annees_list],
                            format_func=lambda x: an_labels[x-1] if x <= len(an_labels) else f"An {x}",
                            key="precompte_annees_pop")
                        p["precompte_annees_actives"] = selected

        if i not in to_delete:
            new_cf[new_label] = new_val

        # Graphique energie retractable sous la ligne
        if is_energie and st.session_state.get(f"_energie_expand_{key_prefix}", False):
            if "_projection_df" in st.session_state:
                _df_e = st.session_state["_projection_df"]
                _cv_en = p.get("cv_hebergement_par_nuitee", {}).get("Energie variable", 10.0)
                # Electricite fixe = Energie fixe (indirect) + Electricite brasserie + Electricite bar
                _elec_brass = p.get("cf_directs_brasserie", {}).get("Electricite", 0) if isinstance(p.get("cf_directs_brasserie"), dict) else 0
                _elec_bar = p.get("cf_directs_bar", {}).get("Electricite", 0) if isinstance(p.get("cf_directs_bar"), dict) else 0
                _fix_total_an = new_val + _elec_brass + _elec_bar
                _annual_e = _df_e.groupby("annee").agg({"nuitees": "sum"}).reset_index()
                _x_e = [str(int(a)) for a in _annual_e["annee"]]
                _var_e = _annual_e["nuitees"] * _cv_en
                _fix_e = [_fix_total_an] * len(_annual_e)
                _tot_e = _var_e + _fix_total_an
                _fig_e = go.Figure()
                _fig_e.add_trace(go.Bar(x=_x_e, y=[vv/1000 for vv in _fix_e],
                    name="Fixe", marker_color="#667eea",
                    hovertemplate="%{x}<br>%{y:,.1f} K\u20ac<extra></extra>"))
                _fig_e.add_trace(go.Bar(x=_x_e, y=[vv/1000 for vv in _var_e],
                    name="Variable", marker_color="#f5576c",
                    hovertemplate="%{x}<br>%{y:,.1f} K\u20ac<extra></extra>"))
                _fig_e.add_trace(go.Scatter(x=_x_e,
                    y=[vv/1000*1.02 for vv in _tot_e], mode="text",
                    text=[f"<b>{vv/1000:,.0f} K\u20ac</b>" for vv in _tot_e],
                    textposition="top center", textfont=dict(size=9, color="black"),
                    showlegend=False, hoverinfo="skip", cliponaxis=False))
                _max_e = max(_tot_e)/1000 if len(_tot_e) > 0 else 1
                _fig_e.update_layout(
                    title=dict(text="Cout total energie par an (fixe + variable)", font=dict(size=14)),
                    height=380, barmode="stack",
                    xaxis=dict(type="category", tickfont=dict(size=10)),
                    yaxis=dict(tickformat=",.0f", title="K\u20ac", range=[0, _max_e*1.15]),
                    showlegend=True, legend=dict(orientation="h", y=1.1, font=dict(size=12),
                        xanchor="center", x=0.5),
                    margin=dict(l=40, r=10, t=40, b=30))
                st.plotly_chart(_fig_e, use_container_width=True, config={"displayModeBar": False})
            else:
                st.caption("Le graphique sera disponible apres le premier calcul de projection.")

    if to_delete:
        # Reconstruire le dict depuis les items originaux (source de verite), en excluant les supprimes
        cleaned = {}
        for j, (orig_k, orig_v) in enumerate(items):
            if j not in to_delete:
                cleaned[orig_k] = orig_v
        st.session_state[state_key] = cleaned
        # Incrementer la version pour que tous les widgets aient de nouvelles cles
        st.session_state[version_key] = _ver + 1
        st.rerun()

    # Bouton ajouter
    _visu = st.session_state.get("_visu_mode", False)
    if st.button(f"\u2795 Ajouter une charge", key=f"{key_prefix}_add", disabled=_visu):
        new_cf["Nouvelle charge"] = 0
        st.session_state[state_key] = new_cf
        st.rerun()

    st.session_state[state_key] = new_cf
    return new_cf


def _fonds_propres_widget(p, key_prefix="fp"):
    """Widget editable pour les fonds propres avec investisseurs + camembert."""
    state_key = f"_fp_{key_prefix}"
    default_investisseurs = p.get("fonds_propres_investisseurs", [
        {"nom": "Investisseur principal", "montant": p.get("fonds_propres_initial", 0)},
    ])
    if state_key not in st.session_state:
        st.session_state[state_key] = list(default_investisseurs)

    current = st.session_state[state_key]
    new_list = []
    to_delete = []

    hcols = st.columns([3, 2, 0.5])
    hcols[0].markdown("**Investisseur**")
    hcols[1].markdown("**Montant (\u20ac)**")

    for i, inv in enumerate(current):
        c1, c2, c3 = st.columns([3, 2, 0.5])
        with c1:
            nom = st.text_input("n", inv["nom"], key=f"{key_prefix}_fpn_{i}", label_visibility="collapsed")
        with c2:
            montant = st.number_input("m", 0, 20_000_000, int(inv["montant"]), step=50_000,
                key=f"{key_prefix}_fpm_{i}", label_visibility="collapsed")
        with c3:
            _visu = st.session_state.get("_visu_mode", False)
            if st.button("\U0001F5D1", key=f"{key_prefix}_fpdel_{i}", help="Supprimer", disabled=_visu):
                to_delete.append(i)
        if i not in to_delete:
            new_list.append({"nom": nom, "montant": montant})

    if to_delete:
        st.session_state[state_key] = new_list
        st.rerun()

    _visu = st.session_state.get("_visu_mode", False)
    if st.button("\u2795 Ajouter un investisseur", key=f"{key_prefix}_fpadd", disabled=_visu):
        new_list.append({"nom": "Nouvel investisseur", "montant": 0})
        st.session_state[state_key] = new_list
        st.rerun()

    st.session_state[state_key] = new_list
    total_fp = sum(inv["montant"] for inv in new_list)
    st.info(f"**Total fonds propres : {total_fp:,.0f} \u20ac**")

    # Camembert repartition
    if len(new_list) > 0 and total_fp > 0:
        labels = [inv["nom"] for inv in new_list if inv["montant"] > 0]
        values = [inv["montant"] for inv in new_list if inv["montant"] > 0]
        colors = ["#4facfe", "#11998e", "#f5576c", "#ffcc00", "#667eea", "#f093fb", "#38ef7d", "#764ba2"]
        fig_fp = go.Figure(data=[go.Pie(
            labels=labels, values=values,
            marker=dict(colors=colors[:len(labels)]),
            textinfo="label+percent+value",
            textfont=dict(size=12),
            texttemplate="%{label}<br>%{value:,.0f} \u20ac<br>(%{percent})",
            hovertemplate="%{label}<br>%{value:,.0f} \u20ac<br>%{percent}<extra></extra>",
        )])
        fig_fp.update_layout(
            height=300, showlegend=False,
            margin=dict(l=10, r=10, t=10, b=10),
        )
        st.plotly_chart(fig_fp, use_container_width=True, config={"displayModeBar": False})

    # Synchroniser le total
    p["fonds_propres_initial"] = total_fp
    p["fonds_propres_investisseurs"] = new_list
    return new_list, total_fp


def tab_hypotheses(p):
    st.header("\U0001F4DD Hypotheses du plan financier")

    # Containers pour les mini-graphiques (remplis apres le calcul)
    charts = {}
    st.session_state["chart_containers"] = charts

    # 2 sous-onglets principaux
    ht_activite, ht5, ht_equilibre = st.tabs([
        "\U0001F4B0 Hypotheses liees a l'activite",
        "\U0001F3D7 Investissements & Financement",
        "\U0001F4CA Equilibre Moyens / Besoins",
    ])

    with ht_activite:
        ht0_gen, ht1, ht2, ht3, ht4, ht6, ht7 = st.tabs([
            "\u2699 Parametres generaux",
            "\U0001F4B0 1. Ventes",
            "\U0001F4C9 2. Frais variables",
            "\U0001F465 3. Frais fixes directs",
            "\U0001F3E2 4. Frais fixes indirects",
            "\U0001F4CA 5. Fiscalite",
            "\u23F1 6. Delais de paiement",
        ])

    # ════════════════════════════════════════════════════════════════════════
    # 0. PARAMETRES GENERAUX
    # ════════════════════════════════════════════════════════════════════════
    with ht0_gen:
        st.subheader("Parametres generaux")
        st.caption("Parametres specifiques a l'activite du Chateau d'Argenteau. "
                   "La date de debut et la duree de projection se definissent sur la page de choix du module.")
        c1, c2, c3 = st.columns(3)
        with c1:
            p["nb_chambres"] = st.number_input("Nombre de chambres", 10, 500, p["nb_chambres"], step=1)
        with c2:
            p["personnes_par_chambre"] = st.number_input("Personnes / chambre", 1.0, 5.0, float(p.get("personnes_par_chambre", 2)), step=0.1, format="%.1f")
        with c3:
            p["inflation_an"] = st.number_input("Inflation annuelle", 0.0, 0.10, p["inflation_an"], step=0.005, format="%.3f")
        st.info(f"L'inflation de **{p['inflation_an']:.1%}** s'applique a toutes les ventes "
                f"et aux charges en \u20ac (par nuitee, par entree, par soin, etc.). "
                f"Les charges en % du CA suivent automatiquement l'evolution du CA.")

    # ════════════════════════════════════════════════════════════════════════
    # 1. HYPOTHESES DE VENTE
    # ════════════════════════════════════════════════════════════════════════
    with ht1:
        st.subheader("Hypotheses de ventes")

        # ── Hebergement ─────────────────────────────────────────────────
        with st.expander("\U0001F6CF **Hebergement**", expanded=True):
            st.markdown("**Taux d'occupation par annee**")
            c1, c2, c3, c4 = st.columns(4)
            for i, (col, lbl) in enumerate(zip([c1, c2, c3, c4], ["An 1", "An 2", "An 3", "Croisiere"])):
                with col:
                    p["taux_occ"][i] = st.number_input(lbl, 0.0, 1.0, p["taux_occ"][i], step=0.01, format="%.2f", key=f"occ_{i}")

            st.markdown("---")
            st.markdown("**Segmentation clientele**")

            seg_names = list(p["segments"].keys())
            nouvelles_parts = {}
            for seg_name in seg_names:
                widget_key = f"seg_p_{seg_name}"
                if widget_key in st.session_state:
                    nouvelles_parts[seg_name] = st.session_state[widget_key]
                else:
                    nouvelles_parts[seg_name] = p["segments"][seg_name]["part"]

            total_actuel = sum(nouvelles_parts.values())
            marge = round(max(0.0, 1.0 - total_actuel), 2)

            cols_seg = st.columns(len(p["segments"]))
            for idx, (seg_name, seg_data) in enumerate(p["segments"].items()):
                val_courante = nouvelles_parts[seg_name]
                max_part = min(1.0, round(val_courante + marge, 2))
                val_defaut = min(seg_data["part"], max_part)
                with cols_seg[idx]:
                    seg_data["part"] = st.number_input(
                        f"{seg_name} (part)", 0.0, max_part,
                        val_defaut, step=0.05, format="%.2f",
                        key=f"seg_p_{seg_name}",
                    )
                    seg_data["prix"] = st.number_input(
                        f"ADR {seg_name} (\u20ac)", 50, 500,
                        int(seg_data["prix"]), step=5,
                        key=f"seg_pr_{seg_name}",
                    )

            total_parts = sum(sd["part"] for sd in p["segments"].values())
            if total_parts > 0:
                prix_moyen = sum(sd["part"] * sd["prix"] for sd in p["segments"].values()) / total_parts
            else:
                prix_moyen = 0

            if abs(total_parts - 1.0) > 0.001:
                reste_pct = (1.0 - total_parts) * 100
                st.warning(f"Total des parts : **{total_parts:.0%}** — il reste **{reste_pct:.0f}%** a attribuer")
            else:
                st.info(f"**Prix moyen pondéré (ADR) : {prix_moyen:,.0f} €** | Total parts : {total_parts:.0%}")

            st.markdown("---")
            st.markdown("**Hausse tarifaire annuelle**")
            c1, c2, c3 = st.columns(3)
            with c1:
                h2 = st.number_input("An 2", 0.0, 0.10, 0.025, step=0.005, format="%.3f", key="hausse_2")
            with c2:
                h3 = st.number_input("An 3", 0.0, 0.10, 0.025, step=0.005, format="%.3f", key="hausse_3")
            with c3:
                h4 = st.number_input("An 4+", 0.0, 0.10, 0.025, step=0.005, format="%.3f", key="hausse_4")
            p["hausse_prix_an"] = [0.0, h2, h3, h4]

            st.markdown("---")
            st.markdown("**Produits divers** (mini-bar, ...)")
            c1, c2 = st.columns(2)
            with c1:
                p["divers_prix_nuitee"] = st.number_input("Consommation par chambre (\u20ac)", 0.0, 50.0, float(p.get("divers_prix_nuitee", 3)), step=0.5, key="divers_prix")
            with c2:
                p["divers_taux"] = st.number_input("Proportion de chambres concernees", 0.0, 1.0, float(p.get("divers_taux", 1.0)), step=0.05, format="%.2f", key="divers_taux")

            st.markdown("---")
            # Saisonnalite avec graphique
            st.markdown("**Saisonnalite** (coefficient mensuel - moyenne = 1.00)")
            mois_noms = ["Janvier", "Fevrier", "Mars", "Avril", "Mai", "Juin",
                         "Juillet", "Aout", "Septembre", "Octobre", "Novembre", "Decembre"]

            # Lire les valeurs actuelles : session_state si deja modifie, sinon p
            _saison_vals = []
            for i in range(12):
                _sk = f"saison_{i}"
                if _sk in st.session_state:
                    _saison_vals.append(st.session_state[_sk])
                else:
                    _saison_vals.append(p["saisonnalite"][i])

            # Securite : si toutes les valeurs sont 0 ou moyenne aberrante, utiliser les defaults
            _moy_check = sum(_saison_vals) / 12 if len(_saison_vals) == 12 else 0
            if abs(_moy_check - 1.0) > 0.15 or all(v == 0 for v in _saison_vals):
                _saison_vals = list(p["saisonnalite"])
                # Si p aussi est invalide, fallback aux defaults
                _moy_p = sum(_saison_vals) / 12 if _saison_vals else 0
                if abs(_moy_p - 1.0) > 0.15 or all(v == 0 for v in _saison_vals):
                    _saison_vals = [0.626, 0.657, 0.86, 1.017, 1.126, 1.22,
                                    1.283, 1.33, 1.173, 1.064, 0.782, 0.86]

            if "saison_prev" not in st.session_state:
                st.session_state["saison_prev"] = list(_saison_vals)

            def _ajuster_saisonnalite():
                prev = st.session_state["saison_prev"]
                curr = [st.session_state[f"saison_{i}"] for i in range(12)]
                changed = next((i for i in range(12) if abs(curr[i] - prev[i]) > 0.001), None)
                if changed is None:
                    return
                target_others = 12.0 - curr[changed]
                sum_others = sum(curr[i] for i in range(12) if i != changed)
                if sum_others > 0:
                    factor = target_others / sum_others
                    for i in range(12):
                        if i != changed:
                            st.session_state[f"saison_{i}"] = round(curr[i] * factor, 2)
                st.session_state["saison_prev"] = [st.session_state[f"saison_{i}"] for i in range(12)]

            col_inputs, col_chart = st.columns([3, 1])
            with col_inputs:
                cols_saison = st.columns(12)
                for i in range(12):
                    with cols_saison[i]:
                        st.number_input(mois_noms[i][:3], min_value=0.0, max_value=3.0,
                                        value=float(_saison_vals[i]),
                                        step=0.05, format="%.2f", key=f"saison_{i}",
                                        on_change=_ajuster_saisonnalite)
            p["saisonnalite"] = [st.session_state[f"saison_{i}"] for i in range(12)]

            with col_chart:
                fig_saison = go.Figure()
                fig_saison.add_trace(go.Scatter(
                    x=[m[:3] for m in mois_noms],
                    y=p["saisonnalite"],
                    mode="lines+markers",
                    line=dict(color="#4facfe", width=2),
                    marker=dict(size=4),
                    fill="tozeroy",
                    fillcolor="rgba(79,172,254,0.15)",
                ))
                fig_saison.add_hline(y=1.0, line_dash="dash", line_color="grey", line_width=1)
                fig_saison.update_layout(
                    height=180, margin=dict(l=0, r=0, t=10, b=0),
                    xaxis=dict(tickfont=dict(size=9), showgrid=False),
                    yaxis=dict(tickfont=dict(size=9), showgrid=True, gridcolor="rgba(0,0,0,0.05)"),
                    plot_bgcolor="rgba(0,0,0,0)",
                    paper_bgcolor="rgba(0,0,0,0)",
                    showlegend=False,
                )
                st.plotly_chart(fig_saison, use_container_width=True, config={"displayModeBar": False})

            st.caption(f"Moyenne : {sum(p['saisonnalite'])/12:.2f}")

            charts["hebergement"] = st.container()

        # ── Brasserie ───────────────────────────────────────────────────
        with st.expander("\U0001F37D **Brasserie / Restaurant**", expanded=False):
            # ── Petit-dejeuner (lie aux nuitees) ──
            st.markdown("**Petit-dejeuner** *(lie au nombre de nuitees)*")
            c1, c2 = st.columns(2)
            with c1:
                p["petit_dej_prix"] = st.number_input("Prix PDJ (\u20ac)", 5.0, 80.0, p["petit_dej_prix"], step=2.5)
            with c2:
                p["petit_dej_taux"] = st.number_input("Taux prise PDJ", 0.0, 1.0, p["petit_dej_taux"], step=0.05)

            st.markdown("---")

            # ── Services midi / soir (lies a la capacite et taux d'occupation) ──
            st.markdown("**Services midi & soir** *(lies a la capacite et au taux d'occupation)*")
            col_cap, _ = st.columns([1, 2])
            with col_cap:
                p["nb_couverts_brasserie"] = st.number_input("Couverts (capacite salle)", 10, 300, int(p.get("nb_couverts_brasserie", 80)), step=5, key="brass_couv")

            c1, c2, c3 = st.columns(3)
            with c1:
                st.markdown("**Service midi**")
                p["brasserie_jours_diner"] = st.number_input("Jours / semaine", 0, 7, int(p.get("brasserie_jours_diner", p.get("brasserie_ouvert_midi", 4))), key="brass_jm")
                p["brasserie_services_diner"] = st.number_input("Nb services / jour", 0.0, 3.0, float(p.get("brasserie_services_diner", 1.5)), step=0.5, format="%.1f", key="brass_sm")
                p["brasserie_prix_diner"] = st.number_input("Prix moyen midi (\u20ac)", 10, 150, int(p.get("brasserie_prix_diner", p.get("brasserie_prix_midi", 45))), step=5, key="brass_pm")
            with c2:
                st.markdown("**Service soir**")
                p["brasserie_jours_souper"] = st.number_input("Jours / semaine", 0, 7, int(p.get("brasserie_jours_souper", p.get("brasserie_ouvert_soir", 4))), key="brass_js")
                p["brasserie_services_souper"] = st.number_input("Nb services / jour", 0.0, 3.0, float(p.get("brasserie_services_souper", 1.0)), step=0.5, format="%.1f", key="brass_ss")
                p["brasserie_prix_souper"] = st.number_input("Prix moyen soir (\u20ac)", 10, 200, int(p.get("brasserie_prix_souper", p.get("brasserie_prix_soir", 75))), step=5, key="brass_ps")
            with c3:
                st.markdown("**Mix nourriture / boissons**")
                p["brasserie_part_nourriture"] = st.number_input("Part nourriture", 0.0, 1.0, float(p.get("brasserie_part_nourriture", 0.60)), step=0.05, format="%.2f", key="brass_food")
                p["brasserie_part_boissons"] = round(1.0 - p["brasserie_part_nourriture"], 2)
                st.caption(f"Part boissons : **{p['brasserie_part_boissons']:.0%}**")

            st.markdown("**Taux d'occupation brasserie par annee**")
            cb1, cb2, cb3, cb4 = st.columns(4)
            taux_brass = p.get("taux_occ_brasserie", [0.35, 0.45, 0.525, 0.60])
            for i, (col, lbl) in enumerate(zip([cb1, cb2, cb3, cb4], ["An 1", "An 2", "An 3", "Croisiere"])):
                with col:
                    taux_brass[i] = st.number_input(lbl, 0.0, 1.0, float(taux_brass[i]), step=0.01, format="%.2f", key=f"occ_brass_{i}")
            p["taux_occ_brasserie"] = taux_brass

            charts["brasserie"] = st.container()

        # ── Bar ──────────────────────────────────────────────────────────
        with st.expander("\U0001F378 **Bar**", expanded=False):
            c1, c2, c3 = st.columns(3)
            with c1:
                st.markdown("**Clients hotel**")
                p["bar_taux_clients_hotel"] = st.number_input("% clients hotel qui vont au bar", 0.0, 1.0, p["bar_taux_clients_hotel"], step=0.05, format="%.2f", key="bar_taux")
                p["bar_conso_moyenne"] = st.number_input("Depense moy. / visite (\u20ac)", 0.0, 100.0, float(p["bar_conso_moyenne"]), step=1.0, key="bar_conso")
            with c2:
                st.markdown("**Clients externes**")
                p["bar_clients_ext_jour"] = st.number_input("Clients ext. / jour", 0, 100, p["bar_clients_ext_jour"], key="bar_ext_nb")
                p["bar_conso_ext_moyenne"] = st.number_input("Depense moy. ext. (\u20ac)", 0.0, 100.0, float(p["bar_conso_ext_moyenne"]), step=1.0, key="bar_ext_conso")
            with c3:
                st.markdown("**Ouverture**")
                p["bar_jours_ouvert_semaine"] = st.number_input("Jours ouvert / semaine", 0, 7, p["bar_jours_ouvert_semaine"], key="bar_jours")

            charts["bar"] = st.container()

        # ── Spa ─────────────────────────────────────────────────────────
        with st.expander("\U0001F9D6 **Spa / Bien-etre**", expanded=False):
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**Clients hotel**")
                p["spa_entree_hotel_prix"] = st.number_input("Prix entree (\u20ac)", 0, 200, int(p.get("spa_entree_hotel_prix", 0)), step=5, key="spa_entree_h")
                p["spa_entree_hotel_taux"] = st.number_input("Taux acces spa (par nuitee)", 0.0, 1.0, float(p.get("spa_entree_hotel_taux", 0.20)), step=0.05, format="%.2f", key="spa_taux_entree_h")
                p["spa_soin_hotel_prix"] = st.number_input("Prix moyen soin (\u20ac)", 0, 500, int(p.get("spa_soin_hotel_prix", 120)), step=5, key="spa_soin_h")
                p["spa_soin_hotel_taux"] = st.number_input("Taux soins (par nuitee)", 0.0, 1.0, float(p.get("spa_soin_hotel_taux", 0.10)), step=0.05, format="%.2f", key="spa_taux_soin_h")
            with c2:
                st.markdown("**Clients externes**")
                p["spa_entree_ext_prix"] = st.number_input("Prix entree ext. (\u20ac)", 0, 200, int(p.get("spa_entree_ext_prix", 55)), step=5, key="spa_entree_ext")
                p["spa_entree_ext_nb_mois"] = st.number_input("Entrees ext. / mois", 0, 500, int(p.get("spa_entree_ext_nb_mois", 25)), step=5, key="spa_nb_entree_ext")
                p["spa_soin_ext_prix"] = st.number_input("Prix moyen soin ext. (\u20ac)", 0, 500, int(p.get("spa_soin_ext_prix", 150)), step=5, key="spa_soin_ext")
                p["spa_soin_ext_nb_mois"] = st.number_input("Soins ext. / mois", 0, 500, int(p.get("spa_soin_ext_nb_mois", 15)), step=5, key="spa_nb_soin_ext")

            charts["spa"] = st.container()

        # ── Salles & Evenements ───────────────────────────────────────
        with st.expander("\U0001F3F0 **Salles & Evenements**", expanded=False):
            c1, c2, c3 = st.columns(3)
            with c1:
                st.markdown('<div class="section-header">\U0001F4BC Seminaires</div>', unsafe_allow_html=True)
                p["seminaire_nb_an"] = st.number_input("Nb seminaires / an", 0, 200, p["seminaire_nb_an"], key="sem_nb")
                p["seminaire_prix_location"] = st.number_input("Location salle (\u20ac)", 0, 50000, p["seminaire_prix_location"], step=50, key="sem_loc")
            with c2:
                st.markdown('<div class="section-header">\U0001F492 Mariages / Evenements</div>', unsafe_allow_html=True)
                p["mariage_nb_an"] = st.number_input("Nb mariages / an", 0, 50, p["mariage_nb_an"], key="mar_nb")
                p["mariage_prix_location"] = st.number_input("Location salle (\u20ac)", 0, 50000, p["mariage_prix_location"], step=100, key="mar_loc")
                p["mariage_nb_convives_moy"] = st.number_input("Convives moyens / mariage", 10, 500, int(p.get("mariage_nb_convives_moy", 120)), key="mar_conv")
                p["mariage_catering_prix_convive"] = st.number_input("Prix catering / convive (\u20ac)", 0, 500, int(p.get("mariage_catering_prix_convive", 95)), step=5, key="mar_cat_prix")
                p["mariage_commission_catering_pct"] = st.number_input("Commission catering (%)", 0.0, 50.0, float(p.get("mariage_commission_catering_pct", 0.10)) * 100, step=1.0, format="%.1f", key="mar_cat_pct") / 100
            with c3:
                st.markdown('<div class="section-header">\U0001F3F0 Salles du chateau</div>', unsafe_allow_html=True)
                p["salles_chateau_nb_an"] = st.number_input("Nb locations / an", 0, 200, int(p.get("salles_chateau_nb_an", 30)), key="chateau_nb")
                p["salles_chateau_prix"] = st.number_input("Prix moyen location (\u20ac)", 0, 50000, int(p.get("salles_chateau_prix", 1500)), step=100, key="chateau_prix")

            charts["seminaires"] = st.container()
            charts["mariages"] = st.container()
            charts["salles_chateau"] = st.container()

        # ── Loyer restaurant gastronomique ──
        with st.expander("\U0001F37D **Location restaurant gastronomique**", expanded=False):
            col_lr, _ = st.columns([1, 2])
            with col_lr:
                p["loyer_restaurant_mensuel"] = st.number_input("Loyer mensuel (\u20ac)", 0, 50000, int(p.get("loyer_restaurant_mensuel", 5800)), step=100, key="loyer_resto")
            st.caption("Revenu fixe mensuel, pas de charges variables associees.")
            charts["loyer_restaurant"] = st.container()

        charts["total"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 2. FRAIS VARIABLES (lies a un service)
    # ════════════════════════════════════════════════════════════════════════
    with ht2:
        st.subheader("Frais variables par service")
        st.caption("Couts proportionnels au chiffre d'affaires ou a l'activite de chaque service")

        with st.expander("\U0001F6CF **Hebergement**", expanded=True):
            # Charges en € par nuitée
            st.markdown("**Charges en \u20ac / nuitee vendue**")
            _default_cv_nuitee = {
                "Linge / Blanchisserie": 5.5,
                "Produits accueil": 3.5,
                "Produits entretien": 2.5,
                "Energie variable": 10.0,
                "Fournitures chambres": 2.5,
            }
            cv_nuitee = p.get("cv_hebergement_par_nuitee", _default_cv_nuitee)
            # Retirer l'ancienne clé commission CB si présente
            cv_nuitee = {k: v for k, v in cv_nuitee.items() if "ommission" not in k}
            new_cv_nuitee = {}
            cols_n = st.columns(len(cv_nuitee))
            for idx, (k, v) in enumerate(cv_nuitee.items()):
                with cols_n[idx]:
                    new_cv_nuitee[k] = st.number_input(k, 0.0, 50.0, float(v), step=0.5, format="%.2f", key=f"cvhn_{k}")
            p["cv_hebergement_par_nuitee"] = new_cv_nuitee

            st.markdown("---")
            # Commission carte de crédit
            st.markdown("**Commission cartes de credit**")
            c1, c2 = st.columns(2)
            with c1:
                p["cv_commission_cb_nuitee"] = st.number_input("Commission (\u20ac / nuitee)", 0.0, 20.0, float(p.get("cv_commission_cb_nuitee", 1.5)), step=0.1, format="%.2f", key="cv_cb_nuitee")
            with c2:
                p["cv_commission_cb_pct_chambres"] = st.number_input("% chambres payees par CB", 0.0, 1.0, float(p.get("cv_commission_cb_pct_chambres", 0.80)), step=0.05, format="%.2f", key="cv_cb_pct")

            st.markdown("---")
            st.markdown("**Consommables produits divers** (lies aux ventes mini-bar, ...)")
            col_div, _ = st.columns([1, 2])
            with col_div:
                p["cv_divers_consommable"] = st.number_input("Cout consommable divers (\u20ac / nuitee)", 0.0, 20.0, float(p.get("cv_divers_consommable", 1.0)), step=0.5, format="%.2f", key="cv_divers_conso")

            st.markdown("---")
            # Commissions OTA
            st.markdown("**Commissions OTA (Booking, Airbnb, Expedia...)**")
            col_taux, _ = st.columns([1, 2])
            with col_taux:
                p["cv_commission_ota_pct"] = st.number_input("Taux commission OTA", 0.0, 0.50, float(p.get("cv_commission_ota_pct", 0.17)), step=0.01, format="%.2f", key="cv_ota_pct")

            # Part OTA par segment
            st.markdown("**Part des reservations via OTA par segment**")
            st.caption("Proportion de clients passant par Booking, Airbnb, Expedia... par type de clientele. "
                       "Impacte le montant des commissions OTA.")
            seg_names = list(p["segments"].keys())
            ota_data = p.get("segments_part_ota", {})
            default_ota = {"Loisirs": 0.75, "Affaires": 0.55, "Groupes": 0.20, "MICE": 0.0, "Evenementiel": 0.0}
            cols_ota = st.columns(len(seg_names))
            new_ota = {}
            for idx, seg_name in enumerate(seg_names):
                with cols_ota[idx]:
                    part_seg = p["segments"][seg_name]["part"]
                    val_ota = float(ota_data.get(seg_name, default_ota.get(seg_name, 0.0)))
                    new_ota[seg_name] = st.number_input(
                        f"{seg_name} ({part_seg:.0%} du CA)",
                        0.0, 1.0, val_ota, step=0.05, format="%.2f",
                        key=f"cv_ota_{seg_name}",
                    )
            p["segments_part_ota"] = new_ota
            # Afficher le poids OTA pondéré
            poids_ota = sum(p["segments"][s]["part"] * new_ota[s] for s in seg_names)
            st.info(f"**Poids OTA pondere : {poids_ota:.2%}** | Commission effective sur CA total : {poids_ota * p.get('cv_commission_ota_pct', 0.17):.2%}")

            st.markdown("---")
            # Franchise en dernier
            st.markdown("**Cout franchise**")
            _modes_options = {"pct": "% du CA hebergement", "nuitee": "\u20ac par nuitee", "forfait": "Forfait par mois"}
            _modes_default = p.get("cv_franchise_modes", ["pct"])
            _modes_sel = st.multiselect("Modalites de la franchise", options=list(_modes_options.keys()),
                                        default=[m for m in _modes_default if m in _modes_options],
                                        format_func=lambda x: _modes_options[x], key="franchise_modes")
            p["cv_franchise_modes"] = _modes_sel

            _fr_cols = st.columns(max(len(_modes_sel), 1))
            _col_idx = 0
            if "pct" in _modes_sel:
                with _fr_cols[_col_idx]:
                    p["cv_franchise_pct"] = st.number_input("Taux franchise (% CA)", 0.0, 0.20, float(p.get("cv_franchise_pct", 0.04)), step=0.01, format="%.2f", key="cv_franchise")
                _col_idx += 1
            if "nuitee" in _modes_sel:
                with _fr_cols[_col_idx]:
                    p["cv_franchise_par_nuitee"] = st.number_input("Franchise (\u20ac/nuitee)", 0.0, 50.0, float(p.get("cv_franchise_par_nuitee", 0)), step=0.5, format="%.2f", key="cv_franchise_nuitee")
                _col_idx += 1
            if "forfait" in _modes_sel:
                with _fr_cols[_col_idx]:
                    p["cv_franchise_forfait_mois"] = st.number_input("Forfait franchise (\u20ac/mois)", 0.0, 50000.0, float(p.get("cv_franchise_forfait_mois", 0)), step=100.0, format="%.0f", key="cv_franchise_forfait")

            charts["cv_detail_heberg"] = st.container()

        st.markdown("---")

        # ── CV Brasserie ──
        with st.expander("\U0001F37D **Brasserie**", expanded=False):
            c1, c2, c3 = st.columns(3)
            with c1:
                p["cv_brasserie_pct"] = st.number_input("Food cost midi & soir (% CA)", 0.0, 1.0, p["cv_brasserie_pct"], step=0.01, format="%.2f", key="cv_brass")
            with c2:
                p["cv_pdj_pct"] = st.number_input("Cout petit-dejeuner (% CA PDJ)", 0.0, 1.0, float(p.get("cv_pdj_pct", p.get("cv_brasserie_pct", 0.35))), step=0.01, format="%.2f", key="cv_pdj")
            with c3:
                p["cv_eau_brasserie_par_client"] = st.number_input("Eau (\u20ac/couvert)", 0.0, 10.0, float(p.get("cv_eau_brasserie_par_client", 0)), step=0.1, format="%.2f", key="cv_eau_brass")
            charts["cv_brasserie"] = st.container()

        # ── CV Bar ──
        with st.expander("\U0001F378 **Bar**", expanded=False):
            c1, c2, c3 = st.columns(3)
            with c1:
                p["cv_bar_pct"] = st.number_input("Cout boissons (% CA)", 0.0, 1.0, p["cv_bar_pct"], step=0.01, format="%.2f", key="cv_bar")
            with c2:
                p["cv_bar_consommable_unite"] = st.number_input("Consommables (serviettes, bougies...) \u20ac/conso", 0.0, 5.0, float(p.get("cv_bar_consommable_unite", 0.20)), step=0.1, format="%.2f", key="cv_bar_conso")
            with c3:
                p["cv_eau_bar_par_client"] = st.number_input("Eau (\u20ac/client)", 0.0, 10.0, float(p.get("cv_eau_bar_par_client", 0)), step=0.1, format="%.2f", key="cv_eau_bar")
            charts["cv_bar"] = st.container()

        # ── CV Spa ──
        with st.expander("\U0001F9D6 **Spa**", expanded=False):
            st.markdown("**Couts par soin** (hotel + externes)")
            c1, c2 = st.columns(2)
            with c1:
                p["cv_spa_soin_cout"] = st.number_input("Prestation soin (\u20ac/soin)", 0.0, 200.0, float(p.get("cv_spa_soin_cout", 50)), step=5.0, format="%.1f", key="cv_spa_soin")
            with c2:
                p["cv_spa_produits_soin"] = st.number_input("Produits soins/cosmetiques (\u20ac/soin)", 0.0, 50.0, float(p.get("cv_spa_produits_soin", 5)), step=1.0, format="%.1f", key="cv_spa_prod")
            st.markdown("**Couts par entree spa** (hotel + externes)")
            c1, c2, c3 = st.columns(3)
            with c1:
                p["cv_spa_consommable_entree"] = st.number_input("Consommables (serviettes, huiles) \u20ac/entree", 0.0, 20.0, float(p.get("cv_spa_consommable_entree", 1.5)), step=0.5, format="%.1f", key="cv_spa_conso")
            with c2:
                p["cv_spa_energie_entree"] = st.number_input("Energie (\u20ac/entree)", 0.0, 20.0, float(p.get("cv_spa_energie_entree", 3.0)), step=0.5, format="%.1f", key="cv_spa_energ")
            with c3:
                p["cv_spa_piscine_entree"] = st.number_input("Produits piscine (\u20ac/entree)", 0.0, 10.0, float(p.get("cv_spa_piscine_entree", 0.5)), step=0.1, format="%.1f", key="cv_spa_pisc")
            charts["cv_spa"] = st.container()

        # ── CV Seminaires ──
        with st.expander("\U0001F4BC **Seminaires**", expanded=False):
            p["seminaire_nb_participants_moy"] = st.number_input("Participants moyens / seminaire", 1, 500, int(p.get("seminaire_nb_participants_moy", 25)), key="sem_part_cv")
            st.markdown("**Couts par participant**")
            c1, c2 = st.columns(2)
            with c1:
                p["cv_seminaire_pause_participant"] = st.number_input("Cafe & collation (\u20ac/participant)", 0.0, 50.0, float(p.get("cv_seminaire_pause_participant", 8)), step=1.0, format="%.1f", key="cv_sem_pause")
            with c2:
                p["cv_seminaire_materiel_participant"] = st.number_input("Consommable (\u20ac/participant)", 0.0, 20.0, float(p.get("cv_seminaire_materiel_participant", 3)), step=0.5, format="%.1f", key="cv_sem_mat")
            st.markdown("**Couts par seminaire**")
            c1, c2, c3 = st.columns(3)
            with c1:
                p["cv_seminaire_equipement"] = st.number_input("Equipement AV (\u20ac/seminaire)", 0.0, 200.0, float(p.get("cv_seminaire_equipement", 15)), step=5.0, format="%.1f", key="cv_sem_equip")
            with c2:
                p["cv_seminaire_energie"] = st.number_input("Energie (\u20ac/seminaire)", 0.0, 500.0, float(p.get("cv_seminaire_energie", 75)), step=5.0, format="%.1f", key="cv_sem_energ")
            with c3:
                p["cv_seminaire_nettoyage"] = st.number_input("Nettoyage (\u20ac/seminaire)", 0.0, 2000.0, float(p.get("cv_seminaire_nettoyage", 500)), step=50.0, format="%.0f", key="cv_sem_nett")
            charts["cv_seminaires"] = st.container()

        # ── CV Mariages ──
        with st.expander("\U0001F492 **Mariages**", expanded=False):
            st.info(f"Convives moyens / mariage : **{int(p.get('mariage_nb_convives_moy', 120))}** (modifiable dans Hypotheses de vente)")
            st.markdown("**Couts par mariage**")
            c1, c2 = st.columns(2)
            with c1:
                p["cv_mariage_energie"] = st.number_input("Energie (\u20ac/mariage)", 0.0, 1000.0, float(p.get("cv_mariage_energie", 150)), step=10.0, format="%.0f", key="cv_mar_energ")
            with c2:
                p["cv_mariage_nettoyage"] = st.number_input("Nettoyage (\u20ac/mariage)", 0.0, 5000.0, float(p.get("cv_mariage_nettoyage", 1000)), step=50.0, format="%.0f", key="cv_mar_nett")
            charts["cv_mariages"] = st.container()

        # ── CV Salles chateau ──
        with st.expander("\U0001F3F0 **Salles du chateau**", expanded=False):
            st.markdown("**Couts par location**")
            c1, c2 = st.columns(2)
            with c1:
                p["cv_salles_chateau_energie"] = st.number_input("Energie (\u20ac/location)", 0.0, 500.0, float(p.get("cv_salles_chateau_energie", 100)), step=10.0, format="%.0f", key="cv_chat_energ")
            with c2:
                p["cv_salles_chateau_nettoyage"] = st.number_input("Nettoyage (\u20ac/location)", 0.0, 2000.0, float(p.get("cv_salles_chateau_nettoyage", 500)), step=50.0, format="%.0f", key="cv_chat_nett")
            charts["cv_salles_chateau"] = st.container()

        st.markdown("---")
        charts["cv_total"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 3. FRAIS FIXES DIRECTS (lies a un service)
    # ════════════════════════════════════════════════════════════════════════
    with ht3:
        st.subheader("Frais fixes directs par service")
        st.caption("Personnel et charges fixes rattaches a un departement / service specifique")

        p["charges_patronales_pct"] = st.number_input("Taux charges patronales", 0.0, 0.60, p["charges_patronales_pct"], step=0.01, format="%.2f")
        cp = p["charges_patronales_pct"]

        # ── Hébergement ──
        with st.expander("\U0001F6CF **Hebergement**", expanded=False):
            st.markdown("**Personnel**")
            _nb_ch = p.get("nb_chambres", 70)
            hcols = st.columns([3, 1.5, 1.0, 2, 1.5, 1.5])
            hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**ETP/10 ch.**"); hcols[3].markdown("**Brut/an**"); hcols[4].markdown("**Cout/an/ETP**"); hcols[5].markdown("**Cout total/an**")
            p["personnel_hebergement"] = _personnel_table(p, p["personnel_hebergement"], "ph", cp, nb_chambres=_nb_ch, allow_delete=False)
            st.markdown("**Autres frais fixes directs**")
            p["cf_directs_hebergement"] = _cf_grid(p, p["cf_directs_hebergement"], "cfdh", 3)

        # ── Brasserie ──
        with st.expander("\U0001F37D **Brasserie**", expanded=False):
            st.markdown("**Personnel**")
            _nb_couv = p.get("nb_couverts_brasserie", 80)
            hcols = st.columns([3, 1.5, 1.0, 2, 1.5, 1.5])
            hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**ETP/10 couv.**"); hcols[3].markdown("**Brut/an**"); hcols[4].markdown("**Cout/an/ETP**"); hcols[5].markdown("**Cout total/an**")
            p["personnel_brasserie"] = _personnel_table(p, p["personnel_brasserie"], "pb", cp, ratio_base=_nb_couv, allow_delete=False)
            st.markdown("**Autres frais fixes directs**")
            p["cf_directs_brasserie"] = _cf_grid(p, p["cf_directs_brasserie"], "cfdb", 2, protected_keys={"Electricite"})

        # ── Bar ──
        with st.expander("\U0001F378 **Bar**", expanded=False):
            st.markdown("**Personnel**")
            hcols = st.columns([3, 1.5, 2, 1.5, 1.5])
            hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**Brut/an**"); hcols[3].markdown("**Cout/an/ETP**"); hcols[4].markdown("**Cout total/an**")
            p["personnel_bar"] = _personnel_table(p, p.get("personnel_bar", []), "pbar", cp, allow_delete=False)
            st.markdown("**Autres frais fixes directs**")
            if "cf_directs_bar" not in p or not isinstance(p["cf_directs_bar"], dict):
                p["cf_directs_bar"] = {"Electricite": 0}
            p["cf_directs_bar"] = _cf_grid(p, p["cf_directs_bar"], "cfdbar", 2, protected_keys={"Electricite"})

        # ── Spa ──
        with st.expander("\U0001F9D6 **Spa**", expanded=False):
            st.markdown("**Personnel**")
            hcols = st.columns([3, 1.5, 2, 1.5, 1.5])
            hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**Brut/an**"); hcols[3].markdown("**Cout/an/ETP**"); hcols[4].markdown("**Cout total/an**")
            p["personnel_spa"] = _personnel_table(p, p["personnel_spa"], "ps", cp, allow_delete=False)
            st.markdown("**Autres frais fixes directs**")
            p["cf_directs_spa"] = _cf_grid(p, p["cf_directs_spa"], "cfds", 2)

        # ── Evenements ──
        with st.expander("\U0001F4BC **Evenements (seminaires, mariages)**", expanded=False):
            st.markdown("**Personnel**")
            hcols = st.columns([3, 1.5, 2, 1.5, 1.5])
            hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**Brut/an**"); hcols[3].markdown("**Cout/an/ETP**"); hcols[4].markdown("**Cout total/an**")
            p["personnel_evenements"] = _personnel_table(p, p.get("personnel_evenements", []), "pe", cp, allow_delete=False)
            st.markdown("**Autres frais fixes directs**")
            p["cf_directs_evenements"] = _cf_grid(p, p["cf_directs_evenements"], "cfde", 2)

        # Totaux
        st.markdown("---")
        all_pers = p["personnel_hebergement"] + p["personnel_brasserie"] + p.get("personnel_bar", []) + p["personnel_spa"] + p.get("personnel_evenements", [])
        total_etp_dir = sum(pe["etp"] for pe in all_pers)
        total_masse_dir = sum(pe["cout_brut"] * (1 + cp) * pe["etp"] for pe in all_pers)
        _cf_bar_dict = p.get("cf_directs_bar", {}) if isinstance(p.get("cf_directs_bar"), dict) else {}
        total_cf_dir = (sum(p["cf_directs_hebergement"].values()) + sum(p["cf_directs_brasserie"].values()) +
                        sum(_cf_bar_dict.values()) +
                        sum(p["cf_directs_spa"].values()) + sum(p["cf_directs_evenements"].values()))
        st.info(f"**Total frais fixes directs : {total_etp_dir:.1f} ETP | "
                f"Personnel : {total_masse_dir:,.0f} \u20ac/an | "
                f"Autres : {total_cf_dir:,.0f} \u20ac/an | "
                f"TOTAL : {total_masse_dir + total_cf_dir:,.0f} \u20ac/an**")

        charts["cf_directs"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 4. FRAIS FIXES INDIRECTS (non lies a un service)
    # ════════════════════════════════════════════════════════════════════════
    with ht4:
        st.subheader("Frais fixes indirects")
        st.caption("Charges non rattachables a un service en particulier")

        # Personnel indirect
        st.markdown('<div class="section-header">\U0001F465 Personnel indirect</div>', unsafe_allow_html=True)
        hcols = st.columns([3, 1.5, 2, 1.5, 1.5])
        hcols[0].markdown("**Poste**"); hcols[1].markdown("**ETP**"); hcols[2].markdown("**Brut/an**"); hcols[3].markdown("**Cout/an/ETP**"); hcols[4].markdown("**Cout total/an**")
        cp = p["charges_patronales_pct"]
        p["personnel_indirect"] = _personnel_table(p, p["personnel_indirect"], "pi", cp)

        total_etp_ind = sum(pe["etp"] for pe in p["personnel_indirect"])
        total_masse_ind = sum(pe["cout_brut"] * (1 + cp) * pe["etp"] for pe in p["personnel_indirect"])
        st.caption(f"Sous-total : {total_etp_ind:.1f} ETP | {total_masse_ind:,.0f} \u20ac/an")

        st.markdown("---")

        # Charges fixes indirectes
        st.markdown('<div class="section-header">\U0001F3E2 Charges fixes indirectes</div>', unsafe_allow_html=True)

        st.markdown("**Charges annuelles**")
        p["charges_fixes_indirectes"] = _cf_grid(p, p["charges_fixes_indirectes"], "cfi", 4, charts=charts)

        st.markdown("---")
        st.markdown("**Loyer (mensuel)** — *defini par l'Immobiliere Rocher*")
        col_loyer, _ = st.columns([1, 2])
        with col_loyer:
            st.number_input("Loyer mensuel (\u20ac)", 0, 500_000, p["loyer_mensuel"], step=500, disabled=True, key="loyer_mens_disabled")
        st.caption(f"Soit **{p['loyer_mensuel'] * 12:,.0f} \u20ac / an** — Modifiable uniquement dans le module Immobiliere Rocher.")

        total_cfi = sum(p["charges_fixes_indirectes"].values()) + p["loyer_mensuel"] * 12
        st.info(f"**Total frais fixes indirects : Personnel {total_masse_ind:,.0f} \u20ac + "
                f"Charges {total_cfi:,.0f} \u20ac = "
                f"**{total_masse_ind + total_cfi:,.0f} \u20ac/an****")

        charts["cf_indirects"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 5. INVESTISSEMENTS & FINANCEMENT
    # ════════════════════════════════════════════════════════════════════════
    with ht5:
        st.subheader("Investissements & Financement")

        st.markdown('<div class="section-header">\U0001F3D7 Investissements & Amortissements</div>', unsafe_allow_html=True)

        # Calcul ETP total et couverts pour les investissements avec multiplicateur
        _all_pers_inv = (p.get("personnel_hebergement", []) + p.get("personnel_brasserie", []) +
                         p.get("personnel_bar", []) + p.get("personnel_spa", []) +
                         p.get("personnel_evenements", []) + p.get("personnel_indirect", []))
        _total_etp = sum(pe["etp"] for pe in _all_pers_inv)
        _nb_couverts = p.get("nb_couverts_brasserie", 80)

        # Definition des investissements supplementaires (avec multiplicateurs)
        inv_supp_cats = [
            ("Bornes electriques & Panneaux photovolt.", "inv_bornes", None, None),
            ("Uniformes travailleurs", "inv_uniformes", "etp", _total_etp),
            ("Equipement brasserie", "inv_equip_brass", None, None),
            ("Vaisselle & couverts brasserie", "inv_vaisselle", "couverts", _nb_couverts),
            ("Equipement salles", "inv_equip_salles", None, None),
        ]
        inv_supp_saved = p.get("investissements_supplementaires", {})

        # En-tetes
        hcols = st.columns([3, 2, 1.5, 1.5, 2])
        hcols[0].markdown("**Categorie**"); hcols[1].markdown("**Montant (\u20ac)**")
        hcols[2].markdown("**Multiplicateur**"); hcols[3].markdown("**Duree amort. (ans)**"); hcols[4].markdown("**Amort. annuel**")

        # ── Investissements classiques (categories fixes) ──
        _fixed_cats = {"Amenagements interieurs", "Mobilier & Equipements", "Branding/Communication", "Imprevus"}
        _supp_cats_names = {c[0] for c in inv_supp_cats}  # noms des investissements supplementaires
        _all_reserved = _fixed_cats | _supp_cats_names
        new_inv = []
        # Separer les investissements fixes des lignes libres (affichees en fin)
        _inv_fixes = [inv for inv in p["investissements"] if inv["categorie"] in _fixed_cats]
        _inv_libres = [inv for inv in p["investissements"] if inv["categorie"] not in _all_reserved]
        # S'assurer qu'il y a exactement 1 ligne libre + 1 "Autre"
        _inv_libres_sans_autre = [inv for inv in _inv_libres if inv["categorie"] != "Autre"]
        _inv_autre = [inv for inv in _inv_libres if inv["categorie"] == "Autre"]
        if len(_inv_libres_sans_autre) == 0:
            _inv_libres_sans_autre = [{"categorie": "", "montant": 0, "duree_amort": 0}]
        if len(_inv_libres_sans_autre) > 1:
            _inv_libres_sans_autre = [_inv_libres_sans_autre[0]]
        if len(_inv_autre) == 0:
            _inv_autre = [{"categorie": "Autre", "montant": 0, "duree_amort": 0}]
        if len(_inv_autre) > 1:
            _inv_autre = [_inv_autre[0]]
        _inv_libres = _inv_libres_sans_autre + _inv_autre

        # Afficher les categories fixes (non modifiables)
        for i, inv in enumerate(_inv_fixes):
            c1, c2, c3, c4, c5 = st.columns([3, 2, 1.5, 1.5, 2])
            with c1:
                st.markdown(f"**{inv['categorie']}**")
            with c2:
                mont = st.number_input("mt", 0, 20_000_000, inv["montant"], step=50_000, key=f"inv_m_{i}", label_visibility="collapsed")
            with c3:
                st.markdown("*-*")
            with c4:
                dur = st.number_input("dur", 0, 30, inv["duree_amort"], key=f"inv_d_{i}", label_visibility="collapsed")
            with c5:
                amort_an = mont / dur if dur > 0 else 0
                st.markdown(f'<div style="text-align:center"><b>{amort_an:,.0f} \u20ac</b></div>' if dur > 0 else '<div style="text-align:center"><i>Non amorti</i></div>', unsafe_allow_html=True)
            new_inv.append({"categorie": inv["categorie"], "montant": mont, "duree_amort": dur})

        # ── Investissements supplementaires (avec multiplicateurs) ──
        inv_supp_new = {}
        for cat_name, key, mult_type, mult_val in inv_supp_cats:
            saved = inv_supp_saved.get(key, {})
            c1, c2, c3, c4, c5 = st.columns([3, 2, 1.5, 1.5, 2])
            with c1:
                st.markdown(f"**{cat_name}**")
            with c2:
                if mult_type is not None:
                    prix_u = st.number_input("prix", 0, 500_000, int(saved.get("prix_unite", 0)),
                                             step=50, key=f"{key}_prix", label_visibility="collapsed")
                else:
                    prix_u = st.number_input("mt", 0, 5_000_000, int(saved.get("prix_unite", 0)),
                                             step=1000, key=f"{key}_prix", label_visibility="collapsed")
            with c3:
                if mult_type == "etp":
                    montant_calc = round(prix_u * _total_etp)
                    st.markdown(f"x {_total_etp:.1f} ETP = **{montant_calc:,.0f} \u20ac**")
                elif mult_type == "couverts":
                    montant_calc = round(prix_u * _nb_couverts)
                    st.markdown(f"x {_nb_couverts} couv. = **{montant_calc:,.0f} \u20ac**")
                else:
                    st.markdown("*-*")
                    montant_calc = prix_u
            with c4:
                dur_s = st.number_input("dur", 0, 30, int(saved.get("duree_amort", 3)),
                                        key=f"{key}_dur", label_visibility="collapsed")
            with c5:
                amort_s = montant_calc / dur_s if dur_s > 0 else 0
                st.markdown(f'<div style="text-align:center"><b>{amort_s:,.0f} \u20ac</b></div>' if dur_s > 0 else '<div style="text-align:center"><i>Non amorti</i></div>', unsafe_allow_html=True)

            inv_supp_new[key] = {"prix_unite": prix_u, "duree_amort": dur_s, "montant": montant_calc}
            if montant_calc > 0:
                new_inv.append({"categorie": cat_name, "montant": montant_calc, "duree_amort": dur_s})

        # ── Lignes libres (intitule modifiable + "Autre") en fin ──
        _offset_libre = len(_inv_fixes) + len(inv_supp_cats)
        for j, inv in enumerate(_inv_libres):
            _is_autre = inv["categorie"] == "Autre"
            c1, c2, c3, c4, c5 = st.columns([3, 2, 1.5, 1.5, 2])
            with c1:
                if _is_autre:
                    st.markdown("**Autre**")
                    cat = "Autre"
                else:
                    cat = st.text_input("cat", inv["categorie"], key=f"inv_c_libre_{j}", label_visibility="collapsed",
                                        placeholder="Intitule libre")
            with c2:
                mont = st.number_input("mt", 0, 20_000_000, inv["montant"], step=50_000, key=f"inv_m_libre_{j}", label_visibility="collapsed")
            with c3:
                st.markdown("*-*")
            with c4:
                dur = st.number_input("dur", 0, 30, inv["duree_amort"], key=f"inv_d_libre_{j}", label_visibility="collapsed")
            with c5:
                amort_an = mont / dur if dur > 0 else 0
                st.markdown(f'<div style="text-align:center"><b>{amort_an:,.0f} \u20ac</b></div>' if dur > 0 else '<div style="text-align:center"><i>Non amorti</i></div>', unsafe_allow_html=True)
            new_inv.append({"categorie": cat, "montant": mont, "duree_amort": dur})

        p["investissements"] = new_inv
        p["investissements_supplementaires"] = inv_supp_new

        total_inv = sum(i["montant"] for i in p["investissements"])
        st.info(f"**Investissement total : {total_inv:,.0f} \u20ac**")

        st.markdown("---")

        # ── Reinvestissements futurs ──
        st.markdown('<div class="section-header">\U0001F504 Reinvestissements annuels</div>', unsafe_allow_html=True)
        st.caption("Pourcentage annuel de l'investissement initial (amenagements + mobilier) "
                   "reinvesti pour maintenir et renouveler les equipements.")

        # Base de reinvestissement = tous les investissements amortissables
        montant_reinvest_base = sum(
            inv["montant"] for inv in p["investissements"]
            if inv.get("montant", 0) > 0 and inv.get("duree_amort", 0) > 0
        )
        st.caption(f"Base de reinvestissement (tous investissements amortissables) : **{montant_reinvest_base:,.0f} \u20ac**")

        pct_list = p.get("reinvest_pct_an", [0.0, 0.005, 0.01, 0.015, 0.02, 0.02, 0.025, 0.025, 0.025, 0.025])
        # S'assurer qu'on a au moins 10 valeurs
        while len(pct_list) < 10:
            pct_list.append(pct_list[-1] if pct_list else 0.02)

        st.caption("Les annees au-dela de l'An 10 reprennent le pourcentage de l'An 10.")

        # Ligne 1 : An 1 à An 5
        cols_r1 = st.columns(5)
        new_pct = []
        for idx in range(5):
            with cols_r1[idx]:
                val = st.number_input(f"An {idx+1}", 0.0, 0.15, float(pct_list[idx]), step=0.005, format="%.3f", key=f"reinv_pct_{idx}")
                new_pct.append(val)
                st.caption(f"{montant_reinvest_base * val:,.0f} \u20ac")

        # Ligne 2 : An 6 à An 10
        cols_r2 = st.columns(5)
        for idx in range(5, 10):
            with cols_r2[idx - 5]:
                val = st.number_input(f"An {idx+1}", 0.0, 0.15, float(pct_list[idx]), step=0.005, format="%.3f", key=f"reinv_pct_{idx}")
                new_pct.append(val)
                st.caption(f"{montant_reinvest_base * val:,.0f} \u20ac")

        p["reinvest_pct_an"] = new_pct

        # Graphique amortissements / an (injecte apres calcul de la projection)
        charts["invest_finance"] = st.container()

        st.markdown("---")

        # Prets
        st.markdown('<div class="section-header">\U0001F3E6 Prets</div>', unsafe_allow_html=True)
        new_prets = []
        pret_cols = st.columns(len(p["prets"]))
        _types_pret = ["Mensualite constante", "Interet seul (in fine)"]
        for i, pret in enumerate(p["prets"]):
            with pret_cols[i]:
                is_rocher = "rocher" in pret.get("nom", "").lower()
                st.markdown(f"**{pret['nom']}**" + (" \U0001F512" if is_rocher else ""))
                if is_rocher:
                    st.caption("Defini par l'Immobiliere Rocher — non modifiable ici.")
                nom = st.text_input("Nom", pret["nom"], key=f"pret_n_{i}", disabled=is_rocher)
                montant = st.number_input("Montant (\u20ac)", 0, 20_000_000, pret["montant"], step=100_000, key=f"pret_m_{i}", disabled=is_rocher)
                taux = st.number_input("Taux annuel", 0.0, 0.15, pret["taux_annuel"], step=0.005, format="%.3f", key=f"pret_t_{i}", disabled=is_rocher)
                duree = st.number_input("Duree (ans)", 1, 30, pret["duree_ans"], key=f"pret_d_{i}", disabled=is_rocher)
                differe = st.number_input("Differe (mois)", 0, 36, pret.get("differe_mois", 0), key=f"pret_df_{i}", disabled=is_rocher)
                # Subside RW (avant le type, car il force le type)
                subside_ch = st.checkbox("Garanti par un subside RW",
                    value=pret.get("subside_rw", False), key=f"pret_sub_{i}",
                    disabled=is_rocher,
                    help="A l'echeance, le subside rembourse le capital restant. "
                         "Force le type a 'Interet seul' (capital rembourse en une fois par le subside).")

                if subside_ch:
                    # Subside RW => forcement interet seul
                    st.selectbox("Type", _types_pret, index=1, key=f"pret_type_{i}", disabled=True)
                    type_val = "interet_seul"
                else:
                    type_idx = 1 if pret.get("type", "annuite") == "interet_seul" else 0
                    type_choisi = st.selectbox("Type", _types_pret, index=type_idx, key=f"pret_type_{i}", disabled=is_rocher)
                    type_val = "interet_seul" if "nteret" in type_choisi else "annuite"

                taux_m = taux / 12
                nb_m = duree * 12
                if type_val == "interet_seul":
                    mens = montant * taux_m if montant > 0 else 0
                else:
                    if taux_m > 0 and montant > 0:
                        mens = montant * taux_m / (1 - (1 + taux_m) ** -nb_m)
                    elif montant > 0:
                        mens = montant / nb_m
                    else:
                        mens = 0

                # Tableau d'amortissement detaille
                if montant > 0:
                    _pret_tmp = {"nom": nom, "montant": montant, "taux_annuel": taux,
                                 "duree_ans": duree, "differe_mois": differe, "type": type_val}
                    _df_amort = calc_tableau_pret(_pret_tmp, p["date_ouverture"], p["nb_mois_projection"])
                    with st.expander("Tableau d'amortissement"):
                        _u_pret = st.radio("Unite", ["\u20ac", "K\u20ac"], horizontal=True, key=f"pu_pret_{i}")
                        _p_pret = st.radio("Periode", ["Annuel", "Mensuel"], horizontal=True, key=f"pp_pret_{i}")
                        _div_p = 1000 if _u_pret.startswith("K") else 1
                        _sfx_p = " K\u20ac" if _u_pret.startswith("K") else " \u20ac"
                        def _fp(v): return f"{v / _div_p:,.0f}{_sfx_p}"
                        if _p_pret == "Annuel":
                            _amort_ann = _df_amort.copy()
                            _amort_ann["annee"] = _amort_ann["date"].apply(lambda x: x.year)
                            _amort_agg = _amort_ann.groupby("annee").agg({
                                "interets": "sum", "capital": "sum", "mensualite": "sum",
                                "capital_restant": "last"
                            }).reset_index()
                            st.dataframe(pd.DataFrame({
                                "Annee": [int(a) for a in _amort_agg["annee"]],
                                "Interets": [_fp(v) for v in _amort_agg["interets"]],
                                "Capital": [_fp(v) for v in _amort_agg["capital"]],
                                "Total": [_fp(v) for v in _amort_agg["mensualite"]],
                                "Capital restant": [_fp(v) for v in _amort_agg["capital_restant"]],
                            }), hide_index=True, use_container_width=True,
                                height=min(600, 35 * len(_amort_agg) + 38))
                        else:
                            st.dataframe(pd.DataFrame({
                                "Mois": [d.strftime("%b %Y") for d in _df_amort["date"]],
                                "Interets": [_fp(v) for v in _df_amort["interets"]],
                                "Capital": [_fp(v) for v in _df_amort["capital"]],
                                "Total": [_fp(v) for v in _df_amort["mensualite"]],
                                "Capital restant": [_fp(v) for v in _df_amort["capital_restant"]],
                            }), hide_index=True, use_container_width=True,
                                height=min(600, 35 * min(len(_df_amort), 15) + 38))

                new_prets.append({"nom": nom, "montant": montant, "taux_annuel": taux, "duree_ans": duree, "differe_mois": differe, "type": type_val, "subside_rw": subside_ch})
        p["prets"] = new_prets

        # Ligne de mensualites resume
        _mens_cols = st.columns(len(new_prets))
        for _mi, _pr in enumerate(new_prets):
            with _mens_cols[_mi]:
                _tm = _pr["taux_annuel"] / 12
                _nm = _pr["duree_ans"] * 12
                if _pr.get("type") == "interet_seul":
                    _ms = _pr["montant"] * _tm if _pr["montant"] > 0 else 0
                    _lbl = f"{_pr['nom']} (interets seuls)"
                else:
                    if _tm > 0 and _pr["montant"] > 0:
                        _ms = _pr["montant"] * _tm / (1 - (1 + _tm) ** -_nm)
                    elif _pr["montant"] > 0:
                        _ms = _pr["montant"] / _nm
                    else:
                        _ms = 0
                    _lbl = _pr["nom"]
                st.metric(_lbl, fmt_eur(_ms) + " /mois")

        st.markdown("---")
        st.markdown('<div class="section-header">\U0001F4B0 Fonds propres</div>', unsafe_allow_html=True)
        _fonds_propres_widget(p, key_prefix="chateau_fp")

        total_emprunte = sum(pr["montant"] for pr in p["prets"])
        if total_inv > 0:
            st.info(f"**Emprunts : {total_emprunte:,.0f} \u20ac | Fonds propres : {p['fonds_propres_initial']:,.0f} \u20ac | "
                    f"Ratio FP/Invest : {p['fonds_propres_initial']/total_inv*100:.1f}%**")

    # ════════════════════════════════════════════════════════════════════════
    # EQUILIBRE MOYENS / BESOINS
    # ════════════════════════════════════════════════════════════════════════
    with ht_equilibre:
        st.subheader("Equilibre Moyens / Besoins")
        st.caption("Comparaison entre les ressources disponibles et les besoins de financement "
                   "(investissements + tresorerie necessaire pendant la montee en puissance). "
                   "Le graphique est calcule apres projection.")

        # Le BFR d'activite sera calcule apres projection - on place un container
        charts["moyens_besoins"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 6. FISCALITE
    # ════════════════════════════════════════════════════════════════════════
    with ht6:
        st.subheader("Hypotheses fiscales")

        col_left, col_right = st.columns(2)

        with col_left:
            st.markdown('<div class="section-header">TVA sur les ventes (par service)</div>', unsafe_allow_html=True)
            new_tva_v = {}
            for k, v in p["tva_ventes"].items():
                new_tva_v[k] = st.number_input(f"TVA {k}", 0.0, 0.30, v, step=0.01, format="%.2f", key=f"tvav_{k}")
            p["tva_ventes"] = new_tva_v

        with col_right:
            st.markdown('<div class="section-header">TVA sur les achats (par type)</div>', unsafe_allow_html=True)
            new_tva_a = {}
            for k, v in p["tva_achats"].items():
                new_tva_a[k] = st.number_input(f"TVA {k}", 0.0, 0.30, v, step=0.01, format="%.2f", key=f"tvaa_{k}")
            p["tva_achats"] = new_tva_a

        st.markdown("---")
        st.markdown('<div class="section-header">Impot & Periodicite</div>', unsafe_allow_html=True)
        c1, c2 = st.columns(2)
        with c1:
            p["isoc"] = st.number_input("ISOC (impot des societes)", 0.0, 0.40, p["isoc"], step=0.01, format="%.2f")
        with c2:
            p["tva_periodicite"] = st.selectbox("Declaration TVA", ["Mensuelle", "Trimestrielle"],
                                                 index=0 if p["tva_periodicite"] == "Mensuelle" else 1)

        charts["fiscalite"] = st.container()

    # ════════════════════════════════════════════════════════════════════════
    # 6. DELAIS DE PAIEMENT
    # ════════════════════════════════════════════════════════════════════════
    with ht7:
        st.subheader("Delais de paiement")
        st.caption("Delai en mois entre la facturation et l'encaissement (clients) "
                   "ou le decaissement (fournisseurs). 0 = paiement comptant. "
                   "Impact sur la tresorerie (BFR) et le bilan.")

        p.setdefault("delais_clients", {})
        p.setdefault("delais_fournisseurs", {})

        st.markdown('<div class="section-header">Cote clients (encaissements)</div>', unsafe_allow_html=True)

        st.markdown("**Hebergement par segment**")
        _seg_cols = st.columns(len(p.get("segments", {})))
        for i, seg in enumerate(p.get("segments", {}).keys()):
            key = f"hebergement_{seg}"
            default = 0 if seg == "Loisirs" else 1
            with _seg_cols[i]:
                p["delais_clients"][key] = st.number_input(
                    seg, 0, 12, int(p["delais_clients"].get(key, default)),
                    step=1, key=f"dc_{key}")

        st.markdown("**Autres revenus**")
        _other_cl = [
            ("brasserie", "Brasserie", 0), ("bar", "Bar", 0),
            ("spa", "Spa", 0), ("salles", "Salles / Evenements", 1),
            ("divers", "Divers", 0), ("loyer_restaurant", "Loyer restaurant", 1),
        ]
        _cl_cols = st.columns(3)
        for i, (key, label, dflt) in enumerate(_other_cl):
            with _cl_cols[i % 3]:
                p["delais_clients"][key] = st.number_input(
                    label, 0, 12, int(p["delais_clients"].get(key, dflt)),
                    step=1, key=f"dc_{key}")

        st.markdown("---")
        st.markdown('<div class="section-header">Cote fournisseurs (decaissements)</div>', unsafe_allow_html=True)

        st.markdown("**Charges variables par nature**")
        _cv_items = [
            ("Linge & blanchisserie", 1), ("Produits & consommables", 1),
            ("Energie variable", 1), ("Nourriture", 1),
            ("Boissons", 1), ("Commissions OTA", 0),
            ("Commissions CB & franchise", 0), ("Soins spa", 1),
            ("Evenements (nettoyage, materiel)", 1),
        ]
        _cv_cols = st.columns(3)
        for i, (key, dflt) in enumerate(_cv_items):
            with _cv_cols[i % 3]:
                p["delais_fournisseurs"][key] = st.number_input(
                    key, 0, 12, int(p["delais_fournisseurs"].get(key, dflt)),
                    step=1, key=f"df_{key.replace(' ', '_')}")

        st.markdown("**Charges fixes (hors personnel)**")
        _cf_items = [
            ("Loyer", 0), ("CF directs departements", 1),
            ("Energie fixe", 1), ("Assurances & taxes", 0),
            ("Maintenance & contrats", 1), ("Marketing", 1),
            ("IT & logiciels", 0), ("Honoraires", 1),
            ("Autres CF indirects", 1),
        ]
        _cf_cols = st.columns(3)
        for i, (key, dflt) in enumerate(_cf_items):
            with _cf_cols[i % 3]:
                p["delais_fournisseurs"][key] = st.number_input(
                    key, 0, 12, int(p["delais_fournisseurs"].get(key, dflt)),
                    step=1, key=f"df_{key.replace(' ', '_')}")

    return p


# ─── Onglet Dashboard ─────────────────────────────────────────────────────────

def tab_dashboard(df, indic, params):
    st.header("\U0001F4CA Dashboard - Vue d'ensemble")

    annees_cal = sorted(indic["Annee calendaire"].unique())
    # Annee de croisiere = 4e annee calendaire (ou derniere si moins de 4 ans)
    idx_croisiere = min(3, len(annees_cal) - 1)
    an_cal = int(annees_cal[idx_croisiere])
    an4 = indic[indic["Annee calendaire"] == an_cal].iloc[0]

    # ── KPIs avec explications ────────────────────────────────────────────
    st.markdown(f'<div class="section-header">Indicateurs cles — Annee de croisiere ({an_cal})</div>',
                unsafe_allow_html=True)
    st.caption("Ces indicateurs correspondent a l'annee de croisiere, "
               "c'est-a-dire le moment ou l'hotel atteint son regime de fonctionnement stable.")

    cols = st.columns(6)
    with cols[0]:
        metric_card("CA Total", fmt_eur(an4["CA Total"]), "metric-blue")
        st.caption("Chiffre d'affaires total toutes activites confondues (hebergement, brasserie, bar, spa, salles, divers).")
    with cols[1]:
        metric_card("EBITDA", fmt_eur(an4["EBITDA"]), "metric-green" if an4["EBITDA"] > 0 else "metric-red")
        st.caption("Excedent brut d'exploitation : capacite de l'hotel a generer du cash avant amortissements, interets et impots.")
    with cols[2]:
        metric_card("Marge EBITDA", fmt_pct(an4["EBITDA %"]), "metric-green" if an4["EBITDA %"] > 20 else "metric-orange")
        st.caption("EBITDA / CA Total. Ref. hotellerie 5* : 25-35% en croisiere. En dessous de 20%, la rentabilite est fragile.")
    with cols[3]:
        metric_card("RevPAR", f"{an4['RevPAR']:.0f} \u20ac", "metric-blue")
        st.caption("Revenue Per Available Room : CA hebergement / nombre de chambres disponibles. Indicateur cle du secteur.")
    with cols[4]:
        metric_card("Taux Occup.", fmt_pct(an4["Taux Occupation"]), "metric-green" if an4["Taux Occupation"] > 60 else "metric-orange")
        st.caption("Pourcentage moyen de chambres occupees. Ref. hotellerie 5* : 55-70% en croisiere.")
    with cols[5]:
        dscr_val = an4["DSCR"]
        dscr_str = f"{dscr_val:.2f}x" if dscr_val < 100 else "\u221e"
        metric_card("DSCR", dscr_str, "metric-green" if dscr_val > 1.3 else "metric-red")
        st.caption("Debt Service Coverage Ratio : EBITDA / service de la dette. Au-dessus de 1.3x, la dette est couverte.")

    st.markdown("---")

    # ── Toggle annuel / mensuel ───────────────────────────────────────────
    vue = st.radio("Affichage", ["Annuel", "Mensuel"], horizontal=True, key="dashboard_vue")

    n_annees = len(indic)
    n_mois = len(df)

    # ── 1. Chiffre d'affaires par segment ─────────────────────────────────
    st.markdown('<div class="section-header">Chiffre d\'affaires par segment</div>', unsafe_allow_html=True)
    st.caption("Decomposition du CA entre hebergement, brasserie et autres activites (bar, spa, salles, divers). "
               "La 1ere annee est partielle si l'ouverture a lieu en cours d'annee.")

    if vue == "Annuel":
        x_labels = [str(int(row['Annee calendaire'])) for _, row in indic.iterrows()]
        fig = go.Figure()
        fig.add_trace(go.Bar(x=x_labels, y=indic["CA Hebergement"], name="Hebergement", marker_color="#4facfe"))
        fig.add_trace(go.Bar(x=x_labels, y=indic["CA Brasserie"], name="Brasserie", marker_color="#f5576c"))
        fig.add_trace(go.Bar(x=x_labels, y=indic["CA Autres"], name="Autres", marker_color="#38ef7d"))
        fig.update_layout(barmode="stack", height=450,
                          xaxis=dict(title="", type="category", tickfont=dict(size=12)),
                          yaxis_title="EUR", yaxis_tickformat=",",
                          legend=dict(orientation="h", y=-0.12))
        _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 450)
    else:
        x_dates = df["date"]
        fig = go.Figure()
        fig.add_trace(go.Bar(x=x_dates, y=df["ca_hebergement"], name="Hebergement", marker_color="#4facfe"))
        fig.add_trace(go.Bar(x=x_dates, y=df["ca_brasserie"], name="Brasserie", marker_color="#f5576c"))
        autres_cols = ["ca_bar", "ca_spa", "ca_salles", "ca_divers"]
        ca_autres = sum(df[c] for c in autres_cols if c in df.columns)
        fig.add_trace(go.Bar(x=x_dates, y=ca_autres, name="Autres", marker_color="#38ef7d"))
        fig.update_layout(barmode="stack", height=450,
                          xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                                     tickfont=dict(size=11), tickangle=-45),
                          yaxis_title="EUR", yaxis_tickformat=",",
                          legend=dict(orientation="h", y=-0.18))
        _scrollable_chart(fig, n_mois, _PX_PAR_MOIS, 450)

    # ── 2. Rentabilite ────────────────────────────────────────────────────
    st.markdown('<div class="section-header">Rentabilite : EBITDA, Resultat Net et Cash Flow</div>', unsafe_allow_html=True)
    st.caption("EBITDA = marge operationnelle brute. Resultat Net = apres amortissements, interets et impots. "
               "Cash Flow = tresorerie generee par l'exploitation. "
               "Un resultat net negatif les premieres annees est courant (amortissements eleves).")

    if vue == "Annuel":
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=x_labels, y=indic["EBITDA"], name="EBITDA",
                                  mode="lines+markers", line=dict(color="#11998e", width=3)))
        fig.add_trace(go.Scatter(x=x_labels, y=indic["Resultat Net"], name="Resultat Net",
                                  mode="lines+markers", line=dict(color="#f5576c", width=3)))
        fig.add_trace(go.Scatter(x=x_labels, y=indic["Cash Flow"], name="Cash Flow",
                                  mode="lines+markers", line=dict(color="#4facfe", width=3, dash="dash")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
        fig.update_layout(height=450,
                          xaxis=dict(title="", type="category", tickfont=dict(size=12)),
                          yaxis_title="EUR", yaxis_tickformat=",",
                          legend=dict(orientation="h", y=-0.12))
        _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 450)
    else:
        x_dates = df["date"]
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=x_dates, y=df["ebitda"], name="EBITDA",
                                  mode="lines", line=dict(color="#11998e", width=2)))
        fig.add_trace(go.Scatter(x=x_dates, y=df["resultat_net"], name="Resultat Net",
                                  mode="lines", line=dict(color="#f5576c", width=2)))
        fig.add_trace(go.Scatter(x=x_dates, y=df["cash_flow_operationnel"], name="Cash Flow",
                                  mode="lines", line=dict(color="#4facfe", width=2, dash="dash")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
        fig.update_layout(height=450,
                          xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                                     tickfont=dict(size=11), tickangle=-45),
                          yaxis_title="EUR", yaxis_tickformat=",",
                          legend=dict(orientation="h", y=-0.18))
        _scrollable_chart(fig, n_mois, _PX_PAR_MOIS, 450)

    # ── 3. Indicateurs hoteliers ──────────────────────────────────────────
    st.markdown('<div class="section-header">Indicateurs hoteliers : Taux d\'occupation, ADR et RevPAR</div>',
                unsafe_allow_html=True)
    st.caption("**Taux d'occupation** : % de chambres vendues. "
               "**ADR** (Average Daily Rate) : prix moyen par chambre vendue. "
               "**RevPAR** (Revenue Per Available Room) : CA hebergement / chambres disponibles = ADR x Taux occup. "
               "Le RevPAR est l'indicateur de performance de reference dans l'hotellerie.")

    col1, col2 = st.columns(2)
    if vue == "Annuel":
        with col1:
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=x_labels, y=indic["Taux Occupation"],
                                      mode="lines+markers", name="Taux Occup.", line=dict(color="#764ba2", width=3)))
            fig.update_layout(title=dict(text="Taux d'occupation moyen (%)", font=dict(size=14)),
                              height=400,
                              xaxis=dict(title="", type="category", tickfont=dict(size=12)),
                              yaxis_title="%", yaxis_range=[0, 100])
            _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 400)

        with col2:
            fig = make_subplots(specs=[[{"secondary_y": True}]])
            fig.add_trace(go.Bar(x=x_labels, y=indic["Prix Moyen (ADR)"],
                                 name="ADR", marker_color="#667eea", opacity=0.7), secondary_y=False)
            fig.add_trace(go.Scatter(x=x_labels, y=indic["RevPAR"],
                                      name="RevPAR", mode="lines+markers", line=dict(color="#f5576c", width=3)),
                          secondary_y=True)
            fig.update_layout(title=dict(text="ADR & RevPAR", font=dict(size=14)),
                              height=400,
                              xaxis=dict(type="category", tickfont=dict(size=12)),
                              legend=dict(orientation="h", y=-0.15))
            fig.update_yaxes(title_text="ADR (\u20ac)", secondary_y=False)
            fig.update_yaxes(title_text="RevPAR (\u20ac)", secondary_y=True)
            _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 400)
    else:
        x_dates = df["date"]
        with col1:
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=x_dates, y=df["taux_occupation"] * 100,
                                      mode="lines", name="Taux Occup.", line=dict(color="#764ba2", width=2)))
            fig.update_layout(title=dict(text="Taux d'occupation mensuel (%)", font=dict(size=14)),
                              height=400,
                              xaxis=dict(title="", dtick="M3", tickformat="%b %Y",
                                         tickfont=dict(size=11), tickangle=-45),
                              yaxis_title="%", yaxis_range=[0, 100])
            _scrollable_chart(fig, n_mois, _PX_PAR_MOIS, 400)

        with col2:
            fig = make_subplots(specs=[[{"secondary_y": True}]])
            fig.add_trace(go.Bar(x=x_dates, y=df["prix_moyen"],
                                 name="ADR", marker_color="#667eea", opacity=0.7), secondary_y=False)
            revpar = df["ca_hebergement"] / (df["nuitees"] / df["taux_occupation"]).replace(0, float("nan"))
            fig.add_trace(go.Scatter(x=x_dates, y=revpar,
                                      name="RevPAR", mode="lines", line=dict(color="#f5576c", width=2)),
                          secondary_y=True)
            fig.update_layout(title=dict(text="ADR & RevPAR (mensuel)", font=dict(size=14)),
                              height=400,
                              xaxis=dict(dtick="M3", tickformat="%b %Y",
                                         tickfont=dict(size=11), tickangle=-45),
                              legend=dict(orientation="h", y=-0.18))
            fig.update_yaxes(title_text="ADR (\u20ac)", secondary_y=False)
            fig.update_yaxes(title_text="RevPAR (\u20ac)", secondary_y=True)
            _scrollable_chart(fig, n_mois, _PX_PAR_MOIS, 400)

    # ── 4. Cash flow cumule & Break-even ──────────────────────────────────
    st.markdown('<div class="section-header">Cash Flow cumule et point mort (break-even)</div>',
                unsafe_allow_html=True)
    st.caption("Le cash flow cumule represente la tresorerie nette generee depuis l'ouverture, "
               "deduction faite du service de la dette. Le break-even est atteint lorsque le cumul "
               "repasse au-dessus de zero, signifiant que l'exploitation a rembourse ses pertes initiales.")

    fig = go.Figure()
    fig.add_trace(go.Scatter(x=df["date"], y=df["cash_flow_cumul"],
                              fill="tozeroy", line=dict(color="#667eea", width=2), name="Cash Flow Cumule"))
    fig.add_hline(y=0, line_dash="dash", line_color="red")
    invest_total = sum(i["montant"] for i in params["investissements"])
    fig.add_hline(y=-invest_total, line_dash="dot", line_color="gray",
                  annotation_text="Investissement initial")
    fig.update_layout(height=400,
                      xaxis=dict(title="", dtick="M6", tickformat="%b %Y",
                                 tickfont=dict(size=11), tickangle=-45),
                      yaxis_title="EUR", yaxis_tickformat=",")
    _scrollable_chart(fig, n_mois, _PX_PAR_MOIS, 400)

    cf_cumul = df["cash_flow_cumul"].values
    break_even_idx = None
    for i in range(1, len(cf_cumul)):
        if cf_cumul[i - 1] < 0 and cf_cumul[i] >= 0:
            break_even_idx = i
            break
    if break_even_idx:
        be_date = df.iloc[break_even_idx]["date"]
        st.success(f"\u2705 **Break-even atteint en {be_date.strftime('%B %Y')}** (mois {break_even_idx} d'exploitation)")
    else:
        st.warning("\u26A0 Break-even non atteint sur la periode de projection")


# ─── Onglet Projection mensuelle ──────────────────────────────────────────────

# ─── Onglet Indicateurs ───────────────────────────────────────────────────────

def tab_indicateurs(df, indic, params):
    st.header("\U0001F3AF Explorateur de donnees")

    # ── Top controls row ──────────────────────────────────────────────────
    c1, c2, c3 = st.columns([2, 1, 1])
    with c1:
        categorie = st.selectbox("Categorie", [
            "Vue d'ensemble", "Hebergement", "Brasserie", "Bar", "Spa",
            "Salles/Evenements", "Charges", "Cash Flow & Dette",
        ], key="indic_cat")
    with c2:
        periode = st.radio("Periode", ["Annuel", "Mensuel"], horizontal=True, key="indic_periode")
    with c3:
        if periode == "Mensuel":
            annees_list = sorted(df["annee"].unique().tolist())
            annee_filtre = st.selectbox("Annee", ["Toutes"] + annees_list, key="indic_annee")
        else:
            annee_filtre = "Toutes"

    # ── Helpers ───────────────────────────────────────────────────────────
    x_labels_an = [str(int(row['Annee calendaire'])) for _, row in indic.iterrows()]
    n_annees = len(indic)
    n_mois = len(df)

    def _filter_monthly(src):
        if annee_filtre != "Toutes":
            return src[src["annee"] == annee_filtre].copy()
        return src.copy()

    def _tick_cfg():
        return dict(dtick="M1" if annee_filtre != "Toutes" else "M3",
                    tickformat="%b %Y", tickfont=dict(size=11), tickangle=-45)

    def _px_pt():
        return _PX_PAR_MOIS if annee_filtre == "Toutes" else 55

    # ── Vue d'ensemble ────────────────────────────────────────────────────
    if categorie == "Vue d'ensemble":
        _agg_cols = {
            "ca_total": "sum", "cv_total": "sum", "marge_brute": "sum",
            "cf_directs_total": "sum", "marge": "sum", "subside_rw": "sum",
            "cf_indirects_total": "sum", "cf_total": "sum",
            "cf_total_cash": "sum",
            "ebitda": "sum", "amortissement": "sum", "ebit": "sum",
            "dette_interets": "sum", "dette_capital": "sum",
            "impot": "sum", "impot_cash": "sum",
            "tva_paiement": "sum", "reinvest_acquisition": "sum",
            "delay_adjustment": "sum",
            "resultat_net": "sum", "cash_flow": "sum", "cash_flow_cumul": "last",
        }

        def _fmt(v):
            return f"{v:,.0f} \u20ac"

        def _pct(num, den):
            return f"{num / den * 100:.1f}%" if den != 0 else "-"

        if periode == "Annuel":
            annual = df.groupby("annee").agg(_agg_cols).reset_index()
            x_lab = [str(int(a)) for a in annual["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CA Total": [_fmt(v) for v in annual["ca_total"]],
                "CV Total": [_fmt(v) for v in annual["cv_total"]],
                "Marge Brute": [_fmt(v) for v in annual["marge_brute"]],
                "Marge Brute %": [_pct(m, c) for m, c in zip(annual["marge_brute"], annual["ca_total"])],
                "CF Directs": [_fmt(v) for v in annual["cf_directs_total"]],
                "Marge service": [_fmt(v) for v in annual["marge"]],
                "Subside RW": [_fmt(v) for v in annual["subside_rw"]],
                "CF Indirects": [_fmt(v) for v in annual["cf_indirects_total"]],
                "EBITDA": [_fmt(v) for v in annual["ebitda"]],
                "EBITDA %": [_pct(e, c) for e, c in zip(annual["ebitda"], annual["ca_total"])],
                "Amortissement": [_fmt(v) for v in annual["amortissement"]],
                "EBIT": [_fmt(v) for v in annual["ebit"]],
                "Interets": [_fmt(v) for v in annual["dette_interets"]],
                "Resultat avant impot": [_fmt(v) for v in (annual["ebit"] - annual["dette_interets"] + annual["subside_rw"])],
                "ISOC (charge)": [_fmt(v) for v in annual["impot"]],
                "Resultat Net": [_fmt(v) for v in annual["resultat_net"]],
                "ISOC (paye)": [_fmt(v) for v in annual["impot_cash"]],
                "Remb. dette": [_fmt(v) for v in annual["dette_capital"]],
                "Reinvestissements": [_fmt(v) for v in annual["reinvest_acquisition"]],
                "TVA a reverser": [_fmt(v) for v in annual["tva_paiement"]],
                "Delta Cash Pers.": [_fmt(v) for v in (annual["cf_total_cash"] - annual["cf_total"])],
                "BFR (impact cash)": [_fmt(v) for v in annual["delay_adjustment"]],
                "Cash Flow": [_fmt(v) for v in annual["cash_flow"]],
                "Cash Flow Cumul": [_fmt(v) for v in annual["cash_flow_cumul"]],
            })
        else:
            df_src = _filter_monthly(df)
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "CA Total": [_fmt(v) for v in df_src["ca_total"]],
                "CV Total": [_fmt(v) for v in df_src["cv_total"]],
                "Marge Brute": [_fmt(v) for v in df_src["marge_brute"]],
                "Marge Brute %": [_pct(m, c) for m, c in zip(df_src["marge_brute"], df_src["ca_total"])],
                "CF Directs": [_fmt(v) for v in df_src["cf_directs_total"]],
                "Marge service": [_fmt(v) for v in df_src["marge"]],
                "Subside RW": [_fmt(v) for v in df_src["subside_rw"]],
                "CF Indirects": [_fmt(v) for v in df_src["cf_indirects_total"]],
                "EBITDA": [_fmt(v) for v in df_src["ebitda"]],
                "EBITDA %": [_pct(e, c) for e, c in zip(df_src["ebitda"], df_src["ca_total"])],
                "Amortissement": [_fmt(v) for v in df_src["amortissement"]],
                "EBIT": [_fmt(v) for v in df_src["ebit"]],
                "Interets": [_fmt(v) for v in df_src["dette_interets"]],
                "ISOC (charge)": [_fmt(v) for v in df_src["impot"]],
                "Resultat Net": [_fmt(v) for v in df_src["resultat_net"]],
                "ISOC (paye)": [_fmt(v) for v in df_src["impot_cash"]],
                "Remb. dette": [_fmt(v) for v in df_src["dette_capital"]],
                "Reinvestissements": [_fmt(v) for v in df_src["reinvest_acquisition"]],
                "TVA a reverser": [_fmt(v) for v in df_src["tva_paiement"]],
                "Delta Cash Pers.": [_fmt(v) for v in (df_src["cf_total_cash"] - df_src["cf_total"])],
                "BFR (impact cash)": [_fmt(v) for v in df_src["delay_adjustment"]],
                "Cash Flow": [_fmt(v) for v in df_src["cash_flow"]],
                "Cash Flow Cumul": [_fmt(v) for v in df_src["cash_flow_cumul"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Hebergement ───────────────────────────────────────────────────────
    elif categorie == "Hebergement":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({
                "nuitees": "sum", "taux_occupation": "mean", "prix_moyen": "mean",
                "ca_hebergement": "sum",
            }).reset_index()
            nb_ch = params["nb_chambres"]
            agg["revpar"] = agg["ca_hebergement"] / (nb_ch * 365)
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "Nuitees": [f"{v:,.0f}" for v in agg["nuitees"]],
                "Taux Occupation": [f"{v*100:.1f}%" for v in agg["taux_occupation"]],
                "Prix Moyen (ADR)": [f"{v:,.0f} \u20ac" for v in agg["prix_moyen"]],
                "CA Hebergement": [f"{v:,.0f} \u20ac" for v in agg["ca_hebergement"]],
                "RevPAR": [f"{v:,.0f} \u20ac" for v in agg["revpar"]],
            })
        else:
            df_src = _filter_monthly(df)
            nb_ch = params["nb_chambres"]
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "Nuitees": [f"{v:,.0f}" for v in df_src["nuitees"]],
                "Taux Occupation": [f"{v*100:.1f}%" for v in df_src["taux_occupation"]],
                "Prix Moyen (ADR)": [f"{v:,.0f} \u20ac" for v in df_src["prix_moyen"]],
                "CA Hebergement": [f"{v:,.0f} \u20ac" for v in df_src["ca_hebergement"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Brasserie ─────────────────────────────────────────────────────────
    elif categorie == "Brasserie":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({"ca_brasserie": "sum", "ca_pdj": "sum"}).reset_index()
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CA Brasserie": [f"{v:,.0f} \u20ac" for v in agg["ca_brasserie"]],
                "CA PDJ": [f"{v:,.0f} \u20ac" for v in agg["ca_pdj"]],
            })
        else:
            df_src = _filter_monthly(df)
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "CA Brasserie": [f"{v:,.0f} \u20ac" for v in df_src["ca_brasserie"]],
                "CA PDJ": [f"{v:,.0f} \u20ac" for v in df_src["ca_pdj"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Bar ───────────────────────────────────────────────────────────────
    elif categorie == "Bar":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({"ca_bar": "sum"}).reset_index()
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CA Bar": [f"{v:,.0f} \u20ac" for v in agg["ca_bar"]],
            })
        else:
            df_src = _filter_monthly(df)
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "CA Bar": [f"{v:,.0f} \u20ac" for v in df_src["ca_bar"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Spa ───────────────────────────────────────────────────────────────
    elif categorie == "Spa":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({"ca_spa": "sum"}).reset_index()
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CA Spa": [f"{v:,.0f} \u20ac" for v in agg["ca_spa"]],
            })
        else:
            df_src = _filter_monthly(df)
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "CA Spa": [f"{v:,.0f} \u20ac" for v in df_src["ca_spa"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Salles/Evenements ─────────────────────────────────────────────────
    elif categorie == "Salles/Evenements":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({
                "ca_seminaires": "sum", "ca_mariages": "sum", "ca_commission_catering": "sum", "ca_salles": "sum",
            }).reset_index()
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CA Seminaires": [f"{v:,.0f} \u20ac" for v in agg["ca_seminaires"]],
                "CA Mariages": [f"{v:,.0f} \u20ac" for v in agg["ca_mariages"]],
                "Commission Catering": [f"{v:,.0f} \u20ac" for v in agg["ca_commission_catering"]],
                "CA Salles Total": [f"{v:,.0f} \u20ac" for v in agg["ca_salles"]],
            })
        else:
            df_src = _filter_monthly(df)
            df_table = pd.DataFrame({
                "Mois": [d.strftime("%b %Y") for d in df_src["date"]],
                "CA Seminaires": [f"{v:,.0f} \u20ac" for v in df_src["ca_seminaires"]],
                "CA Mariages": [f"{v:,.0f} \u20ac" for v in df_src["ca_mariages"]],
                "Commission Catering": [f"{v:,.0f} \u20ac" for v in df_src["ca_commission_catering"]],
                "CA Salles Total": [f"{v:,.0f} \u20ac" for v in df_src["ca_salles"]],
            })
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Charges ───────────────────────────────────────────────────────────
    elif categorie == "Charges":
        if periode == "Annuel":
            agg = df.groupby("annee").agg({
                "cv_hebergement": "sum", "cv_brasserie": "sum", "cv_pdj": "sum",
                "cv_bar": "sum", "cv_spa": "sum", "cv_total": "sum",
                "cf_personnel": "sum", "cf_autres": "sum", "cf_loyer": "sum", "cf_total": "sum",
            }).reset_index()
            x_lab = [str(int(a)) for a in agg["annee"]]
            df_table = pd.DataFrame({
                "Annee": x_lab,
                "CV Hebergement": [f"{v:,.0f} \u20ac" for v in agg["cv_hebergement"]],
                "CV Brasserie": [f"{v:,.0f} \u20ac" for v in agg["cv_brasserie"]],
                "CV PDJ": [f"{v:,.0f} \u20ac" for v in agg["cv_pdj"]],
                "CV Bar": [f"{v:,.0f} \u20ac" for v in agg["cv_bar"]],
                "CV Spa": [f"{v:,.0f} \u20ac" for v in agg["cv_spa"]],
                "CV Total": [f"{v:,.0f} \u20ac" for v in agg["cv_total"]],
                "CF Personnel": [f"{v:,.0f} \u20ac" for v in agg["cf_personnel"]],
                "CF Autres": [f"{v:,.0f} \u20ac" for v in agg["cf_autres"]],
                "CF Loyer": [f"{v:,.0f} \u20ac" for v in agg["cf_loyer"]],
                "CF Total": [f"{v:,.0f} \u20ac" for v in agg["cf_total"]],
            })
        else:
            df_src = _filter_monthly(df)
            charge_cols = ["cv_hebergement", "cv_brasserie", "cv_pdj", "cv_bar",
                           "cv_spa", "cv_total", "cf_personnel", "cf_autres", "cf_loyer", "cf_total"]
            df_table = pd.DataFrame({"Mois": [d.strftime("%b %Y") for d in df_src["date"]]})
            for c in charge_cols:
                df_table[c] = [f"{v:,.0f} \u20ac" for v in df_src[c]]
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))

    # ── Cash Flow & Dette ─────────────────────────────────────────────────
    elif categorie == "Cash Flow & Dette":
        if periode == "Annuel":
            df_table = pd.DataFrame({
                "Annee": x_labels_an,
                "EBITDA": [f"{v:,.0f} \u20ac" for v in indic["EBITDA"]],
                "Amortissement": [f"{v:,.0f} \u20ac" for v in indic["Amortissement"]],
                "Interets": [f"{v:,.0f} \u20ac" for v in indic["Interets"]],
                "Service Dette": [f"{v:,.0f} \u20ac" for v in indic["Service Dette"]],
                "DSCR": [f"{v:.2f}x" if v < 100 else "\u221e" for v in indic["DSCR"]],
                "Resultat Net": [f"{v:,.0f} \u20ac" for v in indic["Resultat Net"]],
                "Cash Flow": [f"{v:,.0f} \u20ac" for v in indic["Cash Flow"]],
                "Cash Flow Cumul": [f"{v:,.0f} \u20ac" for v in indic["Cash Flow Cumul"]],
            })
        else:
            df_src = _filter_monthly(df)
            cols_m = ["date", "ebitda", "amortissement", "dette_mensualite",
                      "dette_interets", "dette_capital", "resultat_net", "cash_flow", "cash_flow_cumul"]
            df_table = df_src[cols_m].copy()
            df_table["date"] = df_table["date"].apply(lambda d: d.strftime("%b %Y"))
            for c in cols_m:
                if c != "date":
                    df_table[c] = df_table[c].apply(lambda v: f"{v:,.0f} \u20ac")
        st.dataframe(df_table, use_container_width=True, hide_index=True,
                      height=min(600, 35 * len(df_table) + 38))



# ─── Onglet Prets ──────────────────────────────────────────────────────────────

def tab_prets(params):
    st.header("\U0001F3E6 Tableaux d'amortissement des prets")

    total_emprunte = sum(p["montant"] for p in params["prets"])
    invest_total = sum(i["montant"] for i in params["investissements"])
    fp = params["fonds_propres_initial"]

    cols = st.columns(4)
    with cols[0]:
        metric_card("Total emprunte", fmt_eur(total_emprunte), "metric-blue")
    with cols[1]:
        metric_card("Fonds propres", fmt_eur(fp), "metric-green")
    with cols[2]:
        metric_card("Investissement total", fmt_eur(invest_total), "metric-orange")
    with cols[3]:
        ratio = fp / invest_total * 100 if invest_total > 0 else 0
        metric_card("Ratio FP/Invest.", fmt_pct(ratio), "metric-blue")

    st.markdown("---")

    for pret in params["prets"]:
        if pret["montant"] == 0:
            continue

        st.subheader(f"{pret['nom']}")
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Montant", fmt_eur(pret["montant"]))
        with col2:
            st.metric("Taux", fmt_pct(pret["taux_annuel"] * 100))
        with col3:
            st.metric("Duree", f"{pret['duree_ans']} ans")
        with col4:
            taux_m = pret["taux_annuel"] / 12
            nb_m = pret["duree_ans"] * 12
            if taux_m > 0:
                mens = pret["montant"] * taux_m / (1 - (1 + taux_m) ** -nb_m)
            else:
                mens = pret["montant"] / nb_m
            st.metric("Mensualite", fmt_eur(mens))

        df_pret = calc_tableau_pret(pret, params["date_ouverture"])

        if not df_pret.empty:
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=df_pret["date"], y=df_pret["capital_restant"],
                                      fill="tozeroy", name="Capital restant", line=dict(color="#667eea")))
            fig.add_trace(go.Bar(x=df_pret["date"], y=df_pret["interets"],
                                  name="Interets", marker_color="#f5576c", opacity=0.6))
            fig.add_trace(go.Bar(x=df_pret["date"], y=df_pret["capital"],
                                  name="Capital rembourse", marker_color="#38ef7d", opacity=0.6))
            fig.update_layout(height=300, yaxis_tickformat=",", barmode="stack")
            st.plotly_chart(fig, use_container_width=True)

            df_pret["annee"] = df_pret["date"].apply(lambda d: d.year)
            annual = df_pret.groupby("annee").agg({
                "mensualite": "sum", "interets": "sum", "capital": "sum",
                "capital_restant": "last"
            }).reset_index()
            annual.columns = ["Annee", "Total Mensualites", "Total Interets", "Capital Rembourse", "Capital Restant"]
            for c in annual.columns[1:]:
                annual[c] = annual[c].apply(lambda v: f"{v:,.0f} \u20ac")
            st.dataframe(annual, use_container_width=True, hide_index=True)

        st.markdown("---")

    st.subheader("Service de la dette consolide")
    all_prets = []
    for pret in params["prets"]:
        if pret["montant"] > 0:
            df_p = calc_tableau_pret(pret, params["date_ouverture"])
            df_p["pret"] = pret["nom"]
            all_prets.append(df_p)

    if all_prets:
        df_all = pd.concat(all_prets)
        df_all["annee"] = df_all["date"].apply(lambda d: d.year)
        consolide = df_all.groupby("annee").agg({"mensualite": "sum", "interets": "sum", "capital": "sum"}).reset_index()

        fig = go.Figure()
        fig.add_trace(go.Bar(x=consolide["annee"], y=consolide["interets"], name="Interets", marker_color="#f5576c"))
        fig.add_trace(go.Bar(x=consolide["annee"], y=consolide["capital"], name="Capital", marker_color="#4facfe"))
        fig.update_layout(title="Service de la dette par annee", barmode="stack", height=350, yaxis_tickformat=",")
        st.plotly_chart(fig, use_container_width=True)


# ─── Onglet Presentation ─────────────────────────────────────────────────────

# Chapter definitions: (key, label, type, has_text, has_image)
_NON_FINANCIAL_CHAPTERS = [
    ("page_de_garde",            "Page de garde",              False, False),
    ("executive_summary",        "Executive Summary",          True,  False),
    ("historique_projet",        "Historique du projet",        True,  True),
    ("site_environnement",       "Le site & l'environnement",  True,  True),
    ("concept_positionnement",   "Concept & Positionnement",   True,  False),
    ("analyse_marche",           "Analyse de marche",           True,  True),
    ("gouvernance_equipe",       "Gouvernance & Equipe",        True,  True),
    ("organigramme",             "Organigramme",                False, True),
    ("strategie_commerciale",    "Strategie commerciale",       True,  False),
    ("planning_realisation",     "Planning de realisation",     True,  True),
]

_FINANCIAL_CHAPTERS = [
    ("invest_financement",   "Investissements & Financement"),
    ("hypotheses_cles",      "Hypotheses cles"),
    ("previsions_ca",        "Previsions de CA"),
    ("structure_couts",      "Structure des couts"),
    ("indicateurs_hoteliers","Indicateurs hoteliers"),
    ("rentabilite_resultat", "Rentabilite & Resultat"),
    ("cashflow_roi",         "Cash Flow & Retour sur investissement"),
    ("bilan",                "Bilan"),
]


def _try_fig_to_png(fig, width=900, height=450):
    """Try to export a plotly figure to PNG bytes. Returns bytes or None."""
    try:
        return fig.to_image(format="png", width=width, height=height)
    except Exception:
        return None


def _init_pres_chapters():
    """Initialise session_state pres_chapters if needed."""
    if "pres_chapters" not in st.session_state:
        chapters = {}
        for key, label, has_text, has_image in _NON_FINANCIAL_CHAPTERS:
            chapters[key] = {
                "label": label,
                "type": "non_financial",
                "included": True,
                "content": "",
                "images": [],
                "has_text": has_text,
                "has_image": has_image,
            }
        for key, label in _FINANCIAL_CHAPTERS:
            chapters[key] = {
                "label": label,
                "type": "financial",
                "included": True,
            }
        st.session_state["pres_chapters"] = chapters


def _make_presentation_charts(indic, params, df):
    """Generate all financial charts and return dict of name -> plotly figure."""
    charts = {}
    x_labels = [str(int(row["Annee calendaire"])) for _, row in indic.iterrows()]

    # 1. CA evolution stacked bar by segment
    fig_ca = go.Figure()
    fig_ca.add_trace(go.Bar(x=x_labels, y=indic["CA Hebergement"],
                            name="Hebergement", marker_color="#4facfe"))
    fig_ca.add_trace(go.Bar(x=x_labels, y=indic["CA Brasserie"],
                            name="Brasserie", marker_color="#f5576c"))
    fig_ca.add_trace(go.Bar(x=x_labels, y=indic["CA Autres"],
                            name="Autres", marker_color="#38ef7d"))
    fig_ca.update_layout(barmode="stack", height=450, width=900,
                         xaxis=dict(type="category"), yaxis_tickformat=",",
                         legend=dict(orientation="h", y=-0.15),
                         title="Evolution du Chiffre d'Affaires")
    charts["ca_evolution"] = fig_ca

    # 2. Rentabilite: EBITDA bar + Resultat Net line
    fig_rent = go.Figure()
    fig_rent.add_trace(go.Bar(x=x_labels, y=indic["EBITDA"],
                              name="EBITDA", marker_color="#11998e"))
    fig_rent.add_trace(go.Scatter(x=x_labels, y=indic["Resultat Net"],
                                  name="Resultat Net", mode="lines+markers",
                                  line=dict(color="#f5576c", width=3)))
    fig_rent.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
    fig_rent.update_layout(height=450, width=900, xaxis=dict(type="category"),
                           yaxis_tickformat=",",
                           legend=dict(orientation="h", y=-0.15),
                           title="Rentabilite")
    charts["rentabilite"] = fig_rent

    # 3. Cash flow cumule filled area
    fig_cf = go.Figure()
    fig_cf.add_trace(go.Scatter(x=df["date"], y=df["cash_flow_cumul"],
                                fill="tozeroy",
                                line=dict(color="#667eea", width=2),
                                name="Cash Flow Cumule"))
    fig_cf.add_hline(y=0, line_dash="dash", line_color="red")
    fig_cf.update_layout(height=450, width=900, yaxis_tickformat=",",
                         title="Cash Flow Cumule")
    charts["cashflow_cumul"] = fig_cf

    # 3b. BFR chart
    if "creances_clients" in df.columns:
        _bfr_ann_ch = df.groupby("annee").agg({
            "creances_clients": "last", "dettes_fournisseurs": "last", "bfr": "last"
        }).reset_index()
        _bfr_x_ch = [str(int(a)) for a in _bfr_ann_ch["annee"]]
        fig_bfr_ch = go.Figure()
        fig_bfr_ch.add_trace(go.Bar(x=_bfr_x_ch, y=_bfr_ann_ch["creances_clients"]/1000,
            name="Creances clients", marker_color="#4facfe"))
        fig_bfr_ch.add_trace(go.Bar(x=_bfr_x_ch, y=-_bfr_ann_ch["dettes_fournisseurs"]/1000,
            name="Dettes fournisseurs", marker_color="#f5576c"))
        fig_bfr_ch.add_trace(go.Scatter(x=_bfr_x_ch, y=_bfr_ann_ch["bfr"]/1000,
            name="BFR net", mode="lines+markers", line=dict(color="#667eea", width=3)))
        fig_bfr_ch.update_layout(height=450, width=900, barmode="relative",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f", title="K\u20ac"),
            title="Besoin en Fonds de Roulement (BFR)")
        fig_bfr_ch.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
        charts["bfr"] = fig_bfr_ch

    # 4. Waterfall chart (reference year = 4th year or last available)
    annees_cal = sorted(indic["Annee calendaire"].unique())
    idx_ref = min(3, len(annees_cal) - 1)
    an_ref = int(annees_cal[idx_ref])
    row_ref = indic[indic["Annee calendaire"] == an_ref].iloc[0]
    fig_wf = go.Figure(go.Waterfall(
        orientation="v",
        measure=["absolute", "relative", "relative", "total",
                 "relative", "total", "relative", "relative", "total"],
        x=["CA Total", "Ch. Var.", "Ch. Fixes", "EBITDA",
           "Amort.", "EBIT", "Interets", "Impot", "Res. Net"],
        y=[row_ref["CA Total"], -row_ref["Charges Variables"],
           -row_ref["Charges Fixes"], 0,
           -row_ref["Amortissement"], 0, -row_ref["Interets"],
           -(row_ref["EBIT"] - row_ref["Interets"] - row_ref["Resultat Net"]), 0],
        increasing={"marker": {"color": "#38ef7d"}},
        decreasing={"marker": {"color": "#f5576c"}},
        totals={"marker": {"color": "#667eea"}},
    ))
    fig_wf.update_layout(height=450, width=900, yaxis_tickformat=",",
                         title=f"Structure des couts ({an_ref})")
    charts["waterfall"] = fig_wf

    # 5. Taux occupation line chart
    fig_occ = go.Figure()
    fig_occ.add_trace(go.Scatter(x=x_labels, y=indic["Taux Occupation"],
                                 mode="lines+markers",
                                 line=dict(color="#667eea", width=3),
                                 name="Taux Occupation (%)"))
    fig_occ.update_layout(height=450, width=900,
                          xaxis=dict(type="category"),
                          yaxis=dict(title="%", ticksuffix="%"),
                          title="Taux d'Occupation")
    charts["occupation"] = fig_occ

    return charts


# ─── Export PPTX ─────────────────────────────────────────────────────────────

def _export_pptx(indic, params, df, chapters):
    """Generate a professional PPTX from the chapter structure."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import io

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    annees_cal = sorted(indic["Annee calendaire"].unique())
    annees5 = annees_cal[:min(5, len(annees_cal))]
    x_labels_an = [str(int(row["Annee calendaire"])) for _, row in indic.iterrows()]
    total_inv = sum(i["montant"] for i in params["investissements"])
    chart_figs = _make_presentation_charts(indic, params, df)

    def _add_title_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        if subtitle_text:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(18)
            p2.alignment = PP_ALIGN.CENTER
        return slide

    def _add_content_slide(title_text, content="", images=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        y_pos = Inches(1.2)
        if content:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12), Inches(3.0))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for line in content.split("\n"):
                pa = tf2.add_paragraph()
                pa.text = line
                pa.font.size = Pt(12)
            y_pos = Inches(4.5)
        if images:
            for img_data in images[:2]:
                img_bytes = img_data.get("bytes")
                if img_bytes:
                    remaining_h = Inches(7.5) - y_pos - Inches(0.3)
                    img_h = min(Inches(3.0), remaining_h)
                    if img_h > Inches(0.5):
                        try:
                            slide.shapes.add_picture(
                                io.BytesIO(img_bytes), Inches(0.5), y_pos,
                                Inches(12), img_h)
                            y_pos += img_h + Inches(0.2)
                        except Exception:
                            pass
        return slide

    def _add_table_slide(title_text, col_headers, rows_data, chart_png=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        n_rows = len(rows_data) + 1
        n_cols = len(col_headers)
        table_height = min(Inches(0.4 * n_rows), Inches(3.0))
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.2),
                                           Inches(12), table_height)
        tbl = tbl_shape.table
        for j, h in enumerate(col_headers):
            cell = tbl.cell(0, j)
            cell.text = str(h)
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(11)
                par.font.bold = True
        for i, row_vals in enumerate(rows_data):
            for j, val in enumerate(row_vals):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                for par in cell.text_frame.paragraphs:
                    par.font.size = Pt(10)
        if chart_png:
            chart_y = Inches(1.2) + table_height + Inches(0.3)
            remaining = Inches(7.5) - chart_y - Inches(0.3)
            chart_h = min(Inches(3.5), remaining)
            if chart_h > Inches(1.0):
                slide.shapes.add_picture(io.BytesIO(chart_png), Inches(0.5), chart_y,
                                         Inches(12), chart_h)
        return slide

    # -- Ordered chapter keys (non-financial first, then financial) --
    nf_keys = [k for k, _, _, _ in _NON_FINANCIAL_CHAPTERS]
    f_keys = [k for k, _ in _FINANCIAL_CHAPTERS]

    for ch_key in nf_keys + f_keys:
        ch = chapters.get(ch_key)
        if not ch or not ch.get("included", False):
            continue

        label = ch["label"]

        # ── Non-financial chapters ──
        if ch.get("type") == "non_financial":
            if ch_key == "page_de_garde":
                _add_title_slide(
                    f"Plan Financier - {params.get('nom_hotel', 'Hotel')}",
                    f"{params['nb_chambres']} chambres | Ouverture : {params['date_ouverture'].strftime('%B %Y')}"
                )
            else:
                _add_content_slide(label, ch.get("content", ""), ch.get("images", []))
            continue

        # ── Financial chapters ──
        if ch_key == "invest_financement":
            inv_headers = ["Categorie", "Montant", "Duree Amort.", "Amort. Annuel"]
            inv_rows = []
            for inv in params["investissements"]:
                amort = inv["montant"] / inv["duree_amort"] if inv["duree_amort"] > 0 else 0
                inv_rows.append([inv["categorie"], f"{inv['montant']:,.0f} \u20ac",
                                 f"{inv['duree_amort']} ans", f"{amort:,.0f} \u20ac"])
            inv_rows.append(["TOTAL", f"{total_inv:,.0f} \u20ac", "", ""])
            _add_table_slide("Investissements", inv_headers, inv_rows)

            fin_headers = ["Source", "Montant", "Taux", "Duree"]
            fin_rows = [["Fonds propres", f"{params['fonds_propres_initial']:,.0f} \u20ac", "-", "-"]]
            for pret in params["prets"]:
                if pret["montant"] > 0:
                    fin_rows.append([pret["nom"], f"{pret['montant']:,.0f} \u20ac",
                                     f"{pret['taux_annuel']*100:.2f}%",
                                     f"{pret['duree_ans']} ans"])
            _add_table_slide("Plan de Financement", fin_headers, fin_rows)

        elif ch_key == "hypotheses_cles":
            hyp_headers = ["Parametre", "Valeur"]
            hyp_rows = [
                ["Nombre de chambres", str(params["nb_chambres"])],
                ["Date d'ouverture", params["date_ouverture"].strftime("%B %Y")],
                ["Investissement total", f"{total_inv:,.0f} \u20ac"],
                ["Fonds propres", f"{params['fonds_propres_initial']:,.0f} \u20ac"],
            ]
            for seg_name, seg_data in params.get("segments", {}).items():
                hyp_rows.append([f"Segment {seg_name}",
                                 f"Part {seg_data['part']*100:.0f}% | ADR {seg_data['prix']:,.0f} \u20ac"])
            taux_occ = params.get("taux_occ", [])
            for i, t in enumerate(taux_occ):
                yr_label = f"Annee {i+1}" if i < len(taux_occ) - 1 else f"Annee {i+1}+"
                hyp_rows.append([f"Taux occupation {yr_label}", f"{t*100:.1f}%"])
            saison = params.get("saisonnalite", [])
            if saison:
                mois_noms = ["Jan", "Fev", "Mar", "Avr", "Mai", "Jun",
                             "Jul", "Aou", "Sep", "Oct", "Nov", "Dec"]
                max_s = max(saison)
                min_s = min(saison)
                max_m = mois_noms[saison.index(max_s)] if len(saison) == 12 else "?"
                min_m = mois_noms[saison.index(min_s)] if len(saison) == 12 else "?"
                hyp_rows.append(["Saisonnalite (pic)", f"{max_m} ({max_s:.3f})"])
                hyp_rows.append(["Saisonnalite (creux)", f"{min_m} ({min_s:.3f})"])
            # Salles & Evenements
            hyp_rows.append(["--- Salles & Evenements ---", ""])
            hyp_rows.append(["Seminaires / an", str(params.get("seminaire_nb_an", 0))])
            hyp_rows.append(["Location salle seminaire", f"{params.get('seminaire_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Mariages / an", str(params.get("mariage_nb_an", 0))])
            hyp_rows.append(["Location salle mariage", f"{params.get('mariage_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Salles chateau / an", str(params.get("salles_chateau_nb_an", 0))])
            hyp_rows.append(["Prix moyen location chateau", f"{params.get('salles_chateau_prix', 0):,.0f} EUR"])
            hyp_rows.append(["CV chateau - Energie", f"{params.get('cv_salles_chateau_energie', 100):,.0f} EUR/loc."])
            hyp_rows.append(["CV chateau - Nettoyage", f"{params.get('cv_salles_chateau_nettoyage', 500):,.0f} EUR/loc."])
            _add_table_slide("Hypotheses Cles", hyp_headers, hyp_rows)

        elif ch_key == "previsions_ca":
            ca_headers = ["Annee"] + [str(int(a)) for a in annees5]
            ca_rows = []
            for field, lbl in [("CA Total", "CA Total"), ("CA Hebergement", "CA Hebergement"),
                                ("CA Brasserie", "CA Brasserie"), ("CA Autres", "CA Autres")]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(f"{d.iloc[0][field]:,.0f} \u20ac" if len(d) > 0 else "")
                ca_rows.append(row_vals)
            ca_png = _try_fig_to_png(chart_figs["ca_evolution"])
            _add_table_slide("Previsions de CA", ca_headers, ca_rows, ca_png)

        elif ch_key == "structure_couts":
            idx_ref = min(3, len(annees_cal) - 1)
            an_ref = int(annees_cal[idx_ref])
            row_ref = indic[indic["Annee calendaire"] == an_ref].iloc[0]
            wf_headers = ["Poste", "Montant"]
            wf_rows = [
                ["CA Total", f"{row_ref['CA Total']:,.0f} \u20ac"],
                ["Charges Variables", f"-{row_ref['Charges Variables']:,.0f} \u20ac"],
                ["Charges Fixes", f"-{row_ref['Charges Fixes']:,.0f} \u20ac"],
                ["EBITDA", f"{row_ref['EBITDA']:,.0f} \u20ac"],
                ["Amortissement", f"-{row_ref['Amortissement']:,.0f} \u20ac"],
                ["EBIT", f"{row_ref['EBIT']:,.0f} \u20ac"],
                ["Interets", f"-{row_ref['Interets']:,.0f} \u20ac"],
                ["Resultat Net", f"{row_ref['Resultat Net']:,.0f} \u20ac"],
            ]
            wf_png = _try_fig_to_png(chart_figs["waterfall"])
            _add_table_slide(f"Structure des Couts ({an_ref})", wf_headers, wf_rows, wf_png)

        elif ch_key == "indicateurs_hoteliers":
            ki_headers = ["Indicateur"] + [str(int(a)) for a in annees5]
            ki_rows = []
            for field, lbl, fmt_fn in [
                ("Taux Occupation", "Taux Occupation", lambda v: f"{v:.1f}%"),
                ("Prix Moyen (ADR)", "ADR", lambda v: f"{v:,.0f} \u20ac"),
                ("RevPAR", "RevPAR", lambda v: f"{v:,.0f} \u20ac"),
                ("GOPPAR", "GOPPAR", lambda v: f"{v:,.0f} \u20ac"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                ki_rows.append(row_vals)
            occ_png = _try_fig_to_png(chart_figs["occupation"])
            _add_table_slide("Indicateurs Hoteliers", ki_headers, ki_rows, occ_png)

        elif ch_key == "rentabilite_resultat":
            prof_headers = ["Indicateur"] + [str(int(a)) for a in annees5]
            prof_rows = []
            for field, lbl, fmt_fn in [
                ("EBITDA", "EBITDA", lambda v: f"{v:,.0f} \u20ac"),
                ("EBITDA %", "Marge EBITDA", lambda v: f"{v:.1f}%"),
                ("Resultat Net", "Resultat Net", lambda v: f"{v:,.0f} \u20ac"),
                ("Cash Flow", "Cash Flow", lambda v: f"{v:,.0f} \u20ac"),
                ("Cash Flow Cumul", "Cash Flow Cumule", lambda v: f"{v:,.0f} \u20ac"),
                ("DSCR", "DSCR", lambda v: f"{v:.2f}x" if v < 100 else "\u221e"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                prof_rows.append(row_vals)
            rent_png = _try_fig_to_png(chart_figs["rentabilite"])
            _add_table_slide("Rentabilite & Resultat", prof_headers, prof_rows, rent_png)

        elif ch_key == "cashflow_roi":
            cf_cumul = df["cash_flow_cumul"].values
            break_even_idx = None
            for i in range(1, len(cf_cumul)):
                if cf_cumul[i - 1] < 0 and cf_cumul[i] >= 0:
                    break_even_idx = i
                    break
            be_text = ""
            if break_even_idx:
                be_date = df.iloc[break_even_idx]["date"]
                be_text = f"Break-even atteint en {be_date.strftime('%B %Y')} (mois {break_even_idx})"
            else:
                be_text = "Break-even non atteint sur la periode de projection"
            cf_png = _try_fig_to_png(chart_figs["cashflow_cumul"])
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
            tf = txBox.text_frame
            pa = tf.paragraphs[0]
            pa.text = "Cash Flow & Retour sur Investissement"
            pa.font.size = Pt(24)
            pa.font.bold = True
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
            tf2 = txBox2.text_frame
            pa2 = tf2.paragraphs[0]
            pa2.text = be_text
            pa2.font.size = Pt(14)
            if cf_png:
                slide.shapes.add_picture(io.BytesIO(cf_png), Inches(0.5), Inches(1.8),
                                         Inches(12), Inches(5))

        elif ch_key == "bilan":
            bfr_png = _try_fig_to_png(chart_figs.get("bfr")) if "bfr" in chart_figs else None
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
            pa = txBox.text_frame.paragraphs[0]
            pa.text = "Bilan"
            pa.font.size = Pt(24)
            pa.font.bold = True
            if bfr_png:
                slide.shapes.add_picture(io.BytesIO(bfr_png), Inches(0.5), Inches(1.2),
                                         Inches(12), Inches(5.5))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─── Export HTML (rapport imprimable) ────────────────────────────────────────

def _export_html(indic, params, df, chapters):
    """Generate a self-contained HTML report with print-optimized CSS and interactive Plotly charts."""
    import base64 as _b64_html
    import plotly.io as pio

    annees_cal = sorted(indic["Annee calendaire"].unique())
    annees5 = annees_cal[:min(5, len(annees_cal))]
    total_inv = sum(i["montant"] for i in params["investissements"])
    chart_figs = _make_presentation_charts(indic, params, df)

    # ── Helper: build HTML table ──
    def _html_table(headers, rows):
        h = "<table><thead><tr>" + "".join(f"<th>{c}</th>" for c in headers) + "</tr></thead><tbody>"
        for row in rows:
            h += "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>"
        h += "</tbody></table>"
        return h

    # ── Helper: embed plotly chart ──
    def _html_chart(fig):
        if fig is None:
            return ""
        fig_copy = fig.full_figure_for_development(warn=False)
        fig_copy.update_layout(width=None, height=400, margin=dict(l=40, r=40, t=50, b=40))
        return (
            '<div class="chart-container">'
            + pio.to_html(fig_copy, full_html=False, include_plotlyjs=False, config={"displayModeBar": False})
            + '</div>'
        )

    # ── Helper: embed user images ──
    def _html_user_images(images):
        html = ""
        for img_data in (images or [])[:2]:
            img_bytes = img_data.get("bytes")
            if img_bytes:
                data_uri = _b64_html.b64encode(img_bytes).decode()
                html += f'<img class="user-image" src="data:image/jpeg;base64,{data_uri}">'
        return html

    # ── Build pages ──
    pages = []
    nf_keys = [k for k, _, _, _ in _NON_FINANCIAL_CHAPTERS]
    f_keys = [k for k, _ in _FINANCIAL_CHAPTERS]

    for ch_key in nf_keys + f_keys:
        ch = chapters.get(ch_key)
        if not ch or not ch.get("included", False):
            continue

        label = ch["label"]

        if ch.get("type") == "non_financial":
            if ch_key == "page_de_garde":
                # Page de garde avec photo
                cover_img = ""
                cover_path = _Path(__file__).parent / "assets" / "chateau_1.jpg"
                if cover_path.exists():
                    with open(str(cover_path), "rb") as _fci:
                        cover_img = (
                            f'<img src="data:image/jpeg;base64,{_b64_html.b64encode(_fci.read()).decode()}" '
                            f'class="cover-image">'
                        )
                pages.append(
                    f'<section class="page cover-page">'
                    f'{cover_img}'
                    f'<h1 class="cover-title">Plan Financier</h1>'
                    f'<h2 class="cover-subtitle">{params.get("nom_hotel", "Hotel")}</h2>'
                    f'<p class="cover-detail">{params["nb_chambres"]} chambres &mdash; '
                    f'Ouverture : {params["date_ouverture"].strftime("%B %Y")}</p>'
                    f'</section>'
                )
            else:
                content_html = ""
                content = ch.get("content", "")
                if content:
                    content_html = "<div class='chapter-text'>" + content.replace("\n", "<br>") + "</div>"
                content_html += _html_user_images(ch.get("images", []))
                pages.append(
                    f'<section class="page">'
                    f'<h2 class="chapter-title">{label}</h2>'
                    f'{content_html}'
                    f'</section>'
                )
            continue

        # ── Financial chapters ──
        body = ""

        if ch_key == "invest_financement":
            body += "<h3>Investissements</h3>"
            inv_rows = []
            for inv in params["investissements"]:
                amort = inv["montant"] / inv["duree_amort"] if inv["duree_amort"] > 0 else 0
                inv_rows.append([inv["categorie"], f'{inv["montant"]:,.0f} \u20ac',
                                 f'{inv["duree_amort"]} ans', f'{amort:,.0f} \u20ac'])
            body += _html_table(["Categorie", "Montant", "Duree Amort.", "Amort. Annuel"], inv_rows)

            body += "<h3>Plan de Financement</h3>"
            fin_rows = [["Fonds propres", f'{params["fonds_propres_initial"]:,.0f} \u20ac', "-", "-"]]
            for pret in params["prets"]:
                if pret["montant"] > 0:
                    fin_rows.append([pret["nom"], f'{pret["montant"]:,.0f} \u20ac',
                                     f'{pret["taux_annuel"]*100:.2f}%', f'{pret["duree_ans"]} ans'])
            body += _html_table(["Source", "Montant", "Taux", "Duree"], fin_rows)

        elif ch_key == "hypotheses_cles":
            hyp_rows = [
                ["Nombre de chambres", str(params["nb_chambres"])],
                ["Date d'ouverture", params["date_ouverture"].strftime("%B %Y")],
                ["Investissement total", f"{total_inv:,.0f} \u20ac"],
                ["Fonds propres", f'{params["fonds_propres_initial"]:,.0f} \u20ac'],
            ]
            for seg_name, seg_data in params.get("segments", {}).items():
                hyp_rows.append([f"Segment {seg_name}",
                                 f'Part {seg_data["part"]*100:.0f}% | ADR {seg_data["prix"]:,.0f} \u20ac'])
            taux_occ = params.get("taux_occ", [])
            for i, t in enumerate(taux_occ):
                yr_label = f"Annee {i+1}" if i < len(taux_occ) - 1 else f"Annee {i+1}+"
                hyp_rows.append([f"Taux occupation {yr_label}", f"{t*100:.1f}%"])
            # Salles & Evenements
            hyp_rows.append(["--- Salles & Evenements ---", ""])
            hyp_rows.append(["Seminaires / an", str(params.get("seminaire_nb_an", 0))])
            hyp_rows.append(["Location salle seminaire", f"{params.get('seminaire_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Mariages / an", str(params.get("mariage_nb_an", 0))])
            hyp_rows.append(["Location salle mariage", f"{params.get('mariage_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Salles chateau / an", str(params.get("salles_chateau_nb_an", 0))])
            hyp_rows.append(["Prix moyen location chateau", f"{params.get('salles_chateau_prix', 0):,.0f} EUR"])
            hyp_rows.append(["CV chateau - Energie", f"{params.get('cv_salles_chateau_energie', 100):,.0f} EUR/loc."])
            hyp_rows.append(["CV chateau - Nettoyage", f"{params.get('cv_salles_chateau_nettoyage', 500):,.0f} EUR/loc."])
            body += _html_table(["Parametre", "Valeur"], hyp_rows)

        elif ch_key == "previsions_ca":
            ca_rows = []
            for field, lbl in [("CA Total", "CA Total"), ("CA Hebergement", "CA Hebergement"),
                                ("CA Brasserie", "CA Brasserie"), ("CA Autres", "CA Autres")]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(f"{d.iloc[0][field]:,.0f} \u20ac" if len(d) > 0 else "")
                ca_rows.append(row_vals)
            body += _html_table(["Poste"] + [str(int(a)) for a in annees5], ca_rows)
            body += _html_chart(chart_figs.get("ca_evolution"))

        elif ch_key == "structure_couts":
            idx_ref = min(3, len(annees_cal) - 1)
            an_ref = int(annees_cal[idx_ref])
            row_ref = indic[indic["Annee calendaire"] == an_ref].iloc[0]
            wf_rows = [
                ["CA Total", f"{row_ref['CA Total']:,.0f} \u20ac"],
                ["Charges Variables", f"-{row_ref['Charges Variables']:,.0f} \u20ac"],
                ["Charges Fixes", f"-{row_ref['Charges Fixes']:,.0f} \u20ac"],
                ["EBITDA", f"{row_ref['EBITDA']:,.0f} \u20ac"],
                ["Amortissement", f"-{row_ref['Amortissement']:,.0f} \u20ac"],
                ["EBIT", f"{row_ref['EBIT']:,.0f} \u20ac"],
                ["Interets", f"-{row_ref['Interets']:,.0f} \u20ac"],
                ["Resultat Net", f"{row_ref['Resultat Net']:,.0f} \u20ac"],
            ]
            body += f"<p class='ref-year'>Annee de reference : {an_ref}</p>"
            body += _html_table(["Poste", "Montant"], wf_rows)
            body += _html_chart(chart_figs.get("waterfall"))

        elif ch_key == "indicateurs_hoteliers":
            ki_rows = []
            for field, lbl, fmt_fn in [
                ("Taux Occupation", "Taux Occupation", lambda v: f"{v:.1f}%"),
                ("Prix Moyen (ADR)", "ADR", lambda v: f"{v:,.0f} \u20ac"),
                ("RevPAR", "RevPAR", lambda v: f"{v:,.0f} \u20ac"),
                ("GOPPAR", "GOPPAR", lambda v: f"{v:,.0f} \u20ac"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                ki_rows.append(row_vals)
            body += _html_table(["Indicateur"] + [str(int(a)) for a in annees5], ki_rows)
            body += _html_chart(chart_figs.get("occupation"))

        elif ch_key == "rentabilite_resultat":
            prof_rows = []
            for field, lbl, fmt_fn in [
                ("EBITDA", "EBITDA", lambda v: f"{v:,.0f} \u20ac"),
                ("EBITDA %", "Marge EBITDA", lambda v: f"{v:.1f}%"),
                ("Resultat Net", "Resultat Net", lambda v: f"{v:,.0f} \u20ac"),
                ("Cash Flow", "Cash Flow", lambda v: f"{v:,.0f} \u20ac"),
                ("Cash Flow Cumul", "Cash Flow Cumule", lambda v: f"{v:,.0f} \u20ac"),
                ("DSCR", "DSCR", lambda v: f"{v:.2f}x" if v < 100 else "inf"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                prof_rows.append(row_vals)
            body += _html_table(["Indicateur"] + [str(int(a)) for a in annees5], prof_rows)
            body += _html_chart(chart_figs.get("rentabilite"))

        elif ch_key == "cashflow_roi":
            cf_cumul = df["cash_flow_cumul"].values
            break_even_idx = None
            for i in range(1, len(cf_cumul)):
                if cf_cumul[i - 1] < 0 and cf_cumul[i] >= 0:
                    break_even_idx = i
                    break
            if break_even_idx:
                be_date = df.iloc[break_even_idx]["date"]
                body += f'<p class="break-even"><strong>Break-even atteint en {be_date.strftime("%B %Y")} (mois {break_even_idx})</strong></p>'
            else:
                body += '<p class="break-even">Break-even non atteint sur la periode de projection</p>'
            body += _html_chart(chart_figs.get("cashflow_cumul"))

        elif ch_key == "bilan":
            body += _html_chart(chart_figs.get("bfr"))

        pages.append(
            f'<section class="page">'
            f'<h2 class="chapter-title">{label}</h2>'
            f'{body}'
            f'</section>'
        )

    # ── Assemble full HTML ──
    all_pages = "\n".join(pages)
    html = f"""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Plan Financier - {params.get('nom_hotel', 'Hotel')}</title>
<script src="https://cdn.plot.ly/plotly-2.35.0.min.js"></script>
<style>
  /* ── General ── */
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    color: #2c3e50; background: #f5f6fa; line-height: 1.6;
  }}
  .page {{
    background: white; max-width: 1100px; margin: 30px auto;
    padding: 50px 60px; border-radius: 8px;
    box-shadow: 0 2px 12px rgba(0,0,0,0.07);
  }}

  /* ── Cover page ── */
  .cover-page {{
    text-align: center; padding: 60px;
    background: linear-gradient(135deg, #1a1a2e 0%, #16213e 50%, #0f3460 100%);
    color: white; min-height: 500px; display: flex; flex-direction: column;
    justify-content: center; align-items: center;
  }}
  .cover-image {{
    width: 350px; max-width: 80%; height: auto; border-radius: 12px;
    margin-bottom: 30px; box-shadow: 0 4px 20px rgba(0,0,0,0.3);
  }}
  .cover-title {{
    font-size: 2.8em; font-weight: 700; margin-bottom: 8px;
    letter-spacing: 1px;
  }}
  .cover-subtitle {{
    font-size: 1.6em; font-weight: 400; color: #a8d8ea;
    margin-bottom: 16px;
  }}
  .cover-detail {{
    font-size: 1.1em; color: rgba(255,255,255,0.7);
  }}

  /* ── Chapter titles ── */
  .chapter-title {{
    font-size: 1.8em; color: #0f3460; border-bottom: 3px solid #4facfe;
    padding-bottom: 10px; margin-bottom: 24px;
  }}
  h3 {{
    font-size: 1.2em; color: #34495e; margin: 24px 0 12px 0;
  }}

  /* ── Tables ── */
  table {{
    width: 100%; border-collapse: collapse; margin: 16px 0 24px 0;
    font-size: 0.92em;
  }}
  thead th {{
    background: #0f3460; color: white; padding: 10px 14px;
    text-align: left; font-weight: 600;
  }}
  tbody td {{
    padding: 8px 14px; border-bottom: 1px solid #e9ecef;
  }}
  tbody tr:nth-child(even) {{ background: #f8f9fa; }}
  tbody tr:hover {{ background: #e8f4fd; }}

  /* ── Charts ── */
  .chart-container {{
    margin: 20px 0; page-break-inside: avoid;
  }}
  .chart-container .plotly-graph-div {{
    width: 100% !important;
  }}

  /* ── Text blocks ── */
  .chapter-text {{
    font-size: 1em; line-height: 1.7; margin-bottom: 20px;
    white-space: pre-wrap;
  }}
  .user-image {{
    max-width: 70%; height: auto; border-radius: 8px;
    margin: 12px 0; display: block;
  }}
  .ref-year {{
    font-style: italic; color: #666; margin-bottom: 8px;
  }}
  .break-even {{
    font-size: 1.15em; padding: 14px 20px; background: #e8f8f0;
    border-left: 4px solid #38ef7d; border-radius: 4px; margin: 12px 0;
  }}

  /* ── Print styles ── */
  @media print {{
    body {{ background: white; }}
    .page {{
      box-shadow: none; border-radius: 0; margin: 0; padding: 35px 40px;
      max-width: 100%; page-break-after: always;
    }}
    .page:last-child {{ page-break-after: auto; }}
    .cover-page {{ min-height: 100vh; }}
    table {{ page-break-inside: avoid; }}
    .chart-container {{ page-break-inside: avoid; }}
    h2.chapter-title {{ page-break-after: avoid; }}
    h3 {{ page-break-after: avoid; }}
  }}
  @page {{
    size: A4 landscape; margin: 15mm;
  }}
</style>
</head>
<body>
{all_pages}
</body>
</html>"""
    return html


# ─── Export PDF ──────────────────────────────────────────────────────────────

def _export_pdf(indic, params, df, chapters):
    """Generate a professional PDF (landscape A4) from the chapter structure."""
    from fpdf import FPDF
    import io

    annees_cal = sorted(indic["Annee calendaire"].unique())
    annees5 = annees_cal[:min(5, len(annees_cal))]
    total_inv = sum(i["montant"] for i in params["investissements"])
    chart_figs = _make_presentation_charts(indic, params, df)

    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)

    def _add_title(text):
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 12, text, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)

    def _add_table(headers, rows, col_widths=None):
        if col_widths is None:
            n = len(headers)
            w = 260 / n
            col_widths = [w] * n
        pdf.set_font("Helvetica", "B", 9)
        for j, h in enumerate(headers):
            pdf.cell(col_widths[j], 7, str(h), border=1)
        pdf.ln()
        pdf.set_font("Helvetica", "", 9)
        for row_vals in rows:
            for j, val in enumerate(row_vals):
                pdf.cell(col_widths[j], 6, str(val), border=1)
            pdf.ln()
        pdf.ln(4)

    def _add_chart_image(fig_png):
        if fig_png:
            try:
                pdf.image(io.BytesIO(fig_png), x=pdf.get_x(), y=pdf.get_y(), w=250)
                pdf.ln(5)
            except Exception:
                pass

    def _add_user_images(images):
        for img_data in (images or [])[:2]:
            img_bytes = img_data.get("bytes")
            if img_bytes:
                try:
                    pdf.image(io.BytesIO(img_bytes), x=pdf.get_x(), y=pdf.get_y(), w=200)
                    pdf.ln(5)
                except Exception:
                    pass

    nf_keys = [k for k, _, _, _ in _NON_FINANCIAL_CHAPTERS]
    f_keys = [k for k, _ in _FINANCIAL_CHAPTERS]

    for ch_key in nf_keys + f_keys:
        ch = chapters.get(ch_key)
        if not ch or not ch.get("included", False):
            continue

        label = ch["label"]

        # ── Non-financial chapters ──
        if ch.get("type") == "non_financial":
            pdf.add_page()
            if ch_key == "page_de_garde":
                pdf.set_font("Helvetica", "B", 28)
                pdf.ln(50)
                pdf.cell(0, 15, f"Plan Financier - {params.get('nom_hotel', 'Hotel')}",
                         new_x="LMARGIN", new_y="NEXT", align="C")
                pdf.set_font("Helvetica", "", 14)
                pdf.cell(0, 10,
                         f"{params['nb_chambres']} chambres | Ouverture : {params['date_ouverture'].strftime('%B %Y')}",
                         new_x="LMARGIN", new_y="NEXT", align="C")
            else:
                _add_title(label)
                content = ch.get("content", "")
                if content:
                    pdf.set_font("Helvetica", "", 11)
                    pdf.multi_cell(0, 6, content)
                    pdf.ln(4)
                _add_user_images(ch.get("images", []))
            continue

        # ── Financial chapters ──
        pdf.add_page()

        if ch_key == "invest_financement":
            _add_title("Investissements")
            inv_headers = ["Categorie", "Montant", "Duree Amort.", "Amort. Annuel"]
            inv_rows = []
            for inv in params["investissements"]:
                amort = inv["montant"] / inv["duree_amort"] if inv["duree_amort"] > 0 else 0
                inv_rows.append([inv["categorie"], f"{inv['montant']:,.0f} EUR",
                                 f"{inv['duree_amort']} ans", f"{amort:,.0f} EUR"])
            _add_table(inv_headers, inv_rows)

            _add_title("Plan de Financement")
            fin_headers = ["Source", "Montant", "Taux", "Duree"]
            fin_rows = [["Fonds propres", f"{params['fonds_propres_initial']:,.0f} EUR", "-", "-"]]
            for pret in params["prets"]:
                if pret["montant"] > 0:
                    fin_rows.append([pret["nom"], f"{pret['montant']:,.0f} EUR",
                                     f"{pret['taux_annuel']*100:.2f}%",
                                     f"{pret['duree_ans']} ans"])
            _add_table(fin_headers, fin_rows)

        elif ch_key == "hypotheses_cles":
            _add_title("Hypotheses Cles")
            hyp_headers = ["Parametre", "Valeur"]
            hyp_rows = [
                ["Nombre de chambres", str(params["nb_chambres"])],
                ["Date d'ouverture", params["date_ouverture"].strftime("%B %Y")],
                ["Investissement total", f"{total_inv:,.0f} EUR"],
                ["Fonds propres", f"{params['fonds_propres_initial']:,.0f} EUR"],
            ]
            for seg_name, seg_data in params.get("segments", {}).items():
                hyp_rows.append([f"Segment {seg_name}",
                                 f"Part {seg_data['part']*100:.0f}% | ADR {seg_data['prix']:,.0f} EUR"])
            taux_occ = params.get("taux_occ", [])
            for i, t in enumerate(taux_occ):
                yr_label = f"Annee {i+1}" if i < len(taux_occ) - 1 else f"Annee {i+1}+"
                hyp_rows.append([f"Taux occupation {yr_label}", f"{t*100:.1f}%"])
            # Salles & Evenements
            hyp_rows.append(["--- Salles & Evenements ---", ""])
            hyp_rows.append(["Seminaires / an", str(params.get("seminaire_nb_an", 0))])
            hyp_rows.append(["Location salle seminaire", f"{params.get('seminaire_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Mariages / an", str(params.get("mariage_nb_an", 0))])
            hyp_rows.append(["Location salle mariage", f"{params.get('mariage_prix_location', 0):,.0f} EUR"])
            hyp_rows.append(["Salles chateau / an", str(params.get("salles_chateau_nb_an", 0))])
            hyp_rows.append(["Prix moyen location chateau", f"{params.get('salles_chateau_prix', 0):,.0f} EUR"])
            hyp_rows.append(["CV chateau - Energie", f"{params.get('cv_salles_chateau_energie', 100):,.0f} EUR/loc."])
            hyp_rows.append(["CV chateau - Nettoyage", f"{params.get('cv_salles_chateau_nettoyage', 500):,.0f} EUR/loc."])
            _add_table(hyp_headers, hyp_rows)

        elif ch_key == "previsions_ca":
            _add_title("Previsions de CA")
            ca_headers = ["Poste"] + [str(int(a)) for a in annees5]
            ca_rows = []
            for field, lbl in [("CA Total", "CA Total"), ("CA Hebergement", "CA Hebergement"),
                                ("CA Brasserie", "CA Brasserie"), ("CA Autres", "CA Autres")]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(f"{d.iloc[0][field]:,.0f} EUR" if len(d) > 0 else "")
                ca_rows.append(row_vals)
            _add_table(ca_headers, ca_rows)
            _add_chart_image(_try_fig_to_png(chart_figs["ca_evolution"]))

        elif ch_key == "structure_couts":
            idx_ref = min(3, len(annees_cal) - 1)
            an_ref = int(annees_cal[idx_ref])
            row_ref = indic[indic["Annee calendaire"] == an_ref].iloc[0]
            _add_title(f"Structure des Couts ({an_ref})")
            wf_headers = ["Poste", "Montant"]
            wf_rows = [
                ["CA Total", f"{row_ref['CA Total']:,.0f} EUR"],
                ["Charges Variables", f"-{row_ref['Charges Variables']:,.0f} EUR"],
                ["Charges Fixes", f"-{row_ref['Charges Fixes']:,.0f} EUR"],
                ["EBITDA", f"{row_ref['EBITDA']:,.0f} EUR"],
                ["Amortissement", f"-{row_ref['Amortissement']:,.0f} EUR"],
                ["EBIT", f"{row_ref['EBIT']:,.0f} EUR"],
                ["Interets", f"-{row_ref['Interets']:,.0f} EUR"],
                ["Resultat Net", f"{row_ref['Resultat Net']:,.0f} EUR"],
            ]
            _add_table(wf_headers, wf_rows)
            _add_chart_image(_try_fig_to_png(chart_figs["waterfall"]))

        elif ch_key == "indicateurs_hoteliers":
            _add_title("Indicateurs Hoteliers")
            ki_headers = ["Indicateur"] + [str(int(a)) for a in annees5]
            ki_rows = []
            for field, lbl, fmt_fn in [
                ("Taux Occupation", "Taux Occupation", lambda v: f"{v:.1f}%"),
                ("Prix Moyen (ADR)", "ADR", lambda v: f"{v:,.0f} EUR"),
                ("RevPAR", "RevPAR", lambda v: f"{v:,.0f} EUR"),
                ("GOPPAR", "GOPPAR", lambda v: f"{v:,.0f} EUR"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                ki_rows.append(row_vals)
            _add_table(ki_headers, ki_rows)
            _add_chart_image(_try_fig_to_png(chart_figs["occupation"]))

        elif ch_key == "rentabilite_resultat":
            _add_title("Rentabilite & Resultat")
            prof_headers = ["Indicateur"] + [str(int(a)) for a in annees5]
            prof_rows = []
            for field, lbl, fmt_fn in [
                ("EBITDA", "EBITDA", lambda v: f"{v:,.0f} EUR"),
                ("EBITDA %", "Marge EBITDA", lambda v: f"{v:.1f}%"),
                ("Resultat Net", "Resultat Net", lambda v: f"{v:,.0f} EUR"),
                ("Cash Flow", "Cash Flow", lambda v: f"{v:,.0f} EUR"),
                ("Cash Flow Cumul", "Cash Flow Cumule", lambda v: f"{v:,.0f} EUR"),
                ("DSCR", "DSCR", lambda v: f"{v:.2f}x" if v < 100 else "inf"),
            ]:
                row_vals = [lbl]
                for a in annees5:
                    d = indic[indic["Annee calendaire"] == a]
                    row_vals.append(fmt_fn(d.iloc[0][field]) if len(d) > 0 else "")
                prof_rows.append(row_vals)
            _add_table(prof_headers, prof_rows)
            _add_chart_image(_try_fig_to_png(chart_figs["rentabilite"]))

        elif ch_key == "cashflow_roi":
            _add_title("Cash Flow & Retour sur Investissement")
            cf_cumul = df["cash_flow_cumul"].values
            break_even_idx = None
            for i in range(1, len(cf_cumul)):
                if cf_cumul[i - 1] < 0 and cf_cumul[i] >= 0:
                    break_even_idx = i
                    break
            if break_even_idx:
                be_date = df.iloc[break_even_idx]["date"]
                pdf.set_font("Helvetica", "B", 12)
                pdf.cell(0, 8, f"Break-even atteint en {be_date.strftime('%B %Y')} (mois {break_even_idx})",
                         new_x="LMARGIN", new_y="NEXT")
            else:
                pdf.set_font("Helvetica", "", 12)
                pdf.cell(0, 8, "Break-even non atteint sur la periode de projection",
                         new_x="LMARGIN", new_y="NEXT")
            pdf.ln(4)
            _add_chart_image(_try_fig_to_png(chart_figs["cashflow_cumul"]))

        elif ch_key == "bilan":
            _add_title("Bilan")
            if "bfr" in chart_figs:
                _add_chart_image(_try_fig_to_png(chart_figs["bfr"]))

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf.getvalue()


# ─── Tab Presentation UI ────────────────────────────────────────────────────

def tab_presentation(df, indic, params):
    st.header("Presentation investisseur")

    annees_cal = sorted(indic["Annee calendaire"].unique())
    annees5 = annees_cal[:min(5, len(annees_cal))]
    n_annees = len(indic)

    if not annees5:
        st.warning("Pas assez de donnees pour la presentation")
        return

    _init_pres_chapters()
    chapters = st.session_state["pres_chapters"]

    # ── Section: Structure du document ────────────────────────────────────
    st.subheader("Structure du document")

    col_edit, col_preview = st.columns([3, 2])

    with col_edit:
        # --- Non-financial chapters ---
        st.markdown('<div class="section-header">Chapitres non-financiers</div>',
                    unsafe_allow_html=True)
        st.caption("Cochez pour inclure, puis remplissez le contenu et ajoutez des images.")

        for ch_key, ch_label, has_text, has_image in _NON_FINANCIAL_CHAPTERS:
            ch = chapters[ch_key]
            ch["included"] = st.checkbox(
                ch_label, value=ch.get("included", True),
                key=f"pres_inc_{ch_key}")

            if ch["included"] and (has_text or has_image):
                with st.expander(f"Contenu : {ch_label}", expanded=False):
                    if has_text:
                        ch["content"] = st.text_area(
                            "Texte", value=ch.get("content", ""),
                            height=150, key=f"pres_txt_{ch_key}",
                            help=f"Contenu pour la section '{ch_label}'")
                    if has_image:
                        uploaded = st.file_uploader(
                            "Image(s)", type=["png", "jpg", "jpeg", "gif", "bmp"],
                            accept_multiple_files=True,
                            key=f"pres_img_{ch_key}",
                            help="Photos, schemas, organigrammes...")
                        if uploaded:
                            ch["images"] = [{"name": f.name, "bytes": f.read()} for f in uploaded]

        st.markdown("---")

        # --- Financial chapters ---
        st.markdown('<div class="section-header">Chapitres financiers (auto-generes)</div>',
                    unsafe_allow_html=True)
        st.caption("Contenu genere automatiquement a partir des donnees du plan financier.")

        for ch_key, ch_label in _FINANCIAL_CHAPTERS:
            ch = chapters[ch_key]
            ch["included"] = st.checkbox(
                ch_label, value=ch.get("included", True),
                key=f"pres_inc_{ch_key}")

    # ── Preview column ───────────────────────────────────────────────────
    with col_preview:
        st.markdown('<div class="section-header">Apercu de la structure</div>',
                    unsafe_allow_html=True)

        nf_keys = [k for k, _, _, _ in _NON_FINANCIAL_CHAPTERS]
        f_keys = [k for k, _ in _FINANCIAL_CHAPTERS]
        num = 1
        for ch_key in nf_keys + f_keys:
            ch = chapters.get(ch_key)
            if ch and ch.get("included", False):
                ch_type = "auto" if ch.get("type") == "financial" else "contenu"
                st.markdown(f"**{num}.** {ch['label']} _({ch_type})_")
                num += 1

        if num == 1:
            st.info("Aucun chapitre selectionne.")

    st.markdown("---")

    # ── Apercu des chapitres financiers (affichage direct) ──────────────
    st.subheader("Apercu de la presentation")
    x_labels_an = [str(int(row["Annee calendaire"])) for _, row in indic.iterrows()]

    if chapters.get("invest_financement", {}).get("included"):
        st.markdown('<div class="section-header">Investissements & Financement</div>', unsafe_allow_html=True)
        total_inv = sum(i["montant"] for i in params["investissements"])
        c1, c2, c3 = st.columns(3)
        with c1:
            metric_card("Investissement", fmt_eur(total_inv), "metric-orange")
        with c2:
            metric_card("Fonds propres", fmt_eur(params["fonds_propres_initial"]), "metric-green")
        with c3:
            total_dette = sum(p["montant"] for p in params["prets"] if p["montant"] > 0)
            metric_card("Dette", fmt_eur(total_dette), "metric-red")
        # Tableau investissements
        inv_data = []
        for inv in params["investissements"]:
            amort = inv["montant"] / inv["duree_amort"] if inv["duree_amort"] > 0 else 0
            inv_data.append({
                "Categorie": inv["categorie"],
                "Montant": f"{inv['montant']:,.0f} \u20ac",
                "Duree amort.": f"{inv['duree_amort']} ans" if inv["duree_amort"] > 0 else "Non amorti",
                "Amort. annuel": f"{amort:,.0f} \u20ac" if amort > 0 else "-",
            })
        st.dataframe(pd.DataFrame(inv_data), use_container_width=True, hide_index=True)
        # Tableau financement
        fin_data = [{"Source": "Fonds propres", "Montant": f"{params['fonds_propres_initial']:,.0f} \u20ac",
                     "Taux": "-", "Duree": "-"}]
        for pret in params["prets"]:
            if pret["montant"] > 0:
                fin_data.append({"Source": pret["nom"], "Montant": f"{pret['montant']:,.0f} \u20ac",
                                 "Taux": f"{pret['taux_annuel']*100:.2f}%", "Duree": f"{pret['duree_ans']} ans"})
        st.dataframe(pd.DataFrame(fin_data), use_container_width=True, hide_index=True)
        # Pie chart répartition investissement
        inv_labels = [i["categorie"] for i in params["investissements"]]
        inv_values = [i["montant"] for i in params["investissements"]]
        fig = px.pie(names=inv_labels, values=inv_values, title="Repartition des investissements",
                     color_discrete_sequence=px.colors.qualitative.Set3)
        fig.update_layout(height=380)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
        st.markdown("---")

    if chapters.get("hypotheses_cles", {}).get("included"):
        st.markdown('<div class="section-header">Hypotheses cles</div>', unsafe_allow_html=True)
        hyp_data = [
            {"Parametre": "Chambres", "Valeur": str(params["nb_chambres"])},
            {"Parametre": "Ouverture", "Valeur": params["date_ouverture"].strftime("%B %Y")},
            {"Parametre": "Personnes par chambre", "Valeur": str(params.get("personnes_par_chambre", 2))},
        ]
        for seg_name, seg_data in params.get("segments", {}).items():
            hyp_data.append({"Parametre": f"Segment {seg_name}",
                             "Valeur": f"Part {seg_data['part']*100:.0f}% | ADR {seg_data['prix']:,.0f} \u20ac"})
        taux_occ = params.get("taux_occ", [])
        for i, t in enumerate(taux_occ):
            yr_label = f"A{i+1}" if i < len(taux_occ) - 1 else f"A{i+1}+"
            hyp_data.append({"Parametre": f"Taux occ. {yr_label}", "Valeur": f"{t*100:.1f}%"})
        st.dataframe(pd.DataFrame(hyp_data), use_container_width=True, hide_index=True)
        st.markdown("---")

    if chapters.get("previsions_ca", {}).get("included"):
        st.markdown('<div class="section-header">Previsions de Chiffre d\'Affaires</div>', unsafe_allow_html=True)
        # Graphique CA stacked
        fig = go.Figure()
        fig.add_trace(go.Bar(x=x_labels_an, y=indic["CA Hebergement"], name="Hebergement", marker_color="#4facfe"))
        fig.add_trace(go.Bar(x=x_labels_an, y=indic["CA Brasserie"], name="Brasserie", marker_color="#f5576c"))
        fig.add_trace(go.Bar(x=x_labels_an, y=indic["CA Autres"], name="Autres", marker_color="#38ef7d"))
        fig.update_layout(barmode="stack", height=450,
                          xaxis=dict(type="category", tickfont=dict(size=12)),
                          yaxis_tickformat=",", legend=dict(orientation="h", y=-0.12))
        _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 450)
        # Tableau CA 5 ans
        ca_data = []
        for a in annees5:
            d = indic[indic["Annee calendaire"] == a].iloc[0]
            ca_data.append({
                "Annee": str(int(a)), "Mois": str(int(d["Mois"])),
                "CA Hebergement": f"{d['CA Hebergement']:,.0f} \u20ac",
                "CA Brasserie": f"{d['CA Brasserie']:,.0f} \u20ac",
                "CA Autres": f"{d['CA Autres']:,.0f} \u20ac",
                "CA Total": f"{d['CA Total']:,.0f} \u20ac",
            })
        st.dataframe(pd.DataFrame(ca_data), use_container_width=True, hide_index=True)
        st.markdown("---")

    if chapters.get("structure_couts", {}).get("included"):
        st.markdown('<div class="section-header">Structure des Couts</div>', unsafe_allow_html=True)
        idx_ref = min(3, len(annees_cal) - 1)
        an_cal_ref = int(annees_cal[idx_ref])
        row_ref = indic[indic["Annee calendaire"] == an_cal_ref].iloc[0]
        col1, col2 = st.columns(2)
        with col1:
            labels_pie = ["Personnel", "Charges fixes autres", "Charges variables"]
            values_pie = [row_ref["Charges Fixes"] * 0.7, row_ref["Charges Fixes"] * 0.3, row_ref["Charges Variables"]]
            fig = px.pie(names=labels_pie, values=values_pie,
                         title=f"Repartition des couts ({an_cal_ref})",
                         color_discrete_sequence=["#667eea", "#764ba2", "#f5576c"])
            fig.update_layout(height=400)
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
        with col2:
            fig = go.Figure(go.Waterfall(
                name="Compte de resultat", orientation="v",
                measure=["absolute", "relative", "relative", "total",
                         "relative", "total", "relative", "relative", "total"],
                x=["CA Total", "Ch. Var.", "Ch. Fixes", "EBITDA",
                   "Amort.", "EBIT", "Interets", "Impot*", "Res. Net"],
                y=[row_ref["CA Total"], -row_ref["Charges Variables"],
                   -row_ref["Charges Fixes"], 0,
                   -row_ref["Amortissement"], 0, -row_ref["Interets"],
                   -(row_ref["EBIT"] - row_ref["Interets"] - row_ref["Resultat Net"]), 0],
                connector={"line": {"color": "#888"}},
                increasing={"marker": {"color": "#38ef7d"}},
                decreasing={"marker": {"color": "#f5576c"}},
                totals={"marker": {"color": "#667eea"}},
            ))
            fig.update_layout(title=f"Waterfall ({an_cal_ref})", height=400, yaxis_tickformat=",")
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
        # Tableau charges sur 5 ans
        ch_data = []
        for a in annees5:
            d = indic[indic["Annee calendaire"] == a].iloc[0]
            ch_data.append({
                "Annee": str(int(a)),
                "Charges Variables": f"{d['Charges Variables']:,.0f} \u20ac",
                "Charges Fixes": f"{d['Charges Fixes']:,.0f} \u20ac",
                "Marge Brute %": f"{d['Marge Brute %']:.1f}%",
                "EBITDA": f"{d['EBITDA']:,.0f} \u20ac",
                "EBITDA %": f"{d['EBITDA %']:.1f}%",
            })
        st.dataframe(pd.DataFrame(ch_data), use_container_width=True, hide_index=True)
        st.markdown("---")

    if chapters.get("indicateurs_hoteliers", {}).get("included"):
        st.markdown('<div class="section-header">Indicateurs Hoteliers</div>', unsafe_allow_html=True)
        # Graphiques cote a cote
        col1, col2 = st.columns(2)
        with col1:
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=x_labels_an, y=indic["Taux Occupation"],
                                      mode="lines+markers", name="Taux Occup. (%)",
                                      line=dict(color="#764ba2", width=3)))
            fig.update_layout(title=dict(text="Taux d'occupation (%)", font=dict(size=14)),
                              height=400, xaxis=dict(type="category", tickfont=dict(size=12)),
                              yaxis=dict(title="%", range=[0, 100]))
            _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 400)
        with col2:
            fig = make_subplots(specs=[[{"secondary_y": True}]])
            fig.add_trace(go.Bar(x=x_labels_an, y=indic["Prix Moyen (ADR)"],
                                 name="ADR", marker_color="#667eea", opacity=0.7), secondary_y=False)
            fig.add_trace(go.Scatter(x=x_labels_an, y=indic["RevPAR"],
                                      name="RevPAR", mode="lines+markers",
                                      line=dict(color="#f5576c", width=3)), secondary_y=True)
            fig.update_layout(title=dict(text="ADR & RevPAR", font=dict(size=14)),
                              height=400, xaxis=dict(type="category", tickfont=dict(size=12)),
                              legend=dict(orientation="h", y=-0.15))
            fig.update_yaxes(title_text="ADR (\u20ac)", secondary_y=False)
            fig.update_yaxes(title_text="RevPAR (\u20ac)", secondary_y=True)
            _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 400)
        # Tableau indicateurs
        ki_data = []
        for a in annees5:
            d = indic[indic["Annee calendaire"] == a].iloc[0]
            ki_data.append({
                "Annee": str(int(a)),
                "Nuitees": f"{d['Nuitees']:,.0f}",
                "Taux Occupation": f"{d['Taux Occupation']:.1f}%",
                "ADR": f"{d['Prix Moyen (ADR)']:,.0f} \u20ac",
                "RevPAR": f"{d['RevPAR']:,.0f} \u20ac",
                "GOPPAR": f"{d['GOPPAR']:,.0f} \u20ac",
            })
        st.dataframe(pd.DataFrame(ki_data), use_container_width=True, hide_index=True)
        st.markdown("---")

    if chapters.get("rentabilite_resultat", {}).get("included"):
        st.markdown('<div class="section-header">Rentabilite & Resultat</div>', unsafe_allow_html=True)
        # Graphique rentabilite
        fig = go.Figure()
        fig.add_trace(go.Bar(x=x_labels_an, y=indic["EBITDA"], name="EBITDA", marker_color="#11998e"))
        fig.add_trace(go.Scatter(x=x_labels_an, y=indic["Resultat Net"], name="Resultat Net",
                                  mode="lines+markers", line=dict(color="#f5576c", width=3)))
        fig.add_trace(go.Scatter(x=x_labels_an, y=indic["Cash Flow"], name="Cash Flow",
                                  mode="lines+markers", line=dict(color="#4facfe", width=3, dash="dash")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
        fig.update_layout(height=450, xaxis=dict(type="category", tickfont=dict(size=12)),
                          yaxis_tickformat=",", legend=dict(orientation="h", y=-0.12))
        _scrollable_chart(fig, n_annees, _PX_PAR_BARRE, 450)
        # Tableau rentabilite
        rent_data = []
        for a in annees5:
            d = indic[indic["Annee calendaire"] == a].iloc[0]
            dscr_str = f"{d['DSCR']:.2f}x" if d["DSCR"] < 100 else "\u221e"
            rent_data.append({
                "Annee": str(int(a)),
                "EBITDA": f"{d['EBITDA']:,.0f} \u20ac",
                "Marge EBITDA": f"{d['EBITDA %']:.1f}%",
                "Resultat Net": f"{d['Resultat Net']:,.0f} \u20ac",
                "Cash Flow": f"{d['Cash Flow']:,.0f} \u20ac",
                "Cash Flow Cumul": f"{d['Cash Flow Cumul']:,.0f} \u20ac",
                "DSCR": dscr_str,
            })
        st.dataframe(pd.DataFrame(rent_data), use_container_width=True, hide_index=True)
        st.markdown("---")

    if chapters.get("cashflow_roi", {}).get("included"):
        st.markdown('<div class="section-header">Cash Flow & Retour sur Investissement</div>', unsafe_allow_html=True)
        # Graphique CF cumule
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=df["date"], y=df["cash_flow_cumul"],
                                  fill="tozeroy", line=dict(color="#667eea", width=2),
                                  name="Cash Flow Cumule"))
        fig.add_hline(y=0, line_dash="dash", line_color="red")
        invest_total = sum(i["montant"] for i in params["investissements"])
        fig.add_hline(y=-invest_total, line_dash="dot", line_color="gray",
                      annotation_text="Investissement initial")
        fig.update_layout(height=450,
                          xaxis=dict(title="", dtick="M6", tickformat="%b %Y",
                                     tickfont=dict(size=11), tickangle=-45),
                          yaxis_title="EUR", yaxis_tickformat=",")
        _scrollable_chart(fig, len(df), _PX_PAR_MOIS, 450)
        # Graphique CF annuel
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(x=x_labels_an, y=indic["Cash Flow"], name="Cash Flow annuel", marker_color="#4facfe"))
        fig2.add_trace(go.Scatter(x=x_labels_an, y=indic["Cash Flow Cumul"], name="CF Cumule",
                                   mode="lines+markers", line=dict(color="#667eea", width=3)))
        fig2.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
        fig2.update_layout(height=400, xaxis=dict(type="category", tickfont=dict(size=12)),
                           yaxis_tickformat=",", legend=dict(orientation="h", y=-0.12))
        _scrollable_chart(fig2, n_annees, _PX_PAR_BARRE, 400)
        # Break-even
        cf_cumul = df["cash_flow_cumul"].values
        break_even_idx = None
        for i in range(1, len(cf_cumul)):
            if cf_cumul[i - 1] < 0 and cf_cumul[i] >= 0:
                break_even_idx = i
                break
        if break_even_idx:
            be_date = df.iloc[break_even_idx]["date"]
            st.success(f"Break-even atteint en {be_date.strftime('%B %Y')} (mois {break_even_idx} d'exploitation)")
        else:
            st.warning("Break-even non atteint sur la periode de projection")

    st.markdown("---")

    # ── Export buttons ────────────────────────────────────────────────────
    st.markdown('<div class="section-header">Exports</div>', unsafe_allow_html=True)
    c_exp0, c_exp1, c_exp2, c_exp3 = st.columns([2, 2, 2, 2])
    with c_exp0:
        html_report = _export_html(indic, params, df, chapters)
        st.download_button("Rapport HTML (recommande)", html_report,
                           "plan_financier.html", "text/html",
                           type="primary", use_container_width=True,
                           help="Ouvrir dans le navigateur puis Ctrl+P pour un PDF parfait")
    with c_exp1:
        try:
            pptx_bytes = _export_pptx(indic, params, df, chapters)
            st.download_button("Exporter en PowerPoint (.pptx)", pptx_bytes,
                               "plan_financier.pptx",
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                               use_container_width=True)
        except ImportError:
            st.warning("Installez python-pptx pour l'export PowerPoint")
    with c_exp2:
        try:
            pdf_bytes = _export_pdf(indic, params, df, chapters)
            st.download_button("Exporter en PDF", pdf_bytes,
                               "plan_financier.pdf", "application/pdf",
                               use_container_width=True)
        except ImportError:
            st.warning("Installez fpdf2 pour l'export PDF")
    with c_exp3:
        csv = indic.to_csv(index=False, sep=";", decimal=",")
        st.download_button("Indicateurs (CSV)", csv, "indicateurs_hotel.csv", "text/csv",
                           use_container_width=True)


# ─── Methodologie ──────────────────────────────────────────────────────────────

def tab_methodologie(df, indic, params):
    st.header("Methodologie & Verification des calculs")
    st.caption("Cette section detaille la construction de chaque indicateur du plan financier. "
               "Utilisez-la pour comprendre et verifier les hypotheses et formules utilisees.")

    categorie = st.selectbox("Categorie", [
        "1. Revenus - Hebergement",
        "2. Revenus - Brasserie & PDJ",
        "3. Revenus - Bar, Spa, Salles",
        "4. Charges variables",
        "5. Charges fixes",
        "6. Amortissements & Reinvestissement",
        "7. Service de la dette",
        "8. Compte de resultat",
        "9. Cash Flow",
        "10. Indicateurs hoteliers (KPI)",
        "11. Inflation",
    ], key="methodo_cat")

    # --- Helper values ---
    nb_ch = params["nb_chambres"]
    pers_ch = params.get("personnes_par_chambre", 2)
    taux_occ = params["taux_occ"]
    saison = params["saisonnalite"]
    segments = params["segments"]
    mois_noms = ["Jan", "Fev", "Mar", "Avr", "Mai", "Jun",
                 "Jul", "Aou", "Sep", "Oct", "Nov", "Dec"]

    # ADR base (annee 1, sans inflation)
    adr_base = sum(seg["part"] * seg["prix"] for seg in segments.values())

    # First year data from indic
    if len(indic) > 0:
        a1 = indic.iloc[0]
    else:
        a1 = None

    # =====================================================================
    # 1. REVENUS - HEBERGEMENT
    # =====================================================================
    if categorie.startswith("1."):
        st.subheader("1. Revenus - Hebergement")

        st.markdown("---")
        st.markdown("#### Chambres disponibles")
        st.markdown(f"""
**Formule** : `Chambres disponibles = nb_chambres x jours_dans_mois`
**Description** : Nombre total de chambres pouvant etre vendues dans le mois.
**Donnees** : nb_chambres = **{nb_ch}**
**Exemple** : Pour un mois de 31 jours : {nb_ch} x 31 = **{nb_ch * 31:,}** chambres disponibles
""")

        st.markdown("---")
        st.markdown("#### Taux d'occupation")
        st.markdown("""
**Formule** : `Taux d'occupation = min(taux_base[annee] x saisonnalite[mois], 100%)`
**Description** : Pourcentage des chambres vendues. Le taux de base augmente progressivement
sur les premieres annees (montee en charge), module par un coefficient de saisonnalite mensuel.
""")
        st.markdown("**Taux de base par annee d'exploitation :**")
        taux_occ_data = {f"A{i+1}{'+'if i==len(taux_occ)-1 else ''}": f"{t:.1%}" for i, t in enumerate(taux_occ)}
        st.dataframe(pd.DataFrame([taux_occ_data]), hide_index=True)

        st.markdown("**Coefficients de saisonnalite (jan-dec) :**")
        saison_data = {mois_noms[i]: f"{s:.3f}" for i, s in enumerate(saison)}
        st.dataframe(pd.DataFrame([saison_data]), hide_index=True)

        st.markdown(f"""
**Verification** : Annee 1, mois de juillet (saisonnalite={saison[6]:.3f}) :
`min({taux_occ[0]:.3f} x {saison[6]:.3f}, 1.0)` = **{min(taux_occ[0] * saison[6], 1.0):.1%}**
""")

        st.markdown("---")
        st.markdown("#### Nuitees (chambres vendues)")
        ex_taux = min(taux_occ[0] * saison[6], 1.0)
        ex_nuitees = nb_ch * 31 * ex_taux
        st.markdown(f"""
**Formule** : `Nuitees = Chambres disponibles x Taux d'occupation`
**Description** : Nombre de chambres vendues dans le mois.
**Verification** : Juillet A1 : {nb_ch} x 31 x {ex_taux:.3f} = **{ex_nuitees:,.0f}** nuitees
""")

        st.markdown("---")
        st.markdown("#### Prix moyen par nuitee (ADR)")
        st.markdown("""
**Formule** : `ADR = Somme(part_segment x prix_segment) x hausse_prix_cumulee x inflation_ventes`
**Description** : Le prix moyen est une moyenne ponderee par la segmentation clientele,
ajustee annuellement par une hausse prix reelle et par l'inflation.
""")
        st.markdown("**Segmentation clientele :**")
        seg_rows = []
        for nom, data in segments.items():
            part_ota = params.get("segments_part_ota", {}).get(nom, 0)
            seg_rows.append({
                "Segment": nom,
                "Part": f"{data['part']:.0%}",
                "Prix (EUR)": f"{data['prix']:.0f}",
                "Part OTA": f"{part_ota:.0%}",
                "Contribution ADR": f"{data['part'] * data['prix']:.1f} EUR",
            })
        st.dataframe(pd.DataFrame(seg_rows), hide_index=True)
        st.markdown(f"**ADR de base (hors inflation)** = **{adr_base:.2f} EUR**")

        hausse = params.get("hausse_prix_an", [0, 0.025, 0.025, 0.025])
        hausse_str = ", ".join(f"A{i+1}={h:.1%}" for i, h in enumerate(hausse))
        st.markdown(f"**Hausse prix annuelle (hors inflation)** : {hausse_str}")

        if a1 is not None:
            st.markdown("---")
            st.markdown("#### Verification annee 1")

            an1_cal = int(a1["Annee calendaire"])
            nb_mois_a1 = int(a1["Mois"])
            ca_h_a1 = a1["CA Hebergement"]
            nuit_a1 = a1["Nuitees"]
            adr_eff = a1["Prix Moyen (ADR)"]

            st.info(f"L'annee {an1_cal} compte **{nb_mois_a1} mois** d'exploitation "
                    f"(ouverture en {params['date_ouverture'].strftime('%B %Y')}).")

            st.markdown(f"""
**ADR effectif** = CA Hebergement / Nuitees = {ca_h_a1:,.0f} / {nuit_a1:,.0f} = **{adr_eff:,.2f} EUR**

> L'ADR effectif est une moyenne **a posteriori**. Il differe de l'ADR de base ({adr_base:.0f} EUR)
> car l'inflation mensuelle ({params.get('inflation_ventes', 0.025):.1%}/an) est appliquee
> mois par mois depuis l'ouverture.

**Verification** : CA = Nuitees × ADR effectif = {nuit_a1:,.0f} × {adr_eff:,.2f} = **{nuit_a1 * adr_eff:,.0f} EUR** ✓
""")

            # Tableau mois par mois
            st.markdown("##### Detail mois par mois")
            df_a1 = df[df["annee"] == an1_cal].copy()
            if len(df_a1) > 0:
                detail_rows = []
                for _, row in df_a1.iterrows():
                    m_nuit = row["nuitees"]
                    m_prix = row["prix_moyen"]
                    m_ca = row["ca_hebergement"]
                    m_taux = row["taux_occupation"]
                    m_verif = m_nuit * m_prix
                    detail_rows.append({
                        "Mois": row["date"].strftime("%b %Y"),
                        "Taux occ.": f"{m_taux:.1%}",
                        "Nuitees": f"{m_nuit:,.0f}",
                        "ADR (EUR)": f"{m_prix:,.1f}",
                        "CA = Nuit. x ADR": f"{m_verif:,.0f}",
                        "CA reel": f"{m_ca:,.0f}",
                        "Ecart": f"{m_ca - m_verif:,.0f}",
                    })
                st.dataframe(pd.DataFrame(detail_rows), hide_index=True, use_container_width=True)
                st.caption("L'ecart est nul : CA = Nuitees × ADR pour chaque mois. "
                           "L'ADR varie d'un mois a l'autre en raison de l'inflation cumulative.")

    # =====================================================================
    # 2. REVENUS - BRASSERIE & PDJ
    # =====================================================================
    elif categorie.startswith("2."):
        st.subheader("2. Revenus - Brasserie & Petit-dejeuner")

        st.markdown("---")
        st.markdown("#### CA Brasserie externe")
        prix_souper = params.get("brasserie_prix_souper", 75)
        jrs_souper = params.get("brasserie_jours_souper", 4)
        srv_souper = params.get("brasserie_services_souper", 1)
        prix_diner = params.get("brasserie_prix_diner", 45)
        jrs_diner = params.get("brasserie_jours_diner", 4)
        srv_diner = params.get("brasserie_services_diner", 1.5)
        nb_couv = params.get("nb_couverts_brasserie", 80)
        taux_occ_brass = params.get("taux_occ_brasserie", taux_occ)

        st.markdown(f"""
**Formule** :
```
CA Brasserie ext = taux_occ_brasserie x saisonnalite
    x (prix_souper x jours_souper/7 x services_souper
     + prix_diner x jours_diner/7 x services_diner)
    x jours_du_mois x nb_couverts
```
**Parametres** :
| Parametre | Valeur |
|---|---|
| Prix souper | {prix_souper} EUR |
| Jours souper / semaine | {jrs_souper} |
| Services souper / jour | {srv_souper} |
| Prix diner | {prix_diner} EUR |
| Jours diner / semaine | {jrs_diner} |
| Services diner / jour | {srv_diner} |
| Nb couverts | {nb_couv} |
""")
        taux_brass_str = ", ".join(f"A{i+1}={t:.1%}" for i, t in enumerate(taux_occ_brass))
        st.markdown(f"**Taux d'occupation brasserie** : {taux_brass_str}")

        st.markdown("---")
        st.markdown("#### CA Petit-dejeuner (PDJ)")
        pdj_prix = params.get("petit_dej_prix", 37.5)
        pdj_taux = params.get("petit_dej_taux", 0.85)
        ex_nuitees = nb_ch * 31 * min(taux_occ[0] * saison[6], 1.0)
        ca_pdj_ex = ex_nuitees * pdj_taux * pdj_prix * pers_ch
        st.markdown(f"""
**Formule** : `CA PDJ = nuitees x taux_prise_pdj x prix_pdj x personnes_par_chambre`
**Parametres** :
| Parametre | Valeur |
|---|---|
| Prix PDJ | {pdj_prix} EUR |
| Taux de prise PDJ | {pdj_taux:.0%} |
| Personnes par chambre | {pers_ch} |

**Exemple** (mois de juillet, A1) : {ex_nuitees:,.0f} nuitees x {pdj_taux:.2f} x {pdj_prix:.1f} x {pers_ch} = **{ca_pdj_ex:,.0f} EUR**
""")

        st.markdown("---")
        st.markdown("#### CA Brasserie total")
        st.markdown("""
**Formule** : `CA Brasserie = (CA Brasserie ext + CA PDJ) x inflation_ventes`
**Description** : Le CA total brasserie combine les revenus externes (couverts restaurant)
et internes (petits-dejeuners des clients hotel), ajuste par l'inflation.
""")

        if a1 is not None:
            st.markdown(f"**CA Brasserie A1** = **{a1['CA Brasserie']:,.0f} EUR**")

    # =====================================================================
    # 3. REVENUS - BAR, SPA, SALLES
    # =====================================================================
    elif categorie.startswith("3."):
        st.subheader("3. Revenus - Bar, Spa, Salles & Divers")

        st.markdown("---")
        st.markdown("#### CA Bar")
        bar_prix = params.get("bar_prix_moyen", 18)
        bar_taux = params.get("bar_taux_clients_hotel", 0.40)
        st.markdown(f"""
**Formule** : `CA Bar = prix_moyen_bar x taux_clients_hotel x nuitees x pers_par_chambre x inflation_ventes`
**Parametres** :
| Parametre | Valeur |
|---|---|
| Prix moyen par personne | {bar_prix} EUR |
| Taux clients hotel au bar | {bar_taux:.0%} |
| Personnes par chambre | {pers_ch} |
""")

        st.markdown("---")
        st.markdown("#### CA Spa")
        st.markdown(f"""
**Formule** :
```
CA Spa = (entrees_hotel + soins_hotel + entrees_ext + soins_ext) x inflation_ventes
```
**Composantes** :
| Composante | Calcul | Parametres |
|---|---|---|
| Entrees hotel | nuitees x taux x prix | taux={params.get('spa_entree_hotel_taux', 0.20):.0%}, prix={params.get('spa_entree_hotel_prix', 0)} EUR |
| Soins hotel | nuitees x taux x prix | taux={params.get('spa_soin_hotel_taux', 0.10):.0%}, prix={params.get('spa_soin_hotel_prix', 120)} EUR |
| Entrees ext | nb/mois x prix | nb={params.get('spa_entree_ext_nb_mois', 25)}/mois, prix={params.get('spa_entree_ext_prix', 55)} EUR |
| Soins ext | nb/mois x prix | nb={params.get('spa_soin_ext_nb_mois', 15)}/mois, prix={params.get('spa_soin_ext_prix', 150)} EUR |
""")

        st.markdown("---")
        st.markdown("#### CA Seminaires & Mariages")
        sem_prix = params.get("seminaire_prix_location", 800)
        sem_nb = params.get("seminaire_nb_an", 50)
        mar_prix = params.get("mariage_prix_location", 2500)
        mar_nb = params.get("mariage_nb_an", 12)
        st.markdown(f"""
**Formule** :
```
CA Seminaires = prix_location_salle x nb_seminaires_par_an / 12 x inflation
CA Mariages   = prix_location_salle x nb_mariages_par_an / 12 x inflation
```
**Parametres** :
| Type | Location salle | Nombre/an | CA annuel (hors inflation) |
|---|---|---|---|
| Seminaires | {sem_prix:,.0f} EUR | {sem_nb} | {sem_prix * sem_nb:,.0f} EUR |
| Mariages | {mar_prix:,.0f} EUR | {mar_nb} | {mar_prix * mar_nb:,.0f} EUR |
| **Total** | | | **{sem_prix * sem_nb + mar_prix * mar_nb:,.0f} EUR/an** |
""")

        st.markdown("---")
        st.markdown("#### CA Divers")
        div_prix = params.get("divers_prix_nuitee", 3)
        div_taux = params.get("divers_taux", 1.0)
        st.markdown(f"""
**Formule** : `CA Divers = prix_nuitee x taux x nuitees x inflation`
**Parametres** : prix = {div_prix} EUR/nuitee, taux = {div_taux:.0%}
""")

        if a1 is not None:
            st.markdown("---")
            st.markdown("#### Verification annee 1")
            st.markdown(f"**CA Autres (Bar+Spa+Salles+Divers) A1** = **{a1['CA Autres']:,.0f} EUR**")

    # =====================================================================
    # 4. CHARGES VARIABLES
    # =====================================================================
    elif categorie.startswith("4."):
        st.subheader("4. Charges variables")

        st.markdown("---")
        st.markdown("#### CV Hebergement")
        st.markdown("""
**Composantes** :
1. **Commission OTA** = poids_OTA_par_segment x CA_hebergement x taux_commission
2. **Couts par nuitee** = somme(cout_unitaire) x nuitees
3. **Franchise** = (CA total - loyer restaurant - spa) x taux_franchise
""")

        comm_ota = params.get("cv_commission_ota_pct", 0.17)
        franchise = params.get("cv_franchise_pct", 0.04)
        st.markdown(f"**Taux commission OTA** : {comm_ota:.0%}")
        st.markdown(f"**Taux franchise** : {franchise:.0%}")

        st.markdown("**Part OTA par segment (pour ponderation) :**")
        ota_data = params.get("segments_part_ota", {})
        ota_rows = []
        for seg, data in segments.items():
            part_ota = ota_data.get(seg, 0)
            ota_rows.append({
                "Segment": seg,
                "Part CA": f"{data['part']:.0%}",
                "Part OTA": f"{part_ota:.0%}",
                "Poids commission": f"{data['part'] * part_ota:.2%}",
            })
        st.dataframe(pd.DataFrame(ota_rows), hide_index=True)
        poids_ota_total = sum(d["part"] * ota_data.get(s, 0) for s, d in segments.items())
        st.markdown(f"**Poids OTA pondere total** = {poids_ota_total:.2%} => Commission effective sur CA = {poids_ota_total * comm_ota:.2%}")

        st.markdown("**Couts par nuitee :**")
        cv_nuitee = params.get("cv_hebergement_par_nuitee", {})
        cv_rows = [{"Poste": k, "Cout/nuitee (EUR)": f"{v:.2f}"} for k, v in cv_nuitee.items()]
        total_cv_nuitee = sum(cv_nuitee.values())
        cv_rows.append({"Poste": "**TOTAL**", "Cout/nuitee (EUR)": f"**{total_cv_nuitee:.2f}**"})
        st.dataframe(pd.DataFrame(cv_rows), hide_index=True)

        st.markdown("---")
        st.markdown("#### CV Brasserie (incl. PDJ)")
        cv_brass = params.get("cv_brasserie_pct", 0.35)
        st.markdown(f"""
**Formule** : `CV Brasserie = (CA PDJ + CA Brasserie ext) x food_cost_pct`
**Food cost** = **{cv_brass:.0%}**
""")

        st.markdown("---")
        st.markdown("#### CV Bar")
        cv_bar_bev = params.get("cv_bar_beverage_pct", 0.30)
        cv_bar_conso = params.get("cv_bar_consommable_unite", 0.20)
        st.markdown(f"""
**Formule** : `CV Bar = beverage_cost x CA Bar + consommable_unite x nb_consommations`
**Beverage cost** = {cv_bar_bev:.0%}, **Consommable/unite** = {cv_bar_conso:.2f} EUR
""")

        st.markdown("---")
        st.markdown("#### CV Spa")
        st.markdown(f"""
**Formule** :
```
CV Spa = (cout_soin + produits_soin) x nb_soins
       + (consommable + energie + piscine) x nb_entrees
```
**Parametres** :
| Poste | Cout unitaire |
|---|---|
| Cout soin | {params.get('cv_spa_soin_cout', 50)} EUR/soin |
| Produits soin | {params.get('cv_spa_produits_soin', 5)} EUR/soin |
| Consommable entree | {params.get('cv_spa_consommable_entree', 1.5)} EUR/entree |
| Energie entree | {params.get('cv_spa_energie_entree', 3.0)} EUR/entree |
| Piscine entree | {params.get('cv_spa_piscine_entree', 0.5)} EUR/entree |
""")

        st.markdown("---")
        st.markdown("#### CV Salles")
        st.markdown(f"""
**CV Seminaires** (par seminaire) :
| Poste | Cout |
|---|---|
| Pause/participant | {params.get('cv_seminaire_pause_participant', 8)} EUR x {params.get('nb_invites_seminaire', 50)} invites |
| Materiel/participant | {params.get('cv_seminaire_materiel_participant', 3)} EUR x {params.get('nb_invites_seminaire', 50)} invites |
| Equipement | {params.get('cv_seminaire_equipement', 15)} EUR |
| Energie | {params.get('cv_seminaire_energie', 75)} EUR |
| Nettoyage | {params.get('cv_seminaire_nettoyage', 500)} EUR |

**CV Mariages** (par mariage) :
| Poste | Cout |
|---|---|
| Energie | {params.get('cv_mariage_energie', 150)} EUR |
| Nettoyage | {params.get('cv_mariage_nettoyage', 1000)} EUR |
""")

        if a1 is not None:
            st.markdown("---")
            st.markdown(f"**Charges Variables totales A1** = **{a1['Charges Variables']:,.0f} EUR**")
            st.markdown(f"**Marge Brute A1** = **{a1['Marge Brute']:,.0f} EUR** ({a1['Marge Brute %']:.1f}%)")

    # =====================================================================
    # 5. CHARGES FIXES
    # =====================================================================
    elif categorie.startswith("5."):
        st.subheader("5. Charges fixes")

        charges_pat = params.get("charges_patronales_pct", 0.35)

        st.markdown("---")
        st.markdown("#### Personnel direct (par departement)")
        st.markdown(f"""
**Formule** :
```
Masse annuelle = cout_brut x (1 + charges_patronales%) x ETP
Mensuel de base = masse / 13.92  (convention belge)
Multiplicateur : juillet = 1.92, decembre = 2.0, autres mois = 1.0
```
**Charges patronales** = {charges_pat:.0%}
""")

        for dept_name, dept_key in [("Hebergement", "personnel_hebergement"),
                                     ("Brasserie", "personnel_brasserie"),
                                     ("Bar", "personnel_bar"),
                                     ("Spa", "personnel_spa")]:
            personnel = params.get(dept_key, [])
            if personnel:
                st.markdown(f"**{dept_name}** :")
                rows = []
                total_masse = 0
                for p in personnel:
                    masse = p["cout_brut"] * (1 + charges_pat) * p["etp"]
                    total_masse += masse
                    rows.append({
                        "Poste": p["poste"],
                        "ETP": p["etp"],
                        "Cout brut": f"{p['cout_brut']:,.0f} EUR",
                        "Masse chargee": f"{masse:,.0f} EUR/an",
                        "Mensuel base": f"{masse / 13.92:,.0f} EUR",
                    })
                rows.append({
                    "Poste": "**TOTAL**",
                    "ETP": sum(p["etp"] for p in personnel),
                    "Cout brut": "",
                    "Masse chargee": f"**{total_masse:,.0f} EUR/an**",
                    "Mensuel base": f"**{total_masse / 13.92:,.0f} EUR**",
                })
                st.dataframe(pd.DataFrame(rows), hide_index=True)

        st.markdown("---")
        st.markdown("#### Personnel indirect")
        personnel_ind = params.get("personnel_indirect", [])
        if personnel_ind:
            rows = []
            total_masse = 0
            for p in personnel_ind:
                masse = p["cout_brut"] * (1 + charges_pat) * p["etp"]
                total_masse += masse
                rows.append({
                    "Poste": p["poste"],
                    "ETP": p["etp"],
                    "Cout brut": f"{p['cout_brut']:,.0f} EUR",
                    "Masse chargee": f"{masse:,.0f} EUR/an",
                    "Mensuel base": f"{masse / 13.92:,.0f} EUR",
                })
            rows.append({
                "Poste": "**TOTAL**",
                "ETP": sum(p["etp"] for p in personnel_ind),
                "Cout brut": "",
                "Masse chargee": f"**{total_masse:,.0f} EUR/an**",
                "Mensuel base": f"**{total_masse / 13.92:,.0f} EUR**",
            })
            st.dataframe(pd.DataFrame(rows), hide_index=True)

        st.markdown("---")
        st.markdown("#### Loyer")
        loyer = params.get("loyer_mensuel", 120000)
        repartition = params.get("loyer_repartition", {})
        st.markdown(f"""
**Loyer mensuel total** = **{loyer:,.0f} EUR/mois**
**Repartition par service** :
""")
        loyer_rows = [{"Service": k.capitalize(), "Part": f"{v:.0%}", "Montant mensuel": f"{loyer * v:,.0f} EUR"}
                      for k, v in repartition.items()]
        st.dataframe(pd.DataFrame(loyer_rows), hide_index=True)

        loyer_rest = params.get("loyer_restaurant_mensuel", 5860)
        st.markdown(f"**Loyer restaurant (reference externe)** = {loyer_rest:,.0f} EUR/mois")

        st.markdown("---")
        st.markdown("#### Charges fixes indirectes")
        cfi = params.get("charges_fixes_indirectes_par_annee", params.get("charges_fixes_indirectes", {}))
        if isinstance(cfi, dict):
            if cfi and isinstance(list(cfi.values())[0], list):
                # Multi-year format
                cf_rows = []
                for poste, vals in cfi.items():
                    row = {"Poste": poste}
                    for i, v in enumerate(vals):
                        row[f"A{i+1}"] = f"{v:,.0f} EUR"
                    cf_rows.append(row)
                # Totaux
                nb_cols = len(list(cfi.values())[0])
                total_row = {"Poste": "**TOTAL**"}
                for i in range(nb_cols):
                    total_row[f"A{i+1}"] = f"**{sum(v[i] for v in cfi.values()):,.0f} EUR**"
                cf_rows.append(total_row)
                st.dataframe(pd.DataFrame(cf_rows), hide_index=True)
            else:
                # Simple format
                cf_rows = [{"Poste": k, "Montant annuel": f"{v:,.0f} EUR", "Mensuel": f"{v/12:,.0f} EUR"}
                           for k, v in cfi.items()]
                total_cfi = sum(cfi.values())
                cf_rows.append({"Poste": "**TOTAL**", "Montant annuel": f"**{total_cfi:,.0f} EUR**",
                                "Mensuel": f"**{total_cfi/12:,.0f} EUR**"})
                st.dataframe(pd.DataFrame(cf_rows), hide_index=True)

        if a1 is not None:
            st.markdown("---")
            st.markdown(f"**Charges Fixes totales A1** = **{a1['Charges Fixes']:,.0f} EUR**")
            st.markdown(f"  dont Directes = {a1['Charges Fixes Directes']:,.0f} EUR, "
                        f"Indirectes = {a1['Charges Fixes Indirectes']:,.0f} EUR")

    # =====================================================================
    # 6. AMORTISSEMENTS & REINVESTISSEMENT
    # =====================================================================
    elif categorie.startswith("6."):
        st.subheader("6. Amortissements & Reinvestissement")

        st.markdown("---")
        st.markdown("#### Amortissement lineaire")
        st.markdown("""
**Formule** : `Amortissement mensuel = montant / (duree_ans x 12)`
**Description** : Amortissement lineaire sur la duree de vie. Les actifs avec duree = 0
(Terrain, Construction) ne sont pas amortis.
""")
        inv = params.get("investissements", [])
        inv_rows = []
        total_amort_an = 0
        for item in inv:
            duree = item.get("duree_amort", 0)
            montant = item.get("montant", 0)
            if duree > 0:
                amort_mens = montant / (duree * 12)
                amort_an = montant / duree
            else:
                amort_mens = 0
                amort_an = 0
            total_amort_an += amort_an
            inv_rows.append({
                "Categorie": item.get("categorie", ""),
                "Montant": f"{montant:,.0f} EUR",
                "Duree (ans)": duree if duree > 0 else "Non amorti",
                "Amort. mensuel": f"{amort_mens:,.0f} EUR" if duree > 0 else "-",
                "Amort. annuel": f"{amort_an:,.0f} EUR" if duree > 0 else "-",
            })
        inv_rows.append({
            "Categorie": "**TOTAL**",
            "Montant": f"**{sum(i.get('montant', 0) for i in inv):,.0f} EUR**",
            "Duree (ans)": "",
            "Amort. mensuel": f"**{total_amort_an/12:,.0f} EUR**",
            "Amort. annuel": f"**{total_amort_an:,.0f} EUR**",
        })
        st.dataframe(pd.DataFrame(inv_rows), hide_index=True)

        st.markdown("---")
        st.markdown("#### Reinvestissement progressif")
        reinv_pct = params.get("reinvest_pct_an", [])
        reinv_cat = params.get("reinvest_categories", [])
        st.markdown(f"""
**Description** : Chaque annee, un pourcentage du CA total est reinvesti
dans les categories : {', '.join(reinv_cat)}.
Les reinvestissements sont amortis sur la duree de leur categorie
(Amenagements = 10 ans, Mobilier = 5 ans).
""")
        st.markdown("**Taux de reinvestissement par annee :**")
        reinv_data = {f"A{i+1}{'+'if i==len(reinv_pct)-1 else ''}": f"{r:.1%}" for i, r in enumerate(reinv_pct)}
        st.dataframe(pd.DataFrame([reinv_data]), hide_index=True)

        if a1 is not None:
            st.markdown(f"**Amortissement total A1** = **{a1['Amortissement']:,.0f} EUR**")

    # =====================================================================
    # 7. SERVICE DE LA DETTE
    # =====================================================================
    elif categorie.startswith("7."):
        st.subheader("7. Service de la dette")

        st.markdown("""
**Types de prets** :
- **Mensualite constante** : mensualite fixe (capital + interets). Formule : `M = P x [r + r / ((1+r)^n - 1)]`
  ou P = principal, r = taux mensuel, n = nombre de mois.
- **Interet seul** : seuls les interets sont payes mensuellement. Formule : `I = capital_restant x taux_mensuel`.
  Le capital est rembourse in fine (a echeance).
""")

        prets = params.get("prets", [])
        for pret in prets:
            nom = pret.get("nom", "")
            montant = pret.get("montant", 0)
            taux = pret.get("taux_annuel", 0)
            duree = pret.get("duree_ans", 0)
            type_pret = pret.get("type", "annuite")
            differe = pret.get("differe_mois", 0)
            taux_mens = taux / 12

            st.markdown(f"---")
            st.markdown(f"#### {nom}")

            if type_pret == "annuite" and taux_mens > 0 and duree > 0:
                n = duree * 12
                mensualite = montant * (taux_mens + taux_mens / ((1 + taux_mens) ** n - 1))
                st.markdown(f"""
| Parametre | Valeur |
|---|---|
| Montant | {montant:,.0f} EUR |
| Taux annuel | {taux:.2%} |
| Duree | {duree} ans ({n} mois) |
| Type | Mensualite constante |
| Differe | {differe} mois |
| **Mensualite** | **{mensualite:,.2f} EUR** |
| Cout total interets | {mensualite * n - montant:,.0f} EUR |
""")
            elif type_pret == "interet_seul":
                interet_mens = montant * taux_mens
                st.markdown(f"""
| Parametre | Valeur |
|---|---|
| Montant | {montant:,.0f} EUR |
| Taux annuel | {taux:.2%} |
| Duree | {duree} ans |
| Type | Interet seul (in fine) |
| Differe | {differe} mois |
| **Interet mensuel** | **{interet_mens:,.2f} EUR** |
| Remboursement capital | A echeance ({duree} ans) |
""")

        if a1 is not None:
            st.markdown("---")
            st.markdown(f"**Service de la dette A1** = **{a1['Service Dette']:,.0f} EUR** "
                        f"(interets : {a1['Interets']:,.0f} EUR)")

    # =====================================================================
    # 8. COMPTE DE RESULTAT
    # =====================================================================
    elif categorie.startswith("8."):
        st.subheader("8. Compte de resultat")

        st.markdown("""
**Structure du compte de resultat** :
```
  CA Total
- Charges Variables
= Marge Brute                     (Marge Brute % = Marge Brute / CA x 100)

  Marge Brute
- Charges Fixes (Directes + Indirectes)
= EBITDA                          (EBITDA % = EBITDA / CA x 100)

  EBITDA
- Amortissements
= EBIT

  EBIT
- Interets dette
= Resultat avant impot

  Resultat avant impot
- ISOC (25% si resultat cumule positif, paye en decembre)
= Resultat Net
```
""")

        isoc = params.get("isoc", 0.25)
        st.markdown(f"**Taux ISOC** = {isoc:.0%}")
        st.markdown("**Note** : L'ISOC est calcule sur le resultat cumule. Il n'est du que lorsque "
                    "le resultat cumule depuis l'ouverture devient positif. Le paiement est comptabilise en decembre.")

        if a1 is not None and len(indic) >= 3:
            st.markdown("---")
            st.markdown("#### Verification sur les premieres annees")
            cdr_cols = ["Annee exploitation", "CA Total", "Charges Variables", "Marge Brute",
                        "Marge Brute %", "Charges Fixes", "EBITDA", "EBITDA %",
                        "Amortissement", "EBIT", "Interets", "Resultat Net"]
            cdr_df = indic[cdr_cols].head(5).copy()
            for col in cdr_df.columns:
                if col not in ["Annee exploitation", "Marge Brute %", "EBITDA %"]:
                    cdr_df[col] = cdr_df[col].apply(lambda x: f"{x:,.0f}" if isinstance(x, (int, float)) else x)
                elif col in ["Marge Brute %", "EBITDA %"]:
                    cdr_df[col] = cdr_df[col].apply(lambda x: f"{x:.1f}%")
            st.dataframe(cdr_df, hide_index=True, use_container_width=True)

    # =====================================================================
    # 9. CASH FLOW
    # =====================================================================
    elif categorie.startswith("9."):
        st.subheader("9. Cash Flow")

        st.markdown("""
**Structure du Cash Flow** :
```
  EBITDA
- Interets dette
= Cash Flow operationnel

  Cash Flow operationnel
- Remboursement capital prets
- ISOC
- Reinvestissement
= Cash Flow mensuel

Cash Flow Cumule = somme cumulative du Cash Flow mensuel
```
**Description** : Le Cash Flow mesure la tresorerie reellement disponible apres
le service de la dette, les impots et les reinvestissements. Le Cash Flow cumule
represente la position de tresorerie depuis l'ouverture.
""")

        if a1 is not None and len(indic) >= 3:
            st.markdown("---")
            st.markdown("#### Evolution sur les premieres annees")
            cf_cols = ["Annee exploitation", "EBITDA", "Interets", "Service Dette",
                       "Cash Flow", "Cash Flow Cumul"]
            cf_df = indic[cf_cols].head(5).copy()
            for col in cf_df.columns:
                if col != "Annee exploitation":
                    cf_df[col] = cf_df[col].apply(lambda x: f"{x:,.0f}" if isinstance(x, (int, float)) else x)
            st.dataframe(cf_df, hide_index=True, use_container_width=True)

    # =====================================================================
    # 10. INDICATEURS HOTELIERS (KPI)
    # =====================================================================
    elif categorie.startswith("10."):
        st.subheader("10. Indicateurs hoteliers (KPI)")

        st.markdown("---")
        st.markdown("#### ADR (Average Daily Rate)")
        st.markdown("""
**Formule** : `ADR = CA Hebergement / Nuitees`
**Description** : Prix moyen effectivement encaisse par chambre vendue.
Inclut les effets de segmentation, hausse prix et inflation.
""")

        st.markdown("---")
        st.markdown("#### RevPAR (Revenue Per Available Room)")
        st.markdown("""
**Formule** : `RevPAR = CA Hebergement / Chambres disponibles = ADR x Taux d'occupation`
**Description** : Revenu genere par chambre disponible (vendue ou non).
Indicateur cle combinant prix et remplissage.
""")

        st.markdown("---")
        st.markdown("#### GOPPAR (Gross Operating Profit Per Available Room)")
        st.markdown("""
**Formule** : `GOPPAR = EBITDA / Chambres disponibles`
**Description** : Profit operationnel brut par chambre disponible.
Mesure la performance operationnelle globale de l'hotel.
""")

        st.markdown("---")
        st.markdown("#### DSCR (Debt Service Coverage Ratio)")
        st.markdown("""
**Formule** : `DSCR = EBITDA / Service de la dette (mensualites totales)`
**Description** : Ratio de couverture du service de la dette.
- DSCR > 1.3x : la dette est confortablement couverte
- DSCR entre 1.0x et 1.3x : couverture limite
- DSCR < 1.0x : l'hotel ne genere pas assez pour couvrir la dette
""")

        if a1 is not None:
            st.markdown("---")
            st.markdown("#### Valeurs sur les premieres annees")
            nb_show = min(len(indic), 5)
            kpi_cols = ["Annee exploitation", "Taux Occupation", "Prix Moyen (ADR)",
                        "RevPAR", "GOPPAR", "DSCR"]
            kpi_df = indic[kpi_cols].head(nb_show).copy()
            kpi_df["Taux Occupation"] = kpi_df["Taux Occupation"].apply(lambda x: f"{x:.1f}%")
            kpi_df["Prix Moyen (ADR)"] = kpi_df["Prix Moyen (ADR)"].apply(lambda x: f"{x:,.2f} EUR")
            kpi_df["RevPAR"] = kpi_df["RevPAR"].apply(lambda x: f"{x:,.2f} EUR")
            kpi_df["GOPPAR"] = kpi_df["GOPPAR"].apply(lambda x: f"{x:,.2f} EUR")
            kpi_df["DSCR"] = kpi_df["DSCR"].apply(lambda x: f"{x:.2f}x" if x < 100 else "N/A")
            st.dataframe(kpi_df, hide_index=True, use_container_width=True)

    # =====================================================================
    # 11. INFLATION
    # =====================================================================
    elif categorie.startswith("11."):
        st.subheader("11. Inflation")

        st.markdown("""
**Formule** : `facteur_inflation = (1 + taux_annuel / 12) ^ mois_depuis_ouverture`
**Description** : L'inflation est appliquee de maniere cumulative mois par mois,
avec des taux differencies par type de charge/revenu.
""")

        infl_items = {
            "Ventes (hebergement, brasserie, bar, spa, salles)": params.get("inflation_ventes", 0.025),
            "Loyer restaurant": params.get("inflation_loyer_restaurant", 0.02),
            "Personnel": params.get("inflation_personnel", 0.025),
            "Loyer": params.get("inflation_loyer", 0.02),
            "Charges variables": params.get("inflation_charges_variables", 0.025),
            "Autres charges fixes": params.get("inflation_autres_charges", 0.02),
        }
        infl_rows = [{"Categorie": k, "Taux annuel": f"{v:.2%}"} for k, v in infl_items.items()]
        st.dataframe(pd.DataFrame(infl_rows), hide_index=True)

        # Exemple
        taux_ex = params.get("inflation_ventes", 0.025)
        mois_ex = 24
        facteur = (1 + taux_ex / 12) ** mois_ex
        st.markdown(f"""
---
#### Exemple de calcul
Apres **{mois_ex} mois** avec un taux d'inflation ventes de **{taux_ex:.2%}** :
`facteur = (1 + {taux_ex:.4f}/12)^{mois_ex} = {facteur:.6f}`
Soit une hausse cumulee de **{(facteur - 1) * 100:.2f}%**

Apres **60 mois** (5 ans) :
`facteur = (1 + {taux_ex:.4f}/12)^60 = {(1 + taux_ex / 12) ** 60:.6f}`
Soit une hausse cumulee de **{((1 + taux_ex / 12) ** 60 - 1) * 100:.2f}%**
""")


# ─── Module Immobiliere Rocher ─────────────────────────────────────────────────

def _module_rocher():
    """Module plan financier de l'Immobiliere Rocher."""
    from pathlib import Path as _Path
    plan_actif = st.session_state.get("plan_actif", "")
    p_chateau = st.session_state.get("params_charges", params_defaut())

    # Initialiser rocher_data dans p_chateau si absent (pour persistance)
    if "rocher_data" not in p_chateau:
        p_chateau["rocher_data"] = {
            "prets": [{"nom": "Pret banque", "creancier": "Banque", "montant": 11250000,
                       "taux_annuel": 0.04, "duree_ans": 15, "differe_mois": 0, "type": "annuite"}],
            "investissements": [
                {"categorie": "Terrain", "montant": 2951000, "duree_amort": 0},
                {"categorie": "Construction", "montant": 16194640, "duree_amort": 15},
            ],
            "fonds_propres_initial": 6750000,
            "fonds_propres_investisseurs": [{"nom": "Investisseur Rocher", "montant": 6750000}],
        }
    rd = p_chateau["rocher_data"]

    # Header
    col_t, col_s = st.columns([4, 1])
    with col_t:
        st.markdown(f'<h2>\U0001F3E2 Immobiliere Rocher <span style="font-size:0.5em; color:#666;">— {plan_actif}</span></h2>',
                    unsafe_allow_html=True)
    with col_s:
        c1, c2 = st.columns(2)
        with c1:
            if st.button("\U0001F4BE Sauver", key="rocher_save", use_container_width=True):
                sauvegarder_plan(plan_actif, p_chateau)
                st.toast("Sauvegarde !")
        with c2:
            if st.button("\U0001F504 Retour", key="rocher_back", use_container_width=True):
                st.session_state.pop("module_actif", None)
                st.rerun()

    # Parametres Rocher (lies au Chateau)
    loyer_mensuel = p_chateau.get("loyer_mensuel", 120000)
    pret_rocher = None
    for pr in p_chateau.get("prets", []):
        if "rocher" in pr.get("nom", "").lower():
            pret_rocher = pr
            break

    if pret_rocher is None:
        pret_rocher = {"nom": "Pret ROCHER", "montant": 2000000, "taux_annuel": 0.04,
                       "duree_ans": 15, "differe_mois": 0, "type": "interet_seul"}
        # Ajouter a la liste des prets du Chateau pour que les modifications persistent
        p_chateau.setdefault("prets", []).append(pret_rocher)

    tab_hyp, tab_res, tab_fin, tab_equil_r = st.tabs([
        "\U0001F4DD Hypotheses", "\U0001F4CA Resultats",
        "\U0001F3E6 Financement", "\U0001F4CA Equilibre Moyens / Besoins"])

    with tab_hyp:
        st.subheader("Parametres de l'Immobiliere Rocher")

        # Revenus lies au Chateau (modifiables ici, repercutes sur le Chateau)
        with st.expander("\U0001F4B0 **Revenus factures au Chateau d'Argenteau**", expanded=True):
            st.markdown("**Loyer mensuel facture au Chateau**")
            new_loyer = st.number_input("Loyer mensuel (\u20ac)", 0, 500_000,
                loyer_mensuel, step=500, key="rocher_loyer_edit")
            if new_loyer != loyer_mensuel:
                p_chateau["loyer_mensuel"] = new_loyer
                st.session_state["params_charges"] = p_chateau
                loyer_mensuel = new_loyer
                _auto_save(p_chateau, plan_actif)
            st.caption(f"Soit **{loyer_mensuel * 12:,.0f} \u20ac/an**. "
                       f"Ce montant est automatiquement repercute dans les charges fixes du Chateau.")

            st.markdown("---")
            st.markdown("**Pret octroye au Chateau**")
            c1, c2, c3, c4, c5 = st.columns(5)
            with c1:
                pr_mont = st.number_input("Montant (\u20ac)", 0, 20_000_000,
                    pret_rocher["montant"], step=100_000, key="rocher_pret_ch_m")
            with c2:
                pr_taux = st.number_input("Taux annuel", 0.0, 0.15,
                    pret_rocher["taux_annuel"], step=0.005, format="%.3f", key="rocher_pret_ch_t")
            with c3:
                pr_duree = st.number_input("Duree (ans)", 1, 30,
                    pret_rocher["duree_ans"], key="rocher_pret_ch_d")
            with c4:
                pr_differe = st.number_input("Differe (mois)", 0, 60,
                    pret_rocher.get("differe_mois", 0), key="rocher_pret_ch_df",
                    help="Pendant le differe, seuls les interets sont payes. "
                         "Duree totale du pret = differe + duree de remboursement.")
            with c5:
                _types_r = ["Mensualite constante", "Interet seul (in fine)"]
                _tidx = 1 if pret_rocher.get("type", "interet_seul") == "interet_seul" else 0
                pr_type_choisi = st.selectbox("Type", _types_r, index=_tidx, key="rocher_pret_ch_type")
                pr_type = "interet_seul" if "nteret" in pr_type_choisi else "annuite"

            # Synchroniser vers les prets du Chateau
            pret_rocher["montant"] = pr_mont
            pret_rocher["taux_annuel"] = pr_taux
            pret_rocher["duree_ans"] = pr_duree
            pret_rocher["differe_mois"] = pr_differe
            pret_rocher["type"] = pr_type
            st.session_state["params_charges"] = p_chateau

            taux_m_r = pr_taux / 12
            interets_recus_mois = pr_mont * taux_m_r
            _duree_totale_mois = pr_duree * 12 + pr_differe
            st.caption(f"Interets recus : **{interets_recus_mois:,.0f} \u20ac/mois**. "
                       f"Duree totale : **{_duree_totale_mois} mois ({_duree_totale_mois/12:.1f} ans)**"
                       f"{f' dont {pr_differe} mois de differe' if pr_differe > 0 else ''}. "
                       f"Les modifications sont automatiquement repercutees dans les prets du Chateau.")

            # Auto-save apres modification du pret Rocher au Chateau
            _auto_save(p_chateau, plan_actif)

        # Dettes Rocher
        with st.expander("\U0001F3E6 **Dettes Rocher**", expanded=True):
            _types_pret_r = ["Mensualite constante", "Interet seul (in fine)"]
            rocher_prets = rd.get("prets", [
                {"nom": "Pret banque", "creancier": "Banque", "montant": 11250000,
                 "taux_annuel": 0.04, "duree_ans": 15, "differe_mois": 0, "type": "annuite"},
            ])

            new_rocher_prets = []
            to_del_rp = []
            for i, pret in enumerate(rocher_prets):
                st.markdown(f"**{pret.get('nom', f'Pret {i+1}')}**")
                c1, c2, c3 = st.columns([2, 2, 0.5])
                with c1:
                    nom_rp = st.text_input("Nom", pret.get("nom", ""), key=f"rp_n_{i}")
                with c2:
                    creancier_rp = st.text_input("Creancier", pret.get("creancier", ""), key=f"rp_cr_{i}")
                with c3:
                    _visu = st.session_state.get("_visu_mode", False)
                    if st.button("\U0001F5D1", key=f"rp_del_{i}", help="Supprimer", disabled=_visu):
                        to_del_rp.append(i)

                c1, c2, c3, c4, c5 = st.columns(5)
                with c1:
                    montant_rp = st.number_input("Montant (\u20ac)", 0, 30_000_000,
                        pret["montant"], step=100_000, key=f"rp_m_{i}")
                with c2:
                    taux_rp = st.number_input("Taux annuel", 0.0, 0.15,
                        pret["taux_annuel"], step=0.005, format="%.3f", key=f"rp_t_{i}")
                with c3:
                    duree_rp = st.number_input("Duree (ans)", 1, 30,
                        pret["duree_ans"], key=f"rp_d_{i}")
                with c4:
                    differe_rp = st.number_input("Differe (mois)", 0, 36,
                        pret.get("differe_mois", 0), key=f"rp_df_{i}")
                # Subside RW (avant le type, car il force le type)
                subside_rw_rp = st.checkbox("Garanti par un subside RW",
                    value=pret.get("subside_rw", False), key=f"rp_sub_{i}",
                    help="A l'echeance, le subside rembourse le capital restant. "
                         "Force le type a 'Interet seul'.")

                with c5:
                    if subside_rw_rp:
                        st.selectbox("Type", _types_pret_r, index=1, key=f"rp_type_{i}", disabled=True)
                        type_val_rp = "interet_seul"
                    else:
                        type_idx_rp = 1 if pret.get("type", "annuite") == "interet_seul" else 0
                        type_choisi_rp = st.selectbox("Type", _types_pret_r, index=type_idx_rp, key=f"rp_type_{i}")
                        type_val_rp = "interet_seul" if "nteret" in type_choisi_rp else "annuite"

                # Mensualite
                taux_m_rp = taux_rp / 12
                nb_m_rp = duree_rp * 12
                if type_val_rp == "interet_seul":
                    mens_rp = montant_rp * taux_m_rp if montant_rp > 0 else 0
                else:
                    if taux_m_rp > 0 and montant_rp > 0:
                        mens_rp = montant_rp * taux_m_rp / (1 - (1 + taux_m_rp) ** -nb_m_rp)
                    elif montant_rp > 0:
                        mens_rp = montant_rp / nb_m_rp
                    else:
                        mens_rp = 0
                st.caption(f"Mensualite : **{mens_rp:,.0f} \u20ac**"
                           + (f" | Subside RW a l'echeance : **{montant_rp:,.0f} \u20ac**" if subside_rw_rp else ""))
                st.markdown("---")

                if i not in to_del_rp:
                    new_rocher_prets.append({
                        "nom": nom_rp, "creancier": creancier_rp, "montant": montant_rp,
                        "taux_annuel": taux_rp, "duree_ans": duree_rp,
                        "differe_mois": differe_rp, "type": type_val_rp,
                        "subside_rw": subside_rw_rp,
                    })

            if to_del_rp:
                rd["prets"] = new_rocher_prets
                st.rerun()

            _visu = st.session_state.get("_visu_mode", False)
            if st.button("\u2795 Ajouter un emprunt", key="rp_add", disabled=_visu):
                new_rocher_prets.append({
                    "nom": "Nouvel emprunt", "creancier": "", "montant": 0,
                    "taux_annuel": 0.04, "duree_ans": 15, "differe_mois": 0, "type": "annuite",
                })
                rd["prets"] = new_rocher_prets
                st.rerun()

            rd["prets"] = new_rocher_prets
            total_dettes_r = sum(pr["montant"] for pr in new_rocher_prets)
            st.info(f"**Total dettes Rocher : {total_dettes_r:,.0f} \u20ac**")

            # Sync rocher_data dans p_chateau
            p_chateau["rocher_data"] = rd
            st.session_state["params_charges"] = p_chateau
            _auto_save(p_chateau, plan_actif)

        # Investissements / Amortissements Rocher
        with st.expander("\U0001F3D7 **Investissements & Amortissements**", expanded=False):
            rocher_inv = rd.get("investissements", [
                {"categorie": "Terrain", "montant": 2951000, "duree_amort": 0},
                {"categorie": "Construction", "montant": 16194640, "duree_amort": 15},
            ])
            new_rinv = []
            hcols = st.columns([4, 2, 1.5, 2])
            hcols[0].markdown("**Categorie**"); hcols[1].markdown("**Montant**")
            hcols[2].markdown("**Duree amort.**"); hcols[3].markdown("**Amort. annuel**")
            for i, inv in enumerate(rocher_inv):
                c1, c2, c3, c4 = st.columns([4, 2, 1.5, 2])
                with c1:
                    cat = st.text_input("c", inv["categorie"], key=f"rinv_c_{i}", label_visibility="collapsed")
                with c2:
                    mont = st.number_input("m", 0, 30_000_000, inv["montant"], step=50_000, key=f"rinv_m_{i}", label_visibility="collapsed")
                with c3:
                    dur = st.number_input("d", 0, 30, inv["duree_amort"], key=f"rinv_d_{i}", label_visibility="collapsed")
                with c4:
                    amort = mont / dur if dur > 0 else 0
                    st.markdown(f"**{amort:,.0f} \u20ac**" if dur > 0 else "*Non amorti*")
                new_rinv.append({"categorie": cat, "montant": mont, "duree_amort": dur})
            rd["investissements"] = new_rinv

        # Fonds propres
        with st.expander("\U0001F4B0 **Fonds propres**", expanded=False):
            rocher_fp_params = {
                "fonds_propres_initial": rd.get("fonds_propres_initial", 6750000),
                "fonds_propres_investisseurs": rd.get("fonds_propres_investisseurs",
                    [{"nom": "Investisseur Rocher", "montant": 6750000}]),
            }
            _fonds_propres_widget(rocher_fp_params, key_prefix="rocher_fp")
            rd["fonds_propres_initial"] = rocher_fp_params["fonds_propres_initial"]
            rd["fonds_propres_investisseurs"] = rocher_fp_params.get("fonds_propres_investisseurs", [])

    # ── Projection Rocher (partagee entre Resultats et Financement) ──
    nb_mois = p_chateau.get("nb_mois_projection", 204)
    date_ouv = p_chateau["date_ouverture"]

    # Projection mensuelle
    rocher_prets_calc = rd.get("prets", [
        {"nom": "Pret banque", "montant": 11250000, "taux_annuel": 0.04,
         "duree_ans": 15, "differe_mois": 0, "type": "annuite"},
    ])
    rocher_inv = rd.get("investissements", [
        {"categorie": "Terrain", "montant": 2951000, "duree_amort": 0},
        {"categorie": "Construction", "montant": 16194640, "duree_amort": 15},
    ])
    amort_mensuel = sum(inv["montant"] / inv["duree_amort"] / 12
                       for inv in rocher_inv if inv["duree_amort"] > 0)

    # Tableaux d'amortissement de toutes les dettes Rocher
    dfs_prets_rocher = {}
    for pr in rocher_prets_calc:
        if pr["montant"] > 0:
            dfs_prets_rocher[pr["nom"]] = calc_tableau_pret(pr, date_ouv, nb_mois)

    # Tableau pret au Chateau (interets/capital recus)
    df_pret_chateau = calc_tableau_pret(pret_rocher, date_ouv, nb_mois)

    rows_r = []
    from dateutil.relativedelta import relativedelta as _rd
    for m in range(nb_mois):
        d = date_ouv + _rd(months=m)
        rev_loyer = loyer_mensuel
        row_ch = df_pret_chateau[df_pret_chateau["date"] == d]
        int_recus = row_ch.iloc[0]["interets"] if not row_ch.empty else 0
        cap_recu = row_ch.iloc[0]["capital"] if not row_ch.empty else 0

        int_payes = 0
        cap_paye = 0
        for _pn, _df_pr in dfs_prets_rocher.items():
            _row = _df_pr[_df_pr["date"] == d]
            if not _row.empty:
                int_payes += _row.iloc[0]["interets"]
                cap_paye += _row.iloc[0]["capital"]

        # Subside RW Rocher : 1/15 du montant par an au resultat, a partir de An 2
        # Pas d'impact cash (neutre)
        annee_idx_r = m // 12
        subside_rw = 0
        if annee_idx_r >= 1:
            for _pr in rocher_prets_calc:
                if _pr.get("subside_rw", False):
                    duree_sub_r = 15  # 15 ans pour le Rocher
                    annee_sub_r = annee_idx_r - 1
                    if annee_sub_r < duree_sub_r:
                        subside_rw += _pr["montant"] / duree_sub_r / 12

        resultat = rev_loyer + int_recus - int_payes - amort_mensuel + subside_rw
        cash = rev_loyer + int_recus + cap_recu - int_payes - cap_paye

        rows_r.append({
            "date": d, "annee": d.year,
            "loyer": rev_loyer, "interets_recus": int_recus, "capital_recu": cap_recu,
            "interets_payes": int_payes, "capital_paye": cap_paye,
            "amortissement": amort_mensuel,
            "subside_rw": subside_rw,
            "resultat": resultat, "cash_flow": cash,
        })

    df_r = pd.DataFrame(rows_r)
    df_r["cash_flow_cumul"] = df_r["cash_flow"].cumsum()
    df_r["resultat_cumul"] = df_r["resultat"].cumsum()

    # Tableau annuel
    annual_r = df_r.groupby("annee").agg({
        "loyer": "sum", "interets_recus": "sum", "capital_recu": "sum",
        "interets_payes": "sum", "capital_paye": "sum",
        "amortissement": "sum", "subside_rw": "sum",
        "resultat": "sum", "cash_flow": "sum",
        "cash_flow_cumul": "last", "resultat_cumul": "last",
    }).reset_index()

    with tab_res:
        st.subheader("Resultats Immobiliere Rocher")
        st.dataframe(pd.DataFrame({
            "Annee": [str(int(a)) for a in annual_r["annee"]],
            "Loyer recu": [f"{v:,.0f} \u20ac" for v in annual_r["loyer"]],
            "Interets recus": [f"{v:,.0f} \u20ac" for v in annual_r["interets_recus"]],
            "Subside RW": [f"{v:,.0f} \u20ac" for v in annual_r["subside_rw"]],
            "Capital recu (Chateau)": [f"{v:,.0f} \u20ac" for v in annual_r["capital_recu"]],
            "Interets payes": [f"{v:,.0f} \u20ac" for v in annual_r["interets_payes"]],
            "Amortissement": [f"{v:,.0f} \u20ac" for v in annual_r["amortissement"]],
            "Resultat": [f"{v:,.0f} \u20ac" for v in annual_r["resultat"]],
            "Resultat cumule": [f"{v:,.0f} \u20ac" for v in annual_r["resultat_cumul"]],
            "Remb. capital paye": [f"{v:,.0f} \u20ac" for v in annual_r["capital_paye"]],
            "Cash Flow": [f"{v:,.0f} \u20ac" for v in annual_r["cash_flow"]],
            "Cash Flow Cumul": [f"{v:,.0f} \u20ac" for v in annual_r["cash_flow_cumul"]],
        }), use_container_width=True, hide_index=True)

        # Graphique
        x_r = [str(int(a)) for a in annual_r["annee"]]
        fig_r = go.Figure()
        fig_r.add_trace(go.Bar(x=x_r, y=annual_r["resultat"], name="Resultat", marker_color="#11998e"))
        fig_r.add_trace(go.Bar(x=x_r, y=annual_r["cash_flow"], name="Cash Flow", marker_color="#4facfe"))
        fig_r.add_trace(go.Scatter(x=x_r, y=annual_r["cash_flow_cumul"],
            name="Cash Flow Cumule", mode="lines+markers", line=dict(color="#f5576c", width=3)))
        fig_r.add_hline(y=0, line_dash="dash", line_color="gray")
        fig_r.update_layout(
            title="Resultat et Cash Flow annuel - Immobiliere Rocher",
            height=500, barmode="group",
            xaxis=dict(type="category", tickfont=dict(size=12)),
            yaxis=dict(tickformat=",.0f"),
            legend=dict(orientation="h", y=-0.12, font=dict(size=13), xanchor="center", x=0.5),
        )
        st.plotly_chart(fig_r, use_container_width=True, config={"displayModeBar": False})

    # ── Onglet Financement ──
    with tab_fin:
        st.subheader("Encours de financement")

        nb_mois_f = p_chateau.get("nb_mois_projection", 204)
        date_ouv_f = p_chateau["date_ouverture"]
        rocher_prets_f = rd.get("prets", [
            {"nom": "Pret banque", "montant": 11250000, "taux_annuel": 0.04,
             "duree_ans": 15, "differe_mois": 0, "type": "annuite"},
        ])
        fonds_propres_r = rd.get("fonds_propres_initial", 6750000)

        # Pret au Chateau
        pret_ch_f = None
        for pr in p_chateau.get("prets", []):
            if "rocher" in pr.get("nom", "").lower():
                pret_ch_f = pr
                break
        if pret_ch_f is None:
            pret_ch_f = {"nom": "Pret ROCHER", "montant": 2000000, "taux_annuel": 0.04,
                         "duree_ans": 15, "differe_mois": 0, "type": "interet_seul"}

        from dateutil.relativedelta import relativedelta as _rd_f

        # Tableaux d'amortissement
        dfs_enc = {}
        for pr in rocher_prets_f:
            if pr["montant"] > 0:
                dfs_enc[pr["nom"]] = calc_tableau_pret(pr, date_ouv_f, nb_mois_f)
        df_pret_ch_f = calc_tableau_pret(pret_ch_f, date_ouv_f, nb_mois_f)

        rows_enc = []
        for m in range(nb_mois_f):
            d = date_ouv_f + _rd_f(months=m)
            enc_dettes = {}
            total_enc_dette = 0
            for pn, df_p in dfs_enc.items():
                row_p = df_p[df_p["date"] == d]
                enc = row_p.iloc[0]["capital_restant"] if not row_p.empty else 0
                enc_dettes[pn] = enc
                total_enc_dette += enc

            row_ch_f = df_pret_ch_f[df_pret_ch_f["date"] == d]
            enc_pret_chateau = row_ch_f.iloc[0]["capital_restant"] if not row_ch_f.empty else 0

            rows_enc.append({
                "date": d, "annee": d.year,
                **{f"enc_{pn}": v for pn, v in enc_dettes.items()},
                "enc_dette_total": total_enc_dette,
                "enc_pret_chateau": enc_pret_chateau,
                "fonds_propres": fonds_propres_r,
            })

        df_enc = pd.DataFrame(rows_enc)
        annual_enc = df_enc.groupby("annee").last().reset_index()
        x_enc = [str(int(a)) for a in annual_enc["annee"]]
        noms_dettes = [pr["nom"] for pr in rocher_prets_f if pr["montant"] > 0]

        # ── PASSIF : Fonds propres + resultats cumules ──
        st.markdown('<div class="section-header">\U0001F7E2 Passif — Fonds propres</div>', unsafe_allow_html=True)
        st.caption("Fonds propres initiaux + resultats cumules de chaque annee")
        investisseurs_r = rd.get("fonds_propres_investisseurs", [{"nom": "Investisseur", "montant": fonds_propres_r}])

        # Resultat cumule par annee (depuis annual_r)
        res_cumul_an = list(annual_r["resultat_cumul"]) if "resultat_cumul" in annual_r.columns else [0] * len(x_enc)

        fp_cols = {"Annee": x_enc}
        for inv_fp in investisseurs_r:
            fp_cols[inv_fp["nom"]] = [f"{inv_fp['montant']:,.0f} \u20ac"] * len(x_enc)
        fp_cols["Fonds propres initiaux"] = [f"{fonds_propres_r:,.0f} \u20ac"] * len(x_enc)
        fp_cols["Resultat cumule"] = [f"{v:,.0f} \u20ac" for v in res_cumul_an]
        capitaux_propres = [fonds_propres_r + rc for rc in res_cumul_an]
        fp_cols["Capitaux propres"] = [f"{v:,.0f} \u20ac" for v in capitaux_propres]
        st.dataframe(pd.DataFrame(fp_cols), use_container_width=True, hide_index=True)

        # Graphique capitaux propres
        fig_fp_r = go.Figure()
        fig_fp_r.add_trace(go.Bar(x=x_enc, y=[fonds_propres_r / 1000] * len(x_enc),
            name="Fonds propres initiaux", marker_color="#38ef7d",
            hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_fp_r.add_trace(go.Bar(x=x_enc, y=[rc / 1000 for rc in res_cumul_an],
            name="Resultat cumule", marker_color="#4facfe" if res_cumul_an[-1] >= 0 else "#f5576c",
            hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_fp_r.add_trace(go.Scatter(x=x_enc, y=[v / 1000 for v in capitaux_propres],
            name="Capitaux propres", mode="lines+markers",
            line=dict(color="#667eea", width=3),
            hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_fp_r.update_layout(
            title="Evolution des capitaux propres (K\u20ac)",
            height=400, barmode="stack",
            xaxis=dict(type="category", tickfont=dict(size=12)),
            yaxis=dict(tickformat=",.0f", title="K\u20ac"),
            legend=dict(orientation="h", y=-0.12, font=dict(size=13), xanchor="center", x=0.5),
        )
        st.plotly_chart(fig_fp_r, use_container_width=True, config={"displayModeBar": False})

        st.markdown("---")

        # ── PASSIF : Encours des dettes ──
        st.markdown('<div class="section-header">\U0001F534 Passif — Encours des dettes</div>', unsafe_allow_html=True)
        st.caption("Capital restant du sur chaque emprunt contracte par l'Immobiliere Rocher.")

        passif_cols = {"Annee": x_enc}
        for pn in noms_dettes:
            col_key = f"enc_{pn}"
            if col_key in annual_enc.columns:
                passif_cols[pn] = [f"{v:,.0f} \u20ac" for v in annual_enc[col_key]]
        passif_cols["Total dettes"] = [f"{v:,.0f} \u20ac" for v in annual_enc["enc_dette_total"]]
        st.dataframe(pd.DataFrame(passif_cols), use_container_width=True, hide_index=True)

        # Graphique passif
        fig_passif = go.Figure()
        colors_p = ["#f5576c", "#ff8c00", "#667eea", "#f093fb", "#764ba2"]
        for idx_p, pn in enumerate(noms_dettes):
            col_key = f"enc_{pn}"
            if col_key in annual_enc.columns:
                fig_passif.add_trace(go.Scatter(x=x_enc, y=annual_enc[col_key] / 1000,
                    name=pn, mode="lines+markers",
                    line=dict(color=colors_p[idx_p % len(colors_p)]),
                    hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_passif.add_trace(go.Scatter(x=x_enc, y=annual_enc["enc_dette_total"] / 1000,
            name="Total dettes", mode="lines+markers",
            line=dict(color="#f5576c", width=3, dash="dash"),
            hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_passif.update_layout(
            title="Passif — Evolution des encours de dettes (K\u20ac)",
            height=450,
            xaxis=dict(type="category", tickfont=dict(size=12)),
            yaxis=dict(tickformat=",.0f", title="K\u20ac"),
            legend=dict(orientation="h", y=-0.12, font=dict(size=13), xanchor="center", x=0.5),
        )
        st.plotly_chart(fig_passif, use_container_width=True, config={"displayModeBar": False})

        st.markdown("---")

        # ── ACTIF : Encours du pret au Chateau ──
        st.markdown('<div class="section-header">\U0001F7E2 Actif — Pret octroye au Chateau d\'Argenteau</div>', unsafe_allow_html=True)
        st.caption(f"Capital restant du par le Chateau d'Argenteau a l'Immobiliere Rocher. "
                   f"Pret de {pret_ch_f['montant']:,.0f} \u20ac a {pret_ch_f['taux_annuel']:.1%} "
                   f"sur {pret_ch_f['duree_ans']} ans ({pret_ch_f.get('type', 'interet_seul')}).")

        actif_cols = {
            "Annee": x_enc,
            "Encours pret au Chateau": [f"{v:,.0f} \u20ac" for v in annual_enc["enc_pret_chateau"]],
        }
        st.dataframe(pd.DataFrame(actif_cols), use_container_width=True, hide_index=True)

        # Graphique actif
        fig_actif = go.Figure()
        fig_actif.add_trace(go.Scatter(x=x_enc, y=annual_enc["enc_pret_chateau"] / 1000,
            name="Pret au Chateau", mode="lines+markers",
            line=dict(color="#11998e", width=3),
            fill="tozeroy", fillcolor="rgba(17,153,142,0.1)",
            hovertemplate="%{x}<br>%{y:,.0f} K\u20ac<extra></extra>"))
        fig_actif.update_layout(
            title="Actif — Encours du pret au Chateau (K\u20ac)",
            height=350,
            xaxis=dict(type="category", tickfont=dict(size=12)),
            yaxis=dict(tickformat=",.0f", title="K\u20ac"),
            showlegend=False,
        )
        st.plotly_chart(fig_actif, use_container_width=True, config={"displayModeBar": False})

    # ── Equilibre Moyens / Besoins Rocher ──
    with tab_equil_r:
        st.subheader("Equilibre Moyens / Besoins — Immobiliere Rocher")

        # Moyens Rocher = dettes + capital
        _rocher_prets = rd.get("prets", [])
        _rocher_fp = rd.get("fonds_propres_initial", 0)
        _total_dettes_r = sum(pr["montant"] for pr in _rocher_prets)
        _total_moyens_r = _total_dettes_r + _rocher_fp

        # Besoins Rocher = investissements + pret au chateau
        _rocher_inv = rd.get("investissements", [])
        _total_inv_r = sum(inv["montant"] for inv in _rocher_inv)
        _pret_au_chateau = pret_rocher["montant"] if pret_rocher else 0
        _total_besoins_r = _total_inv_r + _pret_au_chateau
        _solde_r = _total_moyens_r - _total_besoins_r

        col_m_r, col_b_r = st.columns(2)
        with col_m_r:
            st.markdown("#### Moyens (ressources)")
            for pr in _rocher_prets:
                st.markdown(f"- {pr['nom']} : **{pr['montant']:,.0f} \u20ac**")
            st.markdown(f"- Fonds propres : **{_rocher_fp:,.0f} \u20ac**")
            st.markdown(f"**TOTAL MOYENS : {_total_moyens_r:,.0f} \u20ac**")

        with col_b_r:
            st.markdown("#### Besoins (emplois)")
            for inv in _rocher_inv:
                st.markdown(f"- {inv['categorie']} : **{inv['montant']:,.0f} \u20ac**")
            if _pret_au_chateau > 0:
                st.markdown(f"- Pret octroye au Chateau : **{_pret_au_chateau:,.0f} \u20ac**")
            st.markdown(f"**TOTAL BESOINS : {_total_besoins_r:,.0f} \u20ac**")

        if _solde_r >= 0:
            st.success(f"**Solde : +{_solde_r:,.0f} \u20ac** — Les moyens couvrent les besoins"
                       + (f" avec une marge de {_solde_r/_total_besoins_r*100:.1f}%." if _total_besoins_r > 0 else "."))
        else:
            st.error(f"**Deficit : {_solde_r:,.0f} \u20ac** — Il manque **{abs(_solde_r):,.0f} \u20ac**.")

        # Waterfall Rocher : moyens = prets individuels + FP, besoins = invest + pret chateau
        _wf_prets_r = [{"nom": pr["nom"], "montant": pr["montant"]} for pr in _rocher_prets]
        # Ajouter le pret au chateau comme "emploi" dans les besoins
        # On reconstruit avec investissements detailles + pret chateau
        _wf_labels_r = []
        _wf_values_r = []
        _wf_measures_r = []

        # Moyens
        _nom_c_r = {}
        for pr in _rocher_prets:
            _nom_c_r[pr["nom"]] = _nom_c_r.get(pr["nom"], 0) + 1
        _nom_i_r = {}
        for pr in _rocher_prets:
            _n = pr["nom"]
            if _nom_c_r[_n] > 1:
                _nom_i_r[_n] = _nom_i_r.get(_n, 0) + 1
                lbl = f"{_n} ({_nom_i_r[_n]})"
            else:
                lbl = _n
            _wf_labels_r.append(lbl)
            _wf_values_r.append(pr["montant"])
            _wf_measures_r.append("relative")
        _wf_labels_r.append("Fonds propres")
        _wf_values_r.append(_rocher_fp)
        _wf_measures_r.append("relative")

        _wf_labels_r.append("TOTAL MOYENS")
        _wf_values_r.append(_total_moyens_r)
        _wf_measures_r.append("total")

        # Besoins
        for inv in _rocher_inv:
            _wf_labels_r.append(inv["categorie"])
            _wf_values_r.append(-inv["montant"])
            _wf_measures_r.append("relative")
        if _pret_au_chateau > 0:
            _wf_labels_r.append("Pret au Chateau")
            _wf_values_r.append(-_pret_au_chateau)
            _wf_measures_r.append("relative")

        _wf_labels_r.append("SOLDE")
        _wf_values_r.append(_solde_r)
        _wf_measures_r.append("total")

        fig_wf_r = go.Figure(go.Waterfall(
            x=_wf_labels_r, y=[v / 1000 for v in _wf_values_r],
            measure=_wf_measures_r,
            connector=dict(line=dict(color="rgba(0,0,0,0.3)", width=1)),
            increasing=dict(marker=dict(color="#38ef7d")),
            decreasing=dict(marker=dict(color="#f5576c")),
            totals=dict(marker=dict(color="#667eea")),
            textposition="outside",
            text=[f"<b>{v/1000:,.0f} K\u20ac</b>" for v in _wf_values_r],
            textfont=dict(size=14, color="#1a1a2e"),
            hovertemplate="%{x}<br><b>%{y:,.0f} K\u20ac</b><extra></extra>",
        ))
        fig_wf_r.update_layout(
            title=dict(text="Cascade Moyens / Besoins — Rocher (K\u20ac)", font=dict(size=16, color="#1a1a2e")),
            height=550,
            xaxis=dict(tickfont=dict(size=12, color="#333"), tickangle=-30),
            yaxis=dict(tickformat=",.0f", title="K\u20ac", tickfont=dict(size=12)),
            showlegend=False,
            margin=dict(t=60, b=90),
            uniformtext=dict(minsize=12, mode="show"),
        )
        st.plotly_chart(fig_wf_r, use_container_width=True, config={"displayModeBar": False})

    # Auto-save Rocher
    _auto_save(p_chateau, plan_actif)


# ─── Rapport complet ──────────────────────────────────────────────────────────

def _text_area_height(text, min_h=80, max_h=400, chars_per_line=90, line_h=22):
    """Calcule la hauteur d'un text_area en fonction de la longueur du texte."""
    if not text:
        return min_h
    nb_lines = sum(1 + len(line) // chars_per_line for line in text.split("\n"))
    return max(min_h, min(max_h, nb_lines * line_h + 40))


def _comment_slot(p, plan_nom, slot_id, print_mode=False):
    """Affiche un emplacement de commentaire dans le rapport.
    - Edit mode : bouton + pour ajouter, text_area pour editer, bouton supprimer
    - Visu mode : affiche le commentaire s'il existe (lecture seule)
    - Print mode : affiche le commentaire s'il existe, pas de boutons
    """
    comments = p.setdefault("commentaires", {})
    _is_visu = st.session_state.get("auth_mode") == "visu"
    has_comment = slot_id in comments and comments[slot_id].strip()
    editing_key = f"_editing_comment_{slot_id}"

    if print_mode:
        # Print : afficher le commentaire s'il existe, rien sinon
        if has_comment:
            st.markdown(f'<div style="background:#f8f9fa; border-left:3px solid #667eea; '
                        f'padding:10px 14px; border-radius:4px; margin:8px 0; '
                        f'font-size:0.93em; line-height:1.6; color:#2c3e50; font-style:italic;">'
                        f'{comments[slot_id]}</div>', unsafe_allow_html=True)
        return

    if _is_visu:
        # Visu : afficher le commentaire s'il existe, pas de boutons
        if has_comment:
            st.markdown(f'<div style="background:#f8f9fa; border-left:3px solid #667eea; '
                        f'padding:10px 14px; border-radius:4px; margin:8px 0; '
                        f'font-size:0.93em; line-height:1.6; color:#2c3e50; font-style:italic;">'
                        f'{comments[slot_id]}</div>', unsafe_allow_html=True)
        return

    # Edit mode
    if has_comment or st.session_state.get(editing_key):
        _val = comments.get(slot_id, "")
        _new = st.text_area("Commentaire", value=_val, height=_text_area_height(_val),
                            key=f"ta_cmt_{slot_id}", label_visibility="collapsed")
        _c1, _c2, _ = st.columns([1, 1, 4])
        with _c1:
            if st.button("\U0001F4BE Enregistrer", key=f"save_cmt_{slot_id}",
                         use_container_width=True):
                comments[slot_id] = _new
                p["commentaires"] = comments
                sauvegarder_plan(plan_nom, p)
                st.session_state.pop(editing_key, None)
                st.toast("Commentaire enregistre")
                st.rerun()
        with _c2:
            if st.button("\U0001F5D1 Supprimer", key=f"del_cmt_{slot_id}",
                         use_container_width=True):
                comments.pop(slot_id, None)
                p["commentaires"] = comments
                sauvegarder_plan(plan_nom, p)
                st.session_state.pop(editing_key, None)
                st.toast("Commentaire supprime")
                st.rerun()
    else:
        if st.button("\u270f\ufe0f + Commentaire", key=f"add_cmt_{slot_id}",
                     help="Ajouter un commentaire ici"):
            st.session_state[editing_key] = True
            st.rerun()


def _render_rapport_complet(plan_nom, _Path, print_mode=False):
    """Genere et affiche le rapport complet d'un plan en plein ecran."""

    _color_remap_print = {
        "#38ef7d": "#2dbe68", "#f5576c": "#e04458", "#4facfe": "#3a8fd4",
        "#667eea": "#4a63c8", "#11998e": "#0e7d74", "#ffcc00": "#e0b400",
        "#a0522d": "#8b4726", "#764ba2": "#643e8c", "#f093fb": "#c06ec8",
    }

    _fig_counter = [0]  # compteur mutable pour les slots de graphiques

    # Methodologies par graphique
    _methodologies = {
        "ro_pnl": "**Produits** = Loyer facture au Chateau + Interets recus sur pret intra-groupe.\n\n"
                  "**Charges** = Interets payes sur dettes bancaires + Amortissement des investissements.\n\n"
                  "**Resultat** = Produits - Charges.",
        "ro_cashflow": "**Cash periode** = Loyer recu + Interets recus + Capital rembourse par le Chateau "
                       "- Interets payes aux banques - Capital rembourse aux banques.\n\n"
                       "**Cash cumule** = Somme des cash periodes depuis le debut.",
        "ro_fp": "**FP initiaux** = Capital apporte par les actionnaires.\n\n"
                 "**Capitaux propres** = FP initiaux + Resultat net cumule.",
        "ro_endettement": "**Endettement** = Capital restant du sur chaque emprunt a la fin de chaque annee.\n\n"
                          "Ventile par pret. Pour les prets garantis par un subside RW, le capital "
                          "est rembourse en totalite par le subside a la fin de la duree du pret "
                          "(operation neutre au niveau du cash : perception du subside = remboursement de la dette).",
        "ch_ventes": "**Ventes par service** = CA Hebergement (nuitees x ADR) + CA Brasserie (couverts + PDJ) "
                     "+ CA Bar (nuitees x taux x prix) + CA Spa (soins + entrees) + CA Salles (seminaires + mariages).\n\n"
                     "Tous les CA sont ajustes par l'inflation annuelle.",
        "ch_nuitees": "**Nuitees** = Nombre de chambres x 365 jours x Taux d'occupation x Saisonnalite mensuelle.",
        "ch_occ_adr": "**Taux d'occupation** = Nuitees vendues / Nuitees disponibles.\n\n"
                      "**ADR** (Average Daily Rate) = CA Hebergement / Nuitees vendues.\n\n"
                      "**RevPAR** = CA Hebergement / Nuitees disponibles = ADR x Taux d'occupation.",
        "ch_cv": "**Charges variables** = Couts proportionnels au volume d'activite.\n\n"
                 "Hebergement : cout par nuitee (linge, amenities, energie...) + commission CB.\n"
                 "Brasserie : food cost (% du CA).\n"
                 "Bar : beverage cost (% du CA) + consommables.\n"
                 "Spa : cout des soins + produits.",
        "ch_cf_directs": "**Frais fixes directs** = Masse salariale par departement + charges fixes directes "
                         "(entretien, assurances specifiques...).\n\n"
                         "Repartis proportionnellement entre Hebergement, Brasserie, Spa et Evenements.",
        "ch_marge": "**Marge par service** = CA du service - Charges variables du service.\n\n"
                    "Le subside RW est ajoute comme revenu supplementaire.\n\n"
                    "**Total** = Somme de toutes les marges + Subside RW.",
        "ch_marge_pie": "**Repartition des marges** = Part de chaque service dans le total des marges cumulees "
                        "sur l'ensemble de la periode de projection.",
        "ch_cf_indirects": "**Frais fixes indirects** = Personnel indirect (direction, admin, marketing...) "
                           "+ charges fixes non allouees a un service (loyer, assurances generales, IT...).",
        "ch_rh_pie": "**Repartition du personnel** = Nombre d'ETP (Equivalents Temps Plein) par departement.",
        "ch_ebitda_rn": "**EBITDA** = Marge totale - Frais fixes directs - Frais fixes indirects.\n\n"
                        "**Resultat Net** = EBITDA - Amortissement - Interets sur emprunts - Impot.",
        "ch_cashflow": "**Cash flow** = Resultat Net + Amortissement - Remboursement capital emprunts.\n\n"
                       "**Tresorerie** = Surplus initial (FP + Emprunts - Investissements) + Cash flow cumule.",
        "ch_fp": "**Fonds propres** = Capital initial + Resultat net cumule.\n\n"
                 "**FP + Quasi-FP** = Fonds propres + Pret Rocher (considere comme quasi-fonds propres).",
        "ch_endettement": "**Endettement** = Capital restant du sur chaque emprunt a la fin de chaque annee.\n\n"
                          "Ventile par pret (banque, Rocher, subside RW).",
        "ch_solvabilite": "**Ratio solvabilite** = (FP + Quasi-FP) / Total passif x 100.\n\n"
                          "Numerateur = Capital + Resultat cumule + Pret Rocher.\n"
                          "Denominateur = Toutes dettes (encours) + Capital + Resultat cumule.\n\n"
                          "Seuil de 30% = niveau minimum recommande.",
    }

    def _cmt(slot_id):
        """Raccourci pour inserer un slot de commentaire."""
        _comment_slot(p, plan_nom, slot_id, print_mode)

    def _show_fig(fig, key="", large=False):
        """Affiche un graphique + slot de commentaire apres."""
        _fig_counter[0] += 1
        _slot_name = key if key else f"graph_{_fig_counter[0]}"
        if not print_mode:
            fig.update_layout(margin=dict(l=50, r=30, t=40, b=40))
            # Activer le bouton copie dans la config Plotly
            _cfg_copy = {"displayModeBar": True, "modeBarButtonsToRemove": [
                "zoom2d","pan2d","select2d","lasso2d","zoomIn2d","zoomOut2d",
                "autoScale2d","resetScale2d","hoverClosestCartesian",
                "hoverCompareCartesian","toggleSpikelines"],
                "modeBarButtonsToAdd": ["toImage"],
                "displaylogo": False, "toImageButtonOptions": {
                    "format": "png", "filename": _slot_name, "height": 600, "width": 1000, "scale": 2}}
            st.plotly_chart(fig, use_container_width=True, config=_cfg_copy)
            # Icone methodologie (popup)
            if _slot_name in _methodologies:
                with st.popover("\u2139\ufe0f Methodologie"):
                    st.markdown(f"**{fig.layout.title.text if fig.layout.title and fig.layout.title.text else _slot_name}**")
                    st.markdown(_methodologies[_slot_name])
            _cmt(f"after_{_slot_name}")
            return
        # ── Print mode : largeur fixe pour tenir dans la page PDF ──
        if large:
            _w, _h = 600, 500
        else:
            _w, _h = 700, 360
        fig.update_layout(
            width=_w, height=_h,
            margin=dict(l=50, r=50, t=35, b=55),
            title_font_size=12, font_size=9,
            legend=dict(font_size=8, orientation="h", y=-0.2, x=0.5, xanchor="center"),
            paper_bgcolor="white", plot_bgcolor="white",
        )
        # Centrer le graphique
        _, _print_col, _ = st.columns([1, 4, 1])
        with _print_col:
            st.plotly_chart(fig, use_container_width=False, config=_cfg)
        _cmt(f"after_{_slot_name}")

    if print_mode:
        st.markdown("""<style>
            /* Forcer les couleurs meme hors print */
            *, *::before, *::after {
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
                color-adjust: exact !important;
            }
            iframe { width: 100% !important; }
            @media print {
                *, *::before, *::after {
                    -webkit-print-color-adjust: exact !important;
                    print-color-adjust: exact !important;
                    color-adjust: exact !important;
                }
                body { font-size: 9pt !important; line-height: 1.3 !important; }
                h1 { font-size: 16pt !important; margin: 8px 0 !important; }
                h2 { font-size: 13pt !important; margin: 6px 0 !important;
                     page-break-before: always !important; }
                h3 { font-size: 11pt !important; margin: 4px 0 !important; }
                p, li { margin: 2px 0 !important; font-size: 9pt !important; }
                table { font-size: 8pt !important; }
                td, th { padding: 2px 6px !important; }
                hr { margin: 4px 0 !important; }
                iframe { width: 100% !important; }
                button, [data-testid="stButton"], header, footer,
                [data-testid="stSidebar"], .stDeployButton,
                [data-testid="stAlert"], .stAlert { display: none !important; }
            }
        </style>""", unsafe_allow_html=True)

    try:
        p = charger_plan(plan_nom)
        df_rpt = projection_complete(p)
        _ann = df_rpt.groupby("annee").agg({
            "ca_total": "sum", "ca_hebergement": "sum", "ca_brasserie": "sum",
            "ca_bar": "sum", "ca_spa": "sum", "ca_salles": "sum",
            "ca_loyer_restaurant": "sum", "ca_divers": "sum",
            "cv_total": "sum", "cv_hebergement": "sum", "cv_brasserie": "sum",
            "cv_bar": "sum", "cv_spa": "sum",
            "cf_directs_total": "sum",
            "cf_directs_hebergement": "sum", "cf_directs_brasserie": "sum",
            "cf_directs_bar": "sum", "cf_directs_spa": "sum", "cf_directs_evenements": "sum",
            "cf_indirects_total": "sum", "cf_total": "sum", "cf_total_cash": "sum",
            "marge_brute": "sum", "marge": "sum", "subside_rw": "sum",
            "ebitda": "sum", "amortissement": "sum", "ebit": "sum",
            "dette_interets": "sum", "dette_capital": "sum",
            "impot": "sum", "impot_cash": "sum",
            "tva_paiement": "sum", "reinvest_acquisition": "sum",
            "delay_adjustment": "sum",
            "resultat_net": "sum", "cash_flow": "sum", "cash_flow_cumul": "last",
            "nuitees": "sum", "taux_occupation": "mean",
        }).reset_index()
        _x = [str(int(a)) for a in _ann["annee"]]
        K = lambda v: v / 1000
        _leg = dict(orientation="h", y=-0.15, xanchor="center", x=0.5, font=dict(size=12 if not print_mode else 10))
        _cfg = {"displayModeBar": False}
        _h = lambda h: int(h * 0.7) if print_mode else h  # Hauteur reduite en mode print

        rd = p.get("rocher_data", {})
        fp_ch = p.get("fonds_propres_initial", 0)
        prets_ch = p.get("prets", [])
        fp_ro = rd.get("fonds_propres_initial", 0)
        prets_ro = rd.get("prets", [])
        nb_ch = p["nb_chambres"]

        all_figs = []  # Pour export PDF

        # ════════════════════════════════════════════════════════════════════
        # Scroll-to-top a l'ouverture (efface le hash residuel) + ancre haut
        # ════════════════════════════════════════════════════════════════════
        if not print_mode:
            import streamlit.components.v1 as _scroll_top_comp
            _scroll_top_comp.html(
                """<script>
                // Au chargement, scroller en haut et nettoyer le hash residuel
                // pour ne pas retomber sur la derniere section visitee.
                (function() {
                    try {
                        if (window.parent.location.hash) {
                            window.parent.history.replaceState(null, '',
                                window.parent.location.pathname + window.parent.location.search);
                        }
                    } catch(e) {}
                    try { window.parent.scrollTo({top: 0, behavior: 'instant'}); } catch(e) {}
                    try { window.top.scrollTo({top: 0, behavior: 'instant'}); } catch(e) {}
                })();
                </script>""",
                height=0,
            )
            # Ancre native Streamlit "top" pour le bouton de retour en haut
            st.header(" ", anchor="rapport-top", divider=False)

        # ════════════════════════════════════════════════════════════════════
        # 1. PAGE DE GARDE
        # ════════════════════════════════════════════════════════════════════
        # Page de garde
        import base64 as _b64_cov
        photos = [_Path(__file__).parent / "assets" / f"chateau_{i}.jpg" for i in range(1, 5)]
        ph_exist = [pp for pp in photos if pp.exists()]

        # Bandeau photo panoramique
        if ph_exist:
            with open(str(ph_exist[0]), "rb") as _f_cov:
                _cov_data = _b64_cov.b64encode(_f_cov.read()).decode()
            st.markdown(
                f'<div style="position:relative; border-radius:14px; overflow:hidden; margin:0 0 20px 0;">'
                f'<img src="data:image/jpeg;base64,{_cov_data}" '
                f'style="width:100%; height:250px; object-fit:cover; filter:brightness(0.45);">'
                f'<div style="position:absolute; top:50%; left:50%; transform:translate(-50%,-50%); text-align:center; width:90%;">'
                f'<p style="color:rgba(255,255,255,0.6); font-size:0.9em; letter-spacing:3px; text-transform:uppercase; margin:0 0 8px 0;">Plan Financier</p>'
                f'<h1 style="color:white; margin:0; font-size:2.5em; text-shadow:0 2px 8px rgba(0,0,0,0.5); font-weight:700;">'
                f'{p.get("nom_hotel", plan_nom)}</h1>'
                f'</div></div>',
                unsafe_allow_html=True,
            )

        # Chiffres cles en cartes
        _total_inv_all = sum(i["montant"] for i in p.get("investissements", [])) + sum(i["montant"] for i in rd.get("investissements", []))
        _total_fp_all = fp_ch + fp_ro
        # Endettement ventile par type (exclure pret intra-groupe Rocher → Chateau)
        _dette_bancaire_non_garanti = (
            sum(pr["montant"] for pr in prets_ch if not pr.get("subside_rw") and "rocher" not in pr.get("nom", "").lower())
            + sum(pr["montant"] for pr in prets_ro if not pr.get("subside_rw") and pr.get("creancier", "Banque") == "Banque")
        )
        _dette_garanti_rw = (
            sum(pr["montant"] for pr in prets_ch if pr.get("subside_rw"))
            + sum(pr["montant"] for pr in prets_ro if pr.get("subside_rw"))
        )
        _dette_partenaire = sum(pr["montant"] for pr in prets_ro if not pr.get("subside_rw") and pr.get("creancier", "Banque") != "Banque")
        _total_dette_all = _dette_bancaire_non_garanti + _dette_garanti_rw + _dette_partenaire
        _kpi_style = (
            'text-align:center; padding:14px 8px; background:linear-gradient(135deg,{bg1},{bg2}); '
            'border-radius:10px; border-left:4px solid {clr};'
        )
        _kpi_val = '<div style="font-size:1.5em; font-weight:bold; color:{clr};">{val}</div>'
        _kpi_lbl = '<div style="font-size:0.8em; color:#666; margin-top:2px;">{lbl}</div>'
        _kpi_sub = '<div style="font-size:0.65em; color:{clr}; margin-top:1px;">{txt}</div>'
        k1, k2, k3, k4 = st.columns(4)
        with k1:
            st.markdown(f'<div style="{_kpi_style.format(bg1="#eef2ff", bg2="#e0e7ff", clr="#4f46e5")}">'
                        f'{_kpi_val.format(clr="#4f46e5", val=f"{nb_ch}")}'
                        f'{_kpi_lbl.format(lbl="Chambres")}</div>', unsafe_allow_html=True)
        with k2:
            st.markdown(f'<div style="{_kpi_style.format(bg1="#f0fdf4", bg2="#dcfce7", clr="#16a34a")}">'
                        f'{_kpi_val.format(clr="#16a34a", val=f"{_total_inv_all/1e6:,.1f} M\u20ac")}'
                        f'{_kpi_lbl.format(lbl="Investissement total")}</div>', unsafe_allow_html=True)
        with k3:
            st.markdown(f'<div style="{_kpi_style.format(bg1="#eff6ff", bg2="#dbeafe", clr="#2563eb")}">'
                        f'{_kpi_val.format(clr="#2563eb", val=f"{_total_fp_all/1e6:,.1f} M\u20ac")}'
                        f'{_kpi_lbl.format(lbl="Fonds propres")}</div>', unsafe_allow_html=True)
        with k4:
            _dette_details = (
                f'{_kpi_sub.format(clr="#b91c1c", txt=f"Bancaire non garanti : {_dette_bancaire_non_garanti/1e6:,.1f} M\u20ac")}'
                f'{_kpi_sub.format(clr="#c2410c", txt=f"Garanti RW : {_dette_garanti_rw/1e6:,.1f} M\u20ac")}'
                + (f'{_kpi_sub.format(clr="#7c3aed", txt=f"Partenaire : {_dette_partenaire/1e6:,.1f} M\u20ac")}' if _dette_partenaire > 0 else '')
            )
            st.markdown(f'<div style="{_kpi_style.format(bg1="#fef2f2", bg2="#fee2e2", clr="#dc2626")}">'
                        f'{_kpi_val.format(clr="#dc2626", val=f"{_total_dette_all/1e6:,.1f} M\u20ac")}'
                        f'{_kpi_lbl.format(lbl="Endettement total")}'
                        f'{_dette_details}</div>', unsafe_allow_html=True)

        st.caption(f"Debut d'activite : **{p['date_ouverture'].strftime('%B %Y')}** | "
                   f"Projection : **{p['nb_mois_projection']//12} ans**")

        # ════════════════════════════════════════════════════════════════════
        # SOMMAIRE (TOC) - Navigation rapide entre sections
        # ════════════════════════════════════════════════════════════════════
        if not print_mode:
            _toc_sections = [
                ("sec-1", "1. Montage", "#4facfe"),
                ("sec-2", "2. Capital", "#4facfe"),
                ("sec-3", "3. Injection", "#4facfe"),
                ("sec-4", "4. Invest.", "#4facfe"),
                ("sec-5", "5. M&B", "#4facfe"),
                ("sec-6", "6. Rocher", "#11998e"),
                ("sec-7", "7. Chateau", "#f5576c"),
                ("sec-8", "8. Simulation", "#764ba2"),
            ]
            # Sommaire avec liens HTML natifs dans st.markdown.
            # st.markdown rend dans le DOM principal Streamlit (pas dans un
            # iframe sandboxe), donc les liens <a href="#sec-X"> peuvent
            # naviguer librement vers les ancres dans le meme frame.
            _toc_links = "".join(
                f'<a href="#{aid}" style="display:inline-block; padding:6px 12px; '
                f'margin:3px 4px; background:{clr}15; color:{clr}; border:1px solid {clr}50; '
                f'border-radius:6px; text-decoration:none; font-size:0.85em; font-weight:600;">'
                f'{lbl}</a>'
                for aid, lbl, clr in _toc_sections
            )
            st.markdown(
                f'<div style="background:white; padding:10px 14px; margin:8px 0 16px 0; '
                f'border-radius:10px; box-shadow:0 2px 8px rgba(0,0,0,0.08); '
                f'border:1px solid #e5e7eb;">'
                f'<div style="font-size:0.75em; font-weight:700; color:#6b7280; '
                f'text-transform:uppercase; letter-spacing:1px; margin-bottom:6px;">Sommaire</div>'
                f'<div>{_toc_links}</div>'
                f'</div>',
                unsafe_allow_html=True,
            )
            # Bouton flottant "Haut" - position:fixed via CSS pour rester visible
            # quel que soit le scroll. Pointe vers l'ancre native rapport-top.
            st.markdown(
                '<a href="#rapport-top" title="Retour en haut" '
                'style="position:fixed; bottom:30px; right:30px; '
                'width:48px; height:48px; background:#4f46e5; color:white; '
                'border-radius:50%; text-decoration:none; '
                'display:flex; align-items:center; justify-content:center; '
                'font-size:1.4em; font-weight:bold; '
                'box-shadow:0 4px 12px rgba(0,0,0,0.2); z-index:9999; '
                'transition:all 0.2s;" '
                'onmouseover="this.style.background=\'#3730a3\';this.style.transform=\'scale(1.1)\';" '
                'onmouseout="this.style.background=\'#4f46e5\';this.style.transform=\'scale(1)\';">'
                '↑</a>',
                unsafe_allow_html=True,
            )

        # ════════════════════════════════════════════════════════════════════
        # 2. MONTAGE FINANCIER
        # ════════════════════════════════════════════════════════════════════
        _cmt("before_s1_montage")
        st.header("1. Montage financier", anchor="sec-1", divider="blue")

        _prets_ro_std = [pr for pr in prets_ro if not pr.get("subside_rw")]
        _prets_ro_rw = [pr for pr in prets_ro if pr.get("subside_rw")]
        _prets_ch_std = [pr for pr in prets_ch if not pr.get("subside_rw")]
        _prets_ch_rw = [pr for pr in prets_ch if pr.get("subside_rw")]

        # En-tetes entites
        c1, c2 = st.columns(2)
        with c1:
            st.markdown("### Immobiliere Rocher")
        with c2:
            st.markdown("### Chateau d'Argenteau")

        # Fonds propres (aligne)
        c1, c2 = st.columns(2)
        with c1:
            st.markdown(f"**Fonds propres : {fp_ro:,.0f} \u20ac**")
            for inv in rd.get("fonds_propres_investisseurs", []):
                st.markdown(f"- {inv['nom']} : {inv['montant']:,.0f} \u20ac")
        with c2:
            st.markdown(f"**Fonds propres : {fp_ch:,.0f} \u20ac**")
            for inv in p.get("fonds_propres_investisseurs", []):
                st.markdown(f"- {inv['nom']} : {inv['montant']:,.0f} \u20ac")

        # Dettes bancaires (aligne)
        c1, c2 = st.columns(2)
        with c1:
            if _prets_ro_std:
                st.markdown("**Dettes bancaires :**")
                for pr in _prets_ro_std:
                    st.markdown(f"- {pr['nom']} ({pr.get('creancier','')}) : **{pr['montant']:,.0f} \u20ac** "
                                f"@ {pr['taux_annuel']:.1%} / {pr['duree_ans']} ans")
        with c2:
            if _prets_ch_std:
                st.markdown("**Dettes bancaires :**")
                for pr in _prets_ch_std:
                    st.markdown(f"- {pr['nom']} : **{pr['montant']:,.0f} \u20ac** "
                                f"@ {pr['taux_annuel']:.1%} / {pr['duree_ans']} ans")

        # Dettes garanties RW (aligne)
        c1, c2 = st.columns(2)
        with c1:
            if _prets_ro_rw:
                st.markdown("**Dettes garanties Region Wallonne :**")
                for pr in _prets_ro_rw:
                    st.markdown(f"- \U0001F7E2 {pr['nom']} ({pr.get('creancier','')}) : **{pr['montant']:,.0f} \u20ac** "
                                f"@ {pr['taux_annuel']:.1%} / {pr['duree_ans']} ans — *Subside RW*")
        with c2:
            if _prets_ch_rw:
                st.markdown("**Dettes garanties Region Wallonne :**")
                for pr in _prets_ch_rw:
                    st.markdown(f"- \U0001F7E2 {pr['nom']} : **{pr['montant']:,.0f} \u20ac** "
                                f"@ {pr['taux_annuel']:.1%} / {pr['duree_ans']} ans — *Subside RW*")

        # Schema du montage financier (dans section 1)
        st.markdown("---")
        _nf = lambda v: f"{v:,.0f}".replace(",", " ")
        _fp_ro_m = rd.get("fonds_propres_initial", 0)
        _fp_inv_ro_m = rd.get("fonds_propres_investisseurs", [])
        _fp_inv_ch_m = p.get("fonds_propres_investisseurs", [])
        _pret_rocher_ch_s1 = next((pr for pr in prets_ch if "rocher" in pr.get("nom","").lower()), None)
        _pret_ro_ch_mt = _pret_rocher_ch_s1["montant"] if _pret_rocher_ch_s1 else 0
        _loyer_m_s1 = p.get("loyer_mensuel", 0)

        # Collecter acteurs et prets
        _acteurs_s1 = {}
        for inv in _fp_inv_ro_m:
            _acteurs_s1.setdefault(inv["nom"], {"cap_ro": 0, "cap_ch": 0})
            _acteurs_s1[inv["nom"]]["cap_ro"] += inv["montant"]
        for inv in _fp_inv_ch_m:
            _acteurs_s1.setdefault(inv["nom"], {"cap_ro": 0, "cap_ch": 0})
            _acteurs_s1[inv["nom"]]["cap_ch"] += inv["montant"]
        # Prets bancaires vs partenaires (distinguer par creancier)
        _t_ro_std = sum(pr["montant"] for pr in prets_ro if not pr.get("subside_rw") and pr.get("creancier", "Banque") == "Banque")
        _t_ro_rw = sum(pr["montant"] for pr in prets_ro if pr.get("subside_rw"))
        _prets_ro_partenaires = [(pr["nom"], pr["montant"], pr.get("creancier", "")) for pr in prets_ro
                                  if not pr.get("subside_rw") and pr.get("creancier", "Banque") != "Banque"]
        _t_ch_std = sum(pr["montant"] for pr in prets_ch if not pr.get("subside_rw") and "rocher" not in pr.get("nom","").lower())
        _t_ch_rw = sum(pr["montant"] for pr in prets_ch if pr.get("subside_rw"))

        # Plotly montage schematique — canvas large pour eviter superposition
        _fig_m = go.Figure()
        _fig_m.update_xaxes(range=[-1, 19], showgrid=False, zeroline=False, visible=False, fixedrange=True)
        _fig_m.update_yaxes(range=[-0.5, 10], showgrid=False, zeroline=False, visible=False, fixedrange=True)

        _act_colors = {"LC Cr\u00e9ation": "#e67e22", "Partenaire": "#2980b9"}
        _def_c = ["#e67e22", "#2980b9", "#9b59b6", "#1abc9c"]
        _act_list = list(_acteurs_s1.items())
        _n_act = len(_act_list)

        # Positions acteurs en haut (espaces regulierement)
        _act_positions = {}
        _act_spacing = min(5, 14 / max(_n_act, 1))
        for idx, (nom, _) in enumerate(_act_list):
            if nom not in _act_colors:
                _act_colors[nom] = _def_c[idx % len(_def_c)]
            cx = 2 + idx * _act_spacing
            _act_positions[nom] = cx
            _fig_m.add_shape(type="rect", x0=cx - 1.8, y0=8.2, x1=cx + 1.8, y1=9.5,
                             fillcolor=_act_colors[nom], line=dict(width=0), layer="below")
            _fig_m.add_annotation(x=cx, y=8.85, text=f"<b>{nom}</b>", showarrow=False,
                                  font=dict(size=12, color="white"))

        # Banque (haut droite)
        _bk_x = 15.5
        _fig_m.add_shape(type="rect", x0=_bk_x - 2, y0=8.2, x1=_bk_x + 2, y1=9.5,
                         fillcolor="#1a3a5c", line=dict(width=0), layer="below")
        _fig_m.add_annotation(x=_bk_x, y=8.85, text="<b>BANQUE</b>", showarrow=False,
                              font=dict(size=12, color="white"))

        # ROCHER (bas gauche)
        _ro_cx, _ro_cy = 4, 2
        _fig_m.add_shape(type="rect", x0=0, y0=0.5, x1=8, y1=3.5,
                         fillcolor="#27ae60", line=dict(width=0), layer="below")
        _fig_m.add_annotation(x=_ro_cx, y=2.7, text="<b>ROCHER</b>", showarrow=False,
                              font=dict(size=16, color="white"))
        _fig_m.add_annotation(x=_ro_cx, y=2.1, text="<i>Societe Immobiliere</i>", showarrow=False,
                              font=dict(size=10, color="rgba(255,255,255,0.8)"))
        _fig_m.add_shape(type="rect", x0=0.8, y0=0.65, x1=7.2, y1=1.3,
                         fillcolor="#f1c40f", line=dict(width=0), layer="below")
        _fig_m.add_annotation(x=_ro_cx, y=0.97, text=f"<b>Capital : {_nf(_fp_ro_m)} \u20ac</b>",
                              showarrow=False, font=dict(size=11, color="#333"))

        # ARGENTEAU (bas droite)
        _ag_cx, _ag_cy = 14, 2
        _fig_m.add_shape(type="rect", x0=10, y0=0.5, x1=18, y1=3.5,
                         fillcolor="#8e44ad", line=dict(width=0), layer="below")
        _fig_m.add_annotation(x=_ag_cx, y=2.7, text="<b>ARGENTEAU</b>", showarrow=False,
                              font=dict(size=16, color="white"))
        _fig_m.add_annotation(x=_ag_cx, y=2.1, text="<i>Societe d'Exploitation</i>", showarrow=False,
                              font=dict(size=10, color="rgba(255,255,255,0.8)"))
        _fig_m.add_shape(type="rect", x0=10.8, y0=0.65, x1=17.2, y1=1.3,
                         fillcolor="#f1c40f", line=dict(width=0), layer="below")
        _fig_m.add_annotation(x=_ag_cx, y=0.97, text=f"<b>Capital : {_nf(fp_ch)} \u20ac</b>",
                              showarrow=False, font=dict(size=11, color="#333"))

        # Fleche helper — label place a une position explicite (lx, ly)
        def _arrow(x0, y0, x1, y1, label, montant, color, lx=None, ly=None):
            _fig_m.add_annotation(x=x1, y=y1, ax=x0, ay=y0, xref="x", yref="y", axref="x", ayref="y",
                                  showarrow=True, arrowhead=2, arrowsize=1.5, arrowwidth=2.5,
                                  arrowcolor=color, opacity=0.85, standoff=2, startstandoff=2)
            if lx is None:
                lx = (x0 + x1) / 2
            if ly is None:
                ly = (y0 + y1) / 2
            _fig_m.add_annotation(x=lx, y=ly,
                                  text=f"<b>{label}</b><br>{_nf(montant)} \u20ac",
                                  showarrow=False, font=dict(size=9, color=color),
                                  bgcolor="white", bordercolor=color, borderwidth=1, borderpad=3)

        # Fleches capital acteurs → Rocher
        _ro_idx = 0
        for nom, data in _acteurs_s1.items():
            if data["cap_ro"] > 0:
                c = _act_colors[nom]
                pct = data["cap_ro"] / _fp_ro_m * 100 if _fp_ro_m > 0 else 0
                x_arr = 1.5 + _ro_idx * 3
                sx = _act_positions[nom]
                # Label au 1/3 de la fleche (pres du depart)
                _arrow(sx, 8.2, x_arr, 3.5, f"Capital ({pct:.0f}%)", data["cap_ro"], c,
                       lx=sx + (x_arr - sx) * 0.25, ly=8.2 + (3.5 - 8.2) * 0.25 + 0.4)
                _ro_idx += 1

        # Fleches capital acteurs → Argenteau (labels au 1/3 pres du depart)
        _ag_idx = 0
        for nom, data in _acteurs_s1.items():
            if data["cap_ch"] > 0:
                c = _act_colors[nom]
                pct = data["cap_ch"] / fp_ch * 100 if fp_ch > 0 else 0
                x_arr = 11.5 + _ag_idx * 2.5
                sx = _act_positions[nom]
                _arrow(sx, 8.2, x_arr, 3.5, f"Capital ({pct:.0f}%)", data["cap_ch"], c,
                       lx=sx + (x_arr - sx) * 0.3, ly=8.2 + (3.5 - 8.2) * 0.3 + 0.4)
                _ag_idx += 1

        # Fleches bancaires → Rocher
        if _t_ro_std > 0:
            _sx, _sy, _ex, _ey = _bk_x - 0.5, 8.2, 6, 3.5
            _arrow(_sx, _sy, _ex, _ey, "Pret LT", _t_ro_std, "#e74c3c",
                   lx=(_sx + _ex) / 2 + 0.5, ly=(_sy + _ey) / 2 + 0.4)
        if _t_ro_rw > 0:
            _sx, _sy, _ex, _ey = _bk_x - 1.5, 8.2, 2.5, 3.5
            _arrow(_sx, _sy, _ex, _ey, "Pret garanti RW", _t_ro_rw, "#f39c12",
                   lx=(_sx + _ex) / 2 - 1.5, ly=(_sy + _ey) / 2 + 0.4)

        # Fleches prets partenaires → Rocher
        for _pr_nom, _pr_mt, _pr_creancier in _prets_ro_partenaires:
            if _pr_creancier in _act_positions:
                _sx = _act_positions[_pr_creancier]
                _ex, _ey = 5, 3.5
                _arrow(_sx, 8.2, _ex, _ey, _pr_nom, _pr_mt, _act_colors.get(_pr_creancier, "#2980b9"),
                       lx=_sx + (_ex - _sx) * 0.4, ly=8.2 + (_ey - 8.2) * 0.4 + 0.4)

        # Fleches bancaires → Argenteau (labels proches de leur fleche, decales en y)
        if _t_ch_rw > 0:
            _sx, _sy, _ex, _ey = _bk_x + 1.5, 8.2, 12.5, 3.5
            _arrow(_sx, _sy, _ex, _ey, "Pret garanti RW", _t_ch_rw, "#f39c12",
                   lx=_sx + (_ex - _sx) * 0.6, ly=_sy + (_ey - _sy) * 0.6 + 0.4)
        if _t_ch_std > 0:
            _sx, _sy, _ex, _ey = _bk_x + 0.5, 8.2, 16, 3.5
            _arrow(_sx, _sy, _ex, _ey, "Pret LT", _t_ch_std, "#e74c3c",
                   lx=_sx + (_ex - _sx) * 0.7, ly=_sy + (_ey - _sy) * 0.7 + 0.4)

        # Pret subordonne Rocher → Argenteau
        if _pret_ro_ch_mt > 0:
            _fig_m.add_annotation(x=10, y=2.6, ax=8, ay=2.6, xref="x", yref="y", axref="x", ayref="y",
                                  showarrow=True, arrowhead=2, arrowsize=1.5, arrowwidth=2,
                                  arrowcolor="#27ae60", opacity=0.8, standoff=2, startstandoff=2)
            _fig_m.add_annotation(x=9, y=3.1,
                                  text=f"<b>Pret subordonne</b><br>{_nf(_pret_ro_ch_mt)} \u20ac",
                                  showarrow=False, font=dict(size=9, color="#27ae60"),
                                  bgcolor="#e8f8f5", bordercolor="#27ae60", borderwidth=1, borderpad=3)

        # Bail commercial — ligne pointillee Argenteau → Rocher
        _fig_m.add_shape(type="line", x0=10, y0=1.8, x1=8, y1=1.8,
                         line=dict(color="#c0392b", width=2, dash="dash"), layer="above")
        # Petit coude vers le bas pour symboliser le paiement
        _fig_m.add_shape(type="line", x0=8, y0=1.8, x1=8, y1=1.3,
                         line=dict(color="#c0392b", width=2, dash="dash"), layer="above")
        _fig_m.add_annotation(x=9, y=1.3,
                              text=f"<b>Bail Commercial</b><br>{_nf(_loyer_m_s1)} \u20ac/mois",
                              showarrow=False, font=dict(size=9, color="#c0392b"),
                              bgcolor="#fef9e7", bordercolor="#c0392b", borderwidth=1, borderpad=3)

        _fig_m.update_layout(height=600, margin=dict(l=0, r=0, t=5, b=0),
                             plot_bgcolor="white", paper_bgcolor="white", dragmode=False)
        st.plotly_chart(_fig_m, use_container_width=True, config={"displayModeBar": False})

        # 2. Repartition du capital social
        _cmt("before_s2_capital")
        st.header("2. Repartition du capital social", anchor="sec-2", divider="blue")

        _c_cap1, _c_cap2 = st.columns(2)
        with _c_cap1:
            st.markdown("#### Immobiliere Rocher")
            _fp_inv_ro = rd.get("fonds_propres_investisseurs", [])
            if _fp_inv_ro:
                _cap_labels_ro = [inv["nom"] for inv in _fp_inv_ro]
                _cap_values_ro = [inv["montant"] for inv in _fp_inv_ro]
                _fig_cap_ro = go.Figure(data=[go.Pie(
                    labels=_cap_labels_ro, values=_cap_values_ro,
                    marker=dict(colors=["#e67e22", "#2980b9", "#27ae60", "#8e44ad"][:len(_cap_labels_ro)]),
                    textinfo="label+percent+value",
                    texttemplate="<b>%{label}</b><br><b>%{value:,.0f} \u20ac</b><br>(%{percent})",
                    textfont=dict(size=14),
                    hole=0.3,
                    pull=[0.03] * len(_cap_labels_ro),
                )])
                _fig_cap_ro.update_layout(height=380, showlegend=False, margin=dict(l=10, r=10, t=40, b=10),
                                          title=dict(text=f"Capital : {sum(_cap_values_ro):,.0f} \u20ac", font=dict(size=15)))
                _show_fig(_fig_cap_ro)
            else:
                st.info(f"Fonds propres : **{rd.get('fonds_propres_initial', 0):,.0f} \u20ac**")

        with _c_cap2:
            st.markdown("#### Chateau d'Argenteau")
            _fp_inv_ch = p.get("fonds_propres_investisseurs", [])
            if _fp_inv_ch:
                _cap_labels_ch = [inv["nom"] for inv in _fp_inv_ch]
                _cap_values_ch = [inv["montant"] for inv in _fp_inv_ch]
                _fig_cap_ch = go.Figure(data=[go.Pie(
                    labels=_cap_labels_ch, values=_cap_values_ch,
                    marker=dict(colors=["#e67e22", "#2980b9", "#27ae60", "#8e44ad"][:len(_cap_labels_ch)]),
                    textinfo="label+percent+value",
                    texttemplate="<b>%{label}</b><br><b>%{value:,.0f} \u20ac</b><br>(%{percent})",
                    textfont=dict(size=14),
                    hole=0.3,
                    pull=[0.03] * len(_cap_labels_ch),
                )])
                _fig_cap_ch.update_layout(height=380, showlegend=False, margin=dict(l=10, r=10, t=40, b=10),
                                          title=dict(text=f"Capital : {sum(_cap_values_ch):,.0f} \u20ac", font=dict(size=15)))
                _show_fig(_fig_cap_ch)
            else:
                st.info(f"Fonds propres : **{fp_ch:,.0f} \u20ac**")

        # 3. Injection par acteur
        _cmt("before_s3_injection")
        st.header("3. Injection par acteur", anchor="sec-3", divider="blue")
        st.caption("Recap des montants finances par acteur externe (hors Immobiliere Rocher et Chateau d'Argenteau)")

        # Collecter tous les acteurs et montants
        _injections = {}  # {acteur: montant_total}

        # Fonds propres Rocher
        for inv in rd.get("fonds_propres_investisseurs", []):
            nom_act = inv["nom"]
            _injections[nom_act] = _injections.get(nom_act, 0) + inv["montant"]

        # Fonds propres Chateau
        for inv in p.get("fonds_propres_investisseurs", []):
            nom_act = inv["nom"]
            _injections[nom_act] = _injections.get(nom_act, 0) + inv["montant"]

        # Prets Rocher (creanciers) — distinguer subside RW
        for pr in prets_ro:
            creancier = pr.get("creancier", pr["nom"]) or pr["nom"]
            if pr.get("subside_rw"):
                creancier = f"{creancier} (garanti RW)"
            _injections[creancier] = _injections.get(creancier, 0) + pr["montant"]

        # Prets Chateau (sauf pret Rocher = flux interne) — distinguer subside RW
        for pr in prets_ch:
            if "rocher" in pr.get("nom", "").lower():
                continue  # Flux interne
            creancier = pr.get("creancier", pr["nom"]) or pr["nom"]
            if pr.get("subside_rw"):
                creancier = f"{creancier} (garanti RW)"
            _injections[creancier] = _injections.get(creancier, 0) + pr["montant"]

        # Afficher le tableau
        total_inject = sum(_injections.values())
        for acteur, montant in sorted(_injections.items(), key=lambda x: -x[1]):
            pct = montant / total_inject * 100 if total_inject > 0 else 0
            st.markdown(f"- **{acteur}** : **{montant:,.0f} \u20ac** ({pct:.1f}%)")
        st.markdown(f"\n**Total injections externes : {total_inject:,.0f} \u20ac**")

        # Camembert
        if _injections and total_inject > 0:
            _colors_inj = ["#e67e22", "#2980b9", "#e74c3c", "#f39c12", "#27ae60", "#8e44ad", "#1abc9c", "#f093fb"]
            fig_inj = go.Figure(data=[go.Pie(
                labels=list(_injections.keys()),
                values=list(_injections.values()),
                marker=dict(colors=_colors_inj[:len(_injections)]),
                textinfo="label+percent+value",
                texttemplate="<b>%{label}</b><br><b>%{value:,.0f} \u20ac</b><br>(%{percent})",
                textfont=dict(size=13),
                pull=[0.02] * len(_injections),
            )])
            fig_inj.update_layout(height=500, showlegend=False, margin=dict(l=30, r=30, t=20, b=20))
            _show_fig(fig_inj, large=True)

        # ════════════════════════════════════════════════════════════════════
        # 4. INVESTISSEMENTS INITIAUX
        # ════════════════════════════════════════════════════════════════════
        _pret_rocher_ch = next((pr for pr in prets_ch if "rocher" in pr.get("nom","").lower()), None)
        st.header("4. Investissements initiaux", anchor="sec-4", divider="blue")

        # Regrouper investissements Rocher + Chateau
        _inv_rocher = rd.get("investissements", [])
        _inv_chateau = p.get("investissements", [])

        _tot_inv_ro = sum(inv["montant"] for inv in _inv_rocher)
        _tot_inv_ch = sum(inv["montant"] for inv in _inv_chateau)

        c_inv1, c_inv2 = st.columns(2)
        with c_inv1:
            st.markdown("#### Immobiliere Rocher")
            for inv in _inv_rocher:
                if inv["montant"] > 0:
                    st.markdown(f"- {inv['categorie']} : **{inv['montant']:,.0f} \u20ac** "
                                f"(amort. {inv['duree_amort']} ans)")
        with c_inv2:
            st.markdown("#### Chateau d'Argenteau")
            for inv in _inv_chateau:
                if inv["montant"] > 0:
                    st.markdown(f"- {inv['categorie']} : **{inv['montant']:,.0f} \u20ac** "
                                f"(amort. {inv['duree_amort']} ans)")

        c_t1, c_t2 = st.columns(2)
        with c_t1:
            st.markdown(f"**Total Rocher : {_tot_inv_ro:,.0f} \u20ac**")
        with c_t2:
            st.markdown(f"**Total Chateau : {_tot_inv_ch:,.0f} \u20ac**")

        st.markdown(f"### Total investissements : {_tot_inv_ro + _tot_inv_ch:,.0f} \u20ac")

        # ════════════════════════════════════════════════════════════════════
        # 7. MOYENS & BESOINS
        # ════════════════════════════════════════════════════════════════════
        _cmt("before_s6_moyens_besoins")
        st.header("5. Moyens & Besoins", anchor="sec-5", divider="blue")

        # Donnees Chateau
        _total_inv_ch_mb = sum(i["montant"] for i in p.get("investissements", []))
        _prets_ch_mb = p.get("prets", [])
        _total_prets_ch_mb = sum(pr["montant"] for pr in _prets_ch_mb)
        _fp_ch_mb = p.get("fonds_propres_initial", 0)
        _cf_cum_rpt = df_rpt["cash_flow"].cumsum()
        _besoin_treso_ch = abs(min(0, _cf_cum_rpt.min()))
        _total_moyens_ch_mb = _total_prets_ch_mb + _fp_ch_mb
        _total_besoins_ch_mb = _total_inv_ch_mb + _besoin_treso_ch
        _solde_ch_mb = _total_moyens_ch_mb - _total_besoins_ch_mb

        # Donnees Rocher
        _rocher_prets_mb = rd.get("prets", [])
        _rocher_fp_mb = rd.get("fonds_propres_initial", 0)
        _total_dettes_r_mb = sum(pr["montant"] for pr in _rocher_prets_mb)
        _total_moyens_r_mb = _total_dettes_r_mb + _rocher_fp_mb
        _rocher_inv_mb = rd.get("investissements", [])
        _total_inv_r_mb = sum(inv["montant"] for inv in _rocher_inv_mb)
        _pret_au_ch_mb = _pret_rocher_ch["montant"] if _pret_rocher_ch else 0
        _total_besoins_r_mb = _total_inv_r_mb + _pret_au_ch_mb
        _solde_r_mb = _total_moyens_r_mb - _total_besoins_r_mb

        # Waterfalls par entite (Rocher en premier, puis Chateau)
        _c_mb1, _c_mb2 = st.columns(2)
        with _c_mb1:
            st.markdown("#### Immobiliere Rocher")
            _wf_pr_r = [{"nom": pr["nom"], "montant": pr["montant"]} for pr in _rocher_prets_mb]
            _build_waterfall_moyens_besoins(
                st, go, _wf_pr_r, _rocher_fp_mb, _total_moyens_r_mb,
                _total_inv_r_mb + _pret_au_ch_mb, 0, _solde_r_mb, "rpt_rocher",
                pret_intra=_pret_au_ch_mb, pret_intra_label="Pret a Argenteau"
            )
        with _c_mb2:
            st.markdown("#### Chateau d'Argenteau")
            _build_waterfall_moyens_besoins(
                st, go, _prets_ch_mb, _fp_ch_mb, _total_moyens_ch_mb,
                _total_inv_ch_mb, _besoin_treso_ch, _solde_ch_mb, "rpt_chateau"
            )

        # Waterfall consolide (les 2 entites aggregees)
        # Le pret intra-groupe (Rocher → Chateau) s'annule en conso : exclu des moyens ET des besoins
        st.markdown("#### Vue consolidee")
        # Agreger les prets par nom, en excluant le pret Rocher cote Chateau
        _pret_rocher_nom = _pret_rocher_ch["nom"].lower() if _pret_rocher_ch else ""
        _all_prets_conso = {}
        for pr in _prets_ch_mb:
            if pr["nom"].lower() == _pret_rocher_nom:
                continue
            _n = pr["nom"]
            _all_prets_conso[_n] = _all_prets_conso.get(_n, 0) + pr["montant"]
        for pr in _rocher_prets_mb:
            _n = pr["nom"]
            _all_prets_conso[_n] = _all_prets_conso.get(_n, 0) + pr["montant"]
        _prets_conso = [{"nom": n, "montant": m} for n, m in _all_prets_conso.items()]
        _fp_conso = _fp_ch_mb + _rocher_fp_mb
        _total_moyens_conso = sum(pr["montant"] for pr in _prets_conso) + _fp_conso
        # Besoins consolides : invest Rocher + invest Chateau + besoin treso (sans pret intra-groupe)
        _total_inv_conso = _total_inv_ch_mb + _total_inv_r_mb
        _total_besoins_conso = _total_inv_conso + _besoin_treso_ch
        _solde_conso = _total_moyens_conso - _total_besoins_conso
        _build_waterfall_moyens_besoins(
            st, go, _prets_conso, _fp_conso, _total_moyens_conso,
            _total_inv_conso, _besoin_treso_ch, _solde_conso, "rpt_conso"
        )

        # Rappel reinvestissement annuel
        _reinvest_pct_list = p.get("reinvest_pct_an", [])
        _reinvest_base = sum(inv["montant"] for inv in p.get("investissements", []) if inv.get("duree_amort", 0) > 0)
        if _reinvest_pct_list and _reinvest_base > 0:
            _pct_moyen = sum(_reinvest_pct_list) / len(_reinvest_pct_list) * 100
            _pct_max = max(_reinvest_pct_list) * 100
            st.markdown("---")
            st.markdown(
                f'<div style="background:#f0f4ff; border-left:4px solid #667eea; padding:14px 18px; '
                f'border-radius:6px; margin-top:16px; font-size:0.95em; line-height:1.6; color:#1a1a2e;">'
                f'<b>Reinvestissements prevus</b> : un reinvestissement annuel est prevu, '
                f'correspondant a un pourcentage de l\'investissement initial amortissable '
                f'({_reinvest_base:,.0f} \u20ac). '
                f'Le taux varie de <b>{_reinvest_pct_list[0]*100:.1f}%</b> (annee 1) '
                f'a <b>{_pct_max:.1f}%</b> (a partir de l\'annee {_reinvest_pct_list.index(max(_reinvest_pct_list))+1}), '
                f'soit en moyenne <b>{_pct_moyen:.1f}%</b> par an.</div>',
                unsafe_allow_html=True
            )

        # ════════════════════════════════════════════════════════════════════
        # 7. PLAN FINANCIER ROCHER
        # ════════════════════════════════════════════════════════════════════
        _cmt("before_s7_rocher")
        st.header("6. Plan financier — Immobiliere Rocher", anchor="sec-6", divider="green")

        # Commentaire editable — synthese Rocher
        _comment_key_ro = "commentaire_rocher"
        _default_comment_ro = p.get(_comment_key_ro,
            "L'Immobiliere Rocher porte l'ensemble des investissements immobiliers du projet. "
            "Ses revenus proviennent principalement du loyer facture au Chateau d'Argenteau "
            "et des interets percus sur le pret intra-groupe. "
            "Le resultat est impacte par les charges d'amortissement et le service de la dette externe. "
            "A terme, l'entite degage un resultat positif croissant a mesure que la dette se rembourse."
        )
        st.markdown("### Synthese & constats")
        _is_visu_report = st.session_state.get("auth_mode") == "visu"
        if _is_visu_report:
            st.markdown(f'<div style="background:#f0fdf4; border-left:4px solid #11998e; '
                        f'padding:14px 18px; border-radius:6px; margin-bottom:16px; '
                        f'font-size:0.95em; line-height:1.6; color:#1a3a2a;">'
                        f'{_default_comment_ro}</div>', unsafe_allow_html=True)
        else:
            _new_comment_ro = st.text_area(
                "Commentaire de synthese (editable)",
                value=_default_comment_ro,
                height=_text_area_height(_default_comment_ro, min_h=120),
                key="ta_comment_rocher",
            )
            if _new_comment_ro != p.get(_comment_key_ro):
                p[_comment_key_ro] = _new_comment_ro
                sauvegarder_plan(plan_nom, p, local_only=True)

        # Hypotheses Rocher
        _cmt("before_ro_hypotheses")
        st.markdown("### Hypotheses principales")
        total_inv_ro = sum(inv["montant"] for inv in rd.get("investissements", []))
        amort_an_ro = sum(inv["montant"]/inv["duree_amort"] for inv in rd.get("investissements", []) if inv["duree_amort"]>0)
        total_dettes_ro = sum(pr["montant"] for pr in prets_ro)
        loyer_m = p.get("loyer_mensuel", 0)

        # Projection Rocher (calculee avant le tableau pour avoir les moyennes)
        from dateutil.relativedelta import relativedelta as _rd_rpt
        rocher_inv = rd.get("investissements", [])
        amort_m_ro = sum(inv["montant"]/inv["duree_amort"]/12 for inv in rocher_inv if inv["duree_amort"]>0)
        dfs_ro = {}
        for pr in prets_ro:
            if pr["montant"] > 0:
                dfs_ro[pr["nom"]] = calc_tableau_pret(pr, p["date_ouverture"], p["nb_mois_projection"])
        pret_rocher_ch = _pret_rocher_ch
        df_pret_ch_ro = calc_tableau_pret(pret_rocher_ch, p["date_ouverture"], p["nb_mois_projection"]) if pret_rocher_ch else pd.DataFrame()

        # Calculer interets et capital moyens annuels payes par Rocher
        _total_int_ro = 0
        _total_cap_ro = 0
        for _n, _df in dfs_ro.items():
            _total_int_ro += _df["interets"].sum()
            _total_cap_ro += _df["capital"].sum()
        _nb_ans_proj = max(1, p["nb_mois_projection"] // 12)
        _moy_int_an_ro = _total_int_ro / _nb_ans_proj
        _moy_cap_an_ro = _total_cap_ro / _nb_ans_proj

        _int_rec_an = 0
        if pret_rocher_ch and not df_pret_ch_ro.empty:
            _int_rec_an = df_pret_ch_ro["interets"].sum() / _nb_ans_proj

        _c_hyp1, _c_hyp2 = st.columns(2)
        with _c_hyp1:
            st.markdown("#### Investissements & Dettes")
            st.markdown(f"""
| Poste | Valeur |
|---|---|
| Investissement total | **{total_inv_ro:,.0f} \u20ac** |
| Amortissement annuel | **{amort_an_ro:,.0f} \u20ac** |
| Total dettes | **{total_dettes_ro:,.0f} \u20ac** |
| Interets moyens / an | **{_moy_int_an_ro:,.0f} \u20ac** |
| Capital rembourse moyen / an | **{_moy_cap_an_ro:,.0f} \u20ac** |
""")
        with _c_hyp2:
            st.markdown("#### Revenus")
            _pret_info = f"**{pret_rocher_ch['montant']:,.0f} \u20ac** @ {pret_rocher_ch['taux_annuel']:.1%}" if pret_rocher_ch else "-"
            st.markdown(f"""
| Poste | Valeur |
|---|---|
| Loyer facture au Chateau | **{loyer_m:,.0f} \u20ac/mois** ({loyer_m*12:,.0f} \u20ac/an) |
| Pret octroye au Chateau | {_pret_info} |
| Revenus d'interets moyens / an | **{_int_rec_an:,.0f} \u20ac** |
""")

        rows_ro = []
        for m in range(p["nb_mois_projection"]):
            d = p["date_ouverture"] + _rd_rpt(months=m)
            rev = loyer_m
            int_rec = 0; cap_rec = 0
            if not df_pret_ch_ro.empty:
                r = df_pret_ch_ro[df_pret_ch_ro["date"]==d]
                if not r.empty: int_rec = r.iloc[0]["interets"]; cap_rec = r.iloc[0]["capital"]
            int_pay = 0; cap_pay = 0
            for _n, _df in dfs_ro.items():
                r = _df[_df["date"]==d]
                if not r.empty: int_pay += r.iloc[0]["interets"]; cap_pay += r.iloc[0]["capital"]
            # Subside RW Rocher : 1/15 du montant par an au resultat, a partir de An 2
            # Pas d'impact cash (perception = remboursement de la dette par les ghost rows)
            annee_idx_ro = m // 12
            sub_rw = 0
            if annee_idx_ro >= 1:
                for _pr_sub in prets_ro:
                    if _pr_sub.get("subside_rw", False):
                        duree_sub = 15
                        annee_sub = annee_idx_ro - 1
                        if annee_sub < duree_sub:
                            sub_rw += _pr_sub["montant"] / duree_sub / 12
            resultat_m = rev + int_rec - int_pay - amort_m_ro + sub_rw
            cash_m = rev + int_rec + cap_rec - int_pay - cap_pay
            rows_ro.append({
                "date": d, "annee": d.year,
                "loyer": rev, "interets_recus": int_rec, "capital_recu": cap_rec,
                "interets_payes": int_pay, "capital_paye": cap_pay,
                "amortissement": amort_m_ro, "subside_rw": sub_rw,
                "produits": rev + int_rec + sub_rw,
                "charges": int_pay + amort_m_ro,
                "resultat": resultat_m, "cash": cash_m,
            })
        df_ro = pd.DataFrame(rows_ro)
        df_ro["cash_cumul"] = df_ro["cash"].cumsum()
        df_ro["res_cumul"] = df_ro["resultat"].cumsum()
        ann_ro = df_ro.groupby("annee").agg({
            "loyer": "sum", "interets_recus": "sum", "capital_recu": "sum",
            "interets_payes": "sum", "capital_paye": "sum",
            "amortissement": "sum", "subside_rw": "sum",
            "produits": "sum", "charges": "sum",
            "resultat": "sum", "cash": "sum",
            "cash_cumul": "last", "res_cumul": "last",
        }).reset_index()
        xr = [str(int(a)) for a in ann_ro["annee"]]

        fig = go.Figure()
        fig.add_trace(go.Bar(x=xr, y=K(ann_ro["produits"]), name="Produits", marker_color="#38ef7d"))
        fig.add_trace(go.Bar(x=xr, y=-K(ann_ro["charges"]), name="Charges", marker_color="#f5576c"))
        _res_ro_k = K(ann_ro["resultat"])
        fig.add_trace(go.Scatter(x=xr, y=_res_ro_k, name="Resultat", mode="lines+markers+text",
            line=dict(color="#3b3b98", width=3),
            text=[f"<b>{v:,.0f}</b>" for v in _res_ro_k],
            textposition="top center",
            textfont=dict(size=11, color="#3b3b98")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray")
        fig.update_layout(title="Rocher — Produits, Charges et Resultat (K\u20ac)", height=500, barmode="relative",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ro_pnl")

        fig = go.Figure()
        fig.add_trace(go.Bar(x=xr, y=K(ann_ro["cash"]), name="Cash periode", marker_color="#4facfe"))
        _cash_cum_ro_k = K(ann_ro["cash_cumul"])
        fig.add_trace(go.Scatter(x=xr, y=_cash_cum_ro_k, name="Cash cumule", mode="lines+markers+text",
            line=dict(color="#f5576c", width=3),
            text=[f"<b>{v:,.0f}</b>" for v in _cash_cum_ro_k],
            textposition="top center", textfont=dict(size=11, color="#b71c2e")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray")
        fig.update_layout(title="Rocher — Cash flow (K\u20ac)", height=500, xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ro_cashflow")

        # Detail des chiffres - Rocher (titre a la taille des titres Plotly)
        st.markdown(
            '<div style="font-size:17px; font-weight:700; color:#2a3f5f; '
            'margin:16px 0 8px 0; font-family:Arial, sans-serif;">Details des chiffres</div>',
            unsafe_allow_html=True,
        )
        _fmt_ro = lambda v: f"{v:,.0f} €"
        _df_tbl_ro = pd.DataFrame({
            "Annee": [str(int(a)) for a in ann_ro["annee"]],
            "Loyer recu": [_fmt_ro(v) for v in ann_ro["loyer"]],
            "Interets recus (Chateau)": [_fmt_ro(v) for v in ann_ro["interets_recus"]],
            "Subside RW": [_fmt_ro(v) for v in ann_ro["subside_rw"]],
            "Capital recu (Chateau)": [_fmt_ro(v) for v in ann_ro["capital_recu"]],
            "Interets payes": [_fmt_ro(v) for v in ann_ro["interets_payes"]],
            "Amortissement": [_fmt_ro(v) for v in ann_ro["amortissement"]],
            "Resultat": [_fmt_ro(v) for v in ann_ro["resultat"]],
            "Resultat cumule": [_fmt_ro(v) for v in ann_ro["res_cumul"]],
            "Remb. capital paye": [_fmt_ro(v) for v in ann_ro["capital_paye"]],
            "Cash Flow": [_fmt_ro(v) for v in ann_ro["cash"]],
            "Cash Flow Cumul": [_fmt_ro(v) for v in ann_ro["cash_cumul"]],
        })
        st.dataframe(_df_tbl_ro, use_container_width=True, hide_index=True)

        fig = go.Figure()
        _cp_ro_k = K(fp_ro + ann_ro["res_cumul"])
        fig.add_trace(go.Scatter(x=xr, y=_cp_ro_k, name="Capitaux propres", mode="lines+markers+text",
            line=dict(color="#667eea", width=3),
            fill="tozeroy", fillcolor="rgba(102, 126, 234, 0.15)",
            text=[f"<b>{v:,.0f}</b>" for v in _cp_ro_k],
            textposition="top center", textfont=dict(size=11, color="#3b3b98")))
        fig.update_layout(title="Rocher — Fonds propres (K\u20ac)", height=500, xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ro_fp")

        # Endettement Rocher - ventile par pret
        _enc_ro_by_pret = {}
        for pr in prets_ro:
            if pr["montant"] > 0:
                _pn_ro = pr["nom"]
                if pr.get("subside_rw"):
                    _pn_ro += " (RW)"
                _enc_ro_by_pret[_pn_ro] = []
                _dfp_ro = calc_tableau_pret(pr, p["date_ouverture"], p["nb_mois_projection"])
                for a in ann_ro["annee"]:
                    _d_fin_ro = date(int(a), 12, 31)
                    _r_ro = _dfp_ro[_dfp_ro["date"] <= _d_fin_ro]
                    _enc_ro_by_pret[_pn_ro].append(
                        _r_ro.iloc[-1]["capital_restant"] / 1000 if not _r_ro.empty else 0
                    )
        if _enc_ro_by_pret:
            _enc_ro_colors = ["#f5576c", "#667eea", "#11998e", "#ffcc00", "#764ba2", "#a0522d", "#f093fb"]
            fig = go.Figure()
            for _ci, (_pn_ro, _vals_ro) in enumerate(_enc_ro_by_pret.items()):
                fig.add_trace(go.Bar(x=xr, y=_vals_ro, name=_pn_ro,
                    marker_color=_enc_ro_colors[_ci % len(_enc_ro_colors)]))
            _enc_ro_total_k = [sum(v[i] for v in _enc_ro_by_pret.values()) for i in range(len(xr))]
            fig.add_trace(go.Scatter(x=xr, y=_enc_ro_total_k, mode="markers+text",
                text=[f"<b>{v:,.0f}</b>" for v in _enc_ro_total_k], textposition="top center",
                textfont=dict(size=11, color="#1a1a6e"), marker=dict(size=1, color="rgba(0,0,0,0)"),
                showlegend=False, hoverinfo="skip", cliponaxis=False))
            fig.update_layout(title="Rocher - Evolution endettement (K€)", height=400, barmode="stack",
                xaxis=dict(type="category"),
                yaxis=dict(tickformat=",.0f",
                           range=[0, max(_enc_ro_total_k) * 1.18] if _enc_ro_total_k else None),
                legend=_leg)
            _show_fig(fig, key="ro_endettement")

        # ════════════════════════════════════════════════════════════════════
        # 5. PLAN FINANCIER CHATEAU
        # ════════════════════════════════════════════════════════════════════
        st.header("7. Plan financier — Chateau d'Argenteau", anchor="sec-7", divider="red")

        # Donnees pour commentaire et hypotheses
        total_inv_ch = sum(inv["montant"] for inv in p.get("investissements", []))
        total_dettes_ch = sum(pr["montant"] for pr in prets_ch)
        total_dettes_ch_ext = sum(pr["montant"] for pr in prets_ch if "rocher" not in pr.get("nom", "").lower())
        taux_occ = p.get("taux_occ", [0.42, 0.475, 0.55, 0.62])
        segments = p.get("segments", {})
        adr_pondere = sum(s["part"]*s["prix"] for s in segments.values()) if segments else 0
        loyer_m_ch = p.get("loyer_mensuel", 0)

        # Personnel
        all_pers_ch = (p.get("personnel_hebergement",[]) + p.get("personnel_brasserie",[]) +
            p.get("personnel_spa",[]) + p.get("personnel_evenements",[]) + p.get("personnel_indirect",[]))
        total_etp_ch = sum(pe["etp"] for pe in all_pers_ch)
        cp_ch = p.get("charges_patronales_pct", 0.35)
        masse_ch = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in all_pers_ch)

        # Nuitees et CA en croisiere
        _nuitees_croisiere = int(nb_ch * 365 * taux_occ[-1]) if taux_occ else 0
        _ca_croisiere_est = _nuitees_croisiere * adr_pondere if adr_pondere else 0

        # Segments detail
        _seg_detail = ""
        for seg_name, seg_data in segments.items():
            _seg_detail += f"| Segment {seg_name} | Part : **{seg_data['part']*100:.0f}%** / Prix moyen : **{seg_data['prix']:,.0f} \u20ac** |\n"

        # Charges fixes
        _cf_directs = p.get("cf_directs", {})
        _cf_indirects = p.get("cf_indirects", {})
        _total_cf_dir = sum(_cf_directs.values()) if _cf_directs else 0
        _total_cf_indir = sum(_cf_indirects.values()) if _cf_indirects else 0

        # Commentaire editable — synthese Chateau
        _comment_key_ch = "commentaire_chateau"
        _default_comment_ch = p.get(_comment_key_ch,
            f"Le Chateau d'Argenteau exploite un hotel de {nb_ch} chambres positionne en 5 etoiles. "
            f"Le chiffre d'affaires est genere par l'hebergement, la brasserie, le bar, le spa et la location de salles. "
            f"En regime de croisiere (taux d'occupation de {taux_occ[-1]:.0%}), le CA hebergement est estime a environ {_ca_croisiere_est/1e6:,.1f} M\u20ac. "
            f"La masse salariale represente le poste de charge le plus important avec {total_etp_ch:.1f} ETP pour {masse_ch/1e6:,.2f} M\u20ac/an. "
            f"Le loyer verse a l'Immobiliere Rocher s'eleve a {loyer_m_ch*12:,.0f} \u20ac/an. "
            f"Le plan vise un EBITDA positif des les premieres annees de croisiere et un retour sur investissement progressif."
        )
        st.markdown("### Synthese & constats")
        _is_visu_report = st.session_state.get("auth_mode") == "visu"
        if _is_visu_report:
            st.markdown(f'<div style="background:#fef2f2; border-left:4px solid #f5576c; '
                        f'padding:14px 18px; border-radius:6px; margin-bottom:16px; '
                        f'font-size:0.95em; line-height:1.6; color:#3a1a1a;">'
                        f'{_default_comment_ch}</div>', unsafe_allow_html=True)
        else:
            _new_comment_ch = st.text_area(
                "Commentaire de synthese (editable)",
                value=_default_comment_ch,
                height=_text_area_height(_default_comment_ch, min_h=120),
                key="ta_comment_chateau",
            )
            if _new_comment_ch != p.get(_comment_key_ch):
                p[_comment_key_ch] = _new_comment_ch
                sauvegarder_plan(plan_nom, p, local_only=True)

        # Marge par service
        st.markdown("### Marge par service")
        st.markdown(
            f"Le Chateau d'Argenteau propose **6 sources de revenus** :\n\n"
            f"- Hebergement (5 etoiles)\n"
            f"- Brasserie (petits-dejeuners, midi et soir)\n"
            f"- Bar\n"
            f"- Spa & Bien-etre\n"
            f"- Location de salles (mariages, seminaires, ...)\n"
            f"- Location restaurant gastronomique"
        )

        # Donnees CV
        _cv_nuitee = p.get("cv_hebergement_par_nuitee", {})
        _cv_nuitee_total = sum(_cv_nuitee.values()) if _cv_nuitee else 0
        _cv_cb_nuitee = p.get("cv_commission_cb_nuitee", 1.5)
        _cv_cb_pct = p.get("cv_commission_cb_pct_chambres", 0.80)
        _cv_brass_pct = p.get("cv_brasserie_pct", 0.35)
        _cv_pdj_pct = p.get("cv_pdj_pct", _cv_brass_pct)
        _cv_bar_pct = p.get("cv_bar_pct", 0.30)
        _cv_bar_conso = p.get("cv_bar_consommable_unite", 0.20)
        _cv_spa_soin = p.get("cv_spa_soin_cout", 50)
        _cv_spa_prod = p.get("cv_spa_produits_soin", 5)
        _cv_salles_energie = p.get("cv_salles_chateau_energie", 100)
        _cv_salles_nett = p.get("cv_salles_chateau_nettoyage", 500)

        # Personnel par departement
        _pers_heberg = p.get("personnel_hebergement", [])
        _pers_brass = p.get("personnel_brasserie", [])
        _pers_spa = p.get("personnel_spa", [])
        _pers_event = p.get("personnel_evenements", [])
        _etp_h = sum(pe["etp"] for pe in _pers_heberg)
        _etp_b = sum(pe["etp"] for pe in _pers_brass)
        _etp_s = sum(pe["etp"] for pe in _pers_spa)
        _etp_e = sum(pe["etp"] for pe in _pers_event)
        _masse_h = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_heberg)
        _masse_b = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_brass)
        _masse_s = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_spa)
        _masse_e = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_event)

        # Style commun pour les tableaux de marge
        _tbl = ('<table style="width:100%; border-collapse:collapse; margin:8px 0 16px 0; font-size:0.9em;">'
                '<colgroup><col style="width:33.3%"><col style="width:33.3%"><col style="width:33.3%"></colgroup>'
                '<thead><tr style="background:#f8f9fa; border-bottom:2px solid #dee2e6;">'
                '<th style="padding:8px; text-align:left;">Ventes</th>'
                '<th style="padding:8px; text-align:left;">Frais variables</th>'
                '<th style="padding:8px; text-align:left;">Frais fixes directs</th>'
                '</tr></thead><tbody>')
        _r = lambda c1, c2, c3: (f'<tr><td style="padding:6px 8px; border-bottom:1px solid #eee;">{c1}</td>'
                                  f'<td style="padding:6px 8px; border-bottom:1px solid #eee;">{c2}</td>'
                                  f'<td style="padding:6px 8px; border-bottom:1px solid #eee;">{c3}</td></tr>')

        # Donnees OTA et franchise
        _ota_pct = p.get("cv_commission_ota_pct", 0.17)
        _poids_ota = sum(segments[s]["part"] * p.get("segments_part_ota", {}).get(s, 0) for s in segments)
        _franchise_pct = p.get("cv_franchise_pct", 0.04)
        _franchise_modes = p.get("cv_franchise_modes", ["pct"])
        _franchise_nuitee = p.get("cv_franchise_par_nuitee", 0)
        _franchise_forfait = p.get("cv_franchise_forfait_mois", 0)
        _franchise_desc = []
        if "pct" in _franchise_modes:
            _franchise_desc.append(f"{_franchise_pct:.0%} CA")
        if "nuitee" in _franchise_modes and _franchise_nuitee > 0:
            _franchise_desc.append(f"{_franchise_nuitee:,.1f} \u20ac/nuitee")
        if "forfait" in _franchise_modes and _franchise_forfait > 0:
            _franchise_desc.append(f"{_franchise_forfait:,.0f} \u20ac/mois")

        # 1. Hebergement
        st.markdown(f'**1. Hebergement**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'{nb_ch} chambres', f'CV / nuitee : <b>{_cv_nuitee_total:,.0f} \u20ac</b>', f'Personnel : <b>{_etp_h:.1f} ETP</b>')
            + _r(f'ADR : <b>{adr_pondere:,.0f} \u20ac</b>', f'Commission CB : <b>{_cv_cb_nuitee:,.1f} \u20ac</b> ({_cv_cb_pct:.0%} CB)', f'Masse salariale : <b>{_masse_h:,.0f} \u20ac/an</b>')
            + _r(f'Occ. An1/2/3/Crois. : <b>{taux_occ[0]:.0%} / {taux_occ[1]:.0%} / {taux_occ[2]:.0%} / {taux_occ[3]:.0%}</b>',
                 f'Commission OTA : <b>{_ota_pct:.0%}</b> (poids : {_poids_ota:.0%})', '')
            + _r(f'Saisonnalite mensuelle appliquee', f'Franchise : <b>{" + ".join(_franchise_desc) if _franchise_desc else "Aucune"}</b>', '')
            + '</tbody></table>', unsafe_allow_html=True)

        # Graphique saisonnalite (ligne avec voile)
        _saison = p.get("saisonnalite", [1]*12)
        _mois_labels = ["Jan", "Fev", "Mar", "Avr", "Mai", "Jun", "Jul", "Aou", "Sep", "Oct", "Nov", "Dec"]
        fig_saison = go.Figure()
        fig_saison.add_trace(go.Scatter(
            x=_mois_labels, y=_saison, mode="lines+markers+text",
            line=dict(color="#4facfe", width=3), marker=dict(size=8, color="#4facfe"),
            fill="tozeroy", fillcolor="rgba(79, 172, 254, 0.15)",
            text=[f"<b>{s:.2f}</b>" for s in _saison], textposition="top center",
            textfont=dict(size=11, color="#1a1a2e"),
            hovertemplate="%{x} : %{y:.3f}<extra></extra>"))
        fig_saison.add_hline(y=1.0, line_dash="dash", line_color="#888", opacity=0.5)
        fig_saison.update_layout(title="Coefficients de saisonnalite mensuelle", height=350,
            xaxis=dict(type="category", range=[-0.5, 11.8]),
            yaxis=dict(range=[0, max(_saison)*1.2], title="Coefficient"),
            margin=dict(l=40, r=60, t=40, b=30))
        _show_fig(fig_saison, key="ch_saisonnalite")

        # Graphiques nuitees et taux occupation / ADR / RevPAR
        if "_ann_ch" not in dir():
            _ann["revpar"] = _ann["ca_hebergement"] / (nb_ch * 365)
            _ann["adr"] = _ann["ca_hebergement"] / _ann["nuitees"].replace(0, float("nan"))
            _ann["occ_pct"] = _ann["taux_occupation"] * 100

        fig = go.Figure()
        fig.add_trace(go.Bar(x=_x, y=_ann["nuitees"], name="Nuitees", marker_color="#4facfe",
            text=[f"<b>{v:,.0f}</b>" for v in _ann["nuitees"]], textposition="outside",
            textfont=dict(size=10, color="#1a5e96")))
        fig.update_layout(title="Nuitees par an", height=350, xaxis=dict(type="category"),
            yaxis=dict(tickformat=",.0f", range=[0, max(_ann["nuitees"])*1.15]))
        _show_fig(fig, key="ch_nuitees_hyp")

        st.caption("**ADR** (Average Daily Rate) = prix moyen par chambre vendue | "
                   "**RevPAR** (Revenue Per Available Room) = revenu par chambre disponible = ADR x Taux d'occupation")
        fig = go.Figure()
        fig.add_trace(go.Bar(x=_x, y=_ann["adr"], name="ADR (\u20ac)", marker_color="#667eea",
            text=[f"<b>{v:,.0f}</b>" for v in _ann["adr"]], textposition="outside",
            textfont=dict(size=11, color="#23408f")))
        fig.add_trace(go.Bar(x=_x, y=_ann["revpar"], name="RevPAR (\u20ac)", marker_color="#f5576c",
            text=[f"<b>{v:,.0f}</b>" for v in _ann["revpar"]], textposition="outside",
            textfont=dict(size=11, color="#b71c2e")))
        fig.add_trace(go.Scatter(x=_x, y=_ann["occ_pct"], name="Taux occ. %", mode="lines+markers+text",
            line=dict(color="#11998e", width=3),
            text=[f"<b>{v:.0f}%</b>" for v in _ann["occ_pct"]], textposition="top center",
            textfont=dict(size=12, color="#0b6e66")))
        fig.update_layout(title="Taux occupation, ADR (Average Daily Rate), RevPAR", height=500, barmode="group",
            xaxis=dict(type="category"), yaxis=dict(range=[0, max(_ann["adr"])*1.25]), legend=_leg)
        _show_fig(fig, key="ch_occ_adr_hyp")

        # 2. Brasserie
        _taux_occ_brass = p.get("taux_occ_brasserie", p.get("taux_occ", [0.42]))
        _occ_brass_str = " / ".join(f"{t:.0%}" for t in _taux_occ_brass)
        st.markdown(f'**2. Brasserie & Petit-dejeuner**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'{p.get("nb_couverts_brasserie", 80)} couverts', f'Food cost midi/soir : <b>{_cv_brass_pct:.0%}</b>', f'Personnel : <b>{_etp_b:.1f} ETP</b>')
            + _r(f'Prix midi : <b>{p.get("brasserie_prix_diner", 45):,.0f} \u20ac</b>', f'Food cost PDJ : <b>{_cv_pdj_pct:.0%}</b>', f'Masse salariale : <b>{_masse_b:,.0f} \u20ac/an</b>')
            + _r(f'Prix soir : <b>{p.get("brasserie_prix_souper", 75):,.0f} \u20ac</b>', '', '')
            + _r(f'PDJ : <b>{p.get("petit_dej_prix", 37.5):,.0f} \u20ac</b> (taux prise {p.get("petit_dej_taux", 0.85):.0%})', '', '')
            + _r(f'Occ. brasserie : <b>{_occ_brass_str}</b>', '', '')
            + _r(f'Saisonnalite mensuelle appliquee', '', '')
            + '</tbody></table>', unsafe_allow_html=True)

        # 3. Bar
        _pers_bar = p.get("personnel_bar", [])
        _etp_bar = sum(pe["etp"] for pe in _pers_bar)
        _masse_bar = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_bar)
        _cf_bar_elec = p.get("cf_directs_bar", {}).get("Electricite", 0) if isinstance(p.get("cf_directs_bar"), dict) else 0
        st.markdown(f'**3. Bar**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'Prix client interne : <b>{p.get("bar_conso_moyenne", 22):,.0f} \u20ac</b>', f'Beverage cost : <b>{_cv_bar_pct:.0%}</b>', f'Personnel : <b>{_etp_bar:.1f} ETP</b>')
            + _r(f'Prix client externe : <b>{p.get("bar_conso_ext_moyenne", 25):,.0f} \u20ac</b>', f'Consommables : <b>{_cv_bar_conso:.2f} \u20ac/conso</b>', f'Masse salariale : <b>{_masse_bar:,.0f} \u20ac/an</b>')
            + _r(f'Taux clients hotel : <b>{p.get("bar_taux_clients_hotel", 0.40):.0%}</b>', '', '')
            + _r('', '', f'Electricite : <b>{_cf_bar_elec:,.0f} \u20ac/an</b>' if _cf_bar_elec > 0 else '')
            + '</tbody></table>', unsafe_allow_html=True)

        # 4. Spa
        _cf_spa = p.get("cf_directs_spa", {})
        _cf_spa_details = " | ".join(f"{k} : <b>{v:,.0f} \u20ac</b>" for k, v in _cf_spa.items() if v > 0)
        st.markdown(f'**4. Spa & Bien-etre**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'Entree hotel : <b>{p.get("spa_entree_hotel_prix", 0):,.0f} \u20ac</b> (taux {p.get("spa_entree_hotel_taux", 0.20):.0%})', f'Cout soin : <b>{_cv_spa_soin:,.0f} \u20ac</b>', f'Personnel : <b>{_etp_s:.1f} ETP</b>')
            + _r(f'Soin hotel : <b>{p.get("spa_soin_hotel_prix", 120):,.0f} \u20ac</b> (taux {p.get("spa_soin_hotel_taux", 0.10):.0%})', f'Produits : <b>{_cv_spa_prod:,.0f} \u20ac/soin</b>', f'Masse salariale : <b>{_masse_s:,.0f} \u20ac/an</b>')
            + _r(f'Entree ext. : <b>{p.get("spa_entree_ext_prix", 55):,.0f} \u20ac</b> ({p.get("spa_entree_ext_nb_mois", 25)}/mois)', '', _cf_spa_details if _cf_spa_details else '')
            + _r(f'Soin ext. : <b>{p.get("spa_soin_ext_prix", 150):,.0f} \u20ac</b> ({p.get("spa_soin_ext_nb_mois", 15)}/mois)', '', '')
            + '</tbody></table>', unsafe_allow_html=True)

        # 5. Salles & Evenements
        _cf_events = p.get("cf_directs_evenements", {})
        _cf_events_details = " | ".join(f"{k} : <b>{v:,.0f} \u20ac</b>" for k, v in _cf_events.items() if v > 0)
        _catering_prix = p.get("mariage_catering_prix_convive", 0)
        _catering_pct = p.get("mariage_commission_catering_pct", 0)
        _nb_convives = p.get("mariage_nb_convives_moy", 120)
        _catering_info = f"Commission catering : <b>{_catering_pct:.0%}</b> sur {_nb_convives} convives x {_catering_prix:,.0f} \u20ac" if _catering_pct > 0 else ""
        st.markdown(f'**5. Salles & Evenements**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'Seminaires : <b>{p.get("seminaire_nb_an", 50)}/an</b> a {p.get("seminaire_prix_location", 800):,.0f} \u20ac',
                f'Energie : <b>{p.get("cv_seminaire_energie", 75):,.0f} \u20ac</b> | Nettoyage : <b>{p.get("cv_seminaire_nettoyage", 500):,.0f} \u20ac</b> | Pause : <b>{p.get("cv_seminaire_pause_participant", 8):,.0f} \u20ac/pers</b>',
                _cf_events_details if _cf_events_details else '')
            + _r(f'Mariages : <b>{p.get("mariage_nb_an", 12)}/an</b> a {p.get("mariage_prix_location", 2500):,.0f} \u20ac',
                f'Energie : <b>{p.get("cv_mariage_energie", 150):,.0f} \u20ac</b> | Nettoyage : <b>{p.get("cv_mariage_nettoyage", 1000):,.0f} \u20ac</b>',
                '')
            + (_r(_catering_info, '', '') if _catering_info else '')
            + _r(f'Salles chateau : <b>{p.get("salles_chateau_nb_an", 30)}/an</b> a {p.get("salles_chateau_prix", 1500):,.0f} \u20ac',
                f'Energie : <b>{p.get("cv_salles_chateau_energie", 100):,.0f} \u20ac</b> | Nettoyage : <b>{p.get("cv_salles_chateau_nettoyage", 500):,.0f} \u20ac</b>',
                '')
            + '</tbody></table>', unsafe_allow_html=True)

        # 6. Location restaurant gastronomique
        _loyer_resto = p.get("loyer_restaurant_mensuel", 0)
        st.markdown(f'**6. Location restaurant gastronomique**', unsafe_allow_html=True)
        st.markdown(_tbl
            + _r(f'Loyer mensuel : <b>{_loyer_resto:,.0f} \u20ac/mois</b> ({_loyer_resto*12:,.0f} \u20ac/an)', '', '')
            + '</tbody></table>', unsafe_allow_html=True)

        # ── Graphiques par service ──
        st.markdown("---")

        # Marge par service + Subside (marge = ventes - CV - frais fixes directs)
        fig = go.Figure()
        _ann["marge_heberg"] = _ann["ca_hebergement"] - _ann["cv_hebergement"] - _ann["cf_directs_hebergement"]
        _ann["marge_brass"] = _ann["ca_brasserie"] - _ann["cv_brasserie"] - _ann["cf_directs_brasserie"]
        _ann["marge_bar"] = _ann["ca_bar"] - _ann["cv_bar"] - _ann["cf_directs_bar"]
        _ann["marge_spa"] = _ann["ca_spa"] - _ann["cv_spa"] - _ann["cf_directs_spa"]
        _ann["marge_salles"] = _ann["ca_salles"] - _ann["cf_directs_evenements"]
        _ann["marge_resto"] = _ann.get("ca_loyer_restaurant", 0)
        _marge_services_hyp = [("marge_heberg","Hebergement","#667eea"),("marge_brass","Brasserie","#f5576c"),
            ("marge_bar","Bar","#ffcc00"),("marge_spa","Spa","#11998e"),("marge_salles","Salles","#a0522d"),
            ("marge_resto","Location resto.","#ff8c00")]
        for col, lbl, clr in _marge_services_hyp:
            fig.add_trace(go.Bar(x=_x, y=K(_ann[col]), name=lbl, marker_color=clr))
        _marge_tot_hyp = K(_ann["marge"])
        _max_mh = max(_marge_tot_hyp) if len(_marge_tot_hyp) > 0 else 1
        fig.add_trace(go.Scatter(x=_x, y=[v + _max_mh * 0.04 for v in _marge_tot_hyp], mode="text",
            text=[f"<b>{v:,.0f}</b>" for v in _marge_tot_hyp], textposition="top center",
            textfont=dict(size=12, color="#1a1a6e"), showlegend=False, hoverinfo="skip", cliponaxis=False))
        fig.update_layout(title="Marge par service (K\u20ac)", height=500, barmode="stack",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f", range=[0, _max_mh * 1.5]), legend=_leg)
        _show_fig(fig, key="ch_marge_hyp")

        # Camembert repartition marges (cumul) — marge = ventes - CV - frais fixes directs
        _marge_cumul_hyp = {
            "Hebergement": _ann["marge_heberg"].sum(),
            "Brasserie": _ann["marge_brass"].sum(),
            "Bar": _ann["marge_bar"].sum(),
            "Spa": _ann["marge_spa"].sum(),
            "Salles": _ann["marge_salles"].sum(),
            "Location resto.": _ann["marge_resto"].sum(),
        }
        _marge_cumul_hyp = {k: v for k, v in _marge_cumul_hyp.items() if v > 0}
        if _marge_cumul_hyp:
            _mc_colors = {"Hebergement":"#667eea","Brasserie":"#f5576c","Bar":"#ffcc00",
                          "Spa":"#11998e","Salles":"#a0522d","Location resto.":"#ff8c00"}
            fig = go.Figure(data=[go.Pie(
                labels=list(_marge_cumul_hyp.keys()), values=list(_marge_cumul_hyp.values()),
                marker=dict(colors=[_mc_colors.get(k, "#888") for k in _marge_cumul_hyp]),
                textinfo="label+percent",
                texttemplate="<b>%{label}</b><br>%{percent}",
                textfont=dict(size=13), hole=0.3, pull=[0.02]*len(_marge_cumul_hyp),
            )])
            fig.update_layout(title="Repartition des marges par service (cumul)", height=450,
                              showlegend=False, margin=dict(l=20, r=20, t=40, b=20))
            _show_fig(fig, key="ch_pie_marge_hyp")

        # ── Frais fixes indirects ──
        _pers_indirect = p.get("personnel_indirect", [])
        _etp_ind = sum(pe["etp"] for pe in _pers_indirect)
        _masse_ind = sum(pe["cout_brut"]*(1+cp_ch)*pe["etp"] for pe in _pers_indirect)
        _cfi_dict = p.get("charges_fixes_indirectes", p.get("charges_fixes_indirectes_par_annee", {}))

        st.markdown("### Frais fixes indirects")

        # Tableau personnel indirect
        _hc1, _hc2 = st.columns(2)
        with _hc1:
            st.markdown("**Personnel indirect**")
            _pi_rows = ""
            for pe in _pers_indirect:
                _cout_total = pe["cout_brut"] * (1 + cp_ch) * pe["etp"]
                _pi_rows += f"| {pe['poste']} | {pe['etp']:.1f} | **{_cout_total:,.0f} \u20ac** |\n"
            st.markdown(f"""
| Poste | ETP | Cout total/an |
|---|---|---|
{_pi_rows}| **Total** | **{_etp_ind:.1f}** | **{_masse_ind:,.0f} \u20ac** |
""")

        with _hc2:
            st.markdown("**Autres charges indirectes**")
            _cf_ind_rows = ""
            _total_autres_val = 0
            for _cf_k, _cf_v in _cfi_dict.items():
                _val = _cf_v[0] if isinstance(_cf_v, list) else _cf_v
                if _val > 0:
                    _cf_ind_rows += f"| {_cf_k} | **{_val:,.0f} \u20ac/an** |\n"
                    _total_autres_val += _val
            _cf_ind_rows += f"| Loyer (vers Rocher) | **{loyer_m_ch*12:,.0f} \u20ac/an** ({loyer_m_ch:,.0f} \u20ac/mois) |\n"
            _total_autres_val += loyer_m_ch * 12
            st.markdown(f"""
| Poste | Montant |
|---|---|
{_cf_ind_rows}| **Total** | **{_total_autres_val:,.0f} \u20ac/an** |
""")
        st.caption(f"Inflation annuelle prevue : **{p.get('inflation_an', 0.025):.1%}**")

        # Graphique CF indirects — distinguer salaires vs autres
        _cfi_k = K(_ann["cf_indirects_total"])
        _cf_pers_ind_k = K(_ann.get("cf_personnel_indirect", _ann["cf_indirects_total"] * 0))
        # Estimer la part salaires vs autres si la colonne n'existe pas
        _total_cfi_base = _masse_ind + _total_autres_val
        _pct_sal = _masse_ind / _total_cfi_base if _total_cfi_base > 0 else 0.5
        fig = go.Figure()
        fig.add_trace(go.Bar(x=_x, y=[v * _pct_sal for v in _cfi_k], name="Personnel indirect",
            marker_color="#764ba2"))
        fig.add_trace(go.Bar(x=_x, y=[v * (1 - _pct_sal) for v in _cfi_k], name="Autres charges + Loyer",
            marker_color="#c39bd3"))
        fig.add_trace(go.Scatter(x=_x, y=_cfi_k, mode="markers+text",
            text=[f"<b>{v:,.0f}</b>" for v in _cfi_k], textposition="top center",
            textfont=dict(size=13, color="#4a2470"), marker=dict(size=1, color="rgba(0,0,0,0)"),
            showlegend=False, hoverinfo="skip", cliponaxis=False))
        fig.update_layout(title="Frais fixes indirects (K\u20ac)", height=420, barmode="stack",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f",
            range=[0, max(_cfi_k)*1.25] if len(_cfi_k)>0 else None), legend=_leg)
        _show_fig(fig, key="ch_cf_indirects")

        # Emploi (apres CF indirects)
        _cmt("before_ch_emploi")
        st.markdown("### Emploi")
        pers_dept = {
            "Hebergement": p.get("personnel_hebergement", []),
            "Brasserie": p.get("personnel_brasserie", []),
            "Spa": p.get("personnel_spa", []),
            "Evenements": p.get("personnel_evenements", []),
            "Indirect": p.get("personnel_indirect", []),
        }
        etp_dept = {k: sum(pe["etp"] for pe in v) for k, v in pers_dept.items()}
        total_etp = sum(etp_dept.values())
        cp_rh = p.get("charges_patronales_pct", 0.35)
        masse_tot = sum(pe["cout_brut"]*(1+cp_rh)*pe["etp"] for dept in pers_dept.values() for pe in dept)

        c1, c2, c3, c4 = st.columns(4)
        _card = lambda val, lbl, clr: (
            f'<div style="text-align:center; padding:12px; background:linear-gradient(135deg, {clr}22, {clr}11); '
            f'border-left:4px solid {clr}; border-radius:8px; margin:4px 0;">'
            f'<div style="font-size:1.8em; font-weight:bold; color:{clr};">{val}</div>'
            f'<div style="font-size:0.85em; color:#666;">{lbl}</div></div>')
        with c1:
            st.markdown(_card(f"{total_etp:.1f}", "Total ETP", "#667eea"), unsafe_allow_html=True)
        with c2:
            st.markdown(_card(f"{total_etp/nb_ch:.2f}" if nb_ch>0 else "-", "ETP / chambre", "#11998e"), unsafe_allow_html=True)
        with c3:
            st.markdown(_card(f"{nb_ch/total_etp:.1f}" if total_etp>0 else "-", "Chambres / ETP", "#f5576c"), unsafe_allow_html=True)
        with c4:
            st.markdown(_card(f"{masse_tot/1000:,.0f} K\u20ac", "Masse salariale / an", "#764ba2"), unsafe_allow_html=True)

        _dept_labels = [k for k, v in etp_dept.items() if v > 0]
        _dept_vals = [v for v in etp_dept.values() if v > 0]
        _dept_colors = ["#667eea", "#f5576c", "#11998e", "#a0522d", "#764ba2"]
        fig = go.Figure(data=[go.Pie(labels=_dept_labels, values=_dept_vals,
            marker=dict(colors=_dept_colors[:len(_dept_labels)]),
            textinfo="label+value+percent",
            texttemplate="<b>%{label}</b><br><b>%{value:.1f} ETP</b><br>(%{percent})",
            textfont=dict(size=15),
            pull=[0.02] * len(_dept_labels),
            hole=0.25)])
        fig.update_layout(height=450, showlegend=False, margin=dict(l=20,r=20,t=30,b=20),
                          title=dict(text="Repartition du personnel par departement", font=dict(size=15)))
        _show_fig(fig, key="ch_rh_pie")

        # ── Resultats ──
        st.markdown("### Resultats")

        # EBITDA / Amort / Interets / Subside RW en batonnets groupes + Resultat Net en ligne
        _ebitda_k = K(_ann["ebitda"])
        _amort_k = K(_ann["amortissement"])
        _int_k = K(_ann["dette_interets"])
        _sub_k = K(_ann["subside_rw"])
        _rn_k = K(_ann["resultat_net"])
        fig = go.Figure()
        fig.add_trace(go.Bar(x=_x, y=_ebitda_k, name="EBITDA", marker_color="#11998e"))
        fig.add_trace(go.Bar(x=_x, y=-_amort_k, name="Amortissement", marker_color="#764ba2"))
        fig.add_trace(go.Bar(x=_x, y=-_int_k, name="Interets", marker_color="#ffcc00"))
        fig.add_trace(go.Bar(x=_x, y=_sub_k, name="Subsides RW", marker_color="#f093fb"))
        # Labels EBITDA au dessus des barres
        fig.add_trace(go.Scatter(x=_x, y=_ebitda_k, mode="text",
            text=[f"<b>{v:,.0f}</b>" for v in _ebitda_k], textposition="top center",
            textfont=dict(size=11, color="#0b6e66"), showlegend=False, hoverinfo="skip"))
        # Resultat Net en ligne avec labels en dessous
        fig.add_trace(go.Scatter(x=_x, y=_rn_k, name="Resultat Net", mode="lines+markers+text",
            line=dict(color="#f5576c", width=3),
            text=[f"<b>{v:,.0f}</b>" for v in _rn_k], textposition="bottom center",
            textfont=dict(size=12, color="#b71c2e")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray")
        fig.update_layout(title="Du EBITDA au Resultat Net (K\u20ac)", height=450, barmode="group",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ch_ebitda_rn")

        # Cash flow + tresorerie
        surplus = fp_ch + sum(pr["montant"] for pr in prets_ch) - sum(inv["montant"] for inv in p["investissements"])
        _cf_k = K(_ann["cash_flow"])
        _treso_k = K(_ann["cash_flow_cumul"] + surplus)
        fig = go.Figure()
        fig.add_trace(go.Bar(x=_x, y=_cf_k, name="Cash flow", marker_color="#4facfe",
            text=[f"<b>{v:,.0f}</b>" for v in _cf_k], textposition="outside",
            textfont=dict(size=12, color="#1a5e96")))
        fig.add_trace(go.Scatter(x=_x, y=K(_ann["cash_flow_cumul"]), name="Cumule", mode="lines+markers",
            line=dict(color="#667eea", width=2, dash="dash")))
        # Tresorerie : afficher le montant uniquement si negatif
        _treso_text = [f'<b style="color:red">{v:,.0f}</b>' if v < 0 else "" for v in _treso_k]
        fig.add_trace(go.Scatter(x=_x, y=_treso_k, name="Tresorerie", mode="lines+markers+text",
            line=dict(color="#e04458", width=3),
            text=_treso_text, textposition="bottom center",
            textfont=dict(size=12, color="#b71c2e")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray")
        fig.update_layout(title="Cash flow et Tresorerie (K\u20ac)", height=450,
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ch_cashflow")

        # Detail des chiffres - Chateau (titre a la taille des titres Plotly)
        st.markdown(
            '<div style="font-size:17px; font-weight:700; color:#2a3f5f; '
            'margin:16px 0 8px 0; font-family:Arial, sans-serif;">Details des chiffres</div>',
            unsafe_allow_html=True,
        )
        _fmt_ch = lambda v: f"{v:,.0f} €"
        _pct_ch = lambda num, den: f"{num/den*100:.1f}%" if den != 0 else "-"
        _df_tbl_ch = pd.DataFrame({
            "Annee": [str(int(a)) for a in _ann["annee"]],
            "CA Total": [_fmt_ch(v) for v in _ann["ca_total"]],
            "CV Total": [_fmt_ch(v) for v in _ann["cv_total"]],
            "Marge Brute": [_fmt_ch(v) for v in _ann["marge_brute"]],
            "Marge Brute %": [_pct_ch(m, c) for m, c in zip(_ann["marge_brute"], _ann["ca_total"])],
            "CF Directs": [_fmt_ch(v) for v in _ann["cf_directs_total"]],
            "Marge service": [_fmt_ch(v) for v in _ann["marge"]],
            "Subside RW": [_fmt_ch(v) for v in _ann["subside_rw"]],
            "CF Indirects": [_fmt_ch(v) for v in _ann["cf_indirects_total"]],
            "EBITDA": [_fmt_ch(v) for v in _ann["ebitda"]],
            "EBITDA %": [_pct_ch(e, c) for e, c in zip(_ann["ebitda"], _ann["ca_total"])],
            "Amortissement": [_fmt_ch(v) for v in _ann["amortissement"]],
            "EBIT": [_fmt_ch(v) for v in _ann["ebit"]],
            "Interets": [_fmt_ch(v) for v in _ann["dette_interets"]],
            "Resultat avant impot": [_fmt_ch(v) for v in (_ann["ebit"] - _ann["dette_interets"] + _ann["subside_rw"])],
            "ISOC (charge)": [_fmt_ch(v) for v in _ann["impot"]],
            "Resultat Net": [_fmt_ch(v) for v in _ann["resultat_net"]],
            "ISOC (paye)": [_fmt_ch(v) for v in _ann["impot_cash"]],
            "Remb. dette": [_fmt_ch(v) for v in _ann["dette_capital"]],
            "Reinvestissements": [_fmt_ch(v) for v in _ann["reinvest_acquisition"]],
            "TVA a reverser": [_fmt_ch(v) for v in _ann["tva_paiement"]],
            "Delta Cash Pers.": [_fmt_ch(v) for v in (_ann["cf_total_cash"] - _ann["cf_total"])],
            "BFR (impact cash)": [_fmt_ch(v) for v in _ann["delay_adjustment"]],
            "Cash Flow": [_fmt_ch(v) for v in _ann["cash_flow"]],
            "Cash Flow Cumul": [_fmt_ch(v) for v in _ann["cash_flow_cumul"]],
        })
        st.dataframe(_df_tbl_ch, use_container_width=True, hide_index=True)

        # 5b. Bilan
        st.markdown("### Bilan")
        pret_rocher_mt = next((pr["montant"] for pr in prets_ch if "rocher" in pr.get("nom","").lower()), 0)
        res_cum = _ann["resultat_net"].cumsum()
        # FP = capital + resultats cumules, Quasi-FP = + pret Rocher
        _fp_k = K(fp_ch + res_cum)
        _fp_quasi_k = K(fp_ch + pret_rocher_mt + res_cum)
        fig = go.Figure()
        # Ligne FP (sans quasi) — labels en dessous
        _fp_text = [f'<b>{v:,.0f}</b>' for v in _fp_k]
        fig.add_trace(go.Scatter(x=_x, y=_fp_k, name="Fonds propres", mode="lines+markers+text",
            line=dict(color="#f5576c", width=3),
            text=_fp_text, textposition="bottom center",
            textfont=dict(size=10, color="#b71c2e")))
        # Ligne FP + quasi-FP — labels au dessus
        _fpq_text = [f'<b>{v:,.0f}</b>' for v in _fp_quasi_k]
        fig.add_trace(go.Scatter(x=_x, y=_fp_quasi_k, name="FP + Quasi-FP (Pret Rocher)", mode="lines+markers+text",
            line=dict(color="#667eea", width=3),
            text=_fpq_text, textposition="top center",
            textfont=dict(size=10, color="#23408f")))
        fig.add_hline(y=0, line_dash="dot", line_color="gray")
        fig.update_layout(title="Evolution fonds propres (K\u20ac)", height=500,
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"), legend=_leg)
        _show_fig(fig, key="ch_fp")

        # Endettement — ventile par pret
        _enc_by_pret = {}
        for pr in prets_ch:
            if pr["montant"] > 0:
                _pn = pr["nom"]
                if pr.get("subside_rw"):
                    _pn += " (RW)"
                _enc_by_pret[_pn] = []
                dfp = calc_tableau_pret(pr, p["date_ouverture"], p["nb_mois_projection"])
                for a in _ann["annee"]:
                    d_fin = date(int(a), 12, 31)
                    r = dfp[dfp["date"] <= d_fin]
                    _enc_by_pret[_pn].append(r.iloc[-1]["capital_restant"] / 1000 if not r.empty else 0)
        _enc_colors = ["#f5576c", "#667eea", "#11998e", "#ffcc00", "#764ba2", "#a0522d", "#f093fb"]
        fig = go.Figure()
        for _ci, (_pn, _vals) in enumerate(_enc_by_pret.items()):
            fig.add_trace(go.Bar(x=_x, y=_vals, name=_pn,
                marker_color=_enc_colors[_ci % len(_enc_colors)]))
        _enc_total_k = [sum(v[i] for v in _enc_by_pret.values()) for i in range(len(_x))]
        fig.add_trace(go.Scatter(x=_x, y=_enc_total_k, mode="markers+text",
            text=[f"<b>{v:,.0f}</b>" for v in _enc_total_k], textposition="top center",
            textfont=dict(size=11, color="#1a1a6e"), marker=dict(size=1, color="rgba(0,0,0,0)"),
            showlegend=False, hoverinfo="skip", cliponaxis=False))
        fig.update_layout(title="Evolution endettement (K\u20ac)", height=400, barmode="stack",
            xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f",
            range=[0, max(_enc_total_k)*1.18] if _enc_total_k else None), legend=_leg)
        _show_fig(fig, key="ch_endettement")

        # Ratio solvabilite
        # Numerateur = FP durs (capital + resultat cumule) + Quasi-FP (pret Rocher)
        # Denominateur = Total passif = Dettes totales (encours) + FP + Quasi-FP
        # NB: _enc_by_pret est en K€, convertir en EUR (* 1000)
        ratios = []
        for i in range(len(_x)):
            _enc_i_eur = sum(v[i] for v in _enc_by_pret.values()) * 1000 if _enc_by_pret else 0
            _res_cum_i = res_cum.iloc[i] if i < len(res_cum) else 0
            _fp_quasi = fp_ch + _res_cum_i + pret_rocher_mt  # FP + Quasi-FP
            _total_passif = _enc_i_eur + fp_ch + _res_cum_i + pret_rocher_mt  # Dettes + FP + Quasi-FP
            ratios.append(_fp_quasi / _total_passif * 100 if _total_passif > 0 else 0)
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=_x, y=ratios, name="Ratio solvabilite", mode="lines+markers+text",
            line=dict(color="#667eea", width=3), fill="tozeroy", fillcolor="rgba(102,126,234,0.1)",
            text=[f"<b>{v:.0f}%</b>" for v in ratios], textposition="top center",
            textfont=dict(size=11, color="#23408f")))
        fig.add_hline(y=30, line_dash="dash", line_color="orange", annotation_text="Seuil 30%")
        fig.update_layout(title="Ratio (FP+Quasi-FP) / Total passif (%)", height=400,
            xaxis=dict(type="category"), yaxis=dict(range=[0, 105], tickformat=".0f"))
        _show_fig(fig, key="ch_solvabilite")

        # ════════════════════════════════════════════════════════════════════
        # 8. SIMULATION (curseurs interactifs)
        # ════════════════════════════════════════════════════════════════════
        if not print_mode:
            st.header("8. Simulation", anchor="sec-8", divider="violet")
            st.caption("Faites varier les hypotheses cles et observez l'impact en temps reel sur l'EBITDA, "
                       "le Resultat Net et la tresorerie. Les variations s'appliquent sur **toute la projection**, "
                       "sauf les taux d'occupation qui se reglent par annee.")

            _sim_keys = [
                "sim_cfdh", "sim_prix_heb",
                "sim_occ_heb_0", "sim_occ_heb_1", "sim_occ_heb_2", "sim_occ_heb_3",
                "sim_occ_brass_0", "sim_occ_brass_1", "sim_occ_brass_2", "sim_occ_brass_3",
                "sim_prix_brass", "sim_marge_brass",
                "sim_cf_ind",
            ]
            if st.button("\U0001F504 Reinitialiser tous les curseurs", key="sim_reset"):
                for _k in _sim_keys:
                    st.session_state.pop(_k, None)
                st.rerun()

            sc1, sc2, sc3 = st.columns(3)
            with sc1:
                st.markdown("**\U0001F3E8 Hebergement**")
                d_cf_dir_heb = st.slider("Frais fixes directs (%)", -50, 100, 0, 5, key="sim_cfdh",
                                         help="Variation appliquee aux frais fixes directs hebergement (linge, accueil, amenities...).")
                d_prix_heb = st.slider("Prix moyen (%)", -30, 50, 0, 1, key="sim_prix_heb",
                                       help="Variation des prix de tous les segments hebergement.")
                st.markdown("*Taux d'occupation par annee (pts)*")
                _occ_labels = ["An 1", "An 2", "An 3", "Croisiere (an 4+)"]
                d_occ_heb = []
                for i, lbl in enumerate(_occ_labels):
                    d_occ_heb.append(st.slider(lbl, -20, 20, 0, 1, key=f"sim_occ_heb_{i}",
                                               help=f"Points de pourcentage ajoutes au taux d'occupation {lbl}."))
            with sc2:
                st.markdown("**\U0001F37D Brasserie**")
                d_prix_brass = st.slider("Prix (%)", -30, 50, 0, 1, key="sim_prix_brass",
                                         help="Variation des prix souper/diner/midi/soir.")
                d_marge_brass = st.slider("Marge (pts)", -20, 20, 0, 1, key="sim_marge_brass",
                                          help="Reduction du % de couts variables (food cost). +5 pts de marge = -5 pts de food cost.")
                st.markdown("*Taux d'occupation brasserie par annee (pts)*")
                d_occ_brass = []
                for i, lbl in enumerate(_occ_labels):
                    d_occ_brass.append(st.slider(lbl, -20, 20, 0, 1, key=f"sim_occ_brass_{i}",
                                                 help=f"Points de pourcentage ajoutes au taux d'occupation brasserie {lbl}."))
            with sc3:
                st.markdown("**\U0001F4CA General**")
                d_cf_ind = st.slider("Frais fixes indirects (%)", -50, 100, 0, 5, key="sim_cf_ind",
                                     help="Variation appliquee au personnel indirect ET aux charges fixes indirectes (loyer exclu).")

            # ── Construction des params modifies ──
            p_sim = copy.deepcopy(p)
            # Hebergement
            for _k in list(p_sim.get("cf_directs_hebergement", {}).keys()):
                p_sim["cf_directs_hebergement"][_k] *= (1 + d_cf_dir_heb / 100)
            for _i in range(min(4, len(p_sim.get("taux_occ", [])))):
                p_sim["taux_occ"][_i] = max(0.0, min(1.0, p_sim["taux_occ"][_i] + d_occ_heb[_i] / 100))
            for _seg in p_sim.get("segments", {}).values():
                _seg["prix"] = _seg.get("prix", 0) * (1 + d_prix_heb / 100)
            # Brasserie
            if "taux_occ_brasserie" in p_sim:
                for _i in range(min(4, len(p_sim["taux_occ_brasserie"]))):
                    p_sim["taux_occ_brasserie"][_i] = max(0.0, min(1.0,
                        p_sim["taux_occ_brasserie"][_i] + d_occ_brass[_i] / 100))
            for _kp in ("brasserie_prix_souper", "brasserie_prix_diner",
                        "brasserie_prix_midi", "brasserie_prix_soir"):
                if _kp in p_sim:
                    p_sim[_kp] = p_sim[_kp] * (1 + d_prix_brass / 100)
            if "cv_brasserie_pct" in p_sim:
                p_sim["cv_brasserie_pct"] = max(0.0, p_sim["cv_brasserie_pct"] - d_marge_brass / 100)
            # General : frais fixes indirects (personnel + charges fixes indirectes hors loyer)
            _factor = 1 + d_cf_ind / 100
            for _poste in p_sim.get("personnel_indirect", []):
                _poste["cout_brut"] = _poste.get("cout_brut", 0) * _factor
            if "charges_fixes_indirectes_par_annee" in p_sim:
                p_sim["charges_fixes_indirectes_par_annee"] = {
                    _k: [_v * _factor for _v in _vals]
                    for _k, _vals in p_sim["charges_fixes_indirectes_par_annee"].items()
                }
            if "charges_fixes_indirectes" in p_sim:
                p_sim["charges_fixes_indirectes"] = {
                    _k: _v * _factor
                    for _k, _v in p_sim["charges_fixes_indirectes"].items()
                }

            # ── Recalcul ──
            try:
                df_sim = projection_complete(p_sim)
                ann_sim = df_sim.groupby("annee").agg({
                    "ebitda": "sum", "resultat_net": "sum",
                    "cash_flow": "sum", "cash_flow_cumul": "last",
                }).reset_index()

                # Reference (params actuels)
                ann_ref_eb = list(_ann["ebitda"].values)
                ann_ref_rn = list(_ann["resultat_net"].values)
                surplus_sim = (p_sim.get("fonds_propres_initial", 0)
                              + sum(pr["montant"] for pr in p_sim.get("prets", []))
                              - sum(inv["montant"] for inv in p_sim.get("investissements", [])))
                surplus_ref = (fp_ch + sum(pr["montant"] for pr in prets_ch)
                              - sum(inv["montant"] for inv in p["investissements"]))
                ann_ref_treso = [v + surplus_ref for v in _ann["cash_flow_cumul"].values]
                ann_sim_treso = [v + surplus_sim for v in ann_sim["cash_flow_cumul"].values]

                _x_sim = [str(int(a)) for a in ann_sim["annee"]]

                # ── KPIs deltas (annee de croisiere = derniere) ──
                _last = -1
                _delta_eb = (ann_sim["ebitda"].iloc[_last] - ann_ref_eb[_last]) / 1000
                _delta_rn = (ann_sim["resultat_net"].iloc[_last] - ann_ref_rn[_last]) / 1000
                _delta_tr = (ann_sim_treso[_last] - ann_ref_treso[_last]) / 1000
                _kpi_card = (lambda v, lbl, clr: (
                    f'<div style="background:linear-gradient(135deg,{clr}15,{clr}30); '
                    f'padding:14px; border-radius:10px; text-align:center; border:1px solid {clr}50;">'
                    f'<div style="font-size:1.5em; font-weight:700; color:{clr};">{v:+,.0f} K€</div>'
                    f'<div style="font-size:0.8em; color:#555;">{lbl} (an final)</div></div>'
                ))
                k1, k2, k3 = st.columns(3)
                with k1:
                    st.markdown(_kpi_card(_delta_eb, "EBITDA",
                        "#11998e" if _delta_eb >= 0 else "#dc2626"), unsafe_allow_html=True)
                with k2:
                    st.markdown(_kpi_card(_delta_rn, "Resultat Net",
                        "#11998e" if _delta_rn >= 0 else "#dc2626"), unsafe_allow_html=True)
                with k3:
                    st.markdown(_kpi_card(_delta_tr, "Tresorerie cumulee",
                        "#11998e" if _delta_tr >= 0 else "#dc2626"), unsafe_allow_html=True)

                # ── Graphiques comparatifs ──
                def _comp_chart(title, ref_vals, sim_vals, color_sim, suffix=""):
                    f = go.Figure()
                    f.add_trace(go.Scatter(x=_x_sim, y=[v/1000 for v in ref_vals],
                        name="Reference", mode="lines+markers",
                        line=dict(color="#9ca3af", width=2, dash="dot")))
                    f.add_trace(go.Scatter(x=_x_sim, y=[v/1000 for v in sim_vals],
                        name="Simulation", mode="lines+markers+text",
                        line=dict(color=color_sim, width=3),
                        text=[f"<b>{v/1000:,.0f}</b>" for v in sim_vals],
                        textposition="top center", textfont=dict(size=10, color=color_sim)))
                    f.add_hline(y=0, line_dash="dot", line_color="gray")
                    f.update_layout(title=f"{title} (K€){suffix}", height=350,
                        xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f"),
                        legend=_leg, margin=dict(l=50, r=30, t=40, b=40))
                    return f

                st.plotly_chart(
                    _comp_chart("EBITDA", ann_ref_eb, list(ann_sim["ebitda"].values), "#11998e"),
                    use_container_width=True, key="sim_ebitda")
                st.plotly_chart(
                    _comp_chart("Resultat Net", ann_ref_rn, list(ann_sim["resultat_net"].values), "#f5576c"),
                    use_container_width=True, key="sim_rn")
                st.plotly_chart(
                    _comp_chart("Tresorerie cumulee", ann_ref_treso, ann_sim_treso, "#4facfe",
                                suffix=" — incl. surplus depart"),
                    use_container_width=True, key="sim_treso")
            except Exception as _e_sim:
                st.error(f"Erreur lors du recalcul de la simulation : {_e_sim}")

        _cmt("before_conclusion")
        st.markdown('<div style="margin-top:30px; border-top:2px solid #dee2e6; padding-top:12px; '
                    'text-align:center; color:#888; font-size:0.9em;">'
                    f'Rapport genere pour le plan <b>{plan_nom}</b></div>',
                    unsafe_allow_html=True)

    except Exception as e:
        st.error(f"Erreur lors de la generation du rapport : {e}")
        import traceback
        st.code(traceback.format_exc())


# ─── Main ──────────────────────────────────────────────────────────────────────

def _auto_save(p, plan_name):
    """
    Sauvegarde automatique en local seulement (pas de push GitHub).
    Le push GitHub est reserve aux sauvegardes explicites pour ne pas
    polluer l'historique avec les rerun Streamlit.
    """
    if st.session_state.get("auth_mode") == "visu":
        return  # Pas de sauvegarde en mode lecture
    if plan_name and plan_name != "Defaut":
        sauvegarder_plan(plan_name, p, local_only=True)
        st.session_state["_last_autosave"] = True


@st.dialog("Guide d'utilisation", width="large")
def _show_guide(mode):
    """Guide interactif adapte au mode (edit/visu)."""
    if mode == "edit":
        _steps = [
            ("\U0001F3E0 Page d'accueil",
             "Vous arrivez sur la **page d'accueil** qui liste tous vos plans financiers.\n\n"
             "- **Creer un plan** : entrez un nom et cliquez sur *Creer*\n"
             "- **Renommer** : cliquez sur \u270f\ufe0f a cote du plan\n"
             "- **Supprimer** : cliquez sur \U0001f5d1\ufe0f (une confirmation sera demandee)\n"
             "- **Deconnexion** : bouton en bas de page"),
            ("\U0001F4CA Visualisation du plan",
             "Cliquez sur **Visualisation du plan** pour generer le rapport complet :\n\n"
             "- Page de garde avec chiffres cles\n"
             "- Montage financier (Rocher + Chateau)\n"
             "- Injection par acteur (camembert)\n"
             "- Investissements initiaux\n"
             "- Plans financiers Rocher et Chateau avec graphiques\n"
             "- Indicateurs hoteliers (ADR, RevPAR, occupation)\n\n"
             "\U0001F4DD **Commentaires** : cliquez sur *\u270f\ufe0f + Commentaire* pour ajouter "
             "des annotations a n'importe quel endroit du rapport. Elles sont sauvegardees dans le plan."),
            ("\u2699\ufe0f Hypotheses & Projections",
             "Cliquez sur **Hypotheses, projections & methodologie** pour acceder a l'editeur :\n\n"
             "1. **Choix du module** : Immobiliere Rocher ou Chateau d'Argenteau\n"
             "2. **Onglet Hypotheses** : modifiez tous les parametres (tarifs, occupation, personnel, "
             "charges, investissements, prets...)\n"
             "3. **Onglet Projection** : visualisez les resultats calcules automatiquement\n"
             "4. **Onglet Methodologie** : comprenez les formules utilisees\n\n"
             "\U0001F4BE Les modifications sont **sauvegardees automatiquement**."),
            ("\U0001F4E5 Export PDF",
             "Dans le rapport (Visualisation du plan) :\n\n"
             "1. Cliquez sur **Version PDF** en haut\n"
             "2. Dans la boite de dialogue d'impression (**Ctrl+P**), cochez "
             "**Graphiques d'arriere-plan** dans les options\n"
             "3. Choisissez *Enregistrer au format PDF*\n\n"
             "\u26A0\ufe0f Sans l'option 'Graphiques d'arriere-plan', les couleurs seront delavees."),
        ]
    else:
        _steps = [
            ("\U0001F3E0 Page d'accueil",
             "Vous etes en **mode lecture seule**. Vous pouvez consulter tous les plans "
             "mais ne pouvez pas modifier les donnees.\n\n"
             "- Les champs de saisie sont desactives (grises)\n"
             "- Les boutons de creation, renommage et suppression ne sont pas affiches"),
            ("\U0001F4CA Visualisation du plan",
             "Cliquez sur **Visualisation du plan** pour consulter le rapport complet :\n\n"
             "- Tous les graphiques et tableaux sont accessibles\n"
             "- Les commentaires ajoutes par l'equipe sont visibles\n"
             "- Vous ne pouvez pas ajouter ou modifier de commentaires"),
            ("\u2699\ufe0f Hypotheses & Projections",
             "Cliquez sur **Hypotheses, projections & methodologie** pour consulter :\n\n"
             "- Les hypotheses encodees (en lecture seule)\n"
             "- Les projections calculees\n"
             "- La methodologie de calcul\n\n"
             "Aucune modification ne sera enregistree."),
            ("\U0001F4E5 Export PDF",
             "Vous pouvez exporter le rapport en PDF :\n\n"
             "1. Ouvrez la **Visualisation du plan**\n"
             "2. Cliquez sur **Version PDF**\n"
             "3. **Ctrl+P** > cochez **Graphiques d'arriere-plan** > *Enregistrer au format PDF*"),
        ]

    _step = st.radio("Etape", [s[0] for s in _steps], horizontal=True, label_visibility="collapsed")
    _idx = next(i for i, s in enumerate(_steps) if s[0] == _step)
    st.markdown(f"### {_steps[_idx][0]}")
    st.markdown(_steps[_idx][1])

    # Barre de progression
    st.progress((_idx + 1) / len(_steps))
    st.caption(f"Etape {_idx + 1} / {len(_steps)}")


def main():
    from pathlib import Path as _Path

    # ─── Authentification ──────────────────────────────────────────────
    if "auth_mode" not in st.session_state:
        _, col_login, _ = st.columns([1, 2, 1])
        with col_login:
            st.markdown("")
            st.markdown("")
            _login_photo = _Path(__file__).parent / "assets" / "chateau_1.jpg"
            if _login_photo.exists():
                import base64 as _b64_login
                with open(str(_login_photo), "rb") as _fl:
                    _login_img = _b64_login.b64encode(_fl.read()).decode()
                st.markdown(
                    f'<div style="text-align:center; margin-bottom:20px;">'
                    f'<img src="data:image/jpeg;base64,{_login_img}" '
                    f'style="width:200px; height:200px; object-fit:cover; border-radius:50%; '
                    f'box-shadow:0 4px 15px rgba(0,0,0,0.2);"></div>',
                    unsafe_allow_html=True,
                )
            st.markdown('<h2 style="text-align:center; margin-bottom:4px;">Plan Financier Hotelier</h2>',
                        unsafe_allow_html=True)
            st.markdown('<p style="text-align:center; color:#888; margin-bottom:24px;">Entrez votre code d\'acces</p>',
                        unsafe_allow_html=True)
            _code = st.text_input("Code d'acces", type="password", key="_login_code",
                                  placeholder="Saisissez votre code...")
            if st.button("Se connecter", use_container_width=True, type="primary"):
                if _code == "ArgenteauEdit!":
                    st.session_state["auth_mode"] = "edit"
                    st.rerun()
                elif _code == "ArgenteauVisu!":
                    st.session_state["auth_mode"] = "visu"
                    st.rerun()
                else:
                    st.error("Code incorrect.")
        return

    _is_visu = st.session_state.get("auth_mode") == "visu"

    # ── Mode lecture : desactiver tous les inputs via CSS ──
    if _is_visu:
        st.markdown("""<style>
            /* Desactiver tous les inputs */
            [data-testid="stNumberInput"] input,
            [data-testid="stNumberInput"] button,
            [data-testid="stTextInput"] input,
            [data-testid="stTextArea"] textarea,
            [data-testid="stSelectbox"] > div,
            [data-testid="stMultiSelect"] > div,
            [data-testid="stDateInput"] input,
            [data-testid="stSlider"] > div,
            .stSlider > div,
            [data-testid="stColorPicker"],
            [data-testid="stCheckbox"],
            [data-testid="stFileUploader"] {
                pointer-events: none !important;
                opacity: 0.6 !important;
            }
        </style>""", unsafe_allow_html=True)
        st.warning("\U0001F512 Mode lecture seule — les modifications sont desactivees")
        # Variable globale pour desactiver les boutons de modification
        st.session_state["_visu_mode"] = True

    # ─── Affichage rapport plein ecran ──────────────────────────────────
    if "_rapport_plan" in st.session_state:
        _rpt_nom = st.session_state["_rapport_plan"]
        c_back, c_save_gh, c_dl, c_html = st.columns([1, 1.2, 1, 1.2])
        with c_save_gh:
            _save_disabled = st.session_state.get("auth_mode") == "visu"
            if st.button("\U0001F4BE Sauvegarder sur GitHub",
                         key="btn_save_gh_rpt", use_container_width=True,
                         disabled=_save_disabled,
                         help="Pousse les commentaires et modifications du rapport sur GitHub."):
                try:
                    _p_to_save = charger_plan(_rpt_nom)
                    sauvegarder_plan(_rpt_nom, _p_to_save)
                    st.toast(f"Plan « {_rpt_nom} » sauvegarde sur GitHub")
                except Exception as _e_sv:
                    st.toast(f"Erreur sauvegarde : {_e_sv}", icon="⚠️")
        with c_back:
            if st.button("\u2B05 Retour a l'accueil", key="close_rapport_top"):
                st.session_state.pop("_rapport_plan", None)
                st.rerun()
        with c_dl:
            if st.button("\U0001F4E5 Version PDF", key="btn_pdf_mode", use_container_width=True):
                st.session_state["_rapport_print"] = True
                st.rerun()
        with c_html:
            try:
                from html_export import build_rapport_html as _build_html_rpt
                _p_for_export = charger_plan(_rpt_nom)
                _html_bytes = _build_html_rpt(_rpt_nom, _p_for_export).encode("utf-8")
                st.download_button(
                    "\U0001F4C4 Export HTML interactif",
                    data=_html_bytes,
                    file_name=f"rapport_{_rpt_nom.replace(' ', '_')}.html",
                    mime="text/html",
                    key="btn_html_export",
                    use_container_width=True,
                    help="Telecharge un rapport HTML autonome avec graphiques interactifs.",
                )
            except Exception as _e_html:
                st.button("\U0001F4C4 Export HTML (erreur)",
                          key="btn_html_err", use_container_width=True, disabled=True,
                          help=f"Erreur : {_e_html}")

        _is_print = st.session_state.get("_rapport_print", False)
        if _is_print:
            st.info("\U0001F5A8 Mode PDF actif — Utilisez **Ctrl+P** puis 'Enregistrer au format PDF'.\n\n"
                    "\u26A0\ufe0f **Important** : dans la boite de dialogue d'impression, cochez "
                    "**'Graphiques d'arriere-plan'** (ou 'Background graphics') dans les options "
                    "supplementaires pour que les couleurs des graphiques soient correctement imprimees.")
            import streamlit.components.v1 as _comp_pr
            _comp_pr.html('<script>setTimeout(function(){window.parent.print();}, 2000);</script>', height=0)
            if st.button("Retour au mode normal", key="btn_exit_print"):
                st.session_state.pop("_rapport_print", None)
                st.rerun()
        _render_rapport_complet(_rpt_nom, _Path, print_mode=_is_print)
        return

    # ─── Ecran d'accueil : choix du plan ──────────────────────────────────
    if "plan_actif" not in st.session_state:
        _photos = [_Path(__file__).parent / "assets" / f"chateau_{i}.jpg" for i in range(1, 5)]
        _photos_exist = [p_img for p_img in _photos if p_img.exists()]
        if _photos_exist:
            # Photos reduites avec marges laterales
            _, col_photos, _ = st.columns([1, 4, 1])
            with col_photos:
                cols_ph = st.columns(len(_photos_exist))
                for idx, ph in enumerate(_photos_exist):
                    with cols_ph[idx]:
                        import base64 as _b64_acc
                        with open(str(ph), "rb") as _f_img:
                            _img_data = _b64_acc.b64encode(_f_img.read()).decode()
                        st.markdown(
                            f'<img src="data:image/jpeg;base64,{_img_data}" '
                            f'style="width:100%; border-radius:8px; pointer-events:none;">',
                            unsafe_allow_html=True,
                        )

        st.markdown('<h1 style="text-align:center;">Plan Financier Hotelier</h1>', unsafe_allow_html=True)
        st.markdown('<p style="text-align:center; color:#666;">Selectionnez un plan existant ou creez-en un nouveau.</p>',
                    unsafe_allow_html=True)

        # Bouton guide centre
        _, _col_guide, _ = st.columns([2, 1, 2])
        with _col_guide:
            if st.button("\u2753 Guide d'utilisation", key="btn_guide", use_container_width=True):
                _show_guide("visu" if _is_visu else "edit")

        plans_existants = lister_plans()

        if _is_visu:
            # Mode lecture : plans centres
            _, _col_plans, _ = st.columns([1, 4, 1])
        else:
            # Mode edition : creation + plans sur toute la largeur
            st.markdown("### Nouveau plan")
            _cn1, _cn2, _cn3 = st.columns([3, 1, 4])
            with _cn1:
                nouveau_nom = st.text_input("Nom du plan", placeholder="Ex: Chateau d'Argenteau v2", key="new_plan_name")
            with _cn2:
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("Creer", key="btn_creer_plan", disabled=not nouveau_nom.strip(), use_container_width=True):
                    # Partir des hypotheses du Plan normal s'il existe, sinon des defauts
                    try:
                        p_new = charger_plan("Plan normal")
                    except Exception:
                        p_new = params_defaut()
                    sauvegarder_plan(nouveau_nom.strip(), p_new)
                    st.session_state["params_charges"] = p_new
                    st.session_state["plan_actif"] = nouveau_nom.strip()
                    st.rerun()
            st.markdown("### Ouvrir un plan existant")

        # Contenu des plans (dans _col_plans si visu, sinon pleine largeur)
        _plan_container = _col_plans if _is_visu else st.container()
        with _plan_container:
            if plans_existants:
                for plan_nom in plans_existants:
                    _rename_key = f"_renaming_{plan_nom}"
                    _delete_key = f"_deleting_{plan_nom}"
                    if st.session_state.get(_rename_key, False):
                        # Mode renommage
                        c_input, c_ok, c_cancel = st.columns([4, 1, 1])
                        with c_input:
                            new_name = st.text_input("Nouveau nom", value=plan_nom, key=f"rename_input_{plan_nom}", label_visibility="collapsed")
                        with c_ok:
                            if st.button("\u2714", key=f"rename_ok_{plan_nom}", help="Valider"):
                                new_name = new_name.strip()
                                if new_name and new_name != plan_nom:
                                    if renommer_plan(plan_nom, new_name):
                                        del st.session_state[_rename_key]
                                        st.rerun()
                                    else:
                                        st.toast("Ce nom existe deja.", icon="\u26a0\ufe0f")
                                else:
                                    del st.session_state[_rename_key]
                                    st.rerun()
                        with c_cancel:
                            if st.button("\u2716", key=f"rename_cancel_{plan_nom}", help="Annuler"):
                                del st.session_state[_rename_key]
                                st.rerun()
                    elif st.session_state.get(_delete_key, False):
                        # Mode confirmation suppression
                        st.warning(f"Supprimer le plan **{plan_nom}** ? Cette action est irreversible.")
                        c_yes, c_no = st.columns(2)
                        with c_yes:
                            if st.button("Oui, supprimer", key=f"delete_yes_{plan_nom}", type="primary", use_container_width=True):
                                supprimer_plan(plan_nom)
                                del st.session_state[_delete_key]
                                st.rerun()
                        with c_no:
                            if st.button("Annuler", key=f"delete_no_{plan_nom}", use_container_width=True):
                                del st.session_state[_delete_key]
                                st.rerun()
                    else:
                        # Mode normal — nom du plan + actions
                        if _is_visu:
                            _c_name, _c_actions = st.columns([1, 3])
                        else:
                            _c_name, _c_actions, _c_mgmt = st.columns([1, 3, 1])
                        with _c_name:
                            st.markdown(f"**\U0001F4C2 {plan_nom}**")
                        with _c_actions:
                            _ca1, _ca2 = st.columns(2)
                            with _ca1:
                                if st.button("\U0001F4CA Visualisation du plan", key=f"view_{plan_nom}", use_container_width=True):
                                    st.session_state["_rapport_plan"] = plan_nom
                                    st.rerun()
                            with _ca2:
                                if st.button("\u2699\ufe0f Hypotheses, projections & methodologie", key=f"open_{plan_nom}", use_container_width=True):
                                    st.session_state["params_charges"] = charger_plan(plan_nom)
                                    st.session_state["plan_actif"] = plan_nom
                                    st.rerun()
                        if not _is_visu:
                            with _c_mgmt:
                                _cm1, _cm2 = st.columns(2)
                                with _cm1:
                                    if st.button("\u270f\ufe0f", key=f"rename_{plan_nom}", help="Renommer"):
                                        st.session_state[_rename_key] = True
                                        st.rerun()
                                with _cm2:
                                    if st.button("\U0001f5d1\ufe0f", key=f"delete_{plan_nom}", help="Supprimer"):
                                        st.session_state[_delete_key] = True
                                        st.rerun()
            else:
                st.info("Aucun plan sauvegarde.")

        st.markdown("---")
        _mode_label = "Lecture seule" if _is_visu else "Edition"
        _col_logout, _ = st.columns([1, 3])
        with _col_logout:
            if st.button(f"\U0001F6AA Deconnexion ({_mode_label})", key="btn_logout", use_container_width=True):
                for k in list(st.session_state.keys()):
                    del st.session_state[k]
                st.rerun()

        return  # Stop ici, on n'affiche pas le reste

    # ─── Choix du module ──────────────────────────────────────────────────
    if "module_actif" not in st.session_state:
        # Charger les params pour date/duree
        if "params_charges" in st.session_state:
            _p_mod = st.session_state["params_charges"]
        else:
            _p_mod = params_defaut()

        _plan_name = st.session_state["plan_actif"]
        _photos_mod = [_Path(__file__).parent / "assets" / f"chateau_{i}.jpg" for i in range(1, 5)]
        _photos_mod = [p_img for p_img in _photos_mod if p_img.exists()]
        import base64 as _b64_mod

        # ── Titre du plan (compact) ──
        st.markdown(
            f'<h2 style="text-align:center; margin:0 0 4px 0;">{_plan_name}</h2>'
            f'<p style="text-align:center; color:#888; margin:0 0 20px 0; font-size:0.95em;">Selectionnez un module pour commencer</p>',
            unsafe_allow_html=True,
        )

        # ── Parametres communs (verrouilles par defaut) ──
        st.markdown(
            '<p style="margin: 0 0 2px 0; font-size: 0.85em; font-weight: 600; '
            'color: #667eea; text-transform: uppercase; letter-spacing: 0.05em;">'
            '\u2699\ufe0f  Parametres communs</p>'
            '<hr style="margin: 0 0 10px 0; border: none; border-top: 2px solid #667eea;">',
            unsafe_allow_html=True,
        )

        _date_fin = _p_mod["date_ouverture"].year + _p_mod["nb_mois_projection"] // 12
        _editing_communs = st.session_state.get("_editing_communs", False)

        if not _editing_communs:
            # Mode lecture : afficher les valeurs et un bouton Modifier
            _c1, _c2, _c3 = st.columns([2, 2, 1])
            with _c1:
                st.markdown(
                    f'<p style="margin:0; font-size:0.85em; color:#888;">Date de debut d\'activite</p>'
                    f'<p style="margin:0; font-size:1.1em; font-weight:600;">{_p_mod["date_ouverture"].strftime("%d/%m/%Y")}</p>',
                    unsafe_allow_html=True)
            with _c2:
                st.markdown(
                    f'<p style="margin:0; font-size:0.85em; color:#888;">Duree de projection</p>'
                    f'<p style="margin:0; font-size:1.1em; font-weight:600;">'
                    f'{_p_mod["nb_mois_projection"] // 12} ans ({_p_mod["nb_mois_projection"]} mois)'
                    f' &mdash; jusqu\'en {_date_fin}</p>',
                    unsafe_allow_html=True)
            with _c3:
                if not _is_visu:
                    st.markdown('<div style="margin-top: 4px;"></div>', unsafe_allow_html=True)
                    if st.button("\u270f\ufe0f Modifier", key="btn_edit_communs", use_container_width=True):
                        st.session_state["_editing_communs"] = True
                        st.rerun()
        else:
            # Mode edition
            _c1, _c2 = st.columns(2)
            with _c1:
                _p_mod["date_ouverture"] = st.date_input(
                    "Date de debut d'activite",
                    _p_mod.get("date_ouverture", date(2029, 7, 1)),
                    key="date_ouv_commun",
                )
            with _c2:
                _nb_ans = st.select_slider(
                    "Duree de projection",
                    options=list(range(3, 21)),
                    value=_p_mod["nb_mois_projection"] // 12,
                    format_func=lambda x: f"{x} ans ({x * 12} mois)",
                    key="duree_proj_commun",
                )
                _p_mod["nb_mois_projection"] = _nb_ans * 12

            _date_fin = _p_mod["date_ouverture"].year + _p_mod["nb_mois_projection"] // 12
            st.caption(
                f"Projection : {_p_mod['date_ouverture'].strftime('%B %Y')} \u2192 {_date_fin} "
                f"({_p_mod['nb_mois_projection'] // 12} ans)"
            )

            _bc1, _bc2 = st.columns(2)
            with _bc1:
                if st.button("\U0001F4BE Sauver et fermer", key="btn_save_communs", use_container_width=True, type="primary"):
                    sauvegarder_plan(st.session_state.get("plan_actif", ""), _p_mod)
                    st.session_state["_editing_communs"] = False
                    st.toast("Parametres sauvegardes !")
                    st.rerun()
            with _bc2:
                if st.button("Annuler", key="btn_cancel_communs", use_container_width=True):
                    st.session_state["_editing_communs"] = False
                    st.rerun()

        # ── Separation visuelle ──
        st.markdown(
            '<p style="margin: 16px 0 2px 0; font-size: 0.85em; font-weight: 600; '
            'color: #f5576c; text-transform: uppercase; letter-spacing: 0.05em;">'
            '\U0001F4CB  Plans financiers</p>'
            '<hr style="margin: 0 0 10px 0; border: none; border-top: 2px solid #f5576c;">',
            unsafe_allow_html=True,
        )

        st.session_state["params_charges"] = _p_mod
        _auto_save(_p_mod, st.session_state.get("plan_actif", ""))

        st.markdown("")

        # ── Cartes modules ──
        c1, c2 = st.columns(2, gap="large")
        with c1:
            _rocher_img_html = ""
            if len(_photos_mod) > 1:
                with open(str(_photos_mod[1]), "rb") as _fr:
                    _rocher_img_html = (
                        f'<img src="data:image/jpeg;base64,{_b64_mod.b64encode(_fr.read()).decode()}" '
                        f'style="width:100%; height:100px; object-fit:cover; border-radius:12px 12px 0 0;">'
                    )
            st.markdown(
                f'<div style="border:1px solid #dee2e6; border-radius:12px; overflow:hidden; '
                f'background:white; box-shadow:0 2px 8px rgba(0,0,0,0.06);">'
                f'{_rocher_img_html}'
                f'<div style="padding:16px; text-align:center;">'
                f'<h3 style="margin:0 0 4px 0; color:#11998e; font-size:1.15em;">\U0001F3E2 Immobiliere Rocher</h3>'
                f'<p style="color:#666; font-size:0.85em; margin:0;">Revenus locatifs, prets et amortissements</p>'
                f'</div></div>',
                unsafe_allow_html=True,
            )
            if st.button("Ouvrir Immobiliere Rocher", key="mod_rocher", use_container_width=True, type="primary"):
                st.session_state["module_actif"] = "rocher"
                st.rerun()

        with c2:
            _chateau_img_html = ""
            if len(_photos_mod) > 2:
                with open(str(_photos_mod[2]), "rb") as _fc:
                    _chateau_img_html = (
                        f'<img src="data:image/jpeg;base64,{_b64_mod.b64encode(_fc.read()).decode()}" '
                        f'style="width:100%; height:100px; object-fit:cover; border-radius:12px 12px 0 0;">'
                    )
            st.markdown(
                f'<div style="border:1px solid #dee2e6; border-radius:12px; overflow:hidden; '
                f'background:white; box-shadow:0 2px 8px rgba(0,0,0,0.06);">'
                f'{_chateau_img_html}'
                f'<div style="padding:16px; text-align:center;">'
                f'<h3 style="margin:0 0 4px 0; color:#4facfe; font-size:1.15em;">\U0001F3E8 Activite Chateau d\'Argenteau</h3>'
                f'<p style="color:#666; font-size:0.85em; margin:0;">Plan financier de l\'exploitation hoteliere</p>'
                f'</div></div>',
                unsafe_allow_html=True,
            )
            if st.button("Ouvrir Activite Chateau", key="mod_chateau", use_container_width=True, type="primary"):
                st.session_state["module_actif"] = "chateau"
                st.rerun()

        st.markdown("")

        # ── Retour ──
        if st.button("\u2B05 Revenir a l'accueil", key="btn_back_plan"):
            st.session_state.pop("plan_actif", None)
            st.session_state.pop("params_charges", None)
            st.rerun()

        return

    module = st.session_state["module_actif"]

    # ─── Module Immobiliere Rocher ────────────────────────────────────────
    if module == "rocher":
        _module_rocher()
        return


    # ─── Plan actif : charger les parametres ──────────────────────────────
    if "params_charges" in st.session_state:
        p = st.session_state["params_charges"]
    else:
        p = params_defaut()

    plan_actif = st.session_state.get("plan_actif", "Defaut")

    # ─── Barre superieure : titre + sauvegarde + changer de plan ──────────
    _photo_main = _Path(__file__).parent / "assets" / "chateau_1.jpg"
    _photo_2 = _Path(__file__).parent / "assets" / "chateau_4.jpg"

    col_img1, col_title, col_save, col_img2 = st.columns([0.8, 3, 1, 0.8])
    with col_img1:
        if _photo_main.exists():
            st.image(str(_photo_main), use_container_width=True)
    with col_title:
        st.markdown(
            f'<h2 style="text-align:center; margin:0; padding-top:8px;">'
            f'Plan Financier<br><span style="font-size:0.55em; color:#666;">'
            f'{p["nom_hotel"]} — {plan_actif}</span></h2>',
            unsafe_allow_html=True,
        )
    with col_save:
        st.markdown("<br>", unsafe_allow_html=True)
        c_s1, c_s2 = st.columns(2)
        with c_s1:
            if st.button("\U0001F4BE Sauvegarder", key="btn_save_top", use_container_width=True, disabled=_is_visu):
                sauvegarder_plan(plan_actif, p)
                st.toast(f"Plan \u00ab {plan_actif} \u00bb sauvegarde !")
        with c_s2:
            if st.button("\u2B05 Retour", key="btn_change_plan", use_container_width=True):
                st.session_state.pop("module_actif", None)
                st.rerun()
    with col_img2:
        if _photo_2.exists():
            st.image(str(_photo_2), use_container_width=True)

    if github_sync.is_enabled():
        st.caption(
            "<div style='text-align:center; color:#11998e; font-size:0.8rem;'>"
            "\U0001F7E2 Persistance GitHub active &mdash; vos sauvegardes sont conservees apres redemarrage."
            "</div>",
            unsafe_allow_html=True,
        )
    else:
        st.caption(
            "<div style='text-align:center; color:#999; font-size:0.8rem;'>"
            "\U0001F7E0 Mode local &mdash; sauvegardes uniquement sur la machine actuelle."
            "</div>",
            unsafe_allow_html=True,
        )

    # ─── Onglets principaux ───────────────────────────────────────────────
    tab0, tab_indic, tab_methodo = st.tabs([
        "\U0001F4DD Hypotheses",
        "\U0001F4CA Projection",
        "\U0001F4D6 Methodologie",
    ])

    with tab0:
        p = tab_hypotheses(p)

    # Auto-save apres chaque interaction
    st.session_state["params_charges"] = p
    _auto_save(p, plan_actif)

    st.caption(f"Debut d'activite : {p['date_ouverture'].strftime('%B %Y')} | "
               f"{p['nb_chambres']} chambres | "
               f"Projection : {p['nb_mois_projection'] // 12} ans")

    # Calculs
    with st.spinner("Calcul de la projection..."):
        df = projection_complete(p)
        # Marges mensuelles ventilees par service (ventes - CV - frais fixes directs)
        df["marge_heberg"] = df["ca_hebergement"] - df["cv_hebergement"] - df.get("cf_directs_hebergement", 0)
        df["marge_brass"] = df["ca_brasserie"] - df["cv_brasserie"] - df.get("cf_directs_brasserie", 0)
        df["marge_bar"] = df["ca_bar"] - df["cv_bar"] - df.get("cf_directs_bar", 0)
        df["marge_spa"] = df["ca_spa"] - df["cv_spa"] - df.get("cf_directs_spa", 0)
        df["marge_salles"] = df["ca_salles"] - df["cv_salles"] - df.get("cf_directs_evenements", 0)
        st.session_state["_projection_df"] = df
        indic_expl = indicateurs_annuels(df, p, par_calendaire=False)
        indic = indicateurs_annuels(df, p, par_calendaire=True)

    # Injecter les mini-graphiques dans l'onglet hypotheses
    with tab0:
        _render_hypotheses_charts(df, indic, params_hyp=p)

    with tab_indic:
        st.header("\U0001F4CA Projection")
        sub_resultats, sub_bilan, sub_kpi, sub_chiffres = st.tabs([
            "\U0001F4C8 Resultats",
            "\U0001F3E6 Bilan",
            "\U0001F3AF Indicateurs",
            "\U0001F4CB Chiffres bruts",
        ])

        # Donnees annuelles partagees
        _ann = df.groupby("annee").agg({
            "ca_total": "sum", "ca_hebergement": "sum", "ca_brasserie": "sum",
            "ca_bar": "sum", "ca_spa": "sum", "ca_salles": "sum",
            "ca_loyer_restaurant": "sum",
            "cv_total": "sum", "cv_hebergement": "sum", "cv_brasserie": "sum",
            "cv_bar": "sum", "cv_spa": "sum", "cv_salles": "sum",
            "cf_directs_total": "sum",
            "cf_directs_hebergement": "sum", "cf_directs_brasserie": "sum",
            "cf_directs_bar": "sum", "cf_directs_spa": "sum", "cf_directs_evenements": "sum",
            "cf_indirects_total": "sum",
            "marge_brute": "sum", "marge": "sum", "subside_rw": "sum",
            "ebitda": "sum", "amortissement": "sum", "ebit": "sum",
            "dette_interets": "sum", "dette_capital": "sum",
            "impot": "sum", "resultat_net": "sum",
            "cash_flow": "sum", "cash_flow_cumul": "last",
            "nuitees": "sum", "taux_occupation": "mean",
            "cf_personnel": "sum", "cf_personnel_direct": "sum",
        }).reset_index()
        _x = [str(int(a)) for a in _ann["annee"]]
        # Marges ventilees par service (ventes - CV - frais fixes directs)
        _ann["marge_heberg"] = _ann["ca_hebergement"] - _ann["cv_hebergement"] - _ann["cf_directs_hebergement"]
        _ann["marge_brass"] = _ann["ca_brasserie"] - _ann["cv_brasserie"] - _ann["cf_directs_brasserie"]
        _ann["marge_bar"] = _ann["ca_bar"] - _ann["cv_bar"] - _ann["cf_directs_bar"]
        _ann["marge_spa"] = _ann["ca_spa"] - _ann["cv_spa"] - _ann["cf_directs_spa"]
        _ann["marge_salles"] = _ann["ca_salles"] - _ann["cv_salles"] - _ann["cf_directs_evenements"]
        _ann["marge_resto"] = _ann["ca_loyer_restaurant"]

        def _proj_chart(title, traces_config, key_id, height=450, barmode="stack",
                        show_totals=False, total_col=None, extra_lines=None):
            """Graphique de projection avec controles unite/periode."""
            with st.expander(title, expanded=False):
                _c1, _c2 = st.columns(2)
                with _c1:
                    _u = st.radio("Unite", ["\u20ac", "K\u20ac"], key=f"pu_{key_id}", horizontal=True)
                with _c2:
                    _per = st.radio("Periode", ["Annuel", "Mensuel"], key=f"pp_{key_id}", horizontal=True)
                _div = 1000 if _u.startswith("K") else 1
                _sfx = " K\u20ac" if _u.startswith("K") else " \u20ac"

                fig = go.Figure()
                if _per == "Annuel":
                    for col, lbl, clr in traces_config:
                        vals = _ann[col] / _div
                        fig.add_trace(go.Bar(x=_x, y=vals, name=lbl, marker_color=clr,
                            hovertemplate="%{x} - " + lbl + "<br>%{y:,.0f}" + _sfx + "<extra></extra>"))
                    if show_totals and total_col:
                        _tots = _ann[total_col] / 1000
                        _tots_pos = _ann[total_col] / _div
                        fig.add_trace(go.Scatter(x=_x, y=[v*1.02 for v in _tots_pos],
                            mode="text", text=[f"<b>{v:,.0f} K\u20ac</b>" for v in _tots],
                            textposition="top center", textfont=dict(size=10, color="black"),
                            showlegend=False, hoverinfo="skip", cliponaxis=False))
                    if extra_lines:
                        for col, lbl, clr, dash in extra_lines:
                            fig.add_trace(go.Scatter(x=_x, y=_ann[col]/_div, name=lbl,
                                mode="lines+markers", line=dict(color=clr, width=3, dash=dash or "solid")))
                    _max_y = max(_ann[total_col or traces_config[0][0]]) / _div if total_col or traces_config else 1
                    fig.update_layout(height=height, barmode=barmode,
                        xaxis=dict(type="category", tickfont=dict(size=12)),
                        yaxis=dict(tickformat=",.0f", title=_sfx.strip(),
                                   range=[min(0, fig.data[0].y.min() if hasattr(fig.data[0], 'y') else 0) * 1.1,
                                          _max_y * 1.15] if show_totals else None),
                        legend=dict(orientation="h", y=1.08, xanchor="center", x=0.5, font=dict(size=13)))
                else:
                    for col, lbl, clr in traces_config:
                        fig.add_trace(go.Bar(x=df["date"], y=df[col]/_div, name=lbl, marker_color=clr,
                            hovertemplate="%{x} - " + lbl + "<br>%{y:,.0f}" + _sfx + "<extra></extra>"))
                    if extra_lines:
                        for col, lbl, clr, dash in extra_lines:
                            fig.add_trace(go.Scatter(x=df["date"], y=df[col]/_div, name=lbl,
                                mode="lines", line=dict(color=clr, width=2, dash=dash or "solid")))
                    fig.update_layout(height=height, barmode=barmode,
                        xaxis=dict(dtick="M3", tickformat="%b %Y", tickfont=dict(size=10), tickangle=-45),
                        yaxis=dict(tickformat=",.0f", title=_sfx.strip()),
                        legend=dict(orientation="h", y=1.08, xanchor="center", x=0.5, font=dict(size=13)))
                fig.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
                st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

        # ── 1. RESULTATS ──
        with sub_resultats:
            st.markdown("### Par service")
            _proj_chart("Ventes par service + Subside", [
                ("ca_hebergement","Hebergement","#667eea"),("ca_brasserie","Brasserie","#f5576c"),
                ("ca_bar","Bar","#ffcc00"),("ca_spa","Spa","#11998e"),("ca_salles","Salles","#a0522d"),
                ("ca_loyer_restaurant","Location resto.","#ff8c00"),("subside_rw","Subside RW","#f093fb"),
            ], "ca_svc", show_totals=True, total_col="ca_total")

            _proj_chart("Charges variables par service", [
                ("cv_hebergement","Hebergement","#667eea"),("cv_brasserie","Brasserie","#f5576c"),
                ("cv_bar","Bar","#ffcc00"),("cv_spa","Spa","#11998e"),
            ], "cv_svc", show_totals=True, total_col="cv_total")

            _proj_chart("Frais fixes directs par service", [
                ("cf_directs_hebergement","Hebergement","#667eea"),("cf_directs_brasserie","Brasserie","#f5576c"),
                ("cf_directs_bar","Bar","#ffcc00"),("cf_directs_spa","Spa","#11998e"),
                ("cf_directs_evenements","Evenements","#a0522d"),
            ], "cfd", show_totals=True, total_col="cf_directs_total")

            _proj_chart("Marge par service", [
                ("marge_heberg","Hebergement","#667eea"),("marge_brass","Brasserie","#f5576c"),
                ("marge_bar","Bar","#ffcc00"),("marge_spa","Spa","#11998e"),
                ("marge_salles","Salles","#a0522d"),("marge_resto","Location resto.","#ff8c00"),
            ], "marge_svc")

            st.markdown("### Compte de resultat")
            _proj_chart("EBITDA et resultats", [
                ("marge_brute","Marge Brute","#38ef7d"),
            ], "pl", extra_lines=[
                ("ebitda","EBITDA","#11998e", None),
                ("ebit","EBIT","#667eea","dash"),
                ("resultat_net","Resultat Net","#f5576c", None),
            ])

            st.markdown("### Cash flow")
            _proj_chart("Cash flow par periode", [
                ("cash_flow","Cash flow","#4facfe"),
            ], "cf", extra_lines=[
                ("cash_flow_cumul","Cash flow cumule","#f5576c", None),
            ])

            st.markdown("### Tresorerie")
            surplus_depart = p.get("fonds_propres_initial", 0) + sum(pr["montant"] for pr in p["prets"]) - sum(inv["montant"] for inv in p["investissements"])
            df["_treso"] = df["cash_flow"].cumsum() + surplus_depart
            _ann["_treso"] = _ann["cash_flow_cumul"] + surplus_depart
            _proj_chart("Tresorerie (incl. surplus depart)", [
            ], "treso", extra_lines=[
                ("_treso","Tresorerie","#11998e", None),
            ])

        # ── 2. BILAN ──
        with sub_bilan:
            st.markdown("### Evolution des fonds propres")
            fp_initial = p.get("fonds_propres_initial", 0)
            # Pret Rocher = quasi fonds propres
            pret_rocher_montant = 0
            for pr in p.get("prets", []):
                if "rocher" in pr.get("nom", "").lower():
                    pret_rocher_montant = pr["montant"]
            res_cumul_ann = _ann["resultat_net"].cumsum()

            fig_fp = go.Figure()
            fig_fp.add_trace(go.Bar(x=_x, y=[fp_initial/1000]*len(_x), name="Fonds propres durs", marker_color="#38ef7d"))
            fig_fp.add_trace(go.Bar(x=_x, y=[pret_rocher_montant/1000]*len(_x), name="Quasi-FP (Pret Rocher)", marker_color="#11998e"))
            fig_fp.add_trace(go.Scatter(x=_x, y=(fp_initial + pret_rocher_montant + res_cumul_ann) / 1000,
                name="FP + Quasi-FP + Resultat cumule", mode="lines+markers", line=dict(color="#667eea", width=3)))
            fig_fp.update_layout(height=450, barmode="stack", xaxis=dict(type="category"),
                yaxis=dict(tickformat=",.0f", title="K\u20ac"),
                legend=dict(orientation="h", y=1.08, xanchor="center", x=0.5, font=dict(size=13)))
            st.plotly_chart(fig_fp, use_container_width=True, config={"displayModeBar": False})

            st.markdown("### Evolution de l'endettement")
            # Calculer encours dette par annee
            from dateutil.relativedelta import relativedelta as _rd_b
            enc_dette_ann = []
            for a in _ann["annee"]:
                d_fin = date(int(a), 12, 31)
                enc = 0
                for pr in p.get("prets", []):
                    if pr["montant"] > 0:
                        df_pr = calc_tableau_pret(pr, p["date_ouverture"], p["nb_mois_projection"])
                        row_pr = df_pr[df_pr["date"] <= d_fin]
                        if not row_pr.empty:
                            enc += row_pr.iloc[-1]["capital_restant"]
                enc_dette_ann.append(enc)

            fig_dette = go.Figure()
            fig_dette.add_trace(go.Bar(x=_x, y=[v/1000 for v in enc_dette_ann], name="Encours dettes", marker_color="#f5576c"))
            fig_dette.update_layout(height=400, xaxis=dict(type="category"),
                yaxis=dict(tickformat=",.0f", title="K\u20ac"), showlegend=False)
            st.plotly_chart(fig_dette, use_container_width=True, config={"displayModeBar": False})

            # ── BFR (Besoin en Fonds de Roulement) ──
            if "creances_clients" in df.columns:
                st.markdown("### Besoin en Fonds de Roulement (BFR)")
                _bfr_ann = df.groupby("annee").agg({
                    "creances_clients": "last", "dettes_fournisseurs": "last", "bfr": "last"
                }).reset_index()
                _bfr_x = [str(int(a)) for a in _bfr_ann["annee"]]

                fig_bfr = go.Figure()
                fig_bfr.add_trace(go.Bar(x=_bfr_x, y=_bfr_ann["creances_clients"]/1000,
                    name="Creances clients", marker_color="#4facfe"))
                fig_bfr.add_trace(go.Bar(x=_bfr_x, y=-_bfr_ann["dettes_fournisseurs"]/1000,
                    name="Dettes fournisseurs", marker_color="#f5576c"))
                fig_bfr.add_trace(go.Scatter(x=_bfr_x, y=_bfr_ann["bfr"]/1000,
                    name="BFR net", mode="lines+markers", line=dict(color="#667eea", width=3)))
                fig_bfr.update_layout(height=400, barmode="relative",
                    xaxis=dict(type="category"), yaxis=dict(tickformat=",.0f", title="K\u20ac"),
                    legend=dict(orientation="h", y=1.08, xanchor="center", x=0.5, font=dict(size=13)))
                fig_bfr.add_hline(y=0, line_dash="dot", line_color="gray", opacity=0.5)
                st.plotly_chart(fig_bfr, use_container_width=True, config={"displayModeBar": False})

            # ── Bilan detaille (fin d'annee) ──
            if "dette_tva" in df.columns:
                st.markdown("### Bilan detaille (fin d'annee)")
                _bilan_rows = []
                for idx_b, a in enumerate(_ann["annee"]):
                    _mask_b = df["annee"] == a
                    _last_b = df[_mask_b].iloc[-1]
                    _enc_b = enc_dette_ann[idx_b] if idx_b < len(enc_dette_ann) else 0
                    _res_cum_b = res_cumul_ann.iloc[idx_b] if idx_b < len(res_cumul_ann) else 0
                    _bilan_rows.append({
                        "Annee": int(a),
                        "Creances clients": f"{_last_b['creances_clients']:,.0f} \u20ac",
                        "Dettes fournisseurs": f"{_last_b['dettes_fournisseurs']:,.0f} \u20ac",
                        "BFR": f"{_last_b['bfr']:,.0f} \u20ac",
                        "Dette TVA": f"{_last_b['dette_tva']:,.0f} \u20ac",
                        "Dette ISOC": f"{_last_b['dette_isoc']:,.0f} \u20ac",
                        "Provisions sociales": f"{_last_b['dette_sociale']:,.0f} \u20ac",
                        "Dettes bancaires": f"{_enc_b:,.0f} \u20ac",
                        "Fonds propres": f"{fp_initial + pret_rocher_montant + _res_cum_b:,.0f} \u20ac",
                    })
                st.dataframe(pd.DataFrame(_bilan_rows), use_container_width=True, hide_index=True,
                             height=min(600, 35 * len(_bilan_rows) + 38))

        # ── 3. INDICATEURS ──
        with sub_kpi:
            st.markdown("### Indicateurs hoteliers")
            nb_ch = p["nb_chambres"]
            _ann["revpar"] = _ann["ca_hebergement"] / (nb_ch * 365)
            _ann["adr"] = _ann["ca_hebergement"] / _ann["nuitees"].replace(0, float("nan"))
            _ann["occ_pct"] = _ann["taux_occupation"] * 100

            fig_h = go.Figure()
            fig_h.add_trace(go.Bar(x=_x, y=_ann["nuitees"], name="Nuitees", marker_color="#4facfe"))
            fig_h.update_layout(height=350, xaxis=dict(type="category"),
                yaxis=dict(tickformat=",.0f", title="Nuitees"), showlegend=False,
                title="Nuitees par an")
            st.plotly_chart(fig_h, use_container_width=True, config={"displayModeBar": False})

            fig_occ = go.Figure()
            fig_occ.add_trace(go.Scatter(x=_x, y=_ann["occ_pct"], name="Taux occupation %", mode="lines+markers", line=dict(color="#11998e", width=3)))
            fig_occ.add_trace(go.Bar(x=_x, y=_ann["adr"], name="ADR (\u20ac)", marker_color="#667eea", opacity=0.6))
            fig_occ.add_trace(go.Bar(x=_x, y=_ann["revpar"], name="RevPAR (\u20ac)", marker_color="#f5576c", opacity=0.6))
            fig_occ.update_layout(height=400, barmode="group", xaxis=dict(type="category"),
                yaxis=dict(title=""), title="Taux occupation, ADR et RevPAR",
                legend=dict(orientation="h", y=1.08, xanchor="center", x=0.5, font=dict(size=13)))
            st.plotly_chart(fig_occ, use_container_width=True, config={"displayModeBar": False})

            st.markdown("### Indicateurs RH")
            pers_heberg = p.get("personnel_hebergement", [])
            pers_brass = p.get("personnel_brasserie", [])
            pers_spa = p.get("personnel_spa", [])
            pers_events = p.get("personnel_evenements", [])
            pers_indirect = p.get("personnel_indirect", [])

            etp_heberg = sum(pe["etp"] for pe in pers_heberg)
            etp_brass = sum(pe["etp"] for pe in pers_brass)
            etp_spa = sum(pe["etp"] for pe in pers_spa)
            etp_events = sum(pe["etp"] for pe in pers_events)
            etp_indirect = sum(pe["etp"] for pe in pers_indirect)
            total_etp = etp_heberg + etp_brass + etp_spa + etp_events + etp_indirect

            etp_par_chambre = total_etp / nb_ch if nb_ch > 0 else 0

            # Metriques en colonnes
            c1, c2, c3, c4 = st.columns(4)
            with c1:
                st.markdown(f"""<div style="text-align:center; padding:15px; background:#f0f2f6; border-radius:10px;">
                    <div style="font-size:2em; font-weight:bold; color:#667eea;">{total_etp:.1f}</div>
                    <div style="color:#666;">Total ETP</div></div>""", unsafe_allow_html=True)
            with c2:
                st.markdown(f"""<div style="text-align:center; padding:15px; background:#f0f2f6; border-radius:10px;">
                    <div style="font-size:2em; font-weight:bold; color:#11998e;">{etp_par_chambre:.2f}</div>
                    <div style="color:#666;">ETP / chambre</div></div>""", unsafe_allow_html=True)
            with c3:
                st.markdown(f"""<div style="text-align:center; padding:15px; background:#f0f2f6; border-radius:10px;">
                    <div style="font-size:2em; font-weight:bold; color:#f5576c;">{nb_ch / total_etp:.1f}</div>
                    <div style="color:#666;">Chambres / ETP</div></div>""" if total_etp > 0 else
                    '<div style="text-align:center; padding:15px; background:#f0f2f6; border-radius:10px;"><div>-</div></div>',
                    unsafe_allow_html=True)
            with c4:
                cp_rh = p.get("charges_patronales_pct", 0.35)
                masse_tot = sum(pe["cout_brut"] * (1+cp_rh) * pe["etp"] for pe in
                    pers_heberg + pers_brass + pers_spa + pers_events + pers_indirect)
                st.markdown(f"""<div style="text-align:center; padding:15px; background:#f0f2f6; border-radius:10px;">
                    <div style="font-size:2em; font-weight:bold; color:#764ba2;">{masse_tot/1000:,.0f} K\u20ac</div>
                    <div style="color:#666;">Masse salariale/an</div></div>""", unsafe_allow_html=True)

            # Repartition ETP par service
            st.markdown("")
            fig_etp = go.Figure(data=[go.Pie(
                labels=["Hebergement", "Brasserie", "Spa", "Evenements", "Indirect"],
                values=[etp_heberg, etp_brass, etp_spa, etp_events, etp_indirect],
                marker=dict(colors=["#667eea", "#f5576c", "#11998e", "#a0522d", "#764ba2"]),
                textinfo="label+value+percent",
                textfont=dict(size=13),
                hovertemplate="%{label}<br>%{value:.1f} ETP<br>%{percent}<extra></extra>",
            )])
            fig_etp.update_layout(height=350, showlegend=False, margin=dict(l=10, r=10, t=10, b=10))
            st.plotly_chart(fig_etp, use_container_width=True, config={"displayModeBar": False})

            st.markdown("### Ratio solvabilite")
            st.caption("(Fonds propres + Quasi-FP) / (Dettes + FP + Quasi-FP)")
            fp_total = fp_initial + pret_rocher_montant
            ratios = []
            for i, enc in enumerate(enc_dette_ann):
                fp_cum = fp_total + (res_cumul_ann.iloc[i] if i < len(res_cumul_ann) else 0)
                total_passif = enc + fp_cum
                ratio = fp_cum / total_passif * 100 if total_passif > 0 else 0
                ratios.append(ratio)
            fig_ratio = go.Figure()
            fig_ratio.add_trace(go.Scatter(x=_x, y=ratios, name="Ratio solvabilite %",
                mode="lines+markers", line=dict(color="#667eea", width=3),
                fill="tozeroy", fillcolor="rgba(102,126,234,0.1)"))
            fig_ratio.add_hline(y=30, line_dash="dash", line_color="orange", annotation_text="Seuil 30%")
            fig_ratio.update_layout(height=400, xaxis=dict(type="category"),
                yaxis=dict(tickformat=".0f", title="%", range=[0, 100]),
                showlegend=False, title="Ratio (FP + Quasi-FP) / Total passif")
            st.plotly_chart(fig_ratio, use_container_width=True, config={"displayModeBar": False})

        # ── 4. CHIFFRES BRUTS ──
        with sub_chiffres:
            tab_indicateurs(df, indic, p)
            # Export CSV
            st.markdown("---")
            import io as _io_csv
            df_csv = df.copy()
            df_csv["date"] = df_csv["date"].astype(str)
            csv_data = df_csv.to_csv(index=False)
            st.download_button(
                "\U0001F4E5 Telecharger les donnees en CSV",
                csv_data,
                file_name=f"projection_{p['nom_hotel']}.csv",
                mime="text/csv",
                key="export_csv_brut",
            )

    with tab_methodo:
        tab_methodologie(df, indic, p)


if __name__ == "__main__":
    main()
